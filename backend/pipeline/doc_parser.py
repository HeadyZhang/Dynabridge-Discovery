"""Document parsing module.

Parses uploaded PDF, DOCX, PPTX files into structured text
for AI analysis.
"""
from pathlib import Path

# Maximum characters to extract per document (50K covers most brand docs)
MAX_TEXT_LENGTH = 50000


async def parse_documents(file_paths: list[str]) -> list[dict]:
    """Parse uploaded documents and extract text content.

    Returns:
        [{"filename": str, "text": str, "tables": [...], "images": [...]}]
    """
    results = []

    for fp in file_paths:
        path = Path(fp)
        if not path.exists():
            continue

        suffix = path.suffix.lower()
        text = ""
        tables = []

        try:
            if suffix == ".pdf":
                text, tables = _parse_pdf(path)
            elif suffix in (".docx", ".doc"):
                text, tables = _parse_docx(path)
            elif suffix == ".txt":
                text = path.read_text(encoding="utf-8", errors="ignore")
            elif suffix == ".pptx":
                text = _parse_pptx(path)
            elif suffix in (".csv", ".tsv"):
                text = _parse_csv(path)
            else:
                text = f"[Unsupported file type: {suffix}]"
        except Exception as e:
            print(f"[doc_parser] Error parsing {path.name}: {e}")
            text = f"[Parse error: {e}]"

        results.append({
            "filename": path.name,
            "text": text[:MAX_TEXT_LENGTH],
            "tables": tables,
            "images": [],
        })

    return results


def _parse_pdf(path: Path) -> tuple[str, list]:
    """Extract text and tables from PDF using pdfplumber or PyPDF2."""
    try:
        import pdfplumber
        text_parts = []
        tables = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:100]:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
                for tbl in (page.extract_tables() or []):
                    if tbl and len(tbl) > 1:
                        headers = [str(cell or "").strip() for cell in tbl[0]]
                        rows = [[str(cell or "").strip() for cell in row] for row in tbl[1:]]
                        tables.append({"headers": headers, "rows": rows})
        return "\n\n".join(text_parts), tables
    except ImportError:
        pass

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        text = "\n\n".join(
            page.extract_text() or "" for page in reader.pages[:100]
        )
        return text, []
    except ImportError:
        return "[Install pdfplumber or PyPDF2 to parse PDFs]", []


def _parse_docx(path: Path) -> tuple[str, list]:
    """Extract text and tables from DOCX."""
    try:
        from docx import Document
        doc = Document(str(path))

        # Extract paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)

        # Extract tables
        tables = []
        for tbl in doc.tables:
            rows_data = []
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows_data.append(cells)
            if len(rows_data) > 1:
                tables.append({
                    "headers": rows_data[0],
                    "rows": rows_data[1:],
                })

        return text, tables
    except ImportError:
        return "[Install python-docx to parse DOCX files]", []


def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX slides including notes."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
                # Extract table text
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            texts.append(row_text)
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[Notes: {notes}]")
            if texts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        return "[Install python-pptx to parse PPTX files]"


def _parse_csv(path: Path) -> str:
    """Extract text from CSV/TSV files."""
    import csv
    delimiter = "\t" if path.suffix == ".tsv" else ","
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = list(reader)
        if not rows:
            return ""
        # Format as readable text
        headers = rows[0]
        lines = [" | ".join(headers)]
        lines.append("-" * len(lines[0]))
        for row in rows[1:200]:  # Limit to 200 rows
            lines.append(" | ".join(row))
        return "\n".join(lines)
    except Exception as e:
        return f"[CSV parse error: {e}]"
