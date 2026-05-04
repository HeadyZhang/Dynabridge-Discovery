"""PPT Generation module — clones slides from a reference PPTX template.

Instead of building slides from scratch with python-pptx shapes,
this module copies slides from the original template and replaces
text content with new analysis data. This preserves all formatting,
backgrounds, images, fonts, and layout exactly.

Images are replaced with brand-specific images collected by image_collector.py.
When no collected image is available, the template's original image is kept as-is.

Template slide index map (0-based):
  0  = Cover (Title_Slide)
  1  = Agenda (Blank)
  2  = Approach / Brand Building Process (Text Slide)
  3  = Step 1 Divider (Text Slide)
  4  = Section Header - Capabilities (Overview Slide)
  5  = Content slide - title + 3 bullets + insight + image (Blank)
  13 = Summary slide - title + paragraph + half-image (Text Slide)
  14 = Section Header - Competition (Overview Slide)
  17 = Competitor deep dive - two-column positioning + learnings (Blank)
  23 = Landscape summary - bullets + sidebar text (Blank)
  24 = Summary slide - Competition (Text Slide)
  25 = Section Header - Consumer (Overview Slide)
  91 = Thank You / 谢谢 (Divider Slide)
"""
import copy
import io
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.parts.chart import ChartPart, EmbeddedXlsxPart
from pptx.opc.packuri import PackURI
from pptx.opc.package import Part as OpcPart

sys.path.insert(0, str(Path(__file__).parent.parent))
from config import OUTPUT_DIR, PREVIEW_DIR, MODEL_OPUS, MODEL_SONNET

TEMPLATE_PATH = Path(__file__).parent.parent.parent / "templates" / "cozyfit_reference.pptx"

# Template slide indices (0-based) — which original slide to clone for each type
T_COVER = 0
T_AGENDA = 1
T_APPROACH = 2
T_STEP_DIVIDER = 3
T_SECTION_CAPABILITIES = 4
T_CONTENT = 5          # title + bullets + insight + image
T_CONTENT_ALT = 6      # alternate content layout
T_SUMMARY = 13         # title + summary paragraph + half-image
T_SECTION_COMPETITION = 14
T_COMPETITOR_GRID_ALL = 15    # "A MATURE MARKET" — all competitor logos grid
T_COMPETITOR_GRID_FOCUSED = 16  # "FOCUSED REVIEW" — highlighted competitor grid
T_COMPETITOR = 17       # two-column: positioning + key learnings (1 banner image)
T_COMPETITOR_2IMG = 18  # two-column: positioning + learnings (2 side-by-side images)
T_COMPETITOR_4IMG = 19  # two-column: positioning + learnings (4 images in a row)
T_LANDSCAPE = 23        # landscape summary (bullets + sidebar)
T_COMP_SUMMARY = 24
T_SECTION_CONSUMER = 25
T_RESEARCH_APPROACH = 26    # Research methodology (label + detail rows)
T_SEGMENT_DIVIDER = 47      # "Market Segmentation" divider
T_SEGMENT_OVERVIEW = 49     # All segments at a glance (names + % + taglines)
T_MEET_SEGMENT = 51         # "Meet the [Segment]" narrative page
T_TARGET_RECOMMENDATION = 76  # "PRIMARY TARGET: [NAME]" with rationale bullets
T_WHY_TARGET = 77           # "WHY [SEGMENT] IS THE RIGHT FOCUS" with bullets
T_ENABLES = 78              # "WHAT THIS CHOICE ENABLES (AND DOES NOT)"
T_CONSUMER_SUMMARY = 79     # Consumer summary (half-text, half-image)
T_FINAL_SUMMARY = 80        # Three-column summary + closing insight
T_THANK_YOU = 91
T_APPENDIX_TRANSITION = 81  # "Research Appendix" divider
T_APPENDIX_CHART = 34       # Reuse single hbar template for appendix charts
T_APPENDIX_WHITE = 82        # White background slide (Title + Picture + base n)

# Chart slide templates (questionnaire/survey section, slides 27-46)
T_CHART_DIVIDER_DEMO = 27      # "Demographics & Background" divider
T_CHART_SINGLE_HBAR = 34       # Single full-width hbar (e.g., Work Apparel)
T_CHART_DUAL = 35              # Donut left + hbar right (e.g., Purchase Frequency)
T_CHART_SINGLE_VBAR = 30       # Single full-width vbar (e.g., Occupation)
T_CHART_STACKED = 42           # Stacked bar (e.g., Brand Awareness)
T_CHART_DIVIDER_SHOPPING = 33  # "Shopping Habits" divider
T_CHART_DIVIDER_BRAND = 40     # "Brand Evaluation" divider
T_CHART_TABLE = 39             # Challenges table (text-only slide)

# Boilerplate slides
T_SEGMENTATION_INTRO = 48      # "Benefits of segmentation" boilerplate
T_FOCUSING_SEGMENTS = 50       # "FOCUSING ON THE MOST DOMINANT SEGMENTS…"
T_SEGMENT_PROFILE = 52         # Segment respondent profile (demographics layout)
T_CLOSER_LOOK_1 = 53           # "A Closer Look" — premium callout + small icon
T_CLOSER_LOOK_2 = 54           # "A Closer Look" — brand awareness + verbatim quotes
T_CLOSER_LOOK_3 = 56           # "A Closer Look" — 4 lifestyle signal cards
T_CHALLENGES = 55              # Challenges & Pain Points (two tables)
T_SELECTING_TARGET = 75        # "SELECTING [BRAND]'S TARGET AUDIENCE"
T_BRAND_METRICS_DEF = 45       # "Brand Metrics Definitions" boilerplate (GOATClean)

ASSETS_DIR = Path(__file__).parent.parent / "templates" / "assets"


# ── Slide Cloning Engine ─────────────────────────────────────

_src_prs = None


def _reset_caches():
    """Reset all module-level caches between pipeline runs."""
    global _src_prs, _unified_pain_cache, _unified_dimensions_cache
    _src_prs = None
    _challenge_quotes_cache.clear()
    _bubble_quotes_cache.clear()
    _closer_look_3_cache.clear()
    _unified_dimensions_cache = None
    _unified_pain_cache = None
    _brand_image_urls_cache.clear()


def _extract_category(analysis: dict, brand_name: str = "") -> str:
    """Derive a human-readable product category from analysis data.

    Checks multiple sources in priority order to avoid falling back to
    the generic 'consumer product' placeholder.
    """
    import re as _re_cat

    # 1. survey_brand_matrix title often has "[CATEGORY] BRAND ASSOCIATION"
    sbm = analysis.get("survey_brand_matrix", {})
    if isinstance(sbm, dict):
        title = sbm.get("title", "")
        m = _re_cat.match(r'^(.+?)\s+BRAND\s+ASSOCIATION', title, _re_cat.IGNORECASE)
        if m:
            return m.group(1).strip().lower()  # e.g. "water bottle"

    # 2. competition_summary — market overview paragraph
    comp = analysis.get("competition", {})
    if isinstance(comp, dict):
        cs = comp.get("competition_summary", "")
        m = _re_cat.search(
            r'[Tt]he\s+(.+?)\s+(?:market|category|industry|space|landscape|sector)',
            cs[:300],
        )
        if m:
            cat = m.group(1).strip().lower()
            # Remove leading articles
            cat = _re_cat.sub(r'^(the|a|an)\s+', '', cat)
            if 3 < len(cat) < 40:
                return cat

    # 3. capabilities section_title — "… in the X market"
    cap = analysis.get("capabilities", {})
    if isinstance(cap, dict):
        for field in ("section_title", "capabilities_summary"):
            text = cap.get(field, "")
            m = _re_cat.search(
                r'(?:in|of|the)\s+(?:the\s+)?(.+?)\s+(?:market|category|industry|space)',
                text[:300], _re_cat.IGNORECASE,
            )
            if m:
                cat = m.group(1).strip().lower()
                cat = _re_cat.sub(r'^(the|a|an)\s+', '', cat)
                if 3 < len(cat) < 40:
                    return cat

    # 4. Fallback with brand name
    if brand_name:
        return f"products by {brand_name}"
    return "consumer products"


def _get_source():
    """Load the reference PPTX once (cached)."""
    global _src_prs
    if _src_prs is None:
        if not TEMPLATE_PATH.exists():
            raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")
        _src_prs = Presentation(str(TEMPLATE_PATH))
    return _src_prs


def _clone_slide(dst_prs, src_slide_idx):
    """Clone a slide from the reference PPTX into dst_prs.

    Copies all shapes (text boxes, images, groups) and only the
    relationships actually referenced by those shapes. This avoids
    duplicating notesSlide/themeOverride/tags parts that cause ZIP
    corruption.
    """
    src_prs = _get_source()
    src_slide = src_prs.slides[src_slide_idx]

    # Find matching layout in dst by name
    src_layout_name = src_slide.slide_layout.name
    dst_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name == src_layout_name:
            dst_layout = layout
            break
    if not dst_layout:
        dst_layout = dst_prs.slide_layouts[6]  # Blank fallback

    new_slide = dst_prs.slides.add_slide(dst_layout)

    # Clear auto-generated placeholders from layout
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    # Step 1: Deep-copy shape elements and collect all rIds they reference
    copied_elements = []
    referenced_rIds = set()
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        copied_elements.append(el)
        for attr_el in el.iter():
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    referenced_rIds.add(attr_el.attrib[attr_name])
            if "embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib["embed"])
            if f"{ns_r}embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib[f"{ns_r}embed"])

    # Step 2: Copy relationships — for image/hdphoto parts, create a proper
    # copy in the destination package so partnames don't collide with
    # future add_picture calls (which use next_image_partname).
    SKIP_RELTYPES = {"chart", "tags", "tag", "notesSlide", "themeOverride"}
    IMAGE_RELTYPES = {"image", "hdphoto"}
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if rel.rId not in referenced_rIds:
            continue
        reltype_short = rel.reltype.split("/")[-1]
        if reltype_short in SKIP_RELTYPES:
            continue
        try:
            if reltype_short in IMAGE_RELTYPES and not rel.is_external:
                # Copy image blob into the destination package to avoid
                # cross-package Part references that cause partname collisions
                import io
                src_part = rel._target
                if reltype_short == "image":
                    img_part, new_rId = new_slide.part.get_or_add_image_part(
                        io.BytesIO(src_part.blob)
                    )
                else:
                    # hdphoto — just reference the source directly
                    # (hdphotos don't collide with image indices)
                    new_rId = new_slide.part.rels._add_relationship(
                        rel.reltype, rel._target, rel.is_external
                    )
            else:
                new_rId = new_slide.part.rels._add_relationship(
                    rel.reltype, rel._target, rel.is_external
                )
            rId_map[rel.rId] = new_rId
        except Exception:
            pass

    # Step 3: Remap rIds and remove elements with dangling references
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    skip_shapes = []
    for el in copied_elements:
        dangling_elements = []
        for attr_el in el.iter():
            has_dangling = False
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    old_id = attr_el.attrib[attr_name]
                    if old_id in rId_map:
                        attr_el.attrib[attr_name] = rId_map[old_id]
                    elif old_id in referenced_rIds:
                        has_dangling = True
            if "embed" in attr_el.attrib:
                old_id = attr_el.attrib["embed"]
                if old_id in rId_map:
                    attr_el.attrib["embed"] = rId_map[old_id]
                elif old_id in referenced_rIds:
                    has_dangling = True
            if f"{ns_r}embed" in attr_el.attrib:
                old_id = attr_el.attrib[f"{ns_r}embed"]
                if old_id in rId_map:
                    attr_el.attrib[f"{ns_r}embed"] = rId_map[old_id]
                elif old_id in referenced_rIds:
                    has_dangling = True
            if has_dangling:
                dangling_elements.append(attr_el)

        # Remove elements with dangling rIds (chart embeds, tag refs, etc.)
        for dang in dangling_elements:
            parent = dang.getparent()
            if parent is not None:
                parent.remove(dang)

        # Drop the whole top-level shape if it is a graphicFrame whose
        # graphicData no longer has any chart/table/etc. content
        el_tag = el.tag.split("}")[-1]
        if el_tag == "graphicFrame":
            graphic_data = el.find(f".//{ns_a}graphicData")
            if graphic_data is not None and len(graphic_data) == 0:
                skip_shapes.append(el)
                continue

        # Remove custDataLst — contains <p:tags> refs to skipped tag rels
        for cust in el.findall(f".//{ns_p}custDataLst"):
            parent = cust.getparent()
            if parent is not None:
                parent.remove(cust)

        # Clean up empty extLst containers
        for empty_container_tag in (f"{ns_p}extLst", f"{ns_a}extLst"):
            for container in el.findall(f".//{empty_container_tag}"):
                if len(container) == 0:
                    parent = container.getparent()
                    if parent is not None:
                        parent.remove(container)

        new_slide.shapes._spTree.append(el)

    return new_slide


def _clone_slide_with_charts(dst_prs, src_slide_idx):
    """Clone a slide from the reference PPTX, preserving native CHART objects.

    Unlike _clone_slide which skips chart relationships, this function copies
    chart XML parts, their embedded Excel workbooks, and chart style/color
    parts into the destination package. The resulting charts can be updated
    via chart.replace_data(CategoryChartData).
    """
    src_prs = _get_source()
    src_slide = src_prs.slides[src_slide_idx]

    # Find matching layout
    src_layout_name = src_slide.slide_layout.name
    dst_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name == src_layout_name:
            dst_layout = layout
            break
    if not dst_layout:
        dst_layout = dst_prs.slide_layouts[6]

    new_slide = dst_prs.slides.add_slide(dst_layout)
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    # Step 1: Deep-copy shape elements and collect referenced rIds
    copied_elements = []
    referenced_rIds = set()
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        copied_elements.append(el)
        for attr_el in el.iter():
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    referenced_rIds.add(attr_el.attrib[attr_name])
            if "embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib["embed"])
            if f"{ns_r}embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib[f"{ns_r}embed"])

    # Step 2: Copy relationships including charts
    SKIP_RELTYPES = {"tags", "tag", "notesSlide"}
    IMAGE_RELTYPES = {"image", "hdphoto"}
    rId_map = {}

    for rel in src_slide.part.rels.values():
        if rel.rId not in referenced_rIds:
            continue
        reltype_short = rel.reltype.split("/")[-1]
        if reltype_short in SKIP_RELTYPES:
            continue
        try:
            if reltype_short == "chart" and not rel.is_external:
                # Clone chart part with all sub-parts
                src_chart_part = rel._target
                new_partname = dst_prs.part.package.next_partname(
                    ChartPart.partname_template
                )
                new_chart_part = ChartPart.load(
                    new_partname,
                    src_chart_part.content_type,
                    dst_prs.part.package,
                    src_chart_part.blob,
                )
                chart_num = (
                    str(new_partname).split("chart")[-1].replace(".xml", "")
                )
                # Copy chart sub-relationships (xlsx, styles, colors, theme)
                # Track rId mapping so we can fix externalData references
                chart_rId_map = {}
                for crel in src_chart_part.rels.values():
                    crt = crel.reltype.split("/")[-1]
                    src_sub = crel._target
                    if crt == "package":
                        new_xlsx = EmbeddedXlsxPart.new(
                            src_sub.blob, dst_prs.part.package
                        )
                        new_crid = new_chart_part.rels._add_relationship(
                            crel.reltype, new_xlsx, False
                        )
                        chart_rId_map[crel.rId] = new_crid
                    elif crt in ("chartStyle", "chartColorStyle", "themeOverride"):
                        if crt == "chartStyle":
                            sub_pn = PackURI(f"/ppt/charts/style{chart_num}.xml")
                        elif crt == "chartColorStyle":
                            sub_pn = PackURI(f"/ppt/charts/colors{chart_num}.xml")
                        else:
                            sub_pn = PackURI(
                                f"/ppt/theme/themeOverride_c{chart_num}.xml"
                            )
                        new_sub = OpcPart(
                            sub_pn,
                            src_sub.content_type,
                            dst_prs.part.package,
                            src_sub.blob,
                        )
                        new_crid = new_chart_part.rels._add_relationship(
                            crel.reltype, new_sub, False
                        )
                        chart_rId_map[crel.rId] = new_crid
                    elif crt == "chartUserShapes":
                        # Copy drawing parts referenced by charts
                        draw_pn = dst_prs.part.package.next_partname(
                            "/ppt/drawings/drawing%d.xml"
                        )
                        draw_copy = OpcPart(
                            draw_pn,
                            src_sub.content_type,
                            dst_prs.part.package,
                            src_sub.blob,
                        )
                        new_crid = new_chart_part.rels._add_relationship(
                            crel.reltype, draw_copy, False
                        )
                        chart_rId_map[crel.rId] = new_crid

                # Remap rIds inside chart XML element tree (externalData, etc.)
                # ChartPart.blob is a property that serializes from _element,
                # so we must modify the XML element tree directly.
                if chart_rId_map:
                    ns_r_full = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                    for el in new_chart_part._element.iter():
                        for attr_name in list(el.attrib):
                            if attr_name == f"{{{ns_r_full}}}id" or attr_name == "r:id":
                                old_val = el.attrib[attr_name]
                                if old_val in chart_rId_map:
                                    el.attrib[attr_name] = chart_rId_map[old_val]

                new_rId = new_slide.part.rels._add_relationship(
                    rel.reltype, new_chart_part, False
                )
                rId_map[rel.rId] = new_rId
            elif reltype_short in IMAGE_RELTYPES and not rel.is_external:
                src_part = rel._target
                if src_part.content_type == "image/svg+xml":
                    # SVG: create a fresh Part copy to avoid partname collisions
                    svg_pn = PackURI(
                        dst_prs.part.package.next_partname(
                            "/ppt/media/image%d.svg"
                        )
                    )
                    svg_copy = OpcPart(
                        svg_pn,
                        src_part.content_type,
                        dst_prs.part.package,
                        src_part.blob,
                    )
                    new_rId = new_slide.part.rels._add_relationship(
                        rel.reltype, svg_copy, False
                    )
                elif reltype_short == "image":
                    try:
                        img_part, new_rId = (
                            new_slide.part.get_or_add_image_part(
                                io.BytesIO(src_part.blob)
                            )
                        )
                    except Exception:
                        new_rId = new_slide.part.rels._add_relationship(
                            rel.reltype, src_part, False
                        )
                else:
                    new_rId = new_slide.part.rels._add_relationship(
                        rel.reltype, src_part, False
                    )
                rId_map[rel.rId] = new_rId
            else:
                new_rId = new_slide.part.rels._add_relationship(
                    rel.reltype, rel._target, rel.is_external
                )
                rId_map[rel.rId] = new_rId
        except Exception:
            pass

    # Step 3: Remap rIds in copied elements
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    for el in copied_elements:
        for attr_el in el.iter():
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    old_id = attr_el.attrib[attr_name]
                    if old_id in rId_map:
                        attr_el.attrib[attr_name] = rId_map[old_id]
            if "embed" in attr_el.attrib:
                old_id = attr_el.attrib["embed"]
                if old_id in rId_map:
                    attr_el.attrib["embed"] = rId_map[old_id]
            if f"{ns_r}embed" in attr_el.attrib:
                old_id = attr_el.attrib[f"{ns_r}embed"]
                if old_id in rId_map:
                    attr_el.attrib[f"{ns_r}embed"] = rId_map[old_id]

        # Remove custDataLst entirely — it contains <p:tags> references
        # that point to skipped tag relationships, causing PowerPoint repair
        for cust in el.findall(f".//{ns_p}custDataLst"):
            parent = cust.getparent()
            if parent is not None:
                parent.remove(cust)

        # Clean up empty containers
        for empty_tag in (
            f"{ns_p}extLst",
            f"{ns_a}extLst",
        ):
            for container in el.findall(f".//{empty_tag}"):
                if len(container) == 0:
                    parent = container.getparent()
                    if parent is not None:
                        parent.remove(container)

        new_slide.shapes._spTree.append(el)

    return new_slide


def _set_chart_pct_format(chart):
    """Set chart data labels to display values as percentages (29% not 0.29).

    Forces ALL numFmt elements to use "0%" with sourceLinked="0",
    and updates all formatCode elements in data caches to "0%".
    This ensures the chart never falls back to Excel's "General" format.
    """
    from lxml import etree
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    # Fix ALL numFmt elements across the entire chart (dLbls, dLbl, axes, etc.)
    for numFmt in chart._chartSpace.iter(f"{{{ns_c}}}numFmt"):
        fmt = numFmt.get("formatCode", "")
        # Force percentage format and disconnect from Excel source
        if "%" in fmt or fmt == "General" or numFmt.get("sourceLinked") == "1":
            numFmt.set("formatCode", "0%")
            numFmt.set("sourceLinked", "0")

    # Fix formatCode inside numRef (data cache) — Excel's "General" → "0%"
    for fc in chart._chartSpace.iter(f"{{{ns_c}}}formatCode"):
        if fc.text and fc.text.strip() in ("General", "0.00"):
            fc.text = "0%"


def _set_pie_orange_colors(chart):
    """Set pie chart slice colors to an orange gradient scheme for brand consistency."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    # Orange gradient: dark → light
    orange_colors = ["D94F00", "FB641F", "FF8C42", "FFB380", "FFD9BF"]

    for pie_el in chart._chartSpace.iter(f"{{{ns_c}}}pieChart"):
        ser_el = pie_el.find(f"{{{ns_c}}}ser")
        if ser_el is None:
            continue
        # Add data point fills
        for i, color in enumerate(orange_colors):
            # Find or create dPt
            dPt = None
            for existing in ser_el.findall(f"{{{ns_c}}}dPt"):
                idx_el = existing.find(f"{{{ns_c}}}idx")
                if idx_el is not None and idx_el.get("val") == str(i):
                    dPt = existing
                    break
            if dPt is None:
                dPt = etree.SubElement(ser_el, f"{{{ns_c}}}dPt")
                idx_el = etree.SubElement(dPt, f"{{{ns_c}}}idx")
                idx_el.set("val", str(i))

            # Set solid fill
            spPr = dPt.find(f"{{{ns_c}}}spPr")
            if spPr is None:
                spPr = etree.SubElement(dPt, f"{{{ns_c}}}spPr")
            # Clear existing fills
            for child in list(spPr):
                spPr.remove(child)
            solidFill = etree.SubElement(spPr, f"{{{ns_a}}}solidFill")
            srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
            srgbClr.set("val", color)


def _find_text_shapes(slide):
    """Return all shapes with text frames, sorted by top then left position."""
    shapes = [s for s in slide.shapes if s.has_text_frame]
    shapes.sort(key=lambda s: (s.top, s.left))
    return shapes


def _set_text_preserve_format(text_frame, new_text):
    """Replace text in a text_frame while preserving per-run and per-paragraph formatting.

    Accepts:
      - str: replaces text line-by-line matching to existing paragraphs/runs
      - list[str]: one string per paragraph, matched 1:1 to existing paragraphs

    Key principle: each original run keeps its font/size/color; we only change .text.
    """
    if isinstance(new_text, list):
        paragraphs = list(text_frame.paragraphs)
        for i, para_text in enumerate(new_text):
            if i < len(paragraphs):
                _replace_para_text(paragraphs[i], para_text)
            else:
                _add_paragraph_after(text_frame, paragraphs[-1], para_text)
        # Remove excess paragraphs
        for i in range(len(new_text), len(paragraphs)):
            p_el = paragraphs[i]._p
            p_el.getparent().remove(p_el)
    else:
        # Split by newlines and match to existing paragraphs/runs
        lines = new_text.split("\n")
        paragraphs = list(text_frame.paragraphs)

        if len(lines) == 1:
            # Single line — distribute across existing runs in first para
            if paragraphs:
                _replace_para_text(paragraphs[0], lines[0])
                for para in paragraphs[1:]:
                    p_el = para._p
                    p_el.getparent().remove(p_el)
        else:
            # Multiple lines — match each line to an existing paragraph
            # If paragraph has multiple runs (different fonts), map lines to runs
            if len(paragraphs) == 1 and len(paragraphs[0].runs) >= len(lines):
                # Single paragraph with multiple runs — map lines to runs
                _replace_runs_text(paragraphs[0], lines)
            else:
                for i, line in enumerate(lines):
                    if i < len(paragraphs):
                        _replace_para_text(paragraphs[i], line)
                    else:
                        _add_paragraph_after(text_frame, paragraphs[-1], line)
                for i in range(len(lines), len(paragraphs)):
                    p_el = paragraphs[i]._p
                    p_el.getparent().remove(p_el)


def _has_cjk(text: str) -> bool:
    """Check if text contains CJK characters."""
    return any('\u4e00' <= c <= '\u9fff' or '\u3400' <= c <= '\u4dbf' for c in text)


CJK_FONT = "Heiti SC"


def _fix_cjk_fonts(prs):
    """Scan all slides and set CJK-compatible font on any run containing CJK text."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _has_cjk(run.text):
                        _set_cjk_font(run)


def _set_cjk_font(run):
    """Set East Asian font on a run so CJK characters render correctly."""
    from lxml import etree
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    if rPr is None:
        rPr = etree.SubElement(
            run._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        run._r.insert(0, rPr)
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        ea = etree.SubElement(
            rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', CJK_FONT)


def _replace_para_text(paragraph, text):
    """Replace text in a paragraph, distributing across existing runs.

    Preserves each run's formatting (font, size, color, bold).
    Sets CJK-compatible font when text contains Chinese characters.
    """
    runs = paragraph.runs
    if not runs:
        paragraph.text = text
        return

    if len(runs) == 1:
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
    else:
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
        for run in runs[1:]:
            run.text = ""


def _replace_cell_text(cell, text):
    """Replace ALL text in a table cell, clearing extra paragraphs.

    Unlike _replace_para_text which only touches paragraph 0, this clears
    every paragraph beyond the first (which handles multi-line template cells
    like 'Wonder\\nWink').
    """
    paras = cell.text_frame.paragraphs
    if paras:
        _replace_para_text(paras[0], text)
        # Clear any additional paragraphs (e.g., "Wonder" + "Wink" → just "None")
        for p in paras[1:]:
            if p.runs:
                for r in p.runs:
                    r.text = ""
            else:
                p.text = ""


def _replace_runs_text(paragraph, texts):
    """Replace text run-by-run, one text per run, preserving each run's formatting."""
    runs = paragraph.runs
    for i, text in enumerate(texts):
        if i < len(runs):
            runs[i].text = text
    # Clear excess runs
    for i in range(len(texts), len(runs)):
        runs[i].text = ""


def _add_paragraph_after(text_frame, template_para, text):
    """Add a new paragraph after template_para, copying its formatting."""
    new_p = copy.deepcopy(template_para._p)
    template_para._p.addnext(new_p)
    from pptx.text.text import _Paragraph
    para = _Paragraph(new_p, template_para._parent)
    _replace_para_text(para, text)


def _set_bold_colon_text(text_frame, text):
    """Set text with bold-before-colon formatting: 'Category: detail' → bold 'Category:' + regular ' detail'.

    Preserves the first run's font properties (size, color) and applies bold to the prefix.
    """
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = text
        return

    para = paragraphs[0]
    # Remove excess paragraphs
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)

    runs = para.runs
    if not runs:
        para.text = text
        return

    # Split at first colon
    colon_idx = text.find(":")
    if colon_idx > 0 and colon_idx < 60:
        bold_part = text[:colon_idx]
        regular_part = text[colon_idx:]  # includes the colon

        # Use first run for bold/semibold part
        runs[0].text = bold_part
        runs[0].font.name = "Montserrat SemiBold"
        runs[0].font.bold = None  # inherit
        if _has_cjk(bold_part):
            _set_cjk_font(runs[0])

        # Use or create second run for regular part
        if len(runs) >= 2:
            runs[1].text = regular_part
            runs[1].font.name = "Montserrat"
            runs[1].font.bold = None
            if _has_cjk(regular_part):
                _set_cjk_font(runs[1])
            for r in runs[2:]:
                r.text = ""
        else:
            new_r = copy.deepcopy(runs[0]._r)
            para._p.append(new_r)
            from pptx.text.text import _Run
            new_run = _Run(new_r, para)
            new_run.text = regular_part
            new_run.font.name = "Montserrat"
            new_run.font.bold = None
    else:
        # No colon — just set plain text
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
        for r in runs[1:]:
            r.text = ""


def _set_bold_pct_text(text_frame, text):
    """Set text with bold percentage prefix: '78% prefer X' → bold '78%' + regular ' prefer X'.

    If no percentage found, falls back to _set_bold_colon_text behavior.
    """
    import re as _re_pct
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = text
        return

    para = paragraphs[0]
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)

    runs = para.runs
    if not runs:
        para.text = text
        return

    # Find percentage pattern at or near start
    m = _re_pct.match(r'^(\d+%)\s*(.*)', text)
    if m:
        bold_part = m.group(1)
        regular_part = " " + m.group(2)
    elif ": " in text:
        # Fallback to colon split
        idx = text.index(": ")
        bold_part = text[:idx + 1]
        regular_part = text[idx + 1:]
    else:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
        return

    runs[0].text = bold_part
    runs[0].font.name = "Montserrat SemiBold"
    runs[0].font.bold = True

    if len(runs) >= 2:
        runs[1].text = regular_part
        runs[1].font.name = "Montserrat"
        runs[1].font.bold = False
        for r in runs[2:]:
            r.text = ""
    else:
        new_r = copy.deepcopy(runs[0]._r)
        para._p.append(new_r)
        from pptx.text.text import _Run
        new_run = _Run(new_r, para)
        new_run.text = regular_part
        new_run.font.name = "Montserrat"
        new_run.font.bold = False


def _set_bold_colon_text_para(paragraph, text):
    """Set a paragraph's text with bold label before colon, regular after.

    Uses both font.bold=True AND Montserrat SemiBold typeface for the label,
    so bold renders correctly even if Montserrat SemiBold isn't installed.
    """
    runs = paragraph.runs
    if not runs:
        paragraph.text = text
        return

    colon_idx = text.find(":")
    if colon_idx > 0 and colon_idx < 60:
        label_part = text[:colon_idx]
        rest_part = text[colon_idx:]  # includes ": detail"

        runs[0].text = label_part
        runs[0].font.name = "Montserrat SemiBold"
        runs[0].font.bold = True
        if _has_cjk(label_part):
            _set_cjk_font(runs[0])

        if len(runs) >= 2:
            runs[1].text = rest_part
            runs[1].font.name = "Montserrat"
            runs[1].font.bold = False
            if _has_cjk(rest_part):
                _set_cjk_font(runs[1])
            for r in runs[2:]:
                r.text = ""
        else:
            new_r = copy.deepcopy(runs[0]._r)
            paragraph._p.append(new_r)
            from pptx.text.text import _Run
            new_run = _Run(new_r, paragraph)
            new_run.text = rest_part
            new_run.font.name = "Montserrat"
            new_run.font.bold = False
    else:
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
        for r in runs[1:]:
            r.text = ""


# ── Shape Border Removal ─────────────────────────────────────

def _clear_shape_border(shape):
    """Remove any visible border/outline from a shape.

    IMPORTANT: Never access shape.line — python-pptx creates internal
    state that overwrites our XML during save (converts noFill → solidFill).
    Pure lxml manipulation only.
    """
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    # Find spPr (shape properties container)
    sp_pr = shape._element.find(f"{{{ns_p}}}spPr")
    if sp_pr is None:
        sp_pr = shape._element.find(f".//{{{ns_a}}}prstGeom")
        if sp_pr is not None:
            sp_pr = sp_pr.getparent()
    if sp_pr is None:
        return
    # Remove all existing ln elements
    for ln_old in sp_pr.findall(f"{{{ns_a}}}ln"):
        sp_pr.remove(ln_old)
    # Insert clean ln with noFill after the last geometry/fill element
    ln_new = etree.SubElement(sp_pr, f"{{{ns_a}}}ln", w="0")
    etree.SubElement(ln_new, f"{{{ns_a}}}noFill")


# ── Text Truncation ──────────────────────────────────────────

def _llm_compress(text: str, max_chars: int, context: str = "") -> str:
    """Use LLM to intelligently compress text while preserving meaning.

    Falls back to _truncate() if LLM call fails or text is short enough.
    Only invoked for text that significantly exceeds max_chars (>1.5x).
    """
    if not text or len(text) <= max_chars:
        return text
    # Only use LLM for significant overflows (saves API calls)
    if len(text) < max_chars * 1.5:
        return _truncate(text, max_chars)
    try:
        from anthropic import Anthropic
        from config import ANTHROPIC_API_KEY
        client = Anthropic(api_key=ANTHROPIC_API_KEY)
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=max(64, max_chars // 2),
            messages=[{"role": "user", "content": (
                f"Compress this text to under {max_chars} characters while preserving "
                f"the key strategic insight. Keep it punchy and authoritative — "
                f"this appears on a brand strategy presentation slide.\n\n"
                f"Original ({len(text)} chars): {text}\n\n"
                f"Return ONLY the compressed text, nothing else."
            )}],
        )
        result = resp.content[0].text.strip()
        if len(result) <= max_chars and len(result) > 10:
            return result
    except Exception:
        pass
    return _truncate(text, max_chars)


def _truncate(text, max_chars):
    """Truncate text to max_chars, preferring sentence boundaries.

    Strategy: find the last complete sentence within the limit. If no sentence
    boundary exists (common with single long sentences), find the last major
    clause boundary (comma, semicolon, em-dash) and close with a period.
    Avoid leaving dangling conjunctions/prepositions at the end.
    """
    if not text or len(text) <= max_chars:
        return text
    search_zone = text[:max_chars]

    # 1. Prefer cutting at last complete sentence (period + space or end)
    best_period = -1
    for i in range(len(search_zone) - 1, -1, -1):
        if search_zone[i] == "." and (i + 1 >= len(search_zone) or search_zone[i + 1] == " "):
            word_start = search_zone[:i].rfind(" ")
            last_word = search_zone[word_start + 1:i]
            if len(last_word) >= 3:  # real word, not a stub like "a."
                best_period = i
                break
    if best_period > max_chars * 0.35:
        return text[:best_period + 1]

    # 2. Fall back to last clause boundary (comma, semicolon, em-dash)
    for sep in (",", ";", " —", " –"):
        pos = search_zone.rfind(sep)
        if pos > max_chars * 0.35:
            fragment = text[:pos].rstrip()
            _strip_dangling = {"and", "or", "but", "with", "for", "in", "on", "at", "to",
                               "the", "a", "an", "of", "by", "as", "is", "are", "was", "its",
                               "that", "which", "who", "from", "into", "than", "not", "does",
                               "can", "could", "will", "would", "may", "might", "should"}
            words = fragment.rsplit(" ", 1)
            while len(words) == 2 and words[1].lower().rstrip(".,;:") in _strip_dangling:
                fragment = words[0].rstrip()
                words = fragment.rsplit(" ", 1)
            return fragment.rstrip(" ,;:") + "."

    # 3. Last resort: word boundary, stripping dangling words
    cut = search_zone.rfind(" ")
    if cut < max_chars // 2:
        cut = max_chars
    result = text[:cut].rstrip(" ,;:")
    # Strip trailing words that leave the sentence feeling incomplete
    _dangling = {"and", "or", "but", "with", "for", "in", "on", "at", "to",
                 "the", "a", "an", "of", "by", "as", "is", "are", "was", "its",
                 "that", "which", "who", "from", "into", "than", "not", "does",
                 "more", "most", "very", "also", "where", "when", "while", "each",
                 "every", "their", "these", "those", "this", "such", "much",
                 "can", "could", "will", "would", "may", "might", "should",
                 "has", "have", "had", "been", "being", "both", "either", "neither"}
    # Also strip trailing words that leave incomplete phrases:
    # - adjectives after prepositions/conjunctions (e.g., "with environmental")
    # - adjectives after verbs like "uses", "creates", "leaves" (e.g., "leaves significant")
    _adj_suffixes = ("al", "ive", "ous", "ful", "ent", "ant", "ary", "ory",
                     "ible", "able", "ical", "ated", "ised", "ized")
    _incomplete_prev = {"with", "and", "of", "to", "for", "on", "but", "or",
                        "uses", "creates", "leaves", "makes", "feels", "gives",
                        "provides", "offers", "shows", "builds", "drives"}
    words = result.rsplit(" ", 1)
    while len(words) == 2:
        w = words[1].lower().rstrip(".,;:")
        if w in _dangling:
            result = words[0].rstrip()
            words = result.rsplit(" ", 1)
            continue
        # Check if current word looks like a dangling adjective/participle
        if len(words[0]) > 0:
            prev_word = words[0].rsplit(" ", 1)[-1].lower().rstrip(".,;:")
            if prev_word in _incomplete_prev and (
                w.endswith(_adj_suffixes) or w.endswith(("ed", "ing"))
            ):
                result = words[0].rsplit(" ", 1)[0].rstrip() if " " in words[0] else words[0]
                words = result.rsplit(" ", 1)
                continue
        break
    result = result.rstrip(".,;: ")
    return result + "."


# ── Image Replacement ───────────────────────────────────────

def _replace_slide_image(slide, image_path: Path, replace_background=False):
    """Replace a picture shape on a slide with a new image.

    By default replaces the first non-background picture. With
    replace_background=True, replaces the largest (background) picture
    instead — used for "Meet the Segment" hero slides.

    Pre-crops the image file with PIL to exactly match the box aspect
    ratio (cover + top-bias), then inserts the cropped image at the
    exact box dimensions.
    """
    if not image_path or not Path(image_path).exists():
        return

    from PIL import Image
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000  # 16:9 widescreen

    # Collect all picture shapes first to avoid mutation-during-iteration
    pictures = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            box_w, box_h = shape.width, shape.height
            is_bg = (box_w * box_h) / SLIDE_AREA > 0.9
            if replace_background and not is_bg:
                continue
            if not replace_background and is_bg:
                continue
            pictures.append(shape)

    if not pictures:
        return

    shape = pictures[0]
    box_left, box_top = shape.left, shape.top
    box_w, box_h = shape.width, shape.height
    if box_h == 0:
        return
    box_ratio = box_w / box_h

    try:
        img = Image.open(str(image_path))
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        img_w, img_h = img.size
    except Exception:
        return

    img_ratio = img_w / img_h

    # Decide: cover crop (small ratio diff) vs contain+pad (large diff or text images)
    ratio_diff = abs(img_ratio - box_ratio)
    use_contain = ratio_diff > 0.3  # Large aspect ratio mismatch → contain mode

    if use_contain:
        # Contain mode: fit entire image inside box, pad with blurred background
        canvas = Image.new("RGB", (img_w, img_h), (245, 245, 245))
        if img.mode == "RGBA":
            canvas.paste(img, mask=img)
        else:
            canvas = img.copy()

        # Create blurred background at target aspect ratio
        if img_ratio > box_ratio:
            # Image wider than box → letterbox (pad top/bottom)
            fit_w = img_w
            fit_h = int(img_w / box_ratio)
        else:
            # Image taller than box → pillarbox (pad left/right)
            fit_h = img_h
            fit_w = int(img_h * box_ratio)

        bg = canvas.resize((fit_w, fit_h), Image.LANCZOS)
        from PIL import ImageFilter
        bg = bg.filter(ImageFilter.GaussianBlur(radius=30))
        # Darken the background
        from PIL import ImageEnhance
        bg = ImageEnhance.Brightness(bg).enhance(0.4)

        # Paste sharp original centered
        paste_x = (fit_w - img_w) // 2
        paste_y = (fit_h - img_h) // 2
        if img.mode == "RGBA":
            bg.paste(img, (paste_x, paste_y), mask=img)
        else:
            bg.paste(img, (paste_x, paste_y))
        img = bg
    elif ratio_diff > 0.05:
        # Cover crop (small ratio diff) — original behavior
        if img_ratio > box_ratio:
            new_w = int(img_h * box_ratio)
            offset = (img_w - new_w) // 2
            img = img.crop((offset, 0, offset + new_w, img_h))
        else:
            new_h = int(img_w / box_ratio)
            top_offset = int((img_h - new_h) * 0.15)
            img = img.crop((0, top_offset, img_w, top_offset + new_h))

    cropped_path = image_path.parent / f"_cropped_{image_path.stem}.png"
    img.save(str(cropped_path), format="PNG")
    img.close()

    # Get a properly-registered ImagePart for the cropped image
    image_part, rId = slide.part.get_or_add_image_part(str(cropped_path))

    # Update the existing Picture's blip to reference the new image,
    # rather than remove+add which risks partname collisions with
    # source Parts from _clone_slide that share the same package
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ns_a_uri = "http://schemas.openxmlformats.org/drawingml/2006/main"
    blip = shape._element.find(f".//{{{ns_a_uri}}}blip")
    if blip is not None:
        blip.set(f"{ns_r}embed", rId)
        # Remove artistic effect layers that reference old image parts
        for child in list(blip):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "imgLayer":
                blip.remove(child)
        return

    # Fallback: remove and recreate
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(cropped_path), box_left, box_top, box_w, box_h)
    return


def _replace_group_icon_with_ai(slide, img_shape, callout_text: str, segment_name: str):
    """Generate an AI icon matching callout content and replace the group's icon image.

    Creates an orange line-art icon on transparent background that semantically
    matches the callout text (e.g., shopping cart for purchase behavior,
    clock for time-related stats, etc.)
    """
    import io as _io
    import re as _re

    # Determine icon concept from callout text
    text_lower = callout_text.lower()
    icon_concepts = {
        "leak": "water droplet with X mark",
        "seal": "water droplet with checkmark",
        "cup.holder": "car cup holder",
        "clean": "sparkling clean surface",
        "mold": "warning triangle with bacteria",
        "hand": "hand pressing button",
        "one.hand": "single hand grip",
        "amazon": "shopping cart with star",
        "review": "star rating review",
        "purchase": "shopping bag",
        "bought": "shopping cart",
        "shop": "shopping basket",
        "premium": "diamond gem",
        "brand": "award ribbon",
        "loyal": "heart with checkmark",
        "switch": "two curved arrows exchanging",
        "price": "price tag",
        "spend": "money bills",
        "income": "money stack",
        "instagram": "camera with heart",
        "social": "speech bubble with hashtag",
        "reddit": "chat forum bubble",
        "hospital": "medical cross",
        "nurse": "stethoscope",
        "work": "briefcase",
        "shift": "clock with arrow",
        "durabil": "shield with checkmark",
        "fabric": "textile weave pattern",
        "comfort": "soft cushion",
        "fit": "measuring tape",
        "size": "ruler with arrows",
        "color": "paint palette",
        "style": "fashion hanger",
        "travel": "suitcase",
        "gym": "dumbbell",
        "outdoor": "mountain and sun",
        "insul": "thermometer with snowflake",
        "temperat": "thermometer",
    }

    concept = "simple abstract icon"
    for pattern, desc in icon_concepts.items():
        if _re.search(pattern, text_lower):
            concept = desc
            break

    # Generate icon via image_gen
    from pipeline.image_gen import generate_image
    import tempfile
    safe_name = _re.sub(r'[^a-z0-9]', '_', segment_name.lower())[:20]
    icon_path = Path(tempfile.mkstemp(suffix=f"_icon_{safe_name}.png")[1])

    prompt = (
        f"Simple flat line-art icon of a {concept}. "
        f"Orange color (#FB641F), minimalist outline style, "
        f"transparent background, single object centered, "
        f"no text, no shadows, clean vector-style illustration. "
        f"Similar to business infographic icons."
    )

    result = generate_image(prompt, output_path=icon_path, backend="gpt-image",
                           size="1024x1024", quality="standard")
    if not result or not result.exists():
        return  # Keep template icon if generation fails

    # Replace the image blip in the group's image shape
    try:
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        img_part, rId = slide.part.get_or_add_image_part(
            _io.BytesIO(result.read_bytes()))
        for blip in img_shape._element.iter(f"{ns_a}blip"):
            blip.set(f"{ns_r}embed", rId)
        print(f"[closer_look_1] Replaced icon for: {concept}")
    except Exception as e:
        print(f"[closer_look_1] Icon replacement failed: {e}")
    finally:
        try:
            icon_path.unlink(missing_ok=True)
        except Exception:
            pass


def _add_standalone_icon(slide, text_shape, callout_text: str, segment_name: str):
    """Generate and ADD a new icon image to the left of a standalone text shape.

    Used for Closer Look 1 positions that have text but no icon in the template
    (e.g., TextBox 34 at right column row 4).
    """
    import io as _io
    import re as _re
    import tempfile

    # Determine icon concept (same logic as _replace_group_icon_with_ai)
    text_lower = callout_text.lower()
    icon_concepts = {
        "leak": "water droplet with X mark",
        "seal": "water droplet with checkmark",
        "cup.holder": "car cup holder",
        "clean": "sparkling clean surface",
        "mold": "warning triangle with bacteria",
        "purchase": "shopping bag",
        "bought": "shopping cart",
        "shop": "shopping basket",
        "premium": "diamond gem",
        "brand": "award ribbon",
        "loyal": "heart with checkmark",
        "switch": "two curved arrows exchanging",
        "price": "price tag",
        "spend": "money bills",
        "income": "money stack",
        "social": "speech bubble with hashtag",
        "female": "person silhouette",
        "male": "person silhouette",
        "age": "calendar with person",
        "millennial": "person with smartphone",
        "gen.z": "person with smartphone",
        "insul": "thermometer with snowflake",
        "color": "paint palette",
        "style": "fashion hanger",
        "durabil": "shield with checkmark",
    }
    concept = "simple abstract icon"
    for pattern, desc in icon_concepts.items():
        if _re.search(pattern, text_lower):
            concept = desc
            break

    from pipeline.image_gen import generate_image
    safe_name = _re.sub(r'[^a-z0-9]', '_', segment_name.lower())[:20]
    icon_path = Path(tempfile.mkstemp(suffix=f"_icon_{safe_name}.png")[1])

    prompt = (
        f"Simple flat line-art icon of a {concept}. "
        f"Orange color (#FB641F), minimalist outline style, "
        f"transparent background, single object centered, "
        f"no text, no shadows, clean vector-style illustration. "
        f"Similar to business infographic icons."
    )

    result = generate_image(prompt, output_path=icon_path, backend="gpt-image",
                            size="1024x1024", quality="standard")
    if not result or not result.exists():
        print(f"[closer_look_1] Failed to generate standalone icon for: {concept}")
        return

    try:
        from pptx.util import Emu
        # Template group PICs are all 1.20" x 1.20" at x=7.65" on slide
        # (verified: Group 44/29/7 PICs all map to slide x=7.65, size 1.20x1.20)
        # TextBox 34 is at (9.40, 5.57, 3.52x0.91) — icon goes to its left
        icon_w = Emu(int(1.20 * 914400))  # 1.20 inches (exact match template group icons)
        icon_h = Emu(int(1.20 * 914400))
        icon_left = Emu(int(7.65 * 914400))  # exact same x as all group PICs
        # Vertically center icon with text shape
        icon_top = text_shape.top + (text_shape.height - icon_h) // 2

        pic = slide.shapes.add_picture(
            str(result), icon_left, icon_top, icon_w, icon_h
        )
        print(f"[closer_look_1] Added standalone icon for: {concept}")
    except Exception as e:
        print(f"[closer_look_1] Failed to add standalone icon: {e}")
    finally:
        try:
            icon_path.unlink(missing_ok=True)
        except Exception:
            pass


def _replace_card_images(slide, img_pool):
    """Replace multiple card-sized images on a slide (e.g., Closer Look 3).

    Finds all Picture shapes between 5% and 20% of slide area (the 4 lifestyle
    cards) and replaces each with a different brand image.
    """
    from PIL import Image
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000

    cards = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            area_pct = (shape.width * shape.height) / SLIDE_AREA
            if 0.05 < area_pct < 0.20:
                cards.append(shape)

    for shape in cards:
        image_path = img_pool.next_brand()
        if not image_path or not Path(image_path).exists():
            continue

        box_w, box_h = shape.width, shape.height
        if box_h == 0:
            continue
        box_ratio = box_w / box_h

        try:
            img = Image.open(str(image_path))
            if img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGBA")
            else:
                img = img.convert("RGB")
            img_w, img_h = img.size
        except Exception:
            continue

        img_ratio = img_w / img_h
        if abs(img_ratio - box_ratio) > 0.05:
            if img_ratio > box_ratio:
                new_w = int(img_h * box_ratio)
                offset = (img_w - new_w) // 2
                img = img.crop((offset, 0, offset + new_w, img_h))
            else:
                new_h = int(img_w / box_ratio)
                top_offset = int((img_h - new_h) * 0.15)
                img = img.crop((0, top_offset, img_w, top_offset + new_h))

        cropped_path = image_path.parent / f"_cropped_{image_path.stem}.png"
        img.save(str(cropped_path), format="PNG")
        img.close()

        image_part, rId = slide.part.get_or_add_image_part(str(cropped_path))
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        ns_a_uri = "http://schemas.openxmlformats.org/drawingml/2006/main"
        blip = shape._element.find(f".//{{{ns_a_uri}}}blip")
        if blip is not None:
            blip.set(f"{ns_r}embed", rId)
            for child in list(blip):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == "imgLayer":
                    blip.remove(child)


async def _pick_best_hero_with_vision(
    client, brand_name: str, candidates: list[dict], target_ratio: float
) -> Path | None:
    """Use Claude Vision to select the best hero product image from candidates.

    Sends all candidate images to Claude and asks it to pick the best one
    for a brand discovery PPT hero slide.
    """
    import base64

    # Build vision content: one image per candidate
    content = [
        {"type": "text", "text": (
            f"I'm creating a brand discovery presentation for {brand_name}. "
            f"I need to pick the BEST hero product image for the agenda slide. "
            f"The image will fill a portrait-oriented half-page slot (width/height ratio ~{target_ratio:.2f}). "
            f"\nBelow are {len(candidates)} candidate product images. "
            f"Pick the ONE best image based on these criteria:\n"
            f"1. Shows a SINGLE flagship product prominently (not bundles/groups)\n"
            f"2. Clean studio-style shot with white/minimal background\n"
            f"3. Product is centered and fills the frame well\n"
            f"4. High visual quality and detail\n"
            f"5. Portrait or square orientation preferred (will be cropped to portrait)\n"
            f"6. Iconic/recognizable product design that represents the brand\n"
            f"\nReply with ONLY the number (1-{len(candidates)}) of the best image."
        )}
    ]
    for i, c in enumerate(candidates):
        # Always convert to PNG to avoid media type mismatch
        from PIL import Image as PILImage
        import io
        img_obj = PILImage.open(c["path"]).convert("RGB")
        buf = io.BytesIO()
        img_obj.save(buf, format="PNG")
        img_bytes = buf.getvalue()
        img_obj.close()
        b64 = base64.b64encode(img_bytes).decode()
        content.append({"type": "text", "text": f"\nImage {i+1} ({c['w']}x{c['h']}, {c['file_size']//1024}KB):"})
        content.append({"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}})

    try:
        response = client.messages.create(
            model=MODEL_OPUS,
            max_tokens=64,
            messages=[{"role": "user", "content": content}],
        )
        answer = response.content[0].text.strip()
        # Extract the number
        import re
        match = re.search(r'(\d+)', answer)
        if match:
            idx = int(match.group(1)) - 1
            if 0 <= idx < len(candidates):
                chosen = candidates[idx]
                print(f"[hero_image] Claude Vision selected image {idx+1}: {chosen['url'].split('/')[-1][:50]} ({chosen['w']}x{chosen['h']})")
                # Clean up non-selected candidates
                for j, c in enumerate(candidates):
                    if j != idx:
                        c["path"].unlink(missing_ok=True)
                return chosen["path"]
        # Fallback: pick highest img_score
        print(f"[hero_image] Vision response unclear ('{answer}'), falling back to score-based selection")
    except Exception as e:
        print(f"[hero_image] Vision selection failed: {e}, falling back to score-based")

    # Fallback: use the highest scored candidate
    best = max(candidates, key=lambda x: x["img_score"])
    for c in candidates:
        if c["path"] != best["path"]:
            c["path"].unlink(missing_ok=True)
    return best["path"]


async def _get_hero_product_image(brand_name: str, project_id: int, brand_url: str = "", target_ratio: float = 0.89) -> Path | None:
    """Find or generate the best hero product image for the agenda slide.

    Strategy:
      1. Check existing collected images for a high-quality portrait/square product shot
      2. Scrape brand website for product images (uses brand_url or web_search to find domain)
      3. Claude Vision picks the best from top candidates
      4. Fall back to DALL-E to generate a clean product shot

    Args:
        brand_name: Brand name for search/generation
        project_id: Project ID for output directory
        brand_url: Brand website URL (if known, skips web_search for domain discovery)
        target_ratio: Width/height ratio of the target box (0.89 for agenda half-page)

    Returns:
        Path to the best hero product image, or None
    """
    from PIL import Image
    import httpx
    import re

    output_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
    output_dir.mkdir(parents=True, exist_ok=True)
    hero_path = output_dir / "hero_product.png"

    # If already generated, reuse
    if hero_path.exists() and hero_path.stat().st_size > 50000:
        return hero_path

    # Strategy 1: Find best existing image (portrait/square, high-res, product-like)
    best_existing = None
    best_score = -1
    for img_file in output_dir.iterdir():
        if img_file.name.startswith("_cropped") or img_file.name.startswith("segment_bg") or img_file.name.startswith("persona_"):
            continue
        if img_file.suffix.lower() not in (".png", ".jpg", ".jpeg"):
            continue
        try:
            img = Image.open(img_file)
            w, h = img.size
            img.close()
            if w < 500 or h < 500:
                continue
            ratio = w / h
            # Score: prefer portrait/square, high-res, from brand website (httpx)
            orientation_score = 1.0 - abs(ratio - target_ratio)  # closer to target ratio = better
            resolution_score = min(w * h / (2000 * 2000), 1.0)   # up to 2000x2000
            source_score = 0.3 if "httpx" in img_file.name else 0.1  # prefer website images
            if ratio <= 1.2:  # portrait or square only
                score = orientation_score * 0.5 + resolution_score * 0.3 + source_score * 0.2
                if score > best_score:
                    best_score = score
                    best_existing = img_file
        except Exception:
            continue

    # Strategy 2: Scrape brand website for real product images
    hero_from_search = None
    try:
        client = None
        try:
            from config import ANTHROPIC_API_KEY
            import anthropic
            client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY) if ANTHROPIC_API_KEY else None
        except Exception:
            pass

        # Step 2a: Determine brand domain — use brand_url if available, else web_search
        brand_domain = ""
        if brand_url:
            # Extract domain from provided URL
            domain_match = re.search(r'(?:https?://)?(?:www\.)?([a-zA-Z0-9][-a-zA-Z0-9]*(?:\.[a-zA-Z0-9][-a-zA-Z0-9]*)+)', brand_url)
            if domain_match:
                brand_domain = domain_match.group(1).lower()
                print(f"[hero_image] Using brand_url domain: {brand_domain}")

        if not brand_domain and client:
            try:
                response = client.messages.create(
                    model=MODEL_SONNET,
                    max_tokens=512,
                    tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 2}],
                    messages=[{"role": "user", "content": (
                        f"What is {brand_name}'s official website URL? "
                        f"Just give me the main domain (e.g. nike.com, patagonia.com). "
                        f"Reply with ONLY the URL, nothing else."
                    )}],
                )
                for block in response.content:
                    if hasattr(block, "text"):
                        urls = re.findall(r'(?:https?://)?([a-zA-Z0-9][-a-zA-Z0-9]*(?:\.[a-zA-Z0-9][-a-zA-Z0-9]*)+)', block.text)
                        for u in urls:
                            u_lower = u.lower()
                            if u_lower not in ("google.com", "bing.com", "wikipedia.org") and len(u) < 50:
                                brand_domain = u_lower
                                break
                    if brand_domain:
                        break
            except Exception as e:
                print(f"[hero_image] web_search for domain failed: {e}")

        # Step 2b: Scrape the brand website for product images
        if brand_domain:
            print(f"[hero_image] Scraping {brand_domain} for product images...")
            product_paths = [
                f"https://{brand_domain}/",
                f"https://{brand_domain}/collections/all",
                f"https://{brand_domain}/products",
                f"https://{brand_domain}/shop",
            ]
            all_img_urls = []
            async with httpx.AsyncClient(follow_redirects=True, timeout=15) as http_client:
                for page_url in product_paths:
                    try:
                        resp = await http_client.get(page_url, headers={
                            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
                        })
                        if resp.status_code != 200 or len(resp.text) < 1000:
                            continue
                        html = resp.text
                        cdn_urls = re.findall(r'https?://cdn\.shopify\.com/s/files/[^\"\'\s,)]+', html)
                        img_srcs = re.findall(r'(?:src|data-src|srcset)=["\']([^"\'\s]+\.(?:png|jpg|jpeg|webp)(?:\?[^"\'\s]*)?)', html, re.IGNORECASE)
                        amazon_imgs = re.findall(r'https?://m\.media-amazon\.com/images/[^\"\'\s,)]+', html)
                        for url in cdn_urls + img_srcs + amazon_imgs:
                            base = url.split("?")[0]
                            lower = base.lower()
                            if lower.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                                if base not in all_img_urls:
                                    all_img_urls.append(base)
                        if all_img_urls:
                            print(f"[hero_image] Found {len(all_img_urls)} images from {page_url}")
                            break
                    except Exception:
                        continue

                # Score by filename patterns
                brand_lower = brand_name.lower().replace(" ", "_")
                scored = []
                for url in all_img_urls:
                    name = url.split("/")[-1].lower()
                    score = 0
                    # Positive signals
                    if "_sc" in name or "-sc" in name:  # Studio clean shot
                        score += 3
                    if "product" in name:
                        score += 2
                    if any(kw in name for kw in ("1080", "2000", "4000")):
                        score += 2
                    # Product type keywords — prefer product-related images
                    if "product" in name or "shop" in name or "catalog" in name:
                        score += 3
                    if brand_lower in name:
                        score += 2
                    if brand_lower in name or brand_lower.replace("_", "") in name:
                        score += 1
                    # Negative signals
                    if any(kw in name for kw in ("logo", "icon", "nav", "badge", "blog", "svg", "favicon")):
                        score -= 10
                    if any(kw in name for kw in ("bundle", "set", "group", "kid", "bambino", "accessories",
                                                   "boot", "bag", "pack", "sticker", "cleaning", "tablet",
                                                   "pet", "bowl", "customize", "custom", "gift", "card",
                                                   "banner", "hero_banner", "slider", "carousel")):
                        score -= 5
                    scored.append((score, url))
                scored.sort(key=lambda x: x[0], reverse=True)

                # Download top 5 candidates
                candidates = []
                for i, (score, url) in enumerate(scored[:12]):
                    if score < 0:
                        continue
                    try:
                        resp = await http_client.get(url, headers={
                            "User-Agent": "Mozilla/5.0",
                            "Accept": "image/*,*/*",
                        })
                        if resp.status_code == 200 and len(resp.content) > 20000:
                            search_path = output_dir / f"hero_search_{i}.png"
                            search_path.write_bytes(resp.content)
                            img = Image.open(search_path)
                            w, h = img.size
                            img.close()
                            if w >= 500 and h >= 500:
                                ratio = w / h
                                ratio_fit = 1.0 - abs(ratio - target_ratio)
                                res_score = min(w * h / (2000 * 2000), 1.0)
                                size_score = min(len(resp.content) / 500000, 1.0)
                                # White background detection: sample corner pixels
                                bg_score = 0.0
                                try:
                                    img_check = Image.open(search_path).convert("RGB")
                                    corners = [
                                        img_check.getpixel((5, 5)),
                                        img_check.getpixel((w - 5, 5)),
                                        img_check.getpixel((5, h - 5)),
                                        img_check.getpixel((w - 5, h - 5)),
                                    ]
                                    img_check.close()
                                    white_corners = sum(1 for r, g, b in corners if r > 230 and g > 230 and b > 230)
                                    bg_score = white_corners / 4.0  # 0.0 to 1.0
                                except Exception:
                                    pass
                                img_score = score * 0.25 + ratio_fit * 0.25 + res_score * 0.15 + size_score * 0.15 + bg_score * 0.2
                                candidates.append({
                                    "path": search_path, "url": url,
                                    "w": w, "h": h,
                                    "file_size": len(resp.content),
                                    "name_score": score, "img_score": img_score,
                                })
                                print(f"[hero_image] Candidate {len(candidates)}: {url.split('/')[-1][:50]} ({w}x{h}, score={img_score:.2f})")
                                if len(candidates) >= 5:
                                    break
                            else:
                                search_path.unlink(missing_ok=True)
                    except Exception:
                        continue

                # Use Claude Vision to pick the best (if API available)
                if len(candidates) >= 2 and client:
                    hero_from_search = await _pick_best_hero_with_vision(
                        client, brand_name, candidates, target_ratio
                    )
                elif len(candidates) >= 2:
                    best = max(candidates, key=lambda x: x["img_score"])
                    hero_from_search = best["path"]
                    print(f"[hero_image] Picked by score: {best['url'].split('/')[-1][:50]}")
                    for c in candidates:
                        if c["path"] != hero_from_search:
                            c["path"].unlink(missing_ok=True)
                elif candidates:
                    hero_from_search = candidates[0]["path"]
                    print(f"[hero_image] Only 1 candidate, using it directly")
        else:
            print("[hero_image] Could not determine brand website")
    except Exception as e:
        print(f"[hero_image] Website scraping failed: {e}")

    # Strategy 3: DALL-E generation as fallback
    # Build a product-specific prompt from analysis data if available
    hero_from_gen = None
    try:
        from pipeline.image_gen import generate_image
        from config import OPENAI_API_KEY
        if OPENAI_API_KEY and not hero_from_search:
            # Only generate if web search failed — real product photos are always better
            # Build a category-adaptive prompt from analysis context
            _category_hint = ""
            if analysis and isinstance(analysis, dict):
                _cap_section = analysis.get("capabilities", {})
                if isinstance(_cap_section, dict):
                    _cap_title = _cap_section.get("section_title", "")
                    _cat_m = _re.search(
                        r'(?:shaping|defining|in|of|the)\s+(?:the\s+)?(.+?)\s+(?:market|category|industry|space|landscape|sector)',
                        _cap_title.lower()
                    )
                    if _cat_m:
                        _category_hint = _cat_m.group(1).strip()
                if not _category_hint:
                    # Fallback: extract from any top-level description
                    for _v in _cap_section.values() if isinstance(_cap_section, dict) else []:
                        if isinstance(_v, str) and len(_v) > 20:
                            _category_hint = _v[:60]
                            break
            if not _category_hint:
                _category_hint = f"products by {brand_name}"
            prompt = (
                f"Ultra-realistic professional product photography of a premium "
                f"{_category_hint} product, similar to the brand {brand_name}. "
                f"Shot in a studio with soft directional lighting, clean white background, "
                f"slight shadow and reflection beneath. Portrait orientation, product centered. "
                f"Commercial catalog style, hyper-detailed, 8K quality. "
                f"No text, no logos, no branding, no watermarks."
            )
            gen_path = output_dir / "hero_generated.png"
            result = generate_image(prompt, output_path=gen_path, backend="dalle", size="1024x1792", quality="hd")
            if result and result.exists():
                hero_from_gen = result
                print(f"[hero_image] DALL-E generated hero product image")
    except Exception as e:
        print(f"[hero_image] DALL-E generation failed: {e}")

    # Pick the best result
    candidates = []
    if hero_from_search:
        candidates.append(("search", hero_from_search))
    if best_existing and best_score > 0.4:
        candidates.append(("existing", best_existing))
    if hero_from_gen:
        candidates.append(("generated", hero_from_gen))

    if not candidates:
        return best_existing  # fallback to whatever we have

    # Score candidates: prefer search > generated > existing
    priority = {"search": 3, "generated": 2, "existing": 1}
    best = max(candidates, key=lambda x: priority.get(x[0], 0))

    # Copy best to hero_path
    import shutil
    shutil.copy2(best[1], hero_path)
    print(f"[hero_image] Selected {best[0]} image as hero: {hero_path}")
    return hero_path


class _ImagePool:
    """Manages a pool of collected images, handing them out one at a time.

    Images are categorized: brand images are used first for content slides,
    product images for competitor slides, lifestyle for summaries.
    """

    def __init__(self, images: dict = None):
        self._images = images or {}
        self._brand_idx = 0
        self._product_idx = 0
        self._lifestyle_idx = 0
        self._all_idx = 0

    def next_brand(self) -> Path | None:
        """Get next brand image, cycling through available images."""
        imgs = self._images.get("brand", [])
        if not imgs:
            return self.next_any()
        img = imgs[self._brand_idx % len(imgs)]
        self._brand_idx += 1
        return img

    def next_product(self) -> Path | None:
        """Get next product image."""
        imgs = self._images.get("product", [])
        if not imgs:
            return self.next_brand()
        img = imgs[self._product_idx % len(imgs)]
        self._product_idx += 1
        return img

    def next_lifestyle(self) -> Path | None:
        """Get next lifestyle/stock image."""
        imgs = self._images.get("lifestyle", [])
        if not imgs:
            return self.next_brand()
        img = imgs[self._lifestyle_idx % len(imgs)]
        self._lifestyle_idx += 1
        return img

    def next_any(self) -> Path | None:
        """Get any available image."""
        imgs = self._images.get("all", [])
        if not imgs:
            return None
        img = imgs[self._all_idx % len(imgs)]
        self._all_idx += 1
        return img

    def best_brand(self) -> Path | None:
        """Get the highest-quality brand image (largest landscape-preferred)."""
        from PIL import Image
        imgs = self._images.get("brand", []) or self._images.get("all", [])
        if not imgs:
            return None
        scored = []
        for p in imgs:
            try:
                w, h = Image.open(p).size
                area = w * h
                # Prefer landscape, penalize extreme aspect ratios and screenshots
                if w / h > 2.5:
                    area *= 0.1  # screenshots are bad for summary slides
                elif w >= h:
                    area *= 1.2  # landscape bonus
                scored.append((area, p))
            except Exception:
                scored.append((0, p))
        scored.sort(reverse=True)
        return scored[0][1] if scored else None

    def has_images(self) -> bool:
        return bool(self._images.get("all"))

    def all_paths(self) -> list[Path]:
        """Return all unique image paths across all categories."""
        seen = set()
        result = []
        for imgs in self._images.values():
            for p in imgs:
                if p not in seen:
                    seen.add(p)
                    result.append(p)
        return result


def _pick_clean_product_image(img_pool, exclude: set = None) -> Path | None:
    """Pick the cleanest product image for summary slides.

    Prefers: ecom/product images, high resolution, no text overlays.
    Avoids: screenshots, marketing banners, cropped/segment images.
    """
    from PIL import Image as _PILImg
    exclude = exclude or set()
    candidates = []

    for p in img_pool.all_paths():
        if p in exclude:
            continue
        name = p.name.lower()
        # Skip non-product images
        if any(name.startswith(pfx) for pfx in (
            "_cropped", "segment_bg", "persona_", "hero_", "topic_",
            "comp_summary", "logo_", "chart_",
        )):
            continue
        try:
            with _PILImg.open(p) as img:
                w, h = img.size
        except Exception:
            continue
        if w < 400 or h < 400:
            continue

        score = w * h / 1_000_000  # base score from resolution (megapixels)

        # Prefer ecom/product images (clean white background shots)
        if any(name.startswith(pfx) for pfx in ("ecom_", "amazon_", "product_")):
            score += 5.0

        # Penalize landscape-extreme (likely screenshots or banners with text)
        ratio = w / h
        if ratio > 2.5:
            score -= 20.0  # screenshot
        elif ratio > 1.8:
            score -= 3.0   # likely banner/marketing image with text

        # Penalize images from web scraping (more likely to have text overlays)
        if name.startswith("search_") or name.startswith("httpx_"):
            score -= 2.0

        # Prefer portrait or near-square (typical product photography)
        if 0.6 < ratio < 1.4:
            score += 2.0

        candidates.append((score, p))

    candidates.sort(key=lambda x: x[0], reverse=True)
    return candidates[0][1] if candidates else None


async def _generate_topic_image(
    brand_name: str, slide_title: str, slide_key: str, output_dir: Path,
    source_image: Path | None = None,
) -> Path | None:
    """Generate an AI image tailored to a specific slide topic.

    Strategy:
    - If source_image is provided → Flux img2img (enhance real photo with context)
    - Otherwise → Flux text-to-image (photorealistic) with DALL-E fallback

    Used when no existing image matches the slide's theme.
    """
    import re

    # ── Scene prompts for composite (background only, NO product in scene) ──
    scene_configs = {
        "pricing": {
            "scene": (
                "Clean retail display with three empty white shelves at different heights "
                "against a soft gradient background. Professional studio lighting, "
                "minimalist product photography setup. No products, no text, no logos."
            ),
            "product_position": "center",
            "product_scale": 0.50,
        },
        "review": {
            "scene": (
                "Warm lifestyle desk scene with a smartphone showing a 5-star review "
                "interface and golden star ratings on screen. Shallow depth of field, "
                "warm window light. Leave right side empty for product placement. "
                "No products. Photorealistic."
            ),
            "product_position": "right",
            "product_scale": 0.65,
        },
        "channel": {
            "scene": (
                "Modern desk workspace with a laptop showing a generic e-commerce page "
                "and a half-open shipping box. Natural daylight from left window. "
                "Leave right side empty for product placement. "
                "No products. Photorealistic."
            ),
            "product_position": "right",
            "product_scale": 0.55,
        },
        "sku": {
            "scene": (
                "Clean white studio background with soft gradient lighting, ready for "
                "a product catalog display. Professional product photography setup. "
                "Empty scene, no products, no text."
            ),
            "product_position": "center",
            "product_scale": 0.70,
        },
        "execution": {
            "scene": (
                "Overhead view of a brand strategy workshop desk with color swatches, "
                "design sketches, sticky notes, and a coffee cup. Leave center space "
                "empty for a product. Natural daylight, editorial style. "
                "No products."
            ),
            "product_position": "center",
            "product_scale": 0.45,
        },
        "brand": {
            "scene": (
                "Premium marble desk surface with printed brand style guide pages, "
                "color palette cards, and a notebook. Warm natural light, shallow depth "
                "of field. Leave right side empty for product placement. "
                "No products."
            ),
            "product_position": "right",
            "product_scale": 0.60,
        },
        "challenge": {
            "scene": (
                "Dramatic concrete surface with moody side lighting and slight fog. "
                "Dark cinematic atmosphere suggesting uncertainty. Empty center stage. "
                "No products, no text."
            ),
            "product_position": "center",
            "product_scale": 0.60,
        },
    }

    # Match topic
    key_lower = slide_key.lower()
    title_lower = slide_title.lower()
    config = None
    for topic_key, cfg in scene_configs.items():
        if topic_key in key_lower or topic_key in title_lower:
            config = cfg
            break

    if not config:
        config = {
            "scene": (
                "Clean professional desk with a notebook, pen, and coffee cup. "
                "Warm natural light from left. Shallow depth of field. "
                "Leave right side empty for product. No products."
            ),
            "product_position": "right",
            "product_scale": 0.55,
        }

    try:
        from pipeline.image_gen import generate_image, generate_composite
        from config import OPENAI_API_KEY, REPLICATE_API_TOKEN

        if not OPENAI_API_KEY and not REPLICATE_API_TOKEN:
            return None

        safe_key = re.sub(r'[^a-z0-9]', '_', slide_key.lower())[:30]
        gen_path = output_dir / f"topic_{safe_key}.png"
        if gen_path.exists() and gen_path.stat().st_size > 30000:
            return gen_path

        result = None

        # Strategy 1: Composite — AI scene + real product cutout (best: preserves exact product)
        if source_image:
            backend = "flux" if REPLICATE_API_TOKEN else "dalle"
            # DALL-E 3 only supports 1024x1024, 1024x1792, 1792x1024
            _scene_size = (1344, 768) if backend == "flux" else (1792, 1024)
            result = generate_composite(
                scene_prompt=config["scene"],
                product_image_path=source_image,
                output_path=gen_path,
                backend=backend,
                scene_size=_scene_size,
                product_scale=config.get("product_scale", 0.55),
                product_position=config.get("product_position", "right"),
            )
            if result and result.exists():
                print(f"[image_assign] Composite for '{slide_title[:40]}': {gen_path.name}")
                return result

        # Strategy 2: Text-to-image (no product photo available)
        _t2i_backend = "flux" if REPLICATE_API_TOKEN else "dalle"
        _t2i_size = "1344x768" if _t2i_backend == "flux" else "1792x1024"
        if REPLICATE_API_TOKEN or OPENAI_API_KEY:
            result = generate_image(
                config["scene"],
                output_path=gen_path,
                backend=_t2i_backend,
                size=_t2i_size,
            )
            if result and result.exists():
                print(f"[image_assign] Flux t2i for '{slide_title[:40]}': {gen_path.name}")
                return result

        # Strategy 3: DALL-E HD fallback (only if Strategy 2 used Flux and failed)
        if OPENAI_API_KEY and _t2i_backend != "dalle":
            result = generate_image(
                config["scene"],
                output_path=gen_path,
                backend="dalle",
                size="1792x1024",
                quality="hd",
            )
            if result and result.exists():
                print(f"[image_assign] DALL-E for '{slide_title[:40]}': {gen_path.name}")
                return result

    except Exception as e:
        print(f"[image_assign] Topic image generation failed: {e}")

    return None


async def _assign_images_to_slides(
    img_pool: "_ImagePool",
    slide_topics: list[dict],
    brand_name: str,
    project_id: int = 0,
) -> dict[int, Path]:
    """Use Claude Vision to assign the best image to each capability slide.

    Given a list of slide topics (title + bullets), picks the most relevant
    image from the pool for each slide. Generates DALL-E images for unmatched
    slides. Returns {slide_index: image_path}.

    Strictly enforces no duplicate images across slides.
    """
    import base64
    from PIL import Image
    import io

    all_imgs = img_pool.all_paths()
    if not all_imgs:
        return {}

    # Filter to usable images (skip cropped, segment backgrounds, personas, hero)
    usable = []
    for p in all_imgs:
        name = p.name.lower()
        if any(name.startswith(pfx) for pfx in ("_cropped", "segment_bg", "persona_", "hero_", "topic_")):
            continue
        try:
            img = Image.open(p)
            w, h = img.size
            img.close()
            if w >= 300 and h >= 300:
                usable.append(p)
        except Exception:
            continue

    if not usable:
        return {}

    # Limit to 15 images max for Vision (token budget)
    sample = usable[:15]

    # Try Claude Vision for intelligent assignment
    assignments = {}
    try:
        from config import ANTHROPIC_API_KEY
        import anthropic
        if not ANTHROPIC_API_KEY:
            raise ValueError("No API key")

        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

        content = [{"type": "text", "text": (
            f"You are an art director building a premium brand strategy presentation for {brand_name}. "
            f"Match {len(sample)} available images to {len(slide_topics)} slides. "
            f"Each image must STRENGTHEN the slide's narrative — not just fill space.\n\n"
            f"MATCHING PRINCIPLES:\n"
            f"- 'execution/brand building' → hero product shot, brand identity, or flagship storefront showing how the brand presents itself\n"
            f"- 'product/feature/offer' → close-up product photography showing design, materials, craftsmanship\n"
            f"- 'product foundation/fundamentals' → collection photo showing product range, color variety, or SKU depth\n"
            f"- 'pricing' → use 0 (AI will generate) unless image shows real price tags, shelf pricing, or e-commerce listing with price\n"
            f"- 'channel/distribution' → use 0 unless image shows actual retail shelf, Amazon listing, or DTC storefront\n"
            f"- 'review/social/sentiment' → use 0 unless image shows real customer reviews, ratings, or social media posts\n"
            f"- 'SKU/range' → product lineup, catalog grid, or collection shot showing breadth\n"
            f"- 'website/domain' slides → website screenshot IS the evidence — use it\n"
            f"- 'challenge' slides → use the image that best ILLUSTRATES the specific problem described\n\n"
            f"QUALITY STANDARDS:\n"
            f"- Prefer images with advertising-level composition, lighting, and styling\n"
            f"- A logo alone is NOT a valid match for any content slide\n"
            f"- An unrelated product photo is worse than 0 (AI-generated)\n"
            f"- NEVER reuse an image number — each slide gets a UNIQUE image or 0\n"
            f"- Use 0 generously — a well-generated AI image beats a poorly matched real one\n\n"
            f"SLIDES:\n"
        )}]

        for i, topic in enumerate(slide_topics):
            desc = f"{i+1}. \"{topic['title']}\" — key: {topic.get('key', '')}"
            if topic.get("bullets"):
                desc += f" | bullets: {'; '.join(b[:60] for b in topic['bullets'][:2])}"
            content.append({"type": "text", "text": desc})

        content.append({"type": "text", "text": f"\nIMAGES (numbered 1-{len(sample)}):"})

        for i, img_path in enumerate(sample):
            try:
                img = Image.open(img_path).convert("RGB")
                w, h = img.size
                if w > 400 or h > 400:
                    ratio = min(400/w, 400/h)
                    img = img.resize((int(w*ratio), int(h*ratio)))
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                img.close()
                b64 = base64.b64encode(buf.getvalue()).decode()
                content.append({"type": "text", "text": f"\nImage {i+1} ({img_path.name[:40]}):"})
                content.append({"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}})
            except Exception:
                content.append({"type": "text", "text": f"\nImage {i+1}: [failed to load]"})

        content.append({"type": "text", "text": (
            f"\nFor each slide (1-{len(slide_topics)}), reply with the image number (1-{len(sample)}) or 0.\n"
            f"IMPORTANT: Each number can only appear ONCE. No duplicates.\n"
            f"Format: one number per line.\nExample:\n3\n7\n1\n0\n5"
        )})

        response = client.messages.create(
            model=MODEL_OPUS,
            max_tokens=256,
            messages=[{"role": "user", "content": content}],
        )

        import re
        answer = response.content[0].text.strip()
        numbers = re.findall(r'\d+', answer)
        used = set()
        for slide_idx, num_str in enumerate(numbers[:len(slide_topics)]):
            img_idx = int(num_str) - 1
            if 0 <= img_idx < len(sample) and img_idx not in used:
                assignments[slide_idx] = sample[img_idx]
                used.add(img_idx)
                print(f"[image_assign] Slide '{slide_topics[slide_idx]['title'][:40]}' → {sample[img_idx].name[:40]}")
            else:
                print(f"[image_assign] Slide '{slide_topics[slide_idx]['title'][:40]}' → will generate")

    except Exception as e:
        print(f"[image_assign] Vision failed ({e}), using keyword fallback")

        # Fallback: keyword-based matching (no duplicates)
        topic_kw_map = {
            "execution": ["website", "homepage", "httpx", "brand"],
            "product_offer": ["product", "ecom", "hero", "lineup"],
            "product_fundamentals": ["product", "detail", "close", "feature"],
            "pricing": ["price", "amazon", "listing", "ecom"],
            "channel": ["amazon", "website", "listing", "screenshot"],
            "brand": ["logo", "brand", "website", "homepage"],
            "review": ["review", "amazon", "rating", "star"],
            "sku": ["product", "lineup", "collection", "range"],
        }
        used_paths = set()
        for slide_idx, topic in enumerate(slide_topics):
            key = topic.get("key", "").lower()
            title = topic.get("title", "").lower()
            keywords = []
            for mk, mkws in topic_kw_map.items():
                if mk in key or mk in title:
                    keywords = mkws
                    break
            best_path, best_score = None, -1
            for p in usable:
                if p in used_paths:
                    continue
                name = p.name.lower()
                score = sum(1 for kw in keywords if kw in name)
                if score > best_score:
                    best_score = score
                    best_path = p
            if best_path and best_score > 0:
                assignments[slide_idx] = best_path
                used_paths.add(best_path)

    # Post-process: force AI generation for abstract topics (pricing, channel, review)
    # where real photos are never a good match. Runs after BOTH Vision and fallback paths.
    force_generate_keys = {"pricing", "channel", "review"}
    for si, topic in enumerate(slide_topics):
        key_lower = topic.get("key", "").lower()
        current = assignments.get(si)
        if not current:
            continue
        if any(fk in key_lower for fk in force_generate_keys) and not current.name.startswith("topic_"):
            del assignments[si]
            print(f"[image_assign] FORCE-GEN: '{topic['title'][:40]}' — releasing for AI generation")

    # Post-process: assign website screenshots to evidence slides (brand confusion, etc.)
    # Runs after BOTH Vision and keyword paths.
    evidence_keywords = {"website", "domain", "owned by", ".com"}

    def _is_screenshot(path: Path) -> bool:
        """Detect website screenshots by wide aspect ratio (>2.5:1)."""
        try:
            with Image.open(path) as img:
                w, h = img.size
            return w / h > 2.5
        except Exception:
            return False

    for si, topic in enumerate(slide_topics):
        title_lower = topic.get("title", "").lower()
        bullets_text = " ".join(topic.get("bullets", [])).lower()
        combined = title_lower + " " + bullets_text
        if not any(kw in combined for kw in evidence_keywords):
            continue
        current = assignments.get(si)
        if current and _is_screenshot(current):
            continue

        # Look for a screenshot assigned to a non-evidence slide — swap it
        found = False
        for other_si, other_path in list(assignments.items()):
            if other_si == si:
                continue
            if _is_screenshot(other_path):
                other_title = slide_topics[other_si].get("title", "").lower()
                if not any(kw in other_title for kw in evidence_keywords):
                    assignments[si], assignments[other_si] = other_path, (current if current else None)
                    if assignments[other_si] is None:
                        del assignments[other_si]
                    print(f"[image_assign] SWAP: '{topic['title'][:40]}' ← {other_path.name[:40]} (evidence)")
                    found = True
                    break

        # Search the full pool for unassigned screenshots
        if not found:
            assigned_set = set(assignments.values())
            for p in usable:
                if p not in assigned_set and _is_screenshot(p):
                    assignments[si] = p
                    print(f"[image_assign] EVIDENCE: '{topic['title'][:40]}' ← {p.name[:40]}")
                    found = True
                    break

    # Generate AI images for unmatched slides (Flux img2img > Flux t2i > DALL-E)
    output_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
    # Pick a good product photo as img2img source (first usable non-screenshot)
    source_img = None
    for p in usable:
        name = p.name.lower()
        if not any(x in name for x in ("screenshot", "httpx_", "brand_")):
            try:
                img = Image.open(p)
                w, h = img.size
                img.close()
                if w >= 500 and h >= 500:
                    source_img = p
                    break
            except Exception:
                continue
    for slide_idx, topic in enumerate(slide_topics):
        if slide_idx not in assignments:
            gen_img = await _generate_topic_image(
                brand_name, topic["title"], topic.get("key", ""), output_dir,
                source_image=source_img,
            )
            if gen_img:
                assignments[slide_idx] = gen_img

    return assignments


# ── High-level Slide Builders ────────────────────────────────

def _build_cover(prs, brand_name, date_str):
    """Clone cover slide, replace brand name and date.

    Original cover has one text frame with 2 runs in 1 paragraph:
      Run 0: "CozyFit"  (Montserrat 60pt)
      Run 1: "Brand Discovery" (Montserrat 35pt, preceded by line break)
    We replace run-by-run to preserve each font size.
    """
    slide = _clone_slide(prs, T_COVER)
    shapes = _find_text_shapes(slide)
    if len(shapes) >= 1:
        tf = shapes[0].text_frame
        para = tf.paragraphs[0]
        runs = para.runs
        if len(runs) >= 2:
            # Run 0 = brand name, Run 1 = subtitle
            runs[0].text = brand_name
            runs[1].text = "\nBrand Discovery"
        else:
            _set_text_preserve_format(tf, f"{brand_name}\nBrand Discovery")
    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, date_str)
    return slide


def _build_agenda(prs):
    """Clone agenda slide (no text changes needed — it's generic)."""
    return _clone_slide(prs, T_AGENDA)


def _build_approach(prs):
    """Clone the 'Our Brand Building Process' approach slide."""
    return _clone_slide(prs, T_APPROACH)


def _build_step_divider(prs):
    """Clone the 'Step 1 – Discovery' divider."""
    return _clone_slide(prs, T_STEP_DIVIDER)


def _build_section_header(prs, section_type):
    """Clone a section header. section_type: 'capabilities'|'competition'|'consumer'."""
    idx_map = {
        "capabilities": T_SECTION_CAPABILITIES,
        "competition": T_SECTION_COMPETITION,
        "consumer": T_SECTION_CONSUMER,
    }
    return _clone_slide(prs, idx_map.get(section_type, T_SECTION_CAPABILITIES))


def _build_content_slide(prs, title, bullets, insight_text, template_idx=T_CONTENT):
    """Clone a content slide (title + bullets + insight + image).

    Template shape layout (sorted by position):
      Shape 0 (top): Title — ALL CAPS, orange, Montserrat Bold
      Shape 1 (middle): Bullets — 3 paragraphs, Montserrat, space_before/after
      Shape 2 (bottom): Insight — teal/blue text, single paragraph
      Shape 3: Image (preserved as-is from template)

    Character limits (from original CozyFit template):
      Title: ~55 chars, Bullets: ~100 chars each, Insight: ~90 chars
    """
    slide = _clone_slide(prs, template_idx)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 55))
    if len(shapes) >= 2:
        if isinstance(bullets, list):
            bullets = [_truncate(b, 120) for b in bullets[:3]]
        else:
            bullets = [_truncate(bullets, 120)]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, _llm_compress(insight_text, 120, "strategic insight"))

    return slide


def _build_competitor_banner(prs, competitor_name):
    """Clone an Overview Slide as a competitor banner divider.

    Sample PPTs show "A closer look at [Brand Name]" as a breathing-room
    divider before each competitor deep-dive.
    """
    slide = _clone_slide(prs, T_SECTION_COMPETITION)
    shapes = _find_text_shapes(slide)
    # The section header has text like "A closer look at the competition"
    # Replace with the specific competitor name
    for s in shapes:
        text = s.text_frame.text.strip()
        if "closer look" in text.lower() or "competition" in text.lower():
            _set_text_preserve_format(s.text_frame, f"A closer look at\n{competitor_name}")
        elif text.lower() in ("capabilities", "competition", "consumer"):
            pass  # keep breadcrumb labels
    return slide


def _competitor_grid_position_map(competitors, focused_names=None):
    """Map competitors to consistent grid positions across both overview slides.

    The template 3x3 grid has 9 visual positions (sorted by top, left).
    Focused competitors are placed first so they cluster together;
    the red box dynamically resizes to encompass them.

    Returns: list of 9 competitor names in visual position order,
             where focused competitors occupy the first N positions.
    """
    all_names = [c.get("name", "Competitor") for c in competitors[:9]]

    if focused_names:
        focused_lower = {n.lower() for n in focused_names}
        focused_list = [n for n in all_names if n.lower() in focused_lower]
        non_focused_list = [n for n in all_names if n.lower() not in focused_lower]
    else:
        focused_list = all_names
        non_focused_list = []

    # Fill all 9 positions: focused first, then non-focused
    ordered = (focused_list + non_focused_list)[:9]
    # Pad to 9
    while len(ordered) < 9:
        ordered.append("")

    return ordered


def _build_competitor_overview(prs, competitors, category="", logos: dict = None, focused_names: list = None):
    """Build the 'all competitors' overview slide (template slide 15 pattern).

    Shows all competitor logos/names in a grid with a title and subtitle.
    Uses the same position mapping as the focused slide for consistency.
    """
    slide = _clone_slide(prs, T_COMPETITOR_GRID_ALL)
    shapes = _find_text_shapes(slide)

    cat_label = category.upper() if category else "CATEGORY"
    comp_names = _competitor_grid_position_map(competitors, focused_names)
    title = f"A WELL-DEFINED {cat_label} MARKET"

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 60))
    if len(shapes) >= 2:
        subtitle = (
            f"These brands set clear standards for how {category.lower() or 'products'} "
            f"should look, feel, and perform."
        )
        _set_text_preserve_format(shapes[1].text_frame, _truncate(subtitle, 120))

    # Replace template logos with real competitor logos (visual position order)
    if logos:
        logos_lower = {k.lower(): v for k, v in logos.items()}
        for i, name in enumerate(comp_names):
            if not name:
                continue
            logo_path = logos_lower.get(name.lower())
            if logo_path and Path(logo_path).exists():
                _replace_nth_picture(slide, Path(logo_path), i)

    return slide


def _build_competitor_focused(prs, competitors, category="", logos: dict = None, focused_names: list = None):
    """Build the 'focused review' slide (template slide 16 pattern).

    Uses the same grid position mapping as _build_competitor_overview so
    logos appear in identical positions across both slides. Non-focused
    competitors get a semi-transparent dim overlay.
    """
    slide = _clone_slide(prs, T_COMPETITOR_GRID_FOCUSED)
    shapes = _find_text_shapes(slide)

    cat_label = category.upper() if category else "CATEGORY"
    comp_names = _competitor_grid_position_map(competitors, focused_names)

    # Track text shapes that need border clearing — done LAST to avoid
    # python-pptx re-creating solidFill when we access .line on other shapes
    _border_clear_shapes = []
    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, f"FOCUSED REVIEW OF KEY {cat_label} BRANDS")
        _border_clear_shapes.append(shapes[0])
    # Find the subtitle shape (longer text at bottom)
    for s in shapes[1:]:
        if s.top > 4000000:  # bottom of slide
            subtitle = (
                f"To ground our strategy work, we will take a closer look at a focused set of "
                f"established {category.lower() or 'brands'} that help define the category today."
            )
            _set_text_preserve_format(s.text_frame, _truncate(subtitle, 150))
            _border_clear_shapes.append(s)

    # Replace template logos with real ones
    # Note: _replace_nth_picture sorts pictures by visual position (top, left),
    # so comp_names[i] gets placed at visual position i.
    from pptx.shapes.picture import Picture as _Pic
    SLIDE_AREA = 12192000 * 6858000

    if logos:
        logos_lower = {k.lower(): v for k, v in logos.items()}
        for i, name in enumerate(comp_names):
            if not name:
                continue
            logo_path = logos_lower.get(name.lower())
            if logo_path and Path(logo_path).exists():
                _replace_nth_picture(slide, Path(logo_path), i)

    # Dim non-focused competitor logos and resize red box
    # Use visual order (same as _replace_nth_picture) so indices match comp_names
    if focused_names:
        focused_lower = {n.lower() for n in focused_names}
        vis_pics = sorted(
            [s for s in slide.shapes if isinstance(s, _Pic) and (s.width * s.height) / SLIDE_AREA < 0.9],
            key=lambda s: (s.top, s.left)
        )

        from pptx.enum.shapes import MSO_SHAPE as _MSO
        from lxml import etree as _et
        _ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        focused_pics = []
        for i, pic in enumerate(vis_pics):
            name = comp_names[i] if i < len(comp_names) else ""
            is_focused = name and name.lower() in focused_lower
            if is_focused:
                focused_pics.append(pic)
                print(f"[competitor] BRIGHT: {name} (pos {i})")
            else:
                # Add semi-transparent white overlay to dim non-focused logos
                rect = slide.shapes.add_shape(
                    _MSO.RECTANGLE, pic.left, pic.top, pic.width, pic.height
                )
                fill = rect.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                sp_pr = rect._element.find(f".//{{{_ns_a}}}solidFill")
                if sp_pr is not None:
                    srgb = sp_pr.find(f"{{{_ns_a}}}srgbClr")
                    if srgb is not None:
                        alpha = _et.SubElement(srgb, f"{{{_ns_a}}}alpha")
                        alpha.set("val", "65000")  # 65% opacity — clearly dimmed
                rect.line.fill.background()
                print(f"[competitor] DIM: {name} (pos {i})")

        # Dynamically resize red box to encompass all focused brand logos
        red_box = None
        for s in slide.shapes:
            if isinstance(s, _Pic):
                continue
            # Skip text shapes with content — accessing .line on them creates solidFill
            if s in _border_clear_shapes:
                continue
            try:
                if hasattr(s, 'line') and s.line.color and s.line.color.rgb in (
                    RGBColor(0xFF, 0x00, 0x00), RGBColor(0xE0, 0x00, 0x00),
                    RGBColor(0xC0, 0x00, 0x00), RGBColor(0xFF, 0x33, 0x33)):
                    red_box = s
                    break
            except Exception:
                pass
            try:
                if (abs(s.left - 751063) < 200000 and abs(s.top - 1196502) < 200000
                        and s.width > 5000000):
                    red_box = s
                    break
            except Exception:
                pass

        if red_box and focused_pics:
            margin = 150000
            min_left = min(p.left for p in focused_pics) - margin
            min_top = min(p.top for p in focused_pics) - margin
            max_right = max(p.left + p.width for p in focused_pics) + margin
            max_bottom = max(p.top + p.height for p in focused_pics) + margin
            red_box.left = max(0, min_left)
            red_box.top = max(0, min_top)
            red_box.width = max_right - red_box.left
            red_box.height = max_bottom - red_box.top
            # Ensure uniform line width, style, and sharp corners after resize
            from pptx.util import Pt as _Pt
            from pptx.enum.dml import MSO_LINE_DASH_STYLE as _DASH
            red_box.line.width = _Pt(2.5)
            red_box.line.dash_style = _DASH.SOLID
            red_box.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
            red_box.fill.background()  # no fill inside the box
            # Set sharp miter join to avoid corner gaps
            from lxml import etree as _et
            _ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            ln_el = red_box._element.find(f".//{{{_ns_a}}}ln")
            if ln_el is not None:
                # Remove any existing join element
                for tag in ("miter", "bevel", "round"):
                    old = ln_el.find(f"{{{_ns_a}}}{tag}")
                    if old is not None:
                        ln_el.remove(old)
                _et.SubElement(ln_el, f"{{{_ns_a}}}miter", lim="800000")

            # Bring red box to front so dim overlays don't clip its lines
            sp_tree = red_box._element.getparent()
            sp_tree.remove(red_box._element)
            sp_tree.append(red_box._element)

    # Clear text shape borders LAST — after all shape.line accesses are done
    # (accessing shape.line on ANY shape triggers python-pptx to recreate solidFill)
    for s in _border_clear_shapes:
        _clear_shape_border(s)

    return slide


def _build_competitor_slide(prs, name, positioning_bullets, learnings_bullets, template_idx=None):
    """Clone a competitor deep-dive slide (two-column: positioning + learnings).

    Template variants (matching CozyFit):
      T_COMPETITOR (17): 1 wide banner image
      T_COMPETITOR_2IMG (18): 2 side-by-side images
      T_COMPETITOR_4IMG (19): 4 images in a row

    Template shape layout:
      Shape 0: Title — "DICKIES — POSITIONING & KEY LEARNINGS"
      Shape 1: Left column — "POSITIONING\nbullet1\nbullet2\n..."
      Shape 2: Right column — "KEY LEARNINGS\nbullet1\nbullet2\n..."
      Shape 3+: Images (preserved/replaced)
    """
    if template_idx is None:
        template_idx = T_COMPETITOR
    slide = _clone_slide(prs, template_idx)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(f"{name.upper()} — POSITIONING & KEY LEARNINGS", 60))

    def _fit_competitor_column(tf, header_text, bullets, max_bullets=3):
        """Fill a competitor text column. Controls overflow by truncating text only
        — font size stays consistent across all competitor slides."""
        paragraphs = list(tf.paragraphs)
        if paragraphs:
            _replace_para_text(paragraphs[0], header_text)
            if paragraphs[0].runs:
                paragraphs[0].runs[0].font.name = "Montserrat SemiBold"
                paragraphs[0].runs[0].font.bold = True

        # Progressively truncate until total fits within ~330 chars per column
        # Start at 115 chars per bullet — keeps 3 bullets within slide bounds
        max_chars = 115
        trimmed = [_truncate(b, max_chars) for b in bullets[:max_bullets]]
        total_chars = sum(len(b) for b in trimmed)
        while total_chars > 330 and max_chars > 70:
            max_chars -= 15
            trimmed = [_truncate(b, max_chars) for b in bullets[:max_bullets]]
            total_chars = sum(len(b) for b in trimmed)

        for i, bullet in enumerate(trimmed):
            para_idx = i + 1
            if para_idx < len(paragraphs):
                _set_bold_colon_text_para(paragraphs[para_idx], bullet)
            else:
                _add_paragraph_after(tf, paragraphs[-1], bullet)
                new_paras = list(tf.paragraphs)
                if len(new_paras) > para_idx:
                    _set_bold_colon_text_para(new_paras[para_idx], bullet)

        # Remove excess paragraphs
        all_paras = list(tf.paragraphs)
        for i in range(1 + len(trimmed), len(all_paras)):
            all_paras[i]._p.getparent().remove(all_paras[i]._p)

    # Set positioning column
    if len(shapes) >= 2:
        _fit_competitor_column(shapes[1].text_frame, "POSITIONING", positioning_bullets)

    # Set key learnings column
    if len(shapes) >= 3:
        _fit_competitor_column(shapes[2].text_frame, "KEY LEARNINGS", learnings_bullets)

    return slide


async def _fetch_competitor_logos(
    competitor_names: list[str], output_dir: Path
) -> dict[str, Path]:
    """Download real logos for each competitor.

    Strategy:
    1. Scrape brand's official website for logo/og:image
    2. Fallback: web search per brand for logo image URL
    3. Fallback: generate text-based logo with PIL

    Returns dict of {competitor_name: logo_path}.
    """
    import httpx
    import json
    import re
    from config import ANTHROPIC_API_KEY

    logo_dir = output_dir / "competitor_logos"
    logo_dir.mkdir(parents=True, exist_ok=True)

    # Check for cached logos first
    cached = {}
    for name in competitor_names:
        safe_name = name.lower().replace(" ", "_").replace("'", "").replace("ü", "u")
        for ext in (".png", ".jpg", ".webp"):
            candidate = logo_dir / f"{safe_name}_logo{ext}"
            if candidate.exists() and candidate.stat().st_size > 1000:
                cached[name] = candidate
                break
    if len(cached) == len(competitor_names):
        print(f"[competitor] All {len(cached)} logos cached")
        return cached

    # Domain cache — populated dynamically during competitor image collection
    # (the _collect_competitor_images function discovers domains via HTTP probing)
    brand_domains = {}

    async def _discover_domain(client: httpx.AsyncClient, name: str) -> str:
        """Try to discover a brand's domain via common patterns."""
        slug = name.lower().replace(" ", "").replace("'", "").replace("ü", "u")
        slug_h = name.lower().replace(" ", "-").replace("'", "")
        candidates = [f"{slug}.com", f"{slug_h}.com", f"{slug}usa.com",
                      f"go{slug}.com", f"shop{slug}.com", f"{slug}.co", f"get{slug}.com"]
        if " " not in name:
            candidates.insert(0, f"{name.lower()}.com")
        headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"}
        for cand in candidates:
            try:
                r = await client.head(f"https://www.{cand}", headers=headers, timeout=5, follow_redirects=True)
                if r.status_code < 400:
                    return cand
            except Exception:
                try:
                    r = await client.head(f"https://{cand}", headers=headers, timeout=5, follow_redirects=True)
                    if r.status_code < 400:
                        return cand
                except Exception:
                    continue
        return ""

    async def _extract_logo_from_site(client: httpx.AsyncClient, domain: str) -> str | None:
        """Try to extract logo image URL from brand website HTML."""
        try:
            resp = await client.get(f"https://www.{domain}", timeout=10)
            if resp.status_code != 200:
                resp = await client.get(f"https://{domain}", timeout=10)
            if resp.status_code != 200:
                return None
            html = resp.text
            base_url = str(resp.url)

            # NOTE: Skip og:image — it's often a lifestyle/product hero photo, not the logo.

            # Strategy 1: <img> with "logo" in class/alt/src
            logo_patterns = re.findall(
                r'<img[^>]*(?:class|alt|id)[^>]*logo[^>]*src=["\'](.*?)["\']|<img[^>]*src=["\'](.*?logo[^"\']*)["\']',
                html, re.I
            )
            for groups in logo_patterns:
                url = groups[0] or groups[1]
                if url and not url.endswith(".svg"):
                    if url.startswith("//"):
                        url = "https:" + url
                    elif url.startswith("/"):
                        url = base_url.rstrip("/") + url
                    return url

            # Strategy 3: header area image
            header_match = re.search(r'<header[^>]*>.*?<img[^>]*src=["\'](.*?)["\']', html[:5000], re.I | re.S)
            if header_match:
                url = header_match.group(1)
                if url.startswith("//"):
                    url = "https:" + url
                elif url.startswith("/"):
                    url = base_url.rstrip("/") + url
                return url

        except Exception:
            pass
        return None

    async def _search_logo_via_ai(client_anthropic, name: str) -> str | None:
        """Use Claude web search to find a logo URL for one brand."""
        try:
            response = client_anthropic.messages.create(
                model=MODEL_SONNET,
                max_tokens=500,
                tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 3}],
                messages=[{"role": "user", "content": f'Find one DIRECT IMAGE URL (.png or .jpg) of the official logo for "{name}" brand. Return ONLY the URL, nothing else.'}],
            )
            for block in response.content:
                if hasattr(block, "text") and block.text.strip():
                    # Extract URL from response
                    url_match = re.search(r'https?://[^\s"\'<>]+\.(?:png|jpg|jpeg|webp)', block.text)
                    if url_match:
                        return url_match.group(0)
                    # Maybe the whole text is a URL
                    text = block.text.strip()
                    if text.startswith("http"):
                        return text
        except Exception as e:
            print(f"[competitor] AI logo search failed for {name}: {e}")
        return None

    def _is_likely_logo(file_path: Path) -> bool:
        """Check if a downloaded image is likely a logo vs a lifestyle photo."""
        try:
            from PIL import Image
            # File size check: logos are typically <500KB, photos are >1MB
            if file_path.stat().st_size > 1_000_000:
                return False
            img = Image.open(str(file_path))
            w, h = img.size
            # Very tall portrait images are unlikely logos
            if h > w * 1.5:
                return False
            # Very large images (>2000px) are likely photos
            if w > 2500 and h > 2500:
                return False
            return True
        except Exception:
            return True  # If we can't check, assume it's fine

    def _generate_text_logo(name: str, out_path: Path):
        """Generate a simple text-based logo as last resort."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            w, h = 400, 200
            img = Image.new("RGB", (w, h), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            # Try to use a decent font
            font = None
            for font_name in ["/System/Library/Fonts/Helvetica.ttc", "/System/Library/Fonts/SFNSDisplay.ttf"]:
                try:
                    font = ImageFont.truetype(font_name, 48)
                    break
                except Exception:
                    continue
            if font is None:
                font = ImageFont.load_default()

            bbox = draw.textbbox((0, 0), name, font=font)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            x = (w - tw) // 2
            y = (h - th) // 2
            draw.text((x, y), name, fill=(30, 30, 30), font=font)
            img.save(str(out_path), "PNG")
            return True
        except Exception:
            return False

    # Phase 1: Try scraping brand websites for logos
    dl_headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
    async with httpx.AsyncClient(follow_redirects=True, timeout=15, headers=dl_headers) as dl_client:
        for name in competitor_names:
            if name in cached:
                continue
            safe_name = name.lower().replace(" ", "_").replace("'", "").replace("ü", "u")
            out_path = logo_dir / f"{safe_name}_logo.png"

            domain = brand_domains.get(name, "")
            if not domain:
                domain = await _discover_domain(dl_client, name)
                if domain:
                    brand_domains[name] = domain
            logo_url = None

            # Try site scraping first
            if domain:
                logo_url = await _extract_logo_from_site(dl_client, domain)
                if logo_url:
                    print(f"[competitor] Logo URL from site for {name}: {logo_url[:80]}")

            # Try downloading
            if logo_url:
                try:
                    resp = await dl_client.get(logo_url)
                    if resp.status_code == 200 and len(resp.content) > 1000:
                        out_path.write_bytes(resp.content)
                        if _is_likely_logo(out_path):
                            cached[name] = out_path
                            print(f"[competitor] Logo scraped: {name}")
                            continue
                        else:
                            print(f"[competitor] Rejected (looks like photo): {name}")
                            out_path.unlink(missing_ok=True)
                except Exception:
                    pass

        # Phase 2: AI web search for remaining
        missing = [n for n in competitor_names if n not in cached]
        if missing and ANTHROPIC_API_KEY:
            from anthropic import Anthropic
            ai_client = Anthropic(api_key=ANTHROPIC_API_KEY)
            for name in missing:
                safe_name = name.lower().replace(" ", "_").replace("'", "").replace("ü", "u")
                out_path = logo_dir / f"{safe_name}_logo.png"
                logo_url = await _search_logo_via_ai(ai_client, name)
                if logo_url:
                    try:
                        resp = await dl_client.get(logo_url)
                        if resp.status_code == 200 and len(resp.content) > 500:
                            out_path.write_bytes(resp.content)
                            if _is_likely_logo(out_path):
                                cached[name] = out_path
                                print(f"[competitor] Logo from AI search: {name}")
                                continue
                            else:
                                out_path.unlink(missing_ok=True)
                    except Exception:
                        pass

        # Phase 3: Generate text logos for any still missing
        still_missing = [n for n in competitor_names if n not in cached]
        for name in still_missing:
            safe_name = name.lower().replace(" ", "_").replace("'", "").replace("ü", "u")
            out_path = logo_dir / f"{safe_name}_logo.png"
            if _generate_text_logo(name, out_path):
                cached[name] = out_path
                print(f"[competitor] Text logo generated: {name}")

    return cached


def _is_likely_cutout(img_path: Path) -> bool:
    """Quick heuristic: check if image has uniform corners (product on plain background).

    Works for any color background (white, gray, colored). Detects studio product
    shots where the product is isolated on a clean, uniform background.
    """
    try:
        from PIL import Image
        import numpy as np
        img = Image.open(str(img_path)).convert("RGB")
        arr = np.array(img)
        h, w = arr.shape[:2]
        cs = min(30, h // 8, w // 8)  # corner sample size
        if cs < 5:
            return False
        corners = [arr[:cs, :cs], arr[:cs, -cs:], arr[-cs:, :cs], arr[-cs:, -cs:]]
        all_corner_px = np.concatenate([c.reshape(-1, 3) for c in corners])
        corner_std = float(np.std(all_corner_px))
        if corner_std < 12:
            median_color = np.median(all_corner_px, axis=0)
            bg_mask = np.all(np.abs(arr.astype(float) - median_color) < 30, axis=2)
            bg_ratio = float(np.mean(bg_mask))
            if bg_ratio > 0.35:
                return True
        return False
    except Exception:
        return False


async def _llm_judge_images(
    candidates: list[Path],
    prompt: str,
    pick_count: int = 2,
) -> list[Path]:
    """General-purpose LLM-as-judge for image selection.

    Sends up to 12 candidate images to Claude Haiku vision with the given prompt.
    All images are converted to JPEG before sending (handles GIF/WebP/palette issues).

    Args:
        candidates: List of image file paths to judge.
        prompt: Selection criteria prompt (should end with format instructions).
        pick_count: Max number of images to select.

    Returns:
        List of selected image paths (may be empty if LLM says "none").
    """
    from config import ANTHROPIC_API_KEY
    if not ANTHROPIC_API_KEY or len(candidates) == 0:
        return candidates[:pick_count]

    import anthropic, base64
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    content = []
    valid_candidates = []
    for i, p in enumerate(candidates[:12]):
        try:
            from PIL import Image
            import io
            data = p.read_bytes()
            img = Image.open(io.BytesIO(data))
            img_rgb = img.convert("RGB")
            buf = io.BytesIO()
            img_rgb.save(buf, format="JPEG", quality=85)
            data = buf.getvalue()
            img.close()
            img_rgb.close()
            b64 = base64.b64encode(data).decode()
            content.append({"type": "text", "text": f"Image {i+1}:"})
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
            })
            valid_candidates.append(p)
        except Exception:
            continue

    if len(valid_candidates) == 0:
        return []

    content.append({"type": "text", "text": prompt})

    try:
        resp = client.messages.create(
            model=MODEL_OPUS,
            max_tokens=50,
            messages=[{"role": "user", "content": content}],
        )
        answer = resp.content[0].text.strip()

        if answer.lower() == "none":
            return []

        selected = []
        for part in answer.replace(",", " ").split():
            try:
                idx = int(part.strip()) - 1
                if 0 <= idx < len(valid_candidates):
                    selected.append(valid_candidates[idx])
            except ValueError:
                continue
        return selected[:pick_count]
    except Exception as e:
        print(f"[llm_judge] Selection failed: {e}")
        return valid_candidates[:pick_count]


async def _collect_brand_summary_image(
    brand_name: str, output_dir: Path, category: str = ""
) -> Path | None:
    """Find a high-quality brand advertising image for the competition summary slide.

    Uses Bing Images search + LLM-as-judge to select the best image that
    represents the brand with product elements in an advertising context.
    """
    import httpx, hashlib, base64, re
    from config import ANTHROPIC_API_KEY

    safe_name = brand_name.lower().replace(" ", "_").replace("'", "")
    cache_path = output_dir / f"comp_summary_{safe_name}.jpg"
    # Check cache
    for ext in (".jpg", ".png", ".webp"):
        cp = output_dir / f"comp_summary_{safe_name}{ext}"
        if cp.exists() and cp.stat().st_size > 5000:
            print(f"[comp_summary] {brand_name}: cached → {cp.name}")
            return cp

    dl_headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
    cat_hint = category.lower() if category else "products"
    quoted = f'"{brand_name}"'
    _skip = ["shutterstock", "gettyimages", "istockphoto", "dreamstime", "depositphotos",
             "freepik", "pixabay", "pexels", "unsplash", "mockup", "template",
             "alibaba", "aliexpress", "4imprint", "customink"]
    search_queries = [
        f"{quoted} {cat_hint} official advertising campaign photography",
        f"{quoted} {cat_hint} brand editorial styled product photo",
        f"{quoted} {cat_hint} collection products campaign visual",
        f"{quoted} {cat_hint} product lineup hero image",
        f"site:instagram.com {quoted} {cat_hint} campaign",
    ]

    async with httpx.AsyncClient(follow_redirects=True, timeout=15, headers=dl_headers) as client:
        all_urls = []
        for query in search_queries:
            try:
                from urllib.parse import quote_plus
                bing_url = f"https://www.bing.com/images/search?q={quote_plus(query)}&first=1&count=15"
                resp = await client.get(bing_url, timeout=15)
                if resp.status_code == 200:
                    urls = re.findall(r'murl&quot;:&quot;(https?://[^&]+?)&quot;', resp.text)
                    for u in urls:
                        if u not in all_urls and not any(s in u.lower() for s in _skip):
                            all_urls.append(u)
            except Exception:
                continue

        # Download candidates
        candidates = []
        seen_hashes = set()
        for i, url in enumerate(all_urls):
            if len(candidates) >= 10:
                break
            try:
                resp = await client.get(url, timeout=12)
                if resp.status_code == 200 and len(resp.content) > 10000:
                    content_hash = hashlib.md5(resp.content).hexdigest()
                    if content_hash in seen_hashes:
                        continue
                    seen_hashes.add(content_hash)
                    ext = ".jpg"
                    if ".png" in url.lower(): ext = ".png"
                    elif ".webp" in url.lower(): ext = ".webp"
                    path = output_dir / f"_csummary_cand_{i:02d}{ext}"
                    path.write_bytes(resp.content)
                    from PIL import Image
                    img = Image.open(str(path))
                    w, h = img.size
                    img.close()
                    if w >= 400 and h >= 200:
                        candidates.append(path)
                    else:
                        path.unlink(missing_ok=True)
            except Exception:
                continue

        if not candidates:
            return None

        # LLM selects best image via shared judge utility
        summary_prompt = f"""Select the SINGLE BEST image for a "{brand_name}" competition summary slide in a brand strategy presentation.

The ideal image should:
- Be a high-quality {brand_name} advertising/campaign image featuring their {cat_hint} products
- Show {brand_name}'s products prominently (the actual {cat_hint} must be visible)
- Have advertising-level aesthetics: strong composition, appealing colors, professional photography
- Multiple products together, or products in an appealing scene/context, are preferred
- Feel like something from an official brand campaign or high-end editorial

REJECT: plain uniform-background product cutouts, stock photos, AI-generated images, wrong brand, wrong product category.

Reply with ONLY the single best image number. Example: "3". If none are acceptable, reply "none"."""
        selected = await _llm_judge_images(candidates, summary_prompt, pick_count=1)
        best = selected[0] if selected else None
        if best:
            print(f"[comp_summary] {brand_name}: LLM selected → {best.name}")

        # Rename selected to final path
        result = None
        if best:
            final = output_dir / f"comp_summary_{safe_name}{best.suffix}"
            if best != final:
                best.rename(final)
            result = final
            print(f"[comp_summary] {brand_name}: → {final.name}")

        # Clean up candidates
        for p in output_dir.glob("_csummary_cand_*"):
            p.unlink(missing_ok=True)

        return result


async def _collect_competitor_images(
    competitor_name: str, output_dir: Path, num_images: int = 4, category: str = ""
) -> list[Path]:
    """Collect product/brand images for a specific competitor.

    Strategy:
    1. Scrape the brand's official website for product images
    2. Fallback: use Claude web search to find product page URLs, then scrape those
    3. Filter out tiny/irrelevant images

    Returns list of downloaded image paths.
    """
    import httpx
    import json
    import re
    from config import ANTHROPIC_API_KEY

    safe_dir = competitor_name.lower().replace(" ", "_").replace("'", "").replace("ü", "u")
    comp_dir = output_dir / "competitor_images" / safe_dir
    comp_dir.mkdir(parents=True, exist_ok=True)

    # Check cache — separate hero and product images
    all_cached = [p for p in sorted(comp_dir.glob("*.png")) + sorted(comp_dir.glob("*.jpg")) + sorted(comp_dir.glob("*.webp"))
                  if not p.name.startswith("_")]
    cached_heroes = [p for p in all_cached if p.name.startswith("hero_")]
    cached_products = [p for p in all_cached if not p.name.startswith("hero_")]
    existing = cached_products  # product images for Phase 1/2 dedup

    # Only return from cache if we have 2 hero images (full set)
    if len(cached_heroes) >= 2:
        combined = cached_heroes[:2]
        if len(combined) < num_images:
            combined += cached_products[:num_images - len(combined)]
        if len(combined) >= num_images:
            print(f"[competitor] {competitor_name}: {len(combined)} images cached ({len(cached_heroes)} hero)")
            return combined[:num_images]
    # If only 1 hero, fall through to fetch more (don't return early)

    dl_headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}

    # Auto-discover brand domain by probing likely domain patterns
    brand_domain = ""
    _brand_slug = competitor_name.lower().replace(" ", "").replace("'", "").replace("ü", "u")
    # Generate candidate domains from brand name
    _slug_hyphen = competitor_name.lower().replace(' ', '-').replace(chr(39), '')
    _candidates = [
        f"{_brand_slug}.com",                                          # simplemodern.com
        f"{_slug_hyphen}.com",                                         # simple-modern.com
        f"{_brand_slug}usa.com",                                       # takeyausa.com
        f"go{_brand_slug}.com",                                        # gocontigo.com
        f"shop{_brand_slug}.com",                                      # shopnike.com
        f"{_brand_slug}.co",                                           # brand.co
        f"get{_brand_slug}.com",                                       # getbrand.com
    ]
    # For single-word brands, also try the name directly
    if " " not in competitor_name:
        _candidates.insert(0, f"{competitor_name.lower()}.com")
    try:
        import httpx as _hx
        async with _hx.AsyncClient(follow_redirects=True, timeout=5) as _aclient:
            for _cand in _candidates:
                try:
                    _resp = await _aclient.head(f"https://www.{_cand}", headers=dl_headers)
                    if _resp.status_code < 400:
                        brand_domain = _cand
                        break
                except Exception:
                    try:
                        _resp = await _aclient.head(f"https://{_cand}", headers=dl_headers)
                        if _resp.status_code < 400:
                            brand_domain = _cand
                            break
                    except Exception:
                        continue
        if brand_domain:
            print(f"[competitor] {competitor_name}: discovered domain → {brand_domain}")
        else:
            print(f"[competitor] {competitor_name}: no domain found (tried {len(_candidates)} candidates)")
    except Exception as _e:
        print(f"[competitor] {competitor_name}: domain discovery failed: {_e}")
    # Deduplicate existing images by content hash
    import hashlib
    seen_hashes = set()
    deduped_existing = []
    for p in existing:
        try:
            h = hashlib.md5(p.read_bytes()).hexdigest()
            if h not in seen_hashes:
                seen_hashes.add(h)
                deduped_existing.append(p)
        except Exception:
            deduped_existing.append(p)
    downloaded = list(deduped_existing)
    img_idx = len(existing)  # keep incrementing from original count to avoid name collisions

    # No longer reject images based on aspect ratio heuristics per category —
    # the LLM judge handles relevance filtering more accurately and generalizes
    # across all industries (apparel, beauty, tech, food, drinkware, etc.)

    async def _download_hero_candidate(dl_client, url, comp_dir, idx):
        """Download a hero candidate image. Only basic size/format checks — LLM judges quality."""
        ext = ".jpg"
        if ".png" in url.lower():
            ext = ".png"
        elif ".webp" in url.lower():
            ext = ".webp"
        out_path = comp_dir / f"_candidate_{idx:02d}{ext}"
        try:
            resp = await dl_client.get(url, timeout=12)
            if resp.status_code == 200 and len(resp.content) > 10000:
                ct = resp.headers.get("content-type", "")
                if "image" in ct or url.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                    content_hash = hashlib.md5(resp.content).hexdigest()
                    if content_hash in seen_hashes:
                        return None
                    out_path.write_bytes(resp.content)
                    seen_hashes.add(content_hash)
                    try:
                        from PIL import Image
                        img = Image.open(str(out_path))
                        w, h = img.size
                        img.close()
                        if w < 400 or h < 200:
                            out_path.unlink(missing_ok=True)
                            return None
                    except Exception:
                        out_path.unlink(missing_ok=True)
                        return None
                    return out_path
        except Exception:
            pass
        return None

    async def _select_heroes(candidates: list[Path], brand_name: str, cat: str) -> list[Path]:
        """Select best 2 hero images using shared LLM judge with brand-specific prompt."""
        hero_prompt = f"""You are an art director selecting hero images for a brand strategy deck about "{brand_name}" ({cat}).

Pick the 2 images that would look best on a full-width slide in a premium consulting presentation.

IDEAL (pick these first):
- Campaign/advertising imagery: styled product shots with intentional composition, lighting, and color grading
- Collection shots: multiple {brand_name} products arranged with visual appeal (variety of colors, sizes, or styles)
- Lifestyle context: products in a real environment (retail display, in-use scene, styled flat-lay) with aesthetic quality
- Any image that communicates the brand's visual identity and positioning at a glance

ACCEPTABLE:
- A single product in context (on a surface, in someone's hand, on a shelf)
- Editorial or review photography with decent composition

HARD REJECT (reply "none" if ONLY these remain):
- Products on plain white/gray/solid backgrounds (e-commerce cutouts, studio catalog shots)
- Wrong brand or wrong product category
- Stock photography, watermarked images, AI-generated fakes
- Blurry, pixelated, or extremely low-resolution images

Reply with ONLY the image numbers (1-indexed), comma-separated. Example: "3, 5". If ALL are reject-worthy, reply "none"."""
        selected = await _llm_judge_images(candidates, hero_prompt, pick_count=2)
        print(f"[competitor] {brand_name}: LLM hero judge → {len(selected)} selected")
        return selected

    async def _download_img(dl_client, url, comp_dir, idx):
        """Download one image, return path or None. Rejects tiny/duplicate images."""
        ext = ".jpg"
        if ".png" in url.lower():
            ext = ".png"
        elif ".webp" in url.lower():
            ext = ".webp"
        out_path = comp_dir / f"product_{idx:02d}{ext}"
        if out_path.exists() and out_path.stat().st_size > 5000:
            # Check if duplicate of already-downloaded image
            h = hashlib.md5(out_path.read_bytes()).hexdigest()
            if h in seen_hashes:
                return None
            seen_hashes.add(h)
            return out_path
        try:
            resp = await dl_client.get(url, timeout=12)
            if resp.status_code == 200 and len(resp.content) > 8000:
                ct = resp.headers.get("content-type", "")
                if "image" in ct or url.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                    # Check duplicate before writing
                    content_hash = hashlib.md5(resp.content).hexdigest()
                    if content_hash in seen_hashes:
                        return None
                    out_path.write_bytes(resp.content)
                    seen_hashes.add(content_hash)
                    # Verify dimensions and content quality
                    try:
                        from PIL import Image, ImageStat
                        img = Image.open(str(out_path))
                        w, h = img.size
                        if w < 400 or h < 300:
                            img.close()
                            out_path.unlink(missing_ok=True)
                            return None
                        # Reject small portrait images (wildlife/hero, not products)
                        if w < 600 and h > w * 1.3:
                            img.close()
                            print(f"[competitor] Rejected small portrait: {out_path.name} ({w}x{h})")
                            out_path.unlink(missing_ok=True)
                            return None
                        # Reject near-blank images (>95% white/single color)
                        rgb = img.convert("RGB")
                        stat = ImageStat.Stat(rgb)
                        avg_std = sum(stat.stddev) / 3
                        avg_mean = sum(stat.mean) / 3
                        if avg_std < 10 and avg_mean > 240:
                            img.close()
                            print(f"[competitor] Rejected blank/white: {out_path.name}")
                            out_path.unlink(missing_ok=True)
                            return None
                        # Reject low-content images (badges, icons)
                        if avg_std < 25 and w < 1200:
                            img.close()
                            print(f"[competitor] Rejected low-content: {out_path.name} (std={avg_std:.0f})")
                            out_path.unlink(missing_ok=True)
                            return None
                        # Reject grayscale graphics (logos, badges, text-only images)
                        # Real product photos always have color; badges/icons are often
                        # pure black-and-white or single-hue
                        import numpy as _np
                        arr = _np.array(rgb)
                        r_ch, g_ch, b_ch = arr[:,:,0].astype(int), arr[:,:,1].astype(int), arr[:,:,2].astype(int)
                        color_spread = max(
                            abs(r_ch - g_ch).mean(),
                            abs(r_ch - b_ch).mean(),
                            abs(g_ch - b_ch).mean(),
                        )
                        if color_spread < 3:
                            img.close()
                            print(f"[competitor] Rejected grayscale graphic: {out_path.name}")
                            out_path.unlink(missing_ok=True)
                            return None
                        img.close()
                    except Exception:
                        pass
                    return out_path
        except Exception:
            pass
        return None

    hero_images = []  # Separate list for hero/banner/ad images
    # Brand matching vars (used across all search phases)
    brand_domain_short = brand_domain.replace("www.", "").split(".")[0] if brand_domain else ""
    brand_name_lower = competitor_name.lower().replace(" ", "").replace("'", "").replace("ü", "u")
    # For multi-word brands, build slug variants (e.g., "simplemodern", "simple-modern")
    # Don't use individual common words — they cause false matches
    brand_slugs = []
    if " " in competitor_name:
        words = [w.lower() for w in competitor_name.split() if len(w) >= 3]
        brand_slugs.append("".join(words))         # "simplemodern"
        brand_slugs.append("-".join(words))         # "simple-modern"
        brand_slugs.append("_".join(words))         # "simple_modern"
    _skip_sites = ["qualitylogoproducts", "4imprint", "customink", "alibaba", "mockup",
                   "aliexpress", "shutterstock", "gettyimages", "istockphoto",
                   "dreamstime", "depositphotos", "freepik", "pixabay", "pexels",
                   "dhgate", "wish.com", "freepaani", "adorg", "template",
                   "executiveadvertising", "promotional", "logo-product", "unsplash"]

    async with httpx.AsyncClient(follow_redirects=True, timeout=15, headers=dl_headers) as dl_client:
        # Phase 0: Find brand advertising/campaign images via Bing Images search
        # Target: large lifestyle/advertising shots with taglines (like "FOR WORK THAT MATTERS")
        hero_idx = 0
        cached_heroes = sorted(comp_dir.glob("hero_*.jpg")) + sorted(comp_dir.glob("hero_*.png")) + sorted(comp_dir.glob("hero_*.webp"))
        if len(cached_heroes) >= 2:
            hero_images = cached_heroes[:2]
            print(f"[competitor] {competitor_name}: {len(cached_heroes)} hero images cached")
        elif cached_heroes:
            hero_images = cached_heroes[:1]
            print(f"[competitor] {competitor_name}: 1 hero image cached, will fetch more")
            # Fall through to fetch 1 more hero
        if not hero_images or len(hero_images) < 2:
            cat_hint = category.lower() if category else "brand"
            quoted_name = f'"{competitor_name}"'
            domain = brand_domain
            # Collect candidate URLs from multiple Bing searches
            if domain:
                search_queries = [
                    f"site:{domain} {competitor_name} {cat_hint}",
                    f"{quoted_name} {cat_hint} product photography advertising",
                    f"{quoted_name} {cat_hint} collection lineup campaign",
                    f"{quoted_name} {cat_hint} editorial styled product photo",
                    f"{quoted_name} {cat_hint} brand advertising visual",
                ]
            else:
                search_queries = [
                    f"{quoted_name} {cat_hint} product photography advertising",
                    f"{quoted_name} {cat_hint} collection lineup campaign",
                    f"{quoted_name} {cat_hint} editorial styled product photo",
                    f"{quoted_name} {cat_hint} brand advertising visual",
                ]
            # Gather all candidate URLs across queries
            all_candidate_urls = []
            for query in search_queries:
                try:
                    from urllib.parse import quote_plus
                    bing_url = f"https://www.bing.com/images/search?q={quote_plus(query)}&first=1&count=15"
                    resp = await dl_client.get(bing_url, headers=dl_headers, timeout=15)
                    if resp.status_code == 200:
                        media_urls = re.findall(r'murl&quot;:&quot;(https?://[^&]+?)&quot;', resp.text)
                        for url in media_urls:
                            url_lower = url.lower()
                            if any(skip in url_lower for skip in _skip_sites):
                                continue
                            if url not in [u for u in all_candidate_urls]:
                                all_candidate_urls.append(url)
                except Exception:
                    continue
            print(f"[competitor] {competitor_name}: {len(all_candidate_urls)} candidate URLs from Bing")

            # Download candidates (basic checks only, LLM will judge quality)
            candidate_paths = []
            cand_idx = 0
            for url in all_candidate_urls:
                if len(candidate_paths) >= 12:  # max 12 candidates for LLM
                    break
                path = await _download_hero_candidate(dl_client, url, comp_dir, cand_idx)
                if path:
                    candidate_paths.append(path)
                    cand_idx += 1

            # Pre-filter: separate lifestyle candidates from cutouts
            if candidate_paths:
                lifestyle_candidates = [p for p in candidate_paths if not _is_likely_cutout(p)]
                cutout_count = len(candidate_paths) - len(lifestyle_candidates)
                print(f"[competitor] {competitor_name}: {len(lifestyle_candidates)} lifestyle / {cutout_count} cutout candidates")

                # If we have enough lifestyle candidates, only send those to LLM
                llm_candidates = lifestyle_candidates if len(lifestyle_candidates) >= 2 else candidate_paths
                selected = await _select_heroes(llm_candidates, competitor_name, category or cat_hint)

                # If LLM rejected all OR we had no lifestyle candidates, retry
                if not selected or (len(lifestyle_candidates) == 0):
                    print(f"[competitor] {competitor_name}: all candidates were cutouts, retrying with lifestyle queries")
                    # Clean up old candidates
                    for p in comp_dir.glob("_candidate_*"):
                        p.unlink(missing_ok=True)
                    retry_queries = [
                        f"{quoted_name} {cat_hint} in use real life",
                        f"{quoted_name} {cat_hint} editorial photography magazine",
                        f"{quoted_name} {cat_hint} person using product scene",
                        f"{quoted_name} {cat_hint} store display retail shelf",
                    ]
                    retry_urls = []
                    for query in retry_queries:
                        try:
                            from urllib.parse import quote_plus
                            bing_url = f"https://www.bing.com/images/search?q={quote_plus(query)}&first=1&count=15"
                            resp = await dl_client.get(bing_url, headers=dl_headers, timeout=15)
                            if resp.status_code == 200:
                                media_urls = re.findall(r'murl&quot;:&quot;(https?://[^&]+?)&quot;', resp.text)
                                for url in media_urls:
                                    url_lower = url.lower()
                                    if any(skip in url_lower for skip in _skip_sites):
                                        continue
                                    if url not in all_candidate_urls and url not in retry_urls:
                                        retry_urls.append(url)
                        except Exception:
                            continue
                    print(f"[competitor] {competitor_name}: {len(retry_urls)} retry URLs from lifestyle queries")
                    candidate_paths = []
                    cand_idx = 0
                    for url in retry_urls:
                        if len(candidate_paths) >= 12:
                            break
                        path = await _download_hero_candidate(dl_client, url, comp_dir, cand_idx)
                        if path:
                            candidate_paths.append(path)
                            cand_idx += 1
                    if candidate_paths:
                        selected = await _select_heroes(candidate_paths, competitor_name, category or cat_hint)

                # Rename selected to hero_00, hero_01
                for i, p in enumerate(selected or []):
                    hero_path = comp_dir / f"hero_{i:02d}{p.suffix}"
                    if p != hero_path:
                        p.rename(hero_path)
                    hero_images.append(hero_path)
                    print(f"[competitor] {competitor_name}: LLM selected hero → {hero_path.name}")

            # Clean up unselected candidates
            for p in comp_dir.glob("_candidate_*"):
                p.unlink(missing_ok=True)

        if len(hero_images) < 2:
            print(f"[competitor] {competitor_name}: only {len(hero_images)} hero images found")

        # Phase 1: Amazon search (preferred — white-background product photos)
        try:
            # Include category to get relevant product type (e.g., "BrandX headphones" not "BrandX speakers")
            cat_term = category.lower() if category else ""
            if cat_term:
                search_term = f"{competitor_name} {cat_term}"
            else:
                search_term = competitor_name
            # Use Amazon brand filter (p_89) to restrict to this brand's products
            from urllib.parse import quote_plus
            brand_filter = quote_plus(competitor_name)
            amazon_url = f"https://www.amazon.com/s?k={search_term.replace(' ', '+')}&rh=p_89%3A{brand_filter}"
            resp = await dl_client.get(amazon_url, timeout=15)
            if resp.status_code == 200:
                html = resp.text
                # Extract product image base IDs (e.g., "81qTpMJUfiL")
                raw_urls = re.findall(r'https://m\.media-amazon\.com/images/I/([A-Za-z0-9+]+L)', html)
                seen_ids = set()
                full_urls = []
                for img_id in raw_urls:
                    if img_id in seen_ids:
                        continue
                    seen_ids.add(img_id)
                    # Construct full-size URL (no sizing suffix = original resolution)
                    full_urls.append(f"https://m.media-amazon.com/images/I/{img_id}.jpg")

                for url in full_urls:
                    if len(downloaded) >= num_images:
                        break
                    path = await _download_img(dl_client, url, comp_dir, img_idx)
                    if path:
                        downloaded.append(path)
                        img_idx += 1
                        print(f"[competitor] {competitor_name}: Amazon → {path.name}")
        except Exception as e:
            print(f"[competitor] {competitor_name}: Amazon search failed: {e}")

        # Phase 2: Scrape brand website for additional images if needed
        domain = brand_domain
        if domain and len(downloaded) < num_images:
            try:
                resp = await dl_client.get(f"https://www.{domain}")
                if resp.status_code != 200:
                    resp = await dl_client.get(f"https://{domain}")
                if resp.status_code == 200:
                    html = resp.text
                    base_url = str(resp.url).rstrip("/")
                    # Extract all <img> src URLs
                    img_urls = re.findall(r'<img[^>]*src=["\'](.*?)["\']', html, re.I)
                    # Filter for product-like images (larger images, not icons/logos)
                    candidate_urls = []
                    for url in img_urls:
                        if url.startswith("data:"):
                            continue
                        if url.startswith("//"):
                            url = "https:" + url
                        elif url.startswith("/"):
                            url = base_url + url
                        # Skip tiny icons, SVGs, tracking pixels, UI elements, app store badges
                        if any(skip in url.lower() for skip in [
                            ".svg", "icon", "favicon", "pixel", "1x1", "logo",
                            "badge", "sprite", "arrow", "spinner", "loading",
                            "amazon.com/images/I/", "prime", "fresh",
                            "rating", "star", "button", "banner-ad",
                            "app-store", "appstore", "google-play", "googleplay",
                            "play.google", "apple.com/app", "itunes",
                        ]):
                            continue
                        # Prefer product/collection images
                        if any(kw in url.lower() for kw in ["product", "collection", "hero", "cdn", "shopify", "catalog", "featured"]):
                            candidate_urls.insert(0, url)  # prioritize
                        else:
                            candidate_urls.append(url)

                    for url in candidate_urls[:num_images * 2]:
                        if len(downloaded) >= num_images:
                            break
                        path = await _download_img(dl_client, url, comp_dir, img_idx)
                        if path:
                            downloaded.append(path)
                            img_idx += 1
                            print(f"[competitor] {competitor_name}: scraped {path.name}")
            except Exception as e:
                print(f"[competitor] {competitor_name}: site scrape failed: {e}")

        # Phase 2: AI web search for product page URLs → scrape those
        if len(downloaded) < num_images and ANTHROPIC_API_KEY:
            try:
                from anthropic import Anthropic
                ai_client = Anthropic(api_key=ANTHROPIC_API_KEY)
                response = ai_client.messages.create(
                    model=MODEL_SONNET,
                    max_tokens=1000,
                    tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 3}],
                    messages=[{"role": "user", "content": f'Search the web for "{competitor_name}" {category or "brand"} products. Find 3 PAGE URLs (Amazon listings, official product pages) with product photos. Return ONLY a JSON array of URLs: ["https://...", ...]'}],
                )
                text = ""
                for block in response.content:
                    if hasattr(block, "text"):
                        text += block.text

                page_urls = []
                start = text.find("[")
                end = text.rfind("]") + 1
                if start >= 0 and end > start:
                    try:
                        page_urls = [u for u in json.loads(text[start:end]) if isinstance(u, str) and u.startswith("http")]
                    except json.JSONDecodeError:
                        pass
                # Also extract URLs directly
                for m in re.finditer(r'https?://[^\s"\'<>\]]+', text):
                    url = m.group(0).rstrip(".,;)")
                    if url not in page_urls:
                        page_urls.append(url)

                # Scrape product images from discovered pages
                for page_url in page_urls[:4]:
                    if len(downloaded) >= num_images:
                        break
                    try:
                        resp = await dl_client.get(page_url, timeout=12)
                        if resp.status_code != 200:
                            continue
                        html = resp.text
                        base_url = str(resp.url).rstrip("/")

                        img_urls = re.findall(r'<img[^>]*src=["\'](.*?)["\']', html, re.I)
                        # Also check srcset and data-src
                        img_urls += re.findall(r'data-src=["\'](.*?)["\']', html, re.I)

                        for url in img_urls:
                            if len(downloaded) >= num_images:
                                break
                            if url.startswith("data:"):
                                continue
                            if url.startswith("//"):
                                url = "https:" + url
                            elif url.startswith("/"):
                                url = base_url + url
                            if any(skip in url.lower() for skip in [
                                ".svg", "icon", "favicon", "1x1", "pixel",
                                "badge", "sprite", "arrow", "spinner", "loading",
                                "prime", "fresh", "rating", "star", "button",
                            ]):
                                continue
                            path = await _download_img(dl_client, url, comp_dir, img_idx)
                            if path:
                                downloaded.append(path)
                                img_idx += 1
                    except Exception:
                        continue

            except Exception as e:
                print(f"[competitor] {competitor_name}: AI search failed: {e}")

    # Vision-based relevance filter: reject images that don't show this competitor's products
    if len(downloaded) > 1 and ANTHROPIC_API_KEY:
        try:
            filtered = await _filter_competitor_images(competitor_name, downloaded, category)
            if filtered:
                downloaded = filtered
                print(f"[competitor] {competitor_name}: {len(downloaded)} images after Vision filter")
        except Exception as e:
            print(f"[competitor] {competitor_name}: Vision filter skipped: {e}")

    # Combine: hero images first, then product images
    # This gives the 2-image template: 1 lifestyle/ad + 1 product cutout
    combined = hero_images + downloaded
    # Deduplicate by hash
    final = []
    final_hashes = set()
    for p in combined:
        try:
            h = hashlib.md5(p.read_bytes()).hexdigest()
            if h not in final_hashes:
                final_hashes.add(h)
                final.append(p)
        except Exception:
            final.append(p)
    print(f"[competitor] {competitor_name}: {len(final)} images total ({len(hero_images)} hero + {len(downloaded)} product)")
    return final[:num_images]


async def _filter_competitor_images(
    competitor_name: str, image_paths: list[Path], category: str = ""
) -> list[Path]:
    """Use Claude Vision to filter out images that don't show the competitor's products.

    Rejects images from unrelated brands (e.g., Apple logo on a drinkware competitor
    page, or a random bear/mascot that isn't the brand's product).
    """
    import base64
    from PIL import Image
    import io
    from config import ANTHROPIC_API_KEY
    import anthropic

    if not ANTHROPIC_API_KEY or len(image_paths) < 2:
        return image_paths

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    content = [{"type": "text", "text": (
        f"I'm building a brand report for the competitor '{competitor_name}'"
        f"{f' ({category})' if category else ''}. "
        f"For each image below, reply YES if it shows {competitor_name}'s actual products, "
        f"branding, or relevant content (product photos, lifestyle shots with their products, "
        f"packaging, retail displays). Reply NO if it shows:\n"
        f"- A completely different brand's products or logo\n"
        f"- Generic stock imagery unrelated to {competitor_name}\n"
        f"- Icons, UI elements, or navigation graphics\n"
        f"- Animals, people, or objects with no connection to {competitor_name}\n\n"
        f"Reply with one YES or NO per line, in order."
    )}]

    for i, img_path in enumerate(image_paths):
        try:
            img = Image.open(img_path).convert("RGB")
            w, h = img.size
            if w > 300 or h > 300:
                ratio = min(300/w, 300/h)
                img = img.resize((int(w*ratio), int(h*ratio)))
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            img.close()
            b64 = base64.b64encode(buf.getvalue()).decode()
            content.append({"type": "text", "text": f"\nImage {i+1}:"})
            content.append({"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}})
        except Exception:
            content.append({"type": "text", "text": f"\nImage {i+1}: [failed to load]"})

    try:
        response = client.messages.create(
            model=MODEL_SONNET,
            max_tokens=100,
            messages=[{"role": "user", "content": content}],
        )
        answer = response.content[0].text.strip().upper()
        lines = [line.strip() for line in answer.split("\n") if line.strip()]

        filtered = []
        for i, path in enumerate(image_paths):
            if i < len(lines) and "NO" in lines[i]:
                print(f"[competitor] Vision filter rejected: {path.name} for {competitor_name}")
            else:
                filtered.append(path)

        return filtered  # may return empty if all rejected — caller should handle
    except Exception as e:
        print(f"[competitor] Vision filter failed: {e}")
        return image_paths


def _replace_nth_picture(slide, image_path: Path, index: int = 0, mode: str = "contain"):
    """Replace the Nth picture shape on a slide with a new image.

    Args:
        slide: The slide object.
        image_path: Path to the replacement image.
        index: Which picture shape to replace (0-based, excludes background).
        mode: "contain" = fit with padding (for logos), "cover" = crop to fill (for photos).
    """
    if not image_path or not Path(image_path).exists():
        return

    from PIL import Image
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000
    pictures = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            box_w, box_h = shape.width, shape.height
            is_bg = (box_w * box_h) / SLIDE_AREA > 0.9
            if is_bg:
                continue
            pictures.append(shape)

    # Sort by visual position (top then left) for consistent grid ordering
    pictures.sort(key=lambda s: (s.top, s.left))

    if index >= len(pictures):
        return

    shape = pictures[index]
    box_w, box_h = shape.width, shape.height
    if box_h == 0:
        return
    box_ratio = box_w / box_h

    try:
        img = Image.open(str(image_path))
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        img_w, img_h = img.size
    except Exception:
        return

    img_ratio = img_w / img_h

    if mode == "contain":
        # Contain mode: fit within box with white padding (for logos)
        # Use 150 DPI for crisp logos
        target_w = int(box_w / 914400 * 150)
        target_h = int(box_h / 914400 * 150)
        if target_w < 200:
            target_w = 200
        if target_h < 200:
            target_h = 200

        scale = min(target_w / img_w, target_h / img_h) * 0.75
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        img_resized = img.resize((new_w, new_h), Image.LANCZOS)

        bg = Image.new("RGB", (target_w, target_h), (255, 255, 255))
        x = (target_w - new_w) // 2
        y = (target_h - new_h) // 2
        if img_resized.mode == "RGBA":
            bg.paste(img_resized, (x, y), img_resized)
        else:
            bg.paste(img_resized, (x, y))
        img = bg
    else:
        # Cover mode: fit full image (no cropping!) with edge-color background
        # Never crop — the full image must be visible to preserve information
        target_px_w = max(int(box_w / 914400 * 150), 400)  # 150 DPI
        target_px_h = max(int(box_h / 914400 * 150), 400)

        # Fit entire image within box (contain), fill gaps with edge color
        scale = min(target_px_w / img_w, target_px_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        if img.mode == "RGBA":
            img = img.convert("RGB")
        img_resized = img.resize((new_w, new_h), Image.LANCZOS)

        # Sample edge color from the image borders for seamless background
        import numpy as _np
        arr = _np.array(img_resized)
        # Average the border pixels (top row, bottom row, left col, right col)
        border_pixels = []
        if arr.shape[0] > 2 and arr.shape[1] > 2:
            border_pixels.append(arr[0, :, :])        # top row
            border_pixels.append(arr[-1, :, :])       # bottom row
            border_pixels.append(arr[:, 0, :])         # left col
            border_pixels.append(arr[:, -1, :])        # right col
            all_border = _np.concatenate(border_pixels, axis=0)
            edge_color = tuple(int(c) for c in all_border.mean(axis=0))
        else:
            edge_color = (245, 245, 245)  # light gray fallback

        # If edge color is very dark, use light gray instead
        if sum(edge_color) / 3 < 80:
            edge_color = (245, 245, 245)

        bg = Image.new("RGB", (target_px_w, target_px_h), edge_color)
        x = (target_px_w - new_w) // 2
        y = (target_px_h - new_h) // 2
        bg.paste(img_resized, (x, y))
        img = bg

    cropped_path = image_path.parent / f"_fitted_{image_path.stem}_{index}.png"
    img.save(str(cropped_path), format="PNG")
    img.close()

    image_part, rId = slide.part.get_or_add_image_part(str(cropped_path))
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ns_a_uri = "http://schemas.openxmlformats.org/drawingml/2006/main"
    blip = shape._element.find(f".//{{{ns_a_uri}}}blip")
    if blip is not None:
        blip.set(f"{ns_r}embed", rId)
        # Strip all image effects inherited from template (duotone, grayscale, etc.)
        _effect_tags = {"imgLayer", "duotone", "grayscl", "clrChange", "clrRepl",
                        "lum", "tint", "biLevel", "alphaModFix"}
        for child in list(blip):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in _effect_tags:
                blip.remove(child)
        # Reset srcRect to show full image (template may have cropping)
        blipFill = blip.getparent()
        if blipFill is not None:
            srcRect = blipFill.find(f"{{{ns_a_uri}}}srcRect")
            if srcRect is not None and any(srcRect.get(a, "0") != "0" for a in ("l","t","r","b")):
                for attr in ("l", "t", "r", "b"):
                    if attr in srcRect.attrib:
                        del srcRect.attrib[attr]
        return
    # Fallback
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(cropped_path), shape.left, shape.top, box_w, box_h)


def _add_red_border(slide, picture_index: int, border_width_pt: float = 3.0):
    """Add a red border/box around a specific picture shape on the slide."""
    from pptx.shapes.picture import Picture
    from pptx.enum.shapes import MSO_SHAPE

    SLIDE_AREA = 12192000 * 6858000
    pictures = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            box_w, box_h = shape.width, shape.height
            is_bg = (box_w * box_h) / SLIDE_AREA > 0.9
            if is_bg:
                continue
            pictures.append(shape)

    # Sort by visual position for consistent grid ordering
    pictures.sort(key=lambda s: (s.top, s.left))

    if picture_index >= len(pictures):
        return

    pic = pictures[picture_index]
    # Draw red border slightly outside the picture so it frames it cleanly
    outset = Emu(20000)  # ~0.02 inch outset
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        pic.left - outset, pic.top - outset,
        pic.width + outset * 2, pic.height + outset * 2,
    )
    rect.fill.background()  # no fill (transparent)
    rect.line.color.rgb = RGBColor(0xE0, 0x20, 0x20)  # Red
    rect.line.width = Pt(border_width_pt)


def _dim_unfocused(slide, focused_indices: list[int]):
    """Add a semi-transparent white overlay on non-focused pictures.

    Uses visual position order (sorted by top, left) to match the grid layout,
    not DOM order which can differ between template slides.
    """
    from pptx.shapes.picture import Picture
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree

    SLIDE_AREA = 12192000 * 6858000
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pictures = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            if (shape.width * shape.height) / SLIDE_AREA < 0.9:
                pictures.append(shape)

    # Sort by visual position (top then left) to match grid layout
    pictures.sort(key=lambda s: (s.top, s.left))

    for i, pic in enumerate(pictures):
        if i not in focused_indices:
            # Add semi-transparent white overlay to dim
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, pic.left, pic.top, pic.width, pic.height
            )
            fill = rect.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            # Set transparency via XML — find srgbClr and add alpha child
            sp_pr = rect._element.find(f".//{{{ns_a}}}solidFill")
            if sp_pr is not None:
                srgb = sp_pr.find(f"{{{ns_a}}}srgbClr")
                if srgb is not None:
                    alpha = etree.SubElement(srgb, f"{{{ns_a}}}alpha")
                    alpha.set("val", "50000")  # 50% opacity
            rect.line.fill.background()  # no border


def _build_landscape_slide(prs, title, bullets, sidebar_text):
    """Clone the landscape summary slide (slide 24 pattern).

    Matches CozyFit reference: SemiBold label before colon, regular text after.
    Sidebar (Shape 2) is a 3.8"×2.0" callout box at the right side.
    """
    from pptx.util import Pt

    slide = _clone_slide(prs, T_LANDSCAPE)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 60))

    if len(shapes) >= 2:
        if isinstance(bullets, list):
            bullet_texts = [_truncate(b, 130) for b in bullets[:4]]
        else:
            bullet_texts = [_truncate(bullets, 130)]
        tf = shapes[1].text_frame
        paragraphs = list(tf.paragraphs)
        for i, bullet in enumerate(bullet_texts):
            if i < len(paragraphs):
                _set_bold_colon_text_para(paragraphs[i], bullet)
            else:
                _add_paragraph_after(tf, paragraphs[-1], bullet)
                new_paras = list(tf.paragraphs)
                if len(new_paras) > i:
                    _set_bold_colon_text_para(new_paras[i], bullet)
        all_paras = list(tf.paragraphs)
        for i in range(len(bullet_texts), len(all_paras)):
            all_paras[i]._p.getparent().remove(all_paras[i]._p)

    # Sidebar callout — split into header + body, format clearly
    if len(shapes) >= 3 and sidebar_text:
        # Move sidebar higher — align top with the bullet list area
        from pptx.util import Inches
        if len(shapes) >= 2:
            # Align sidebar top with the bullet list shape top
            shapes[2].top = shapes[1].top
        else:
            shapes[2].top = Inches(1.5)
        tf = shapes[2].text_frame
        # Split "Header:\nBody" into header and body
        if "\n" in sidebar_text:
            header, body = sidebar_text.split("\n", 1)
        else:
            header, body = "", sidebar_text
        body = _truncate(body.strip(), 280)

        paragraphs = list(tf.paragraphs)
        if header and paragraphs:
            _replace_para_text(paragraphs[0], header.strip())
            if paragraphs[0].runs:
                paragraphs[0].runs[0].font.name = "Montserrat SemiBold"
                paragraphs[0].runs[0].font.bold = True
                paragraphs[0].runs[0].font.size = Pt(22)
        # Set body text
        if len(paragraphs) > 1:
            _replace_para_text(paragraphs[1], body)
            if paragraphs[1].runs:
                paragraphs[1].runs[0].font.size = Pt(18)
        else:
            _add_paragraph_after(tf, paragraphs[-1], body)
            new_paras = list(tf.paragraphs)
            if len(new_paras) > 1 and new_paras[1].runs:
                new_paras[1].runs[0].font.size = Pt(18)
        # Remove excess paragraphs
        all_paras = list(tf.paragraphs)
        keep = 2 if header else 1
        for i in range(keep, len(all_paras)):
            all_paras[i]._p.getparent().remove(all_paras[i]._p)

    return slide


def _build_summary_slide(prs, title, summary_text, template_idx=T_SUMMARY):
    """Clone a summary slide (title + flowing paragraph + half-image)."""
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    from pptx.util import Pt as _PtFont

    slide = _clone_slide(prs, template_idx)
    shapes = _find_text_shapes(slide)

    # Summary slide has 2 text shapes: paragraph body and title
    # They may be in different order depending on position sort
    title_shape = None
    body_shape = None
    for s in shapes:
        text = s.text_frame.text.strip().upper()
        if "SUMMARY" in text or len(text) < 40:
            title_shape = s
        else:
            body_shape = s

    if title_shape:
        _set_text_preserve_format(title_shape.text_frame, _truncate(title, 40))
    if body_shape:
        # Vertically center the body text within its frame
        from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR
        body_shape.text_frame.word_wrap = True
        try:
            body_shape.text_frame.auto_size = None  # disable auto-size so anchor works
            # Set vertical anchor to middle via XML (more reliable than property)
            from lxml import etree as _et_sum
            _ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            body_pr = body_shape.text_frame._txBody.find(f"{{{_ns_a}}}bodyPr")
            if body_pr is not None:
                body_pr.set("anchor", "ctr")
        except Exception:
            pass
        # Truncate to a length that reads well at the template's default font size
        _set_text_preserve_format(body_shape.text_frame, _llm_compress(summary_text, 600, "section summary"))

    return slide


def _build_thank_you(prs):
    """Clone the Thank You slide."""
    return _clone_slide(prs, T_THANK_YOU)


# ── Contract-Required Analysis Slides ────────────────────────


def _build_clarity_scoring_slide(prs, clarity_data: dict):
    """Build a Clarity Scoring slide showing 5-dimension brand clarity assessment.

    Layout: Title + 5 horizontal bar rows (dimension name, score bar, evidence) + headline.
    Uses T_CONTENT template, replaces text shapes.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    headline = clarity_data.get("headline", "BRAND CLARITY ASSESSMENT")
    overall = clarity_data.get("overall_score", 0)
    dimensions = clarity_data.get("dimensions", [])
    strongest = clarity_data.get("strongest_zone", "")
    weakest = clarity_data.get("weakest_zone", "")

    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    # Title
    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, f"BRAND CLARITY SCORE: {overall}/100")

    # Build dimension summary as bullets
    dim_lines = []
    for d in dimensions[:5]:
        name = d.get("name", "")
        score = d.get("score", 0)
        max_s = d.get("max", 10)
        bar = "█" * score + "░" * (max_s - score)
        dim_lines.append(f"{name}: {bar} {score}/{max_s}")

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, dim_lines[:3])
    if len(shapes) >= 3:
        # Use insight shape for strongest/weakest zones
        insight = f"Strongest: {strongest[:50]} | Weakest: {weakest[:50]}"
        _set_text_preserve_format(shapes[2].text_frame, _truncate(insight, 120))

    # Add remaining dimensions + evidence as a table overlay
    _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    from pptx.util import Inches as _Inches, Pt as _Pt
    from lxml import etree

    # Create table for detailed scoring
    rows = len(dimensions) + 1  # header + dimensions
    cols = 3  # Dimension | Score | Evidence
    tbl_width = Emu(8229600)   # ~8.6 inches
    tbl_height = Emu(2743200)  # ~2.85 inches
    left = Emu(457200)         # 0.48 inches
    top = Emu(3200400)         # 3.33 inches

    try:
        table_shape = slide.shapes.add_table(rows, cols, left, top, tbl_width, tbl_height)
        tbl = table_shape.table

        # Column widths
        tbl.columns[0].width = Emu(2286000)  # Dimension name
        tbl.columns[1].width = Emu(1143000)  # Score
        tbl.columns[2].width = Emu(4800600)  # Evidence

        # Header row
        header_color = RGBColor(0xE8, 0x6C, 0x00)  # Orange
        for ci, label in enumerate(["DIMENSION", "SCORE", "EVIDENCE"]):
            cell = tbl.cell(0, ci)
            cell.text = label
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color

        # Data rows
        for ri, d in enumerate(dimensions[:5]):
            row_idx = ri + 1
            score = d.get("score", 0)
            max_s = d.get("max", 10)

            tbl.cell(row_idx, 0).text = d.get("name", "")
            tbl.cell(row_idx, 1).text = f"{score}/{max_s}"
            tbl.cell(row_idx, 2).text = _truncate(d.get("evidence", ""), 120)

            # Style cells
            for ci in range(3):
                cell = tbl.cell(row_idx, ci)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                # Color-code score
                if ci == 1:
                    if score >= 7:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xD4, 0xED, 0xDA)  # green
                    elif score <= 4:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF8, 0xD7, 0xDA)  # red
    except Exception as e:
        print(f"[clarity] Table creation failed: {e}")

    return slide


def _build_conflict_matrix_slide(prs, conflict_data: dict):
    """Build a Conflict Matrix slide showing pairwise segment conflicts.

    Shows a triangular matrix of segment conflicts with severity coloring.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    seg_names = conflict_data.get("segments", [])
    conflicts = conflict_data.get("conflicts", [])
    implication = conflict_data.get("strategic_implication", "")

    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "INTER-SEGMENT CONFLICT MATRIX")

    # Build conflict lookup
    conflict_map = {}
    for c in conflicts:
        key = (c.get("segment_a", ""), c.get("segment_b", ""))
        conflict_map[key] = c
        conflict_map[(key[1], key[0])] = c  # symmetric

    # Build table
    n = min(len(seg_names), 5)
    if n < 2:
        if len(shapes) >= 2:
            _set_text_preserve_format(shapes[1].text_frame, ["Insufficient segments for conflict analysis"])
        return slide

    rows = n + 1  # header row + segment rows
    cols = n + 1  # header col + segment cols
    tbl_width = Emu(8229600)
    tbl_height = Emu(3200400)
    left = Emu(457200)
    top = Emu(2400000)

    try:
        table_shape = slide.shapes.add_table(rows, cols, left, top, tbl_width, tbl_height)
        tbl = table_shape.table

        # Set column widths
        col_w = Emu(8229600 // (n + 1))
        for ci in range(cols):
            tbl.columns[ci].width = col_w

        # Header row and column
        header_color = RGBColor(0x2D, 0x3A, 0x4A)  # Dark blue
        tbl.cell(0, 0).text = ""
        for i, name in enumerate(seg_names[:n]):
            # Short name for header (first word only if long)
            short = name.split()[0] if len(name) > 12 else name
            tbl.cell(0, i + 1).text = short
            tbl.cell(i + 1, 0).text = short
            for ci in [0]:
                cell = tbl.cell(i + 1, ci)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(7)
                    p.font.bold = True
            cell = tbl.cell(0, i + 1)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color

        # Fill in matrix cells
        severity_colors = {
            "high": RGBColor(0xF8, 0xD7, 0xDA),    # Red
            "medium": RGBColor(0xFF, 0xF3, 0xCD),   # Yellow
            "low": RGBColor(0xD4, 0xED, 0xDA),      # Green
        }

        for ri in range(n):
            for ci in range(n):
                cell = tbl.cell(ri + 1, ci + 1)
                if ri == ci:
                    cell.text = "—"
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
                elif ri < ci:  # Upper triangle — show conflicts
                    key = (seg_names[ri], seg_names[ci])
                    conflict = conflict_map.get(key)
                    if conflict:
                        severity = conflict.get("severity", "low")
                        desc = conflict.get("description", "")
                        cell.text = _truncate(desc, 80)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = severity_colors.get(severity, RGBColor(0xF0, 0xF0, 0xF0))
                    else:
                        cell.text = "Low"
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = severity_colors["low"]
                else:  # Lower triangle — mirror
                    cell.text = ""
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(6)
                cell.text_frame.word_wrap = True

        # Header corner
        tbl.cell(0, 0).fill.solid()
        tbl.cell(0, 0).fill.fore_color.rgb = header_color

    except Exception as e:
        print(f"[conflict] Table creation failed: {e}")

    # Strategic implication in insight shape
    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, [_truncate(implication, 200)])
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, "Conflicts drive target selection — choosing one segment means accepting trade-offs with others.")

    return slide


def _build_hypothesis_validation_slide(prs, validations: list):
    """Build a Hypothesis Validation slide showing confirmed/partial/refuted status.

    Shows a status table with color-coded validation results.
    """
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor

    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "HYPOTHESIS VALIDATION SCORECARD")

    if not validations:
        if len(shapes) >= 2:
            _set_text_preserve_format(shapes[1].text_frame, ["No hypotheses available for validation"])
        return slide

    # Build table
    rows = min(len(validations), 6) + 1  # header + data rows
    cols = 4  # ID | Hypothesis | Status | Evidence
    tbl_width = Emu(8229600)
    tbl_height = Emu(3200400)
    left = Emu(457200)
    top = Emu(2200000)

    try:
        table_shape = slide.shapes.add_table(rows, cols, left, top, tbl_width, tbl_height)
        tbl = table_shape.table

        tbl.columns[0].width = Emu(571500)    # ID
        tbl.columns[1].width = Emu(2857500)   # Hypothesis
        tbl.columns[2].width = Emu(1143000)   # Status
        tbl.columns[3].width = Emu(3657600)   # Evidence

        # Header
        header_color = RGBColor(0x2D, 0x3A, 0x4A)
        for ci, label in enumerate(["#", "HYPOTHESIS", "STATUS", "EVIDENCE"]):
            cell = tbl.cell(0, ci)
            cell.text = label
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color

        # Status colors and icons
        status_config = {
            "confirmed": (RGBColor(0xD4, 0xED, 0xDA), "CONFIRMED"),
            "partially_supported": (RGBColor(0xFF, 0xF3, 0xCD), "PARTIAL"),
            "refuted": (RGBColor(0xF8, 0xD7, 0xDA), "REFUTED"),
        }

        for ri, h in enumerate(validations[:6]):
            row_idx = ri + 1
            tbl.cell(row_idx, 0).text = h.get("id", f"H{ri+1}")
            tbl.cell(row_idx, 1).text = _truncate(h.get("statement", ""), 100)

            status = h.get("status", "partially_supported")
            color, label = status_config.get(status, (RGBColor(0xFF, 0xF3, 0xCD), "PARTIAL"))
            status_cell = tbl.cell(row_idx, 2)
            status_cell.text = label
            status_cell.fill.solid()
            status_cell.fill.fore_color.rgb = color

            tbl.cell(row_idx, 3).text = _truncate(h.get("evidence", ""), 120)

            for ci in range(4):
                for p in tbl.cell(row_idx, ci).text_frame.paragraphs:
                    p.font.size = Pt(7)
                tbl.cell(row_idx, ci).text_frame.word_wrap = True

    except Exception as e:
        print(f"[hypothesis] Table creation failed: {e}")

    # Use remaining shapes for implication summary
    if len(shapes) >= 2:
        confirmed = sum(1 for h in validations if h.get("status") == "confirmed")
        partial = sum(1 for h in validations if h.get("status") == "partially_supported")
        refuted = sum(1 for h in validations if h.get("status") == "refuted")
        summary = f"{confirmed} confirmed, {partial} partially supported, {refuted} refuted out of {len(validations)} hypotheses tested"
        _set_text_preserve_format(shapes[1].text_frame, [summary])
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, "Consumer data validates strategic direction while revealing key assumption gaps.")

    return slide


def _build_evidence_plan_slide(prs, evidence_plan: dict):
    """Build an Evidence Collection Plan slide summarizing the research design."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor

    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "EVIDENCE COLLECTION PLAN")

    qs = evidence_plan.get("questionnaire_summary", {})
    gaps = evidence_plan.get("coverage_gaps", [])
    hyps = evidence_plan.get("hypotheses_to_validate", [])

    # Build summary bullets
    bullets = []
    if qs:
        target = qs.get("target_respondent", "Category purchasers 18+")
        bullets.append(f"Target: {_truncate(target, 100)}")
        total_q = qs.get("total_questions", "~22")
        duration = qs.get("estimated_duration", "10 minutes")
        sections = ", ".join(qs.get("sections", [])[:4])
        bullets.append(f"Survey: {total_q} questions, {duration} | {sections}")
    if hyps:
        bullets.append(f"Hypotheses to validate: {len(hyps)} from Phase 1-2 findings")

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, bullets[:3])

    # Coverage gaps as insight
    if len(shapes) >= 3 and gaps:
        gap_text = "Gaps: " + "; ".join(_truncate(g, 60) for g in gaps[:2])
        _set_text_preserve_format(shapes[2].text_frame, _truncate(gap_text, 120))

    # Add hypothesis validation targets as table
    if hyps:
        rows = min(len(hyps), 5) + 1
        cols = 3  # Hypothesis | Method | Sample
        tbl_width = Emu(8229600)
        tbl_height = Emu(2000000)
        left = Emu(457200)
        top = Emu(3600000)

        try:
            table_shape = slide.shapes.add_table(rows, cols, left, top, tbl_width, tbl_height)
            tbl = table_shape.table
            tbl.columns[0].width = Emu(3657600)
            tbl.columns[1].width = Emu(2286000)
            tbl.columns[2].width = Emu(2286000)

            header_color = RGBColor(0xE8, 0x6C, 0x00)
            for ci, label in enumerate(["HYPOTHESIS", "METHOD", "SAMPLE"]):
                cell = tbl.cell(0, ci)
                cell.text = label
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color

            for ri, h in enumerate(hyps[:5]):
                tbl.cell(ri + 1, 0).text = _truncate(h.get("hypothesis", ""), 90)
                tbl.cell(ri + 1, 1).text = _truncate(h.get("collection_method", ""), 60)
                tbl.cell(ri + 1, 2).text = _truncate(h.get("sample_target", ""), 50)
                for ci in range(3):
                    for p in tbl.cell(ri + 1, ci).text_frame.paragraphs:
                        p.font.size = Pt(7)
                    tbl.cell(ri + 1, ci).text_frame.word_wrap = True
        except Exception as e:
            print(f"[evidence_plan] Table creation failed: {e}")

    return slide


# ── Consumer Slide Builders ─────────────────────────────────


def _generate_research_approach(brand_name: str, analysis: dict, date_str: str) -> list[dict]:
    """Generate research approach methodology using Claude Opus.

    Uses the full brand analysis context (capabilities, competition, consumer)
    to produce a precise, category-appropriate quantitative survey methodology.
    Falls back to heuristic generation if API is unavailable.
    """
    # Try LLM-powered generation first
    try:
        from config import ANTHROPIC_API_KEY
        if ANTHROPIC_API_KEY:
            result = _generate_research_approach_llm(brand_name, analysis, date_str)
            if result and len(result) >= 4:
                return result
    except Exception as e:
        print(f"[research_approach] LLM generation failed: {e}")

    # Fallback: heuristic generation
    return _generate_research_approach_heuristic(brand_name, analysis, date_str)


def _generate_research_approach_llm(brand_name: str, analysis: dict, date_str: str) -> list[dict]:
    """Use Claude Opus to generate research methodology from analysis context."""
    import json as _json
    from anthropic import Anthropic as _Anthropic
    from config import ANTHROPIC_API_KEY

    # Build concise context from analysis
    context_parts = []
    cap = analysis.get("capabilities", {})
    if isinstance(cap, dict):
        es = cap.get("execution_summary", {})
        if isinstance(es, dict):
            context_parts.append(f"Brand positioning: {es.get('title', '')}")
        context_parts.append(f"Capabilities: {cap.get('capabilities_summary', '')[:300]}")

    comp = analysis.get("competition", {})
    if isinstance(comp, dict):
        mo = comp.get("market_overview", {})
        if isinstance(mo, dict):
            context_parts.append(f"Market: {mo.get('title', '')}")
        context_parts.append(f"Competition: {comp.get('competition_summary', '')[:300]}")
        # Include competitor names for brand awareness question design
        analyses = comp.get("competitor_analyses", [])
        comp_names = [c.get("name", "") for c in analyses if isinstance(c, dict)]
        if comp_names:
            context_parts.append(f"Key competitors: {', '.join(comp_names[:8])}")

    consumer = analysis.get("consumer", {})
    if isinstance(consumer, dict):
        context_parts.append(f"Consumer: {consumer.get('overview', '')[:300]}")

    context = "\n".join(context_parts)

    prompt = f"""You are a senior quantitative research director at a top-tier brand consulting firm (like Interbrand, Landor, or Kantar). Design the methodology for a consumer survey for this brand discovery project.

Brand: {brand_name}
Fielding date: {date_str}

## Analysis Context
{context}

## Task
Generate the "Participants" section for the QUANTITATIVE RESEARCH APPROACH slide. This must be realistic and precisely calibrated to this brand's category, consumer base, and competitive landscape.

You must determine:
1. **Sample size** (typically 200-300 for brand discovery; choose a specific number like 201, 253, etc.)
2. **Age minimum** (18+ for general consumer, 21+ for alcohol/parenting, 22+ for professional categories)
3. **Gender quota** — based on the actual consumer demographics of this specific category:
   - If the category consumer base skews female (e.g., beauty, yoga, lifestyle drinkware, fashion), use Female (60-70%) Male (30-40%)
   - If it skews male (e.g., power tools, gaming, tactical gear), use Male (60-70%) Female (30-40%)
   - If relatively balanced (e.g., tech, food, general home), use Female (50%) Male (50%)
   - Base this on real market knowledge of WHO actually buys in this category
4. **Screener criteria** — 2-3 bullets that precisely qualify respondents for THIS category:
   - Professional categories: specify the profession and role requirements
   - Consumer categories: specify purchase recency and category familiarity
   - Niche categories: specify relevant life stage or usage context
   - Always end with "Primary or shared decision-maker when purchasing"

Return ONLY a JSON object:
{{
  "sample_size": 253,
  "age_min": "18+ years old",
  "gender_quota": "Female (60%) Male (40%)",
  "screener_bullets": [
    "First screener criterion specific to this category",
    "Second screener criterion",
    "Primary or shared decision-maker when purchasing their own [category]"
  ]
}}"""

    client = _Anthropic(api_key=ANTHROPIC_API_KEY)
    try:
        response = client.messages.create(
            model=MODEL_OPUS,
            max_tokens=500,
            messages=[{"role": "user", "content": prompt}],
        )
        text = response.content[0].text if response.content else ""
    except Exception as e:
        print(f"[research_approach] LLM generation failed: {e}")
        return None

    start = text.find("{")
    end = text.rfind("}") + 1
    if start >= 0 and end > start:
        data = _json.loads(text[start:end])
        sample_n = data.get('sample_size', 200)
        age_min = data.get('age_min', '18+ years old')
        gender_q = data.get('gender_quota', 'Female (50%) Male (50%)')
        participants_detail = (
            f"Total of {sample_n} US Consumers:\n"
            f"US resident; {age_min}\n"
            f"Quota on Gender: {gender_q}\n"
            + "\n".join(data.get("screener_bullets", []))
        )
        print(f"[research_approach] LLM generated: n={sample_n}, gender={gender_q}")
        return [
            {"label": "Format", "detail": "Methodology:\tOnline survey\nLength:\t\t10 minutes\nBranding:\tUnbranded"},
            {"label": "Participants", "detail": participants_detail},
            {"label": "Analysis", "detail": "Demographics & Background\nShopping Habits, Category Usage and Ownership\nBrand Evaluation & Competitor Analysis\nMarket Segmentation"},
            {"label": "Timing", "detail": f"Fielding: {date_str}"},
        ]

    return []


def _generate_research_approach_heuristic(brand_name: str, analysis: dict, date_str: str) -> list[dict]:
    """Heuristic fallback for research approach when API is unavailable."""
    import re as _re
    import random

    # Extract category from analysis
    category = ""
    _cat_pattern = r'(?:shaping|defining|in|of|the)\s+(?:the\s+)?(.+?)\s+(?:market|category|industry|space|landscape|sector)'
    for source in [
        analysis.get("capabilities", {}).get("section_title", ""),
        (analysis.get("competition", {}).get("market_overview", {}) or {}).get("title", ""),
        (analysis.get("capabilities", {}).get("execution_summary", {}) or {}).get("title", ""),
    ]:
        cat_m = _re.search(_cat_pattern, source.lower()) if source else None
        if cat_m:
            category = cat_m.group(1).strip()
            break
    if not category:
        category = "consumer products"

    sample_n = random.choice([201, 203, 205, 208, 211, 215, 250, 253])

    participants_detail = (
        f"Total of {sample_n} US Consumers:\n"
        f"US resident; 18+ years old\n"
        f"Quota on Gender: Female (50%) Male (50%)\n"
        f"Have purchased {category} at least once in the past 12 months\n"
        f"Can name at least 2 brands in the {category} category (aided or unaided)\n"
        f"Primary or shared decision-maker when purchasing their own {category}"
    )

    return [
        {"label": "Format", "detail": "Methodology:\tOnline survey\nLength:\t\t10 minutes\nBranding:\tUnbranded"},
        {"label": "Participants", "detail": participants_detail},
        {"label": "Analysis", "detail": "Demographics & Background\nShopping Habits, Category Usage and Ownership\nBrand Evaluation & Competitor Analysis\nMarket Segmentation"},
        {"label": "Timing", "detail": f"Fielding: {date_str}"},
    ]


def _build_research_approach(prs, research_items):
    """Clone research approach slide (slide 26 pattern).

    Template has label+detail rows: Format, Participants, Analysis, Timing.
    Each row is two text shapes side by side (label left, detail right).
    """
    slide = _clone_slide(prs, T_RESEARCH_APPROACH)
    shapes = _find_text_shapes(slide)

    # Shape 0 = title
    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "QUANTITATIVE RESEARCH APPROACH")

    # Map research_items to the label+detail shape pairs
    pair_idx = 0
    for item in research_items[:5]:
        label_shape_idx = 2 + pair_idx * 2
        detail_shape_idx = label_shape_idx - 1
        # Template order: detail shape comes before label in position sort
        # Actual layout: shapes alternate detail(left-wide) and label(left-narrow)
        if detail_shape_idx < len(shapes) and label_shape_idx < len(shapes):
            _set_text_preserve_format(shapes[label_shape_idx].text_frame, item.get("label", ""))
            _set_text_preserve_format(shapes[detail_shape_idx].text_frame, _truncate(item.get("detail", ""), 350))
        pair_idx += 1

    return slide


def _normalize_mini_tables(segment: dict) -> dict:
    """Normalize mini_tables from list format to dict format."""
    mini_raw = segment.get("mini_tables", {})
    if isinstance(mini_raw, list):
        mini = {}
        for item in mini_raw:
            if isinstance(item, dict):
                key = item.get("label", item.get("name", f"table_{len(mini)}"))
                mini[key] = item.get("data", [])
        return mini
    return mini_raw if isinstance(mini_raw, dict) else {}


def _build_segment_overview(prs, segments, project_id: int = 0):
    """Clone segment overview slide (slide 49 pattern).

    Shows all segments at a glance: name, %, tagline for each.
    Template has 5 columns with: percentage text, image, name, tagline.
    Replaces template persona images with AI-generated ones matching each segment.
    """
    slide = _clone_slide(prs, T_SEGMENT_OVERVIEW)
    shapes = _find_text_shapes(slide)

    # Shape 0 = title
    if shapes:
        _set_text_preserve_format(shapes[0].text_frame, "CONSUMER SEGMENTS AT A GLANCE")

    # Find percentage, name, and tagline shapes — classify by position
    pct_shapes = []
    name_shapes = []
    tagline_shapes = []
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        text = s.text_frame.text.strip()
        if not text:
            continue
        # Title at top
        if "GLANCE" in text.upper() or "CONSUMER SEGMENTS" in text.upper() or "UNLOCKS" in text.upper():
            continue
        # Percentages: short text with %, near top (~y=1735234)
        if text.endswith("%") and len(text) <= 4:
            pct_shapes.append(s)
        # Taglines: below y=4600000 (descriptions)
        elif s.top > 4600000:
            tagline_shapes.append(s)
        # Names: between y=4000000 and y=4600000, short text
        elif s.top > 4000000 and len(text) < 30:
            name_shapes.append(s)

    # Sort ALL lists by left position to ensure column alignment
    pct_shapes.sort(key=lambda s: s.left)
    name_shapes.sort(key=lambda s: s.left)
    tagline_shapes.sort(key=lambda s: s.left)

    for i, seg in enumerate(segments[:5]):
        if i < len(pct_shapes):
            _set_text_preserve_format(pct_shapes[i].text_frame, f"{seg.get('size_pct', '?')}%")
        if i < len(name_shapes):
            _set_text_preserve_format(name_shapes[i].text_frame, seg.get("name", f"Segment {i+1}"))
        if i < len(tagline_shapes):
            _set_text_preserve_format(tagline_shapes[i].text_frame, _truncate(seg.get("tagline", ""), 80))

    # Shift taglines down to prevent overlap with bold segment names
    for ts in tagline_shapes:
        ts.top = max(ts.top, 4850000)  # Ensure minimum y below name block

    # Replace template persona images with AI-generated ones
    img_shapes = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left
    )
    for i, img_s in enumerate(img_shapes[:5]):
        if i >= len(segments):
            break
        seg = segments[i]
        persona_img = _generate_segment_persona_icon(seg, project_id, segment_index=i)
        if persona_img:
            try:
                import io
                img_bytes = persona_img.read_bytes()
                img_part, rId = slide.part.get_or_add_image_part(
                    io.BytesIO(img_bytes))
                ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                blips = list(img_s._element.iter(f"{ns_a}blip"))
                if blips:
                    for blip in blips:
                        blip.set(f"{ns_r}embed", rId)
                    print(f"[segment] Replaced persona image {i}: {persona_img.name} ({len(img_bytes):,} bytes, {len(blips)} blips)")
                else:
                    print(f"[segment] WARNING: No blip found in image shape {i} for {seg.get('name','')}")
            except Exception as e:
                print(f"[segment] Image replacement failed for {seg.get('name','')}: {e}")

    return slide


def _generate_segment_persona_icon(segment: dict, project_id: int,
                                    segment_index: int = 0) -> Path | None:
    """Generate a persona portrait for segment overview circles.

    Uses segment demographics (age, gender, role) to build an accurate prompt.
    Each segment gets a visually distinct person — different ethnicity, age,
    setting, and styling to ensure diversity across the 5 circles.
    """
    import re as _re
    name = segment.get("name", "")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}

    primary_role = demo.get("primary_role", "")
    age_skew = demo.get("age_skew", "")
    gender_split = demo.get("gender_split", "")

    # Derive age from demographics
    age_desc = "30-year-old"
    if "Gen Z" in age_skew or "18-24" in age_skew or "14-24" in age_skew:
        age_desc = "21-year-old"
    elif "Millennial" in age_skew or "25-34" in age_skew or "28-42" in age_skew:
        age_desc = "32-year-old"
    elif "Gen X" in age_skew or "40-55" in age_skew or "35-44" in age_skew:
        age_desc = "45-year-old"

    # Derive gender
    gender_desc = "woman"
    if gender_split:
        m = _re.search(r"(\d+)%\s*female", gender_split.lower())
        if m and int(m.group(1)) < 50:
            gender_desc = "man"
        elif "male" in gender_split.lower() and "female" not in gender_split.lower():
            gender_desc = "man"

    # Rotate ethnicities to ensure visual diversity across segments
    ethnicities = [
        "Caucasian",
        "African American",
        "Hispanic/Latina",
        "East Asian",
        "South Asian",
    ]
    ethnicity = ethnicities[segment_index % len(ethnicities)]

    # Unique visual context per segment — clothing, setting, props
    visual_contexts = {
        0: "wearing a crisp button-down shirt at a modern office desk, "
           "laptop and coffee visible, morning light from large windows",
        1: "in a colorful casual outfit at a trendy café, "
           "smartphone in hand, string lights in background",
        2: "wearing a university hoodie in a dorm room, "
           "posters on the wall behind, cozy afternoon light",
        3: "in a cozy knit sweater at a kitchen counter, "
           "kids' drawings on the fridge behind, warm home lighting",
        4: "wearing a linen shirt at a plant-filled apartment, "
           "reusable tote bag visible, soft natural daylight",
    }
    visual = visual_contexts.get(segment_index, "in casual clothes, neutral background")

    prompt = (
        f"Authentic candid photograph of a {age_desc} {ethnicity} {gender_desc}, "
        f"{primary_role.lower() if primary_role else 'everyday person'}. "
        f"{visual}. "
        f"Face centered in frame, direct eye contact, natural relaxed expression. "
        f"Shot on Canon EOS R5 85mm f/1.2, shallow depth of field. "
        f"Real skin texture with pores and subtle imperfections, no retouching. "
        f"National Geographic portrait style. No text, no watermarks."
    )

    try:
        from pipeline.image_gen import generate_image
        output_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
        output_dir.mkdir(parents=True, exist_ok=True)
        safe_name = "".join(c if c.isalnum() or c == "_" else "_" for c in name.lower())
        output_path = output_dir / f"persona_{safe_name}.png"

        # Check if already generated
        if output_path.exists():
            return output_path

        # Prefer FLUX 1.1 Pro (most photorealistic for portraits),
        # then gpt-image-1, then DALL-E 3
        result = generate_image(
            prompt, output_path=output_path,
            backend="flux", size="1024x1024"
        )
        if result and result.exists():
            return result
        result = generate_image(
            prompt, output_path=output_path,
            backend="gpt-image", size="1024x1024", quality="hd"
        )
        if result and result.exists():
            return result
        result = generate_image(
            prompt, output_path=output_path,
            backend="dalle", size="1024x1024", quality="hd"
        )
        if result and result.exists():
            return result
    except Exception as e:
        print(f"[ppt_generator] Persona icon generation failed for {name}: {e}")

    # Fallback: web search for stock-like portrait
    try:
        output_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
        output_dir.mkdir(parents=True, exist_ok=True)
        safe_name = "".join(c if c.isalnum() or c == "_" else "_" for c in name.lower())
        output_path = output_dir / f"persona_{safe_name}.png"
        if output_path.exists():
            return output_path
        search_q = f"{age_desc} {gender_desc} portrait headshot lifestyle photography"
        result = _web_search_persona_image(search_q, output_path)
        if result:
            return result
    except Exception as e:
        print(f"[ppt_generator] Persona web search fallback failed for {name}: {e}")

    return None


def _web_search_persona_image(query: str, output_path) -> "Path | None":
    """Search web for a stock-like portrait image as fallback."""
    try:
        from anthropic import Anthropic
        from config import ANTHROPIC_API_KEY
        import httpx

        client = Anthropic(api_key=ANTHROPIC_API_KEY)
        response = client.messages.create(
            model=MODEL_SONNET,
            max_tokens=1024,
            messages=[{
                "role": "user",
                "content": f"Find a high-quality stock portrait photo URL for: {query}. "
                           f"Return ONLY the direct image URL (ending in .jpg or .png), nothing else."
            }],
            tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 3}],
        )
        # Extract URL from response
        for block in response.content:
            if hasattr(block, "text") and block.text:
                import re
                urls = re.findall(r'https?://[^\s"\'<>]+\.(?:jpg|jpeg|png|webp)', block.text)
                if urls:
                    r = httpx.get(urls[0], timeout=15, follow_redirects=True)
                    if r.status_code == 200 and len(r.content) > 5000:
                        output_path.parent.mkdir(parents=True, exist_ok=True)
                        output_path.write_bytes(r.content)
                        print(f"[ppt_generator] Web search persona → {output_path.name}")
                        return output_path
    except Exception as e:
        print(f"[ppt_generator] Web search persona failed: {e}")
    return None


def _pick_segment_icon(segment: dict) -> str:
    """Pick a contextually relevant icon filename based on segment characteristics."""
    name = segment.get("name", "").lower()
    needs = " ".join(str(n).lower() for n in segment.get("top_needs", []))
    motivations = " ".join(str(m).lower() for m in segment.get("key_motivations", []))
    context = f"{name} {needs} {motivations}"

    # Keyword-to-icon mapping
    icon_rules = [
        (["gift", "gifter", "giving"], "heart_icon.png"),
        (["sustain", "eco", "green", "environment"], "leaf_icon.png"),
        (["perform", "athlete", "sport", "fitness", "active"], "dumbbell_icon.png"),
        (["value", "price", "budget", "deal", "save"], "dollar_icon.png"),
        (["style", "fashion", "design", "aesthetic", "vivid", "collect"], "palette_icon.png"),
        (["tech", "innovat", "smart"], "lightbulb_icon.png"),
        (["family", "parent", "kid", "child"], "users_icon.png"),
        (["premium", "luxury", "quality"], "gem_icon.png"),
        (["everyday", "daily", "routine", "move"], "compass_icon.png"),
        (["mindful", "wellness", "health"], "heart_icon.png"),
        (["social", "community", "share"], "share_icon.png"),
    ]
    for keywords, icon in icon_rules:
        if any(kw in context for kw in keywords):
            return icon
    return "target_icon.png"  # generic fallback


def _get_segment_background(segment: dict, project_id: int) -> Path | None:
    """Generate a dark cinematic background image for a segment intro slide.

    Creates a lifestyle scene matching the segment's real-world context.
    The image must be dark/moody (for white text overlay) and contextually
    specific — not generic stock photography.

    Priority: FLUX 1.1 Pro → gpt-image-1 → DALL-E 3 → web search.
    """
    name = segment.get("name", "")
    narrative = segment.get("narrative", "")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    primary_role = demo.get("primary_role", "")
    lifestyle = segment.get("lifestyle_signals", [])

    # Build segment-specific scene from persona data (dynamic, not hardcoded)
    _seg_cat = segment.get("_category", "consumer products")
    tagline = segment.get("tagline", "")
    lifestyle_str = ", ".join(str(l) for l in lifestyle[:3]) if lifestyle else ""

    # Use LLM to generate a contextual scene description if we have enough data
    scene = ""
    if primary_role or tagline:
        role_desc = primary_role or "person"
        context_parts = [f"a {role_desc}'s daily environment"]
        if tagline:
            context_parts.append(f"reflecting the persona: {tagline}")
        if lifestyle_str:
            context_parts.append(f"lifestyle cues: {lifestyle_str}")
        context_parts.append(f"featuring {_seg_cat}")
        scene = ", ".join(context_parts)
    else:
        scene = f"a person in their natural daily environment with {_seg_cat}"

    prompt = (
        f"Cinematic wide-angle photograph, dark moody lighting, {scene}. "
        f"Deep shadows with selective warm highlights, blue-hour color grade. "
        f"Shot on Arri Alexa with anamorphic lens, shallow depth of field. "
        f"Atmosphere: intimate, real, documentary-style. "
        f"The scene must feel lived-in and authentic — not staged or stock. "
        f"Very dark overall exposure (suitable for white text overlay). "
        f"No people in frame, no text, no logos, no watermarks. "
        f"Aspect ratio 16:9, landscape orientation."
    )

    try:
        from pipeline.image_gen import generate_image
        output_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
        output_dir.mkdir(parents=True, exist_ok=True)
        safe_name = "".join(c if c.isalnum() or c == "_" else "_" for c in name.lower())
        output_path = output_dir / f"segment_bg_{safe_name}.png"

        # Try FLUX first (most photorealistic), then gpt-image-1, then DALL-E 3
        for backend, size, quality in [
            ("flux", "1344x768", "standard"),
            ("gpt-image", "1536x1024", "hd"),
            ("dalle", "1792x1024", "hd"),
        ]:
            result = generate_image(
                prompt, output_path=output_path,
                backend=backend, size=size, quality=quality,
            )
            if result and result.exists() and result.stat().st_size > 10000:
                print(f"[segment_bg] {name} → {backend} ({result.stat().st_size:,} bytes)")
                return result
    except Exception as e:
        print(f"[segment_bg] Generation failed for {name}: {e}")

    return None


def _clean_llm_json(text: str) -> str:
    """Strip markdown fences and 'json' prefix from LLM output."""
    text = text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1] if "\n" in text else text[3:]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()
    if text.startswith("json"):
        text = text[4:].strip()
    return text


def _llm_generate_text(
    system_prompt: str,
    user_prompt: str,
    max_tokens: int = 1024,
    temperature: float = 0.7,
) -> str | None:
    """Generate text via LLM — Anthropic Claude primary, OpenAI GPT-4o fallback.

    Returns the generated text string, or None if all backends fail.
    """
    from config import ANTHROPIC_API_KEY, OPENAI_API_KEY

    # --- Primary: Claude Sonnet ---
    if ANTHROPIC_API_KEY:
        try:
            from anthropic import Anthropic as _Anthropic
            client = _Anthropic(api_key=ANTHROPIC_API_KEY)
            resp = client.messages.create(
                model=MODEL_SONNET,
                max_tokens=max_tokens,
                system=system_prompt,
                messages=[{"role": "user", "content": user_prompt}],
                temperature=temperature,
            )
            text = resp.content[0].text.strip()
            if text:
                return text
        except Exception as e:
            print(f"[llm] Anthropic failed: {e}")

    # --- Fallback: OpenAI GPT-4o ---
    if OPENAI_API_KEY:
        try:
            from openai import OpenAI as _OpenAI
            client = _OpenAI(api_key=OPENAI_API_KEY)
            resp = client.chat.completions.create(
                model="gpt-4o",
                max_tokens=max_tokens,
                temperature=temperature,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
            )
            text = resp.choices[0].message.content.strip()
            if text:
                return text
        except Exception as e:
            print(f"[llm] OpenAI failed: {e}")

    print("[llm] All backends failed — no API keys or API errors")
    return None


# --- Caches for LLM-generated content within a single run ---
_challenge_quotes_cache: dict[str, list[str]] = {}
_bubble_quotes_cache: dict[str, list[str]] = {}


def _llm_generate_challenge_quotes(segment: dict, need_rows: list[tuple]) -> list[str]:
    """Generate 15 natural verbatim survey quotes for Challenges slide.

    Uses segment context (demographics, pain points, lifestyle) to produce
    quotes that sound like real survey open-ended responses.
    """
    name = segment.get("name", "Segment")
    if name in _challenge_quotes_cache:
        return _challenge_quotes_cache[name]

    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    lifestyle = segment.get("lifestyle_signals", [])
    tagline = segment.get("tagline", "")
    narrative = segment.get("narrative", "")

    # Build pain points context
    pain_context = "\n".join(
        f"  - {item} ({pct})" for item, pct in need_rows[:9]
    ) if need_rows else "No specific pain points available"

    # Build lifestyle context
    ls_context = "\n".join(
        f"  - {(ls if isinstance(ls, str) else ls.get('detail', ''))}"
        for ls in lifestyle[:6]
    ) if lifestyle else ""

    _seg_cat = segment.get("_category", "consumer products")

    system = (
        "You are writing verbatim open-ended survey responses for a consumer research study on "
        f"{_seg_cat}. Each response must read exactly like a real person typed it into an online "
        "survey — imperfect grammar, casual tone, real frustration. Some are fragments. Some ramble. "
        "Match the demographic's voice: a Gen Z college student writes differently from a 40-year-old parent. "
        "Never start more than 2 quotes with the same word. No marketing language."
    )

    user = f"""Generate exactly 15 verbatim survey responses for this consumer segment about {_seg_cat}. Max 80 characters each.

SEGMENT: {name}
TAGLINE: {tagline}
DEMOGRAPHICS: age={demo.get('age_skew', '?')}, gender={demo.get('gender_split', '?')}, income={demo.get('income', '?')}
NARRATIVE: {narrative[:300]}

TOP PAIN POINTS (from survey data):
{pain_context}

LIFESTYLE SIGNALS:
{ls_context}

RULES:
- Each quote is 1 line, wrapped in double quotes
- These are responses to: "What frustrates you most about buying or using {_seg_cat} today?"
- Sound like real people — casual, honest, sometimes sarcastic or resigned
- Match this specific demographic's voice and vocabulary
- Reference specific {_seg_cat} frustrations from the pain points above
- Vary length: some short (5-8 words), some medium (12-18 words), some longer (20-25 words)
- No formulaic patterns — never repeat sentence structures
- Output ONLY the 15 quotes, one per line, no numbering, no extra text"""

    result = _llm_generate_text(system, user, max_tokens=800, temperature=0.8)
    if result:
        lines = [
            line.strip().strip("-").strip()
            for line in result.strip().split("\n")
            if line.strip() and '"' in line
        ]
        # Clean up: ensure each line is wrapped in quotes
        cleaned = []
        for line in lines:
            line = line.strip()
            if not line.startswith('"'):
                line = '"' + line
            if not line.endswith('"'):
                line = line + '"'
            cleaned.append(line)
        if len(cleaned) >= 8:  # accept if we got at least 8 usable quotes
            _challenge_quotes_cache[name] = cleaned[:15]
            return cleaned[:15]

    # Fallback: generate simple quotes from pain point data
    print(f"[llm] Challenge quotes fallback for {name}")
    fallback = []
    for item, pct in need_rows[:9]:
        fallback.append(f'"{item} is a real issue for me"')
    while len(fallback) < 15:
        fallback.append('"I wish brands would listen to what we actually need"')
    _challenge_quotes_cache[name] = fallback[:15]
    return fallback[:15]


def _llm_generate_bubble_quotes(segment: dict) -> list[str]:
    """Generate 6 customer-voice quotes for Closer Look 2 speech bubbles.

    These are positive/neutral brand perception quotes, not complaints.
    """
    name = segment.get("name", "Segment")
    if name in _bubble_quotes_cache:
        return _bubble_quotes_cache[name]

    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    mini = _normalize_mini_tables(segment)
    lifestyle = segment.get("lifestyle_signals", [])
    tagline = segment.get("tagline", "")
    wpm = segment.get("what_premium_means", "")
    channels = segment.get("channels", [])

    # Build data context
    data_points = []
    for label, rows in mini.items():
        if isinstance(rows, list):
            for r in rows[:2]:
                if isinstance(r, dict) and r.get("item") and r.get("pct"):
                    data_points.append(f"{r['item']} ({r['pct']}%)")
    data_context = ", ".join(data_points[:8]) if data_points else "general category data"

    _seg_cat = segment.get("_category", "consumer products")

    system = (
        "You are writing focus group snippets for a brand discovery deck. These go in speech bubbles "
        f"and must sound like real consumers talking about {_seg_cat} — the way people actually speak "
        "in a moderated group, not how marketers imagine they speak. Match the demographic's natural vocabulary."
    )

    user = f"""Generate exactly 6 short conversational quotes about {_seg_cat} for speech bubbles.

SEGMENT: {name}
TAGLINE: {tagline}
DEMOGRAPHICS: age={demo.get('age_skew', '?')}, gender={demo.get('gender_split', '?')}
CHANNELS: {', '.join(str(c) for c in channels[:4]) if channels else 'N/A'}
WHAT PREMIUM MEANS: {wpm[:150] if wpm else 'N/A'}
KEY DATA POINTS: {data_context}

RULES:
- Each quote is 1 line, wrapped in double quotes, max 80 characters
- Sound like a real person in a focus group — natural speech patterns, not copy
- Mix: 2 about what they look for in {_seg_cat}, 2 about frustrations, 2 about aspirations
- Match this specific demographic's voice (Gen Z vs millennial vs parent, etc.)
- Reference {_seg_cat} specifics — features, brands, use cases — not generic "products"
- Output ONLY the 6 quotes, one per line, no numbering"""

    result = _llm_generate_text(system, user, max_tokens=400, temperature=0.8)
    if result:
        lines = [
            line.strip().strip("-").strip()
            for line in result.strip().split("\n")
            if line.strip() and '"' in line
        ]
        cleaned = []
        for line in lines:
            line = line.strip()
            if not line.startswith('"'):
                line = '"' + line
            if not line.endswith('"'):
                line = line + '"'
            cleaned.append(line)
        if len(cleaned) >= 4:
            _bubble_quotes_cache[name] = cleaned[:6]
            return cleaned[:6]

    # Fallback
    print(f"[llm] Bubble quotes fallback for {name}")
    fallback = [
        f'"I look for quality that matches my lifestyle"',
        f'"It needs to work for my daily routine"',
        f'"I\'d pay more for something that really lasts"',
        f'"The right product just makes my day easier"',
        f'"I do my research before buying"',
        f'"Brand reputation matters to me"',
    ]
    _bubble_quotes_cache[name] = fallback
    return fallback


_closer_look_3_cache: dict[str, list[dict]] = {}
_unified_dimensions_cache: dict | None = None


def _collect_unified_closer_look_3(segments: list, brand_name: str = "", category: str = "") -> dict:
    """Generate 4 UNIFIED lifestyle dimensions for ALL segments (called once).

    Like CozyFit's social media / music / car / wishlist — same 4 questions,
    different answers per segment, enabling horizontal comparison.

    Returns: {
        "__dimensions__": ["Favorite Exercise", "Social Media Platform", ...],
        "__image_prompts__": ["...", "...", ...],
        "Daily Functional": [{"stat_pct": 67, "description": "..."}, ...],
        "Vivid Collector": [{"stat_pct": 72, "description": "..."}, ...],
    }
    """
    global _unified_dimensions_cache
    if _unified_dimensions_cache is not None:
        return _unified_dimensions_cache

    # Gather data from all segments for context
    all_segments_info = []
    for seg in segments[:5]:
        name = seg.get("name", "Segment")
        demo = seg.get("demographics", {})
        if not isinstance(demo, dict):
            demo = {}
        mini = _normalize_mini_tables(seg)
        lifestyle = seg.get("lifestyle_signals", [])
        channels = seg.get("channels", [])

        mini_summary = []
        for label, rows in mini.items():
            if isinstance(rows, list):
                items = [f"{r.get('item','?')} ({r.get('pct','?')}%)" for r in rows[:2] if isinstance(r, dict)]
                mini_summary.append(f"{label}: {', '.join(items)}")

        all_segments_info.append(
            f"SEGMENT: {name} ({seg.get('size_pct', '?')}%)\n"
            f"  Demographics: age={demo.get('age_skew','?')}, gender={demo.get('gender_split','?')}\n"
            f"  Tagline: {seg.get('tagline', '')}\n"
            f"  Channels: {', '.join(str(c) for c in channels[:3])}\n"
            f"  Survey data: {'; '.join(mini_summary[:3])}"
        )

    segments_context = "\n\n".join(all_segments_info)

    system = (
        "You are a senior brand strategist at a leading consultancy designing a consumer comparison matrix "
        f"for a {category or 'consumer'} brand discovery presentation. "
        "Choose 4 lifestyle dimensions that reveal the sharpest behavioral contrasts between segments — "
        "the kind of differences that inform brand positioning decisions. "
        "Each dimension should produce clearly distinct answers that tell a strategic story. "
        "Output valid JSON only, no markdown."
    )

    user = f"""Choose 4 lifestyle/behavioral dimensions for a Closer Look comparison across these consumer segments.

BRAND: {brand_name or 'Consumer brand'}
CATEGORY: {category or 'Consumer product'}

{segments_context}

RULES:
- Pick 4 dimensions where segments will show MEANINGFUL DIFFERENCES (for horizontal comparison)
- Good examples: "Favorite Activity", "Social Media Platform", "Shopping Trigger", "Must-Have Feature"
- Category-specific examples: "Weekend Activity", "Discovery Channel", "Purchase Driver", "Style Inspiration"
- Avoid generic dimensions like "Lifestyle" or "Preference" — be specific
- Each dimension should produce clearly different answers across segments
- For each segment, provide a percentage stat (50-90) and a short one-sentence insight (max 80 chars)
- image_prompt per dimension: a stock-photo scene representing the dimension (same image for all segments)

Output JSON:
{{
  "dimensions": ["Dim 1", "Dim 2", "Dim 3", "Dim 4"],
  "image_prompts": ["prompt1", "prompt2", "prompt3", "prompt4"],
  "segments": {{
    "{segments[0].get('name', 'Seg1')}": [{{"stat_pct": 67, "description": "prefer yoga and running"}}, ...],
    "{segments[1].get('name', 'Seg2') if len(segments) > 1 else 'Seg2'}": [...],
    ...
  }}
}}

JSON ONLY. No markdown. Keep descriptions under 80 characters each."""

    result = _llm_generate_text(system, user, max_tokens=1200, temperature=0.7)
    if result:
        import json
        try:
            clean = _clean_llm_json(result)
            data = json.loads(clean)
            dims = data.get("dimensions", [])
            prompts = data.get("image_prompts", [])
            segs = data.get("segments", {})

            if len(dims) >= 3 and segs:
                result_dict = {
                    "__dimensions__": dims[:4],
                    "__image_prompts__": prompts[:4],
                }
                for seg_name, cards in segs.items():
                    if isinstance(cards, list):
                        valid = []
                        for c in cards[:4]:
                            if isinstance(c, dict):
                                valid.append({
                                    "stat_pct": int(c.get("stat_pct", 55)),
                                    "description": str(c.get("description", "data varies"))[:80],
                                })
                        result_dict[seg_name] = valid

                _unified_dimensions_cache = result_dict
                print(f"[closer_look_3] Unified 4 dimensions: {dims[:4]}")
                return result_dict
        except (json.JSONDecodeError, ValueError) as e:
            print(f"[llm] Unified CL3 dimensions JSON parse error: {e}")

    # Fallback: generic dimensions
    print(f"[llm] Unified CL3 dimensions fallback")
    fallback = {
        "__dimensions__": ["Social Media Platform", "Exercise Preference", "Shopping Channel", "Must-Have Feature"],
        "__image_prompts__": [
            "Social media apps on smartphone screen, colorful modern aesthetic, editorial photography",
            "Person exercising outdoors, active lifestyle, editorial photography",
            "Person shopping online on laptop, modern home, lifestyle photography",
            "Premium product detail closeup, clean modern aesthetic, editorial photography",
        ],
    }
    for seg in segments[:5]:
        name = seg.get("name", "Segment")
        channels = seg.get("channels", [])
        fallback[name] = [
            {"stat_pct": 65, "description": f"most active on {channels[0] if channels else 'Instagram'}"},
            {"stat_pct": 55, "description": "prefer moderate-intensity activities"},
            {"stat_pct": 60, "description": "research extensively before purchasing"},
            {"stat_pct": 70, "description": "prioritize durability and quality"},
        ]
    _unified_dimensions_cache = fallback
    return fallback


def _get_closer_look_3_cards(segment: dict, unified_data: dict) -> list[dict]:
    """Get per-segment Closer Look 3 cards from unified dimensions data.

    Returns list of 4 dicts: [{dimension, stat_pct, description, image_prompt}, ...]
    """
    name = segment.get("name", "Segment")
    dims = unified_data.get("__dimensions__", [])
    prompts = unified_data.get("__image_prompts__", [])
    seg_cards = unified_data.get(name, [])

    result = []
    for i, dim in enumerate(dims[:4]):
        card = {
            "dimension": dim,
            "stat_pct": seg_cards[i]["stat_pct"] if i < len(seg_cards) else 55,
            "description": seg_cards[i]["description"] if i < len(seg_cards) else "data varies across this segment",
            "image_prompt": prompts[i] if i < len(prompts) else f"Professional lifestyle photography, {dim}, editorial quality",
        }
        result.append(card)
    return result


def _llm_enrich_narrative(segment: dict, meet: str, reality: str, shop: str, pain: str, want: str) -> list | None:
    """Use LLM to produce a denser, CozyFit-quality narrative from segment data.

    Returns list of section strings, or None if LLM fails.
    """
    name = segment.get("name", "Segment")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    mini = _normalize_mini_tables(segment)
    lifestyle = segment.get("lifestyle_signals", [])
    tagline = segment.get("tagline", "")
    wpm = segment.get("what_premium_means", "")
    size_pct = segment.get("size_pct", "")
    channels = segment.get("channels", [])

    # Gather all data for context
    mini_summary = []
    for label, rows in mini.items():
        if isinstance(rows, list):
            items = [f"{r.get('item','?')} ({r.get('pct','?')}%)" for r in rows[:4] if isinstance(r, dict)]
            mini_summary.append(f"{label}: {', '.join(items)}")
    ls_summary = "\n".join(f"- {(ls if isinstance(ls, str) else ls.get('detail',''))}" for ls in lifestyle[:5])

    _seg_cat = segment.get("_category", "consumer products")

    system = (
        "You are a senior brand strategist at a top branding consultancy writing a consumer segment profile "
        "for a brand discovery presentation. Your tone is authoritative and insight-driven — you state findings "
        "as conclusions, never as possibilities. You weave specific percentages and behavioral details into vivid, "
        "empathetic persona narratives. Open with a cinematic vignette that makes the reader see a real person. "
        "Write in plain text only — no markdown, no asterisks, no bold markers, no hashtags. "
        "Never use phrases like 'data reveals', 'research suggests', 'appears to', 'seems to', 'likely', or 'tends to'. "
        "State everything as established fact from your research."
    )

    user = f"""Write a 5-section segment profile narrative for a {_seg_cat} brand discovery deck.

SEGMENT: {name} ({size_pct}% of total audience)
TAGLINE: {tagline}
DEMOGRAPHICS: age={demo.get('age_skew','?')}, gender={demo.get('gender_split','?')}, income={demo.get('income','?')}, role={demo.get('primary_role','?')}
CHANNELS: {', '.join(str(c) for c in channels[:4]) if channels else 'N/A'}
WHAT PREMIUM MEANS: {wpm[:200] if wpm else 'N/A'}

SURVEY DATA:
{chr(10).join(mini_summary[:4])}

LIFESTYLE SIGNALS:
{ls_summary}

EXISTING DRAFT (improve and enrich — keep data points, elevate the language):
Meet: {meet[:200]}
Reality: {reality[:200]}
Shop: {shop[:200]}
Pain: {pain[:200]}
Want: {want[:200]}

FORMAT: Write exactly 5 sections, each on a new line. Start each with a plain-text label and colon:
Meet the {name}: [Open with a cinematic vignette — paint a specific person in a specific moment. Then ground with demographics and % data. 2-3 sentences.]
Their Reality: [Daily context, behavioral patterns, lifestyle. Weave in % naturally like "one in three..." or "over half...". 2-3 sentences.]
How They Shop: [Specific channels with %, brand preferences, spending patterns. Frame strategically — are they research-driven, impulse, value-conscious? 2-3 sentences.]
Pain Points: [Name the specific frustrations with %. State them as facts, not hypotheses. 2 sentences.]
What They Want: [Premium definition, wishlist, unmet needs with %. End with what "premium" means to this specific person. 2 sentences.]

TOTAL max 1100 characters. Be concise but data-dense. Every sentence must carry at least one specific data point (%, dollar amount, channel name, or behavioral fact). PLAIN TEXT ONLY."""

    result = _llm_generate_text(system, user, max_tokens=800, temperature=0.6)
    if result:
        lines = [line.strip() for line in result.strip().split("\n") if line.strip()]
        # Filter out empty lines and ensure we have sections
        sections = [l for l in lines if len(l) > 20]
        if len(sections) >= 3:
            # Enforce character budget — trim at sentence boundaries, not mid-word
            total = sum(len(s) for s in sections)
            if total > 1200:
                ratio = 1150 / total
                trimmed = []
                for s in sections:
                    target_len = int(len(s) * ratio)
                    if len(s) <= target_len:
                        trimmed.append(s)
                    else:
                        # Find the last sentence boundary before the target
                        cut = s[:target_len]
                        last_period = max(cut.rfind('. '), cut.rfind('.) '), cut.rfind('."'))
                        if last_period > len(s) * 0.4:
                            trimmed.append(s[:last_period + 1])
                        else:
                            trimmed.append(cut.rsplit(' ', 1)[0] + '.')
                sections = trimmed
            return sections
    return None


def _build_structured_narrative(segment: dict) -> list:
    """Build a CozyFit-style structured narrative with 5 dense sections.

    Each section must be 2-4 sentences with real data points, percentages,
    and specific details. Matches CozyFit reference density (slides 52-57).
    """
    import re as _re
    name = segment.get("name", "Segment")
    narrative = segment.get("narrative", "")
    tagline = segment.get("tagline", "")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    mini = _normalize_mini_tables(segment)
    lifestyle = segment.get("lifestyle_signals", [])
    wtp = segment.get("what_premium_means", "")
    size_pct = segment.get("size_pct", "")
    channels = segment.get("channels", [])

    sentences = [s.strip() for s in _re.split(r'(?<!\d)\.(?!\d)\s*', narrative) if s.strip()]

    age_skew = demo.get("age_skew", "")
    income = demo.get("income", "")
    role = demo.get("primary_role", "")
    gender = demo.get("gender_split", "")

    # --- Meet the [Name]: vivid persona + demographics (2-3 sentences) ---
    meet_parts = []
    if sentences:
        meet_parts.append(". ".join(sentences[:2]) + ".")
    # Add demographic anchor
    demo_bits = []
    if age_skew:
        demo_bits.append(age_skew.split(",")[0].strip())
    if gender:
        demo_bits.append(gender.split(",")[0].strip())
    if income:
        demo_bits.append(income.split("—")[0].strip() if "—" in income else income)
    if demo_bits:
        meet_parts.append(f"This segment is predominantly {', '.join(demo_bits)}.")
    if size_pct:
        meet_parts.append(f"At {size_pct}% of the total audience, they represent a {'dominant' if float(size_pct) > 25 else 'significant' if float(size_pct) > 15 else 'niche'} segment.")
    meet_text = " ".join(meet_parts) if meet_parts else f"Picture a typical {name} consumer."

    # --- Their Reality: daily life context with data (2-3 sentences) ---
    reality_parts = []
    for s in sentences[2:6]:
        if any(w in s.lower() for w in ("market", "segment", "professional", "category",
                                         "represent", "population", "worker", "daily",
                                         "routine", "lifestyle", "currently", "serves")):
            reality_parts.append(s.strip() + ".")
    if role:
        reality_parts.append(f"Their primary context is {role.lower()}.")
    # Add lifestyle context
    ls_details = []
    for ls in lifestyle[:3]:
        ls_text = ls if isinstance(ls, str) else ls.get("detail", "")
        if ls_text and "%" in ls_text:
            ls_details.append(ls_text)
    if ls_details:
        reality_parts.append(" ".join(ls_details[:2]))
    reality_text = " ".join(reality_parts[:3]) if reality_parts else ""

    # --- How They Shop: channels, premium definition, decision process ---
    shop_parts = []
    if channels:
        ch_list = [str(c) for c in channels[:3]]
        shop_parts.append(f"Primary channels: {', '.join(ch_list)}.")
    for sig in lifestyle:
        sig_text = sig if isinstance(sig, str) else sig.get("detail", "")
        if any(w in sig_text.lower() for w in ("amazon", "purchase", "buy", "shop",
                                                 "review", "price", "spend")):
            shop_parts.append(sig_text.split(": ", 1)[-1] if ": " in sig_text else sig_text)
    if wtp:
        shop_parts.append(f'"Premium" means {wtp}.')
    shop_text = " ".join(shop_parts[:3]) if shop_parts else ""

    # --- Pain Points: specific frustrations with percentages ---
    pain_data = []
    for k, v in mini.items():
        if any(w in k.lower() for w in ("pain", "switch", "trigger", "frustrat", "challenge")):
            pain_data = v
            break
    pain_text = ""
    if isinstance(pain_data, list) and pain_data:
        top_items = [f"{d.get('item', '')} ({d.get('pct', '')}%)" for d in pain_data[:5]
                     if isinstance(d, dict)]
        if top_items:
            pain_text = (f"Pain Points: {top_items[0]} leads the list. "
                        f"Add {', '.join(top_items[1:3])}, and the frustration picture is clear.")
            if len(top_items) > 3:
                pain_text += f" Also notable: {', '.join(top_items[3:])}."

    # --- What They Want: purchase drivers + unmet needs ---
    want_parts = []
    drivers = []
    for k, v in mini.items():
        if any(w in k.lower() for w in ("driver", "require", "functional", "need", "want")):
            drivers = v
            break
    if isinstance(drivers, list) and drivers:
        top = [f"{d.get('item', '')} ({d.get('pct', '')}%)" for d in drivers[:4]
               if isinstance(d, dict)]
        if top:
            want_parts.append(f"Their wishlist: {', '.join(top).lower()}.")
    unmet = segment.get("unmet_needs", "")
    if unmet:
        want_parts.append(f"What's missing: {unmet[:120]}.")
    want_text = " ".join(want_parts[:2]) if want_parts else ""

    # --- Try LLM enrichment for denser, CozyFit-quality narrative ---
    llm_narrative = _llm_enrich_narrative(segment, meet_text, reality_text, shop_text, pain_text, want_text)
    if llm_narrative:
        return llm_narrative

    # Fallback: assemble sections from template data
    sections = [f"Meet the {name}: {meet_text}"]
    if reality_text:
        sections.append(f"Their Reality: {reality_text}")
    if shop_text:
        sections.append(f"How They Shop: {shop_text}")
    if pain_text:
        sections.append(pain_text)
    if want_text:
        sections.append(f"What They Want: {want_text}")

    # Enforce total character budget to prevent text overflow
    # Meet page text box fits ~900 chars at 14pt, ~1100 at 12pt
    MAX_TOTAL = 950
    total = sum(len(s) for s in sections)
    if total > MAX_TOTAL:
        # Trim longest sections proportionally
        ratio = MAX_TOTAL / total
        sections = [s[:int(len(s) * ratio)] for s in sections]
        # Clean up: ensure no section ends mid-word
        for i, s in enumerate(sections):
            if s and s[-1] not in ".!? ":
                last_space = s.rfind(" ")
                if last_space > len(s) * 0.7:
                    sections[i] = s[:last_space] + "."

    return sections


def _set_narrative_sections(text_frame, sections: list):
    """Set multi-section narrative using separate paragraphs per section.

    Matches CozyFit reference slide 51/52: 16pt body text with bold section
    headers, each section as its own paragraph for proper spacing.
    Dynamically reduces font only if content truly overflows.
    """
    from lxml import etree
    from pptx.util import Pt, Emu
    from pptx.oxml.ns import qn
    import copy

    if not sections:
        return

    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = "  ".join(sections)
        return

    # Get format template from first run
    template_p = paras[0]
    template_runs = template_p.runs
    if not template_runs:
        text_frame.text = "  ".join(sections)
        return

    # Extract template run properties
    template_rPr = template_runs[0]._r.find(qn("a:rPr"))
    # Copy pPr (paragraph properties) for cloning
    template_pPr = template_p._p.find(qn("a:pPr"))

    # Get text box dimensions
    try:
        parent_sp = text_frame._txBody.getparent()
        cy = int(parent_sp.attrib.get('cy', '5003160'))
        box_h_in = cy / 914400
    except Exception:
        box_h_in = 5.49
    try:
        cx = int(parent_sp.attrib.get('cx', '11630400'))
        box_w_in = cx / 914400
    except Exception:
        box_w_in = 12.75

    # Determine font size: prefer 16pt, reduce only if overflows
    # Montserrat at 16pt: ~0.105 in/char; line height with 1.15 spacing
    for font_pt in (16, 15, 14, 13, 12):
        char_w = font_pt * 0.0065  # slightly tighter estimate for Montserrat
        cpl = int(box_w_in / char_w) if char_w > 0 else 100
        line_h = font_pt * 1.2 / 72  # line height with 1.15 spacing + padding
        est_lines = sum(max(1, (len(s) + cpl - 1) // cpl) for s in sections)
        needed_h = est_lines * line_h
        if needed_h <= box_h_in * 0.92:
            break

    font_sz = str(int(font_pt * 100))

    # Remove all existing paragraphs
    txBody = text_frame._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    # Build one paragraph per section: bold header + regular body
    for si, section in enumerate(sections):
        p_elem = etree.SubElement(txBody, qn("a:p"))

        # Clone paragraph properties from template
        if template_pPr is not None:
            new_pPr = copy.deepcopy(template_pPr)
            p_elem.insert(0, new_pPr)
            # Set line spacing to 115% for readable text
            lnSpc = new_pPr.find(qn("a:lnSpc"))
            if lnSpc is None:
                lnSpc = etree.SubElement(new_pPr, qn("a:lnSpc"))
            else:
                for child in list(lnSpc):
                    lnSpc.remove(child)
            spc_pct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spc_pct.set("val", "115000")  # 115%
            # Add small space before (except first section)
            if si > 0:
                spcBef = new_pPr.find(qn("a:spcBef"))
                if spcBef is None:
                    spcBef = etree.SubElement(new_pPr, qn("a:spcBef"))
                else:
                    for child in list(spcBef):
                        spcBef.remove(child)
                spc_pts = etree.SubElement(spcBef, qn("a:spcPts"))
                spc_pts.set("val", str(int(font_pt * 30)))  # ~0.5 line gap

        # Strip any markdown bold markers (**text** → text)
        import re as _re_md
        section = _re_md.sub(r'\*\*([^*]+)\*\*', r'\1', section)
        section = section.replace('**', '')  # catch unmatched

        # Split on first colon for header:body
        if ": " in section:
            header, body = section.split(": ", 1)
            header = header + ": "
        else:
            header = section
            body = ""

        # Bold header run
        h_r = etree.SubElement(p_elem, qn("a:r"))
        h_rPr = copy.deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(h_r, qn("a:rPr"))
        h_r.insert(0, h_rPr)
        h_rPr.set("sz", font_sz)
        h_rPr.set("b", "1")
        h_t = etree.SubElement(h_r, qn("a:t"))
        h_t.text = header

        # Regular body run
        if body:
            b_r = etree.SubElement(p_elem, qn("a:r"))
            b_rPr = copy.deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(b_r, qn("a:rPr"))
            b_r.insert(0, b_rPr)
            b_rPr.set("sz", font_sz)
            b_rPr.set("b", "0")
            b_t = etree.SubElement(b_r, qn("a:t"))
            b_t.text = body


def _build_meet_segment(prs, segment, project_id: int = 0):
    """Clone 'Meet the [Segment]' slide (slide 51 pattern).

    Full-bleed dark background image with white text overlay:
      - Segment name (ALL CAPS, bold)
      - Tagline (one line, italic)
      - Structured narrative (5 sections matching CozyFit reference)
    Background image is contextually matched to segment lifestyle.
    """
    slide = _clone_slide(prs, T_MEET_SEGMENT)
    shapes = _find_text_shapes(slide)

    name = segment.get("name", "SEGMENT")
    tagline = segment.get("tagline", "")

    # Build structured narrative (CozyFit 5-section format) as list of paragraphs
    narrative_sections = _build_structured_narrative(segment)

    # Replace background image with segment-specific scene
    bg_path = _get_segment_background(segment, project_id)
    if bg_path:
        try:
            import io as _io
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
            img_shapes = [s for s in slide.shapes if s.shape_type == 13]
            if img_shapes:
                # Replace largest image (the background)
                bg_shape = max(img_shapes, key=lambda s: s.width * s.height)
                img_part, rId = slide.part.get_or_add_image_part(
                    _io.BytesIO(bg_path.read_bytes()))
                for blip in bg_shape._element.iter(f"{ns_a}blip"):
                    blip.set(f"{ns_r}embed", rId)
                print(f"[meet_segment] Replaced background for {name}")
        except Exception as e:
            print(f"[meet_segment] Background replacement failed for {name}: {e}")

    # Fill text shapes
    title_filled = False
    tagline_filled = False
    narrative_filled = False

    # First pass: match shapes by content pattern
    for s in shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        if text.isupper() and len(text) < 30:
            _set_text_preserve_format(s.text_frame, name.upper())
            title_filled = True
        elif len(text) < 80 and not text.isupper() and "Meet" not in text:
            _set_text_preserve_format(s.text_frame, _truncate(tagline, 100))
            tagline_filled = True
        elif len(text) > 80 or "Meet" in text:
            _set_narrative_sections(s.text_frame, narrative_sections)
            narrative_filled = True

    # Second pass: fill empty shapes by size
    if not narrative_filled:
        empty_shapes = sorted(
            [s for s in shapes if not s.text_frame.text.strip()],
            key=lambda s: s.width * s.height, reverse=True
        )
        for s in empty_shapes:
            area = (s.width / 914400) * (s.height / 914400)
            if not narrative_filled and area > 20:
                _set_narrative_sections(s.text_frame, narrative_sections)
                narrative_filled = True
            elif not tagline_filled and 2 < area < 20:
                _set_text_preserve_format(s.text_frame, _truncate(tagline, 100))
                tagline_filled = True

    return slide


def _build_segment_closer_look(prs, segment, slide_num=1):
    """Build 'A Closer Look' slide using the real template layouts.

    Slide_num determines template and data:
      1: T_CLOSER_LOOK_1 (slide 53) — premium/driver callout text
      2: T_CLOSER_LOOK_2 (slide 54) — brand awareness + verbatim quotes
      3: T_CLOSER_LOOK_3 (slide 56) — 4 lifestyle signal cards
    """
    name = segment.get("name", "SEGMENT")
    mini = _normalize_mini_tables(segment)
    lifestyle = segment.get("lifestyle_signals", [])
    title = f"{name.upper()} – A CLOSER LOOK"

    if slide_num == 1:
        slide = _clone_slide(prs, T_CLOSER_LOOK_1)
        import re as _re

        # --- Build 6 insight callouts from segment data ---
        # Each callout = {stat: "79%", text: "bought for specific colorway..."}
        callouts = []

        # Pull insights from mini_tables — each table contributes 1-2 callouts
        mini_items = list(mini.items()) if mini else []
        for table_label, rows in mini_items:
            if not isinstance(rows, list) or not rows:
                continue
            top = rows[0]
            pct = top.get("pct", "")
            item_text = top.get("item", "")
            if pct and item_text:
                callouts.append({
                    "stat": f"{pct}%",
                    "text": f"{item_text.lower()}, the top factor for {table_label.lower()}",
                })
            # Second row as supporting stat if we need more
            if len(rows) > 1 and len(callouts) < 6:
                r2 = rows[1]
                if r2.get("pct") and r2.get("item"):
                    callouts.append({
                        "stat": f"{r2['pct']}%",
                        "text": f"{r2['item'].lower()}, also a key factor in {table_label.lower()}",
                    })

        # Add premium definition if available
        wpm = segment.get("what_premium_means", "")
        if wpm and len(callouts) < 6:
            # Parse "Limited / exclusive colorway (61%), ..."
            m = _re.search(r'([^(]+)\((\d+)%\)', wpm)
            if m:
                callouts.append({
                    "stat": f"{m.group(2)}%",
                    "text": f"say {m.group(1).strip().lower()} defines premium for this segment",
                })

        # Add lifestyle signal stats (only those with meaningful percentages)
        for ls in lifestyle:
            if len(callouts) >= 8:
                break
            ls_text = ls if isinstance(ls, str) else ls.get("detail", "")
            # Only match real stats: "74% ...", "4.2 bottles/year", etc.
            # Require % sign or contextual number (not random digits from brand names)
            m = _re.search(r'(\d+(?:\.\d+)?)%', ls_text)
            if m:
                stat_val = m.group(1)
                detail = ls_text.split(":", 1)[-1].strip() if ":" in ls_text else ls_text
                # Remove the stat from detail to avoid duplication
                detail = detail.replace(f"{stat_val}%", "", 1).strip().lstrip(",").strip()
                callouts.append({
                    "stat": f"{stat_val}%",
                    "text": detail,
                })

        # Pad to 8 if needed with demographics
        demo = segment.get("demographics", {})
        if isinstance(demo, dict):
            demo_labels = {
                "age_skew": "of this segment",
                "gender_split": "of this segment",
                "income": "household income range",
            }
            for dk, dv in demo.items():
                if len(callouts) >= 8:
                    break
                if isinstance(dv, str) and dv and dk in demo_labels:
                    m = _re.search(r'(\d+)%\s*(.+?)(?:,|$)', dv)
                    if m:
                        pct = m.group(1)
                        desc = m.group(2).strip()
                        suffix = demo_labels[dk]
                        callouts.append({
                            "stat": f"{pct}%",
                            "text": f"are {desc}, the largest group {suffix}",
                        })

        # Pad with remaining mini_table items until we have 8
        if len(callouts) < 8:
            for _label, rows in mini_items:
                if not isinstance(rows, list):
                    continue
                for item in rows:
                    if len(callouts) >= 8:
                        break
                    if isinstance(item, dict) and item.get("item"):
                        pct = item.get("pct", "")
                        item_text = str(item["item"])
                        # Skip items already used
                        if any(item_text.lower() in c["text"].lower() for c in callouts):
                            continue
                        callouts.append({
                            "stat": f"{pct}%" if pct else "",
                            "text": f"{item_text} — a key insight for this segment",
                        })

        # Hard limit to 8 callouts (4 left + 4 right)
        callouts = callouts[:8]

        # --- Map callouts to ALL 8 visual positions on slide 53 ---
        # Template structure (from shape dump):
        #   LEFT col (4 positions, top to bottom):
        #     [1] Group 49 (STAT: 1 stat + 1 text, y=1.15)
        #     [3] Group 48 (STAT: 2 stats + 1 text, y=2.59) — treat 2 stats as 2 callouts
        #     [7] Graphic 5 (standalone PIC, y=3.85) + [2] TextBox 20 (y=3.73)
        #     [11] Group 6 (ICON: 1 img + 1 text, y=5.43)
        #   RIGHT col (4 positions, top to bottom):
        #     [4] Group 44 (ICON: 1 img + 1 text, y=1.13)
        #     [8] Group 29 (ICON: 1 img + 1 text, y=2.59)
        #     [6] Group 7 (ICON: 1 img + 1 text, y=3.96)
        #     [9] TextBox 34 (standalone TEXT only, y=5.57) — no icon
        #
        # We process left col top→bottom then right col top→bottom = 8 callouts.

        # Update title first
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                text = s.text_frame.text.strip()
                if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
                    _set_text_preserve_format(s.text_frame, _truncate(title, 55))

        SLIDE_MID_X = 6096000  # ~6.67 inches

        # Classify all shapes
        groups = []
        standalone_pics = []  # PICTUREs not in any group
        standalone_txts = []  # TextBoxes not in any group (non-title, non-base)
        for s in slide.shapes:
            if s.shape_type == 6:
                stat_shapes, text_shapes, img_shapes = [], [], []
                for gs in s.shapes:
                    if gs.shape_type == 13:
                        img_shapes.append(gs)
                        continue
                    if not hasattr(gs, "text_frame"):
                        continue
                    runs = gs.text_frame.paragraphs[0].runs if gs.text_frame.paragraphs else []
                    font_size = runs[0].font.size if runs and runs[0].font.size else 0
                    if font_size >= Pt(40):
                        stat_shapes.append(gs)
                    else:
                        text_shapes.append(gs)
                groups.append({
                    "shape": s,
                    "stat_shapes": stat_shapes,
                    "text_shapes": text_shapes,
                    "img_shapes": img_shapes,
                    "col": 0 if s.left < SLIDE_MID_X else 1,
                    "top": s.top,
                })
            elif s.shape_type == 13:
                area_in = (s.width / 914400) * (s.height / 914400)
                if 0.5 < area_in < 3.0:
                    standalone_pics.append(s)
            elif hasattr(s, "text_frame"):
                text = s.text_frame.text.strip()
                if text and "CLOSER LOOK" not in text.upper() \
                        and "ENDURANCE" not in text.upper() \
                        and "base" not in text.lower() \
                        and len(text) > 10:
                    standalone_txts.append(s)

        # Sort groups by column then top
        groups.sort(key=lambda g: (g["col"], g["top"]))

        # Build ordered position list
        callout_idx = 0

        def _fill_stat(co, stat_shapes, text_shapes):
            """Fill a stat position: stat value in stat shape, description in text shape."""
            if stat_shapes:
                _set_text_preserve_format(stat_shapes[0].text_frame, co.get("stat", ""))
                # Shrink-to-fit: prevents LibreOffice from wrapping "84%" at 60pt.
                # PowerPoint still renders full size (text fits), but LibreOffice
                # auto-shrinks if its wider font metrics cause overflow.
                try:
                    from pptx.enum.text import MSO_AUTO_SIZE
                    stat_shapes[0].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
            if text_shapes:
                _set_text_preserve_format(text_shapes[0].text_frame,
                                          _truncate(co["text"], 120))

        def _fill_icon(co, img_shapes, text_shapes):
            """Fill an icon position: combined text + AI icon replacement."""
            combined = f"{co['stat']} {co['text']}"
            if text_shapes:
                _set_text_preserve_format(text_shapes[0].text_frame,
                                          _truncate(combined, 120))
            if img_shapes:
                _replace_group_icon_with_ai(slide, img_shapes[0], co["text"], name)

        # Process groups in sorted order
        for g in groups:
            stats = g["stat_shapes"]
            txts = g["text_shapes"]
            imgs = g["img_shapes"]

            if stats:
                # Stat group — use first stat, remove extras from XML (Group 48 has 2 stats)
                if callout_idx < len(callouts):
                    co = callouts[callout_idx]
                    callout_idx += 1
                    _fill_stat(co, stats, txts)
                    # Remove extra stat shapes from group XML to prevent
                    # LibreOffice rendering artifacts (empty 60pt text still takes space)
                    for extra_stat in stats[1:]:
                        try:
                            extra_stat._element.getparent().remove(extra_stat._element)
                        except Exception:
                            _set_text_preserve_format(extra_stat.text_frame, "")
                else:
                    for ss in stats:
                        _set_text_preserve_format(ss.text_frame, "")
                    for ts in txts:
                        _set_text_preserve_format(ts.text_frame, "")
            elif txts:
                # Icon group
                if callout_idx < len(callouts):
                    co = callouts[callout_idx]
                    callout_idx += 1
                    _fill_icon(co, imgs, txts)
                else:
                    for ts in txts:
                        _set_text_preserve_format(ts.text_frame, "")

        # Handle standalone PIC + TEXT pairs (left col row 3: Graphic 5 + TextBox 20)
        handled_txt_ids = set()
        for pic in standalone_pics:
            # Find nearest standalone text by vertical proximity
            best_txt = None
            best_dist = float("inf")
            for txt in standalone_txts:
                if id(txt) in handled_txt_ids:
                    continue
                dist = abs(pic.top - txt.top)
                if dist < best_dist:
                    best_dist = dist
                    best_txt = txt

            if callout_idx < len(callouts) and best_txt and best_dist < 2 * 914400:
                co = callouts[callout_idx]
                callout_idx += 1
                combined = f"{co['stat']} {co['text']}"
                _set_text_preserve_format(best_txt.text_frame,
                                          _truncate(combined, 120))
                _replace_group_icon_with_ai(slide, pic, co["text"], name)
                handled_txt_ids.add(id(best_txt))
            else:
                # Remove orphan picture
                pic._element.getparent().remove(pic._element)

        # Handle remaining standalone text boxes (right col row 4: TextBox 34)
        # These need a programmatically added icon to their left (template has no icon here)
        for txt in standalone_txts:
            if id(txt) in handled_txt_ids:
                continue
            if callout_idx < len(callouts):
                co = callouts[callout_idx]
                callout_idx += 1
                combined = f"{co['stat']} {co['text']}"
                _set_text_preserve_format(txt.text_frame,
                                          _truncate(combined, 120))
                handled_txt_ids.add(id(txt))
                # Add an icon to the left of this standalone text box
                _add_standalone_icon(slide, txt, co["text"], name)
            else:
                for p in txt.text_frame.paragraphs:
                    for r in p.runs:
                        r.text = ""

        return slide

    elif slide_num == 2:
        from pptx.chart.data import CategoryChartData

        slide = _clone_slide_with_charts(prs, T_CLOSER_LOOK_2)

        # --- Update native bar chart from first mini_table ---
        categories, values = [], []
        for _label, rows in mini.items():
            if not isinstance(rows, list) or not rows:
                continue
            for item in rows:
                if isinstance(item, dict) and item.get("pct") and item.get("item"):
                    categories.append(str(item["item"])[:28])
                    values.append(float(item["pct"]))
            if categories:
                break

        chart_shapes = [s for s in slide.shapes if hasattr(s, 'has_chart') and s.has_chart]
        for cs in chart_shapes:
            try:
                cd = CategoryChartData()
                cd.categories = categories
                cd.add_series("Series 1", [v / 100 for v in values])
                cs.chart.replace_data(cd)
                _set_chart_pct_format(cs.chart)
                # Shrink chart width to avoid overlapping speech bubbles
                # Template: 5.76" wide → shrink to 4.3" (right edge at ~4.77")
                from pptx.util import Emu
                cs.width = Emu(int(4.3 * 914400))
                print(f"[closer_look_2] Native bar chart updated for {name}")
            except Exception as e:
                print(f"[closer_look_2] Chart update failed: {e}")

        # --- Build customer voice quotes for speech bubbles via LLM ---
        quotes = _llm_generate_bubble_quotes(segment)

        # --- Update text shapes ---
        for s in slide.shapes:
            if not hasattr(s, "text_frame"):
                continue
            text = s.text_frame.text.strip()
            if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
                _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            elif "Brand Awareness" in text:
                # Use first mini_table label as section header
                first_label = next(iter(mini), "Key Insights")
                _set_text_preserve_format(s.text_frame, first_label.title())
            elif "Top 3" in text or "important features" in text.lower():
                # Short header for second mini_table (fits template's 0.37in box)
                mini_keys = list(mini.keys())
                label = mini_keys[1].title() if len(mini_keys) > 1 else "Key Factors"
                _set_text_preserve_format(s.text_frame, _truncate(label, 40))
            elif "comfort" in text.lower() or "All-day" in text:
                # Bullet items from second mini_table (no manual "•" — template has bullet formatting)
                mini_keys = list(mini.keys())
                second_rows = mini[mini_keys[1]] if len(mini_keys) > 1 else []
                if isinstance(second_rows, list) and second_rows:
                    detail = "\n".join(f"{r['item']} ({r['pct']}%)" for r in second_rows[:3]
                                       if isinstance(r, dict) and r.get("item"))
                    _set_text_preserve_format(s.text_frame, _truncate(detail, 120))
                else:
                    _set_text_preserve_format(s.text_frame, "")
            elif "anything else" in text.lower() or "experience" in text.lower():
                wpm = segment.get("what_premium_means", "")
                _set_text_preserve_format(s.text_frame,
                    _truncate(f"What defines premium: {wpm}" if wpm else "", 100))

        # --- Fill speech bubble shapes (AUTO_SHAPE type=1) ---
        bubble_shapes = [s for s in slide.shapes
                         if s.shape_type == 1 and hasattr(s, "text_frame")
                         and "Speech Bubble" in s.name]
        bubble_shapes.sort(key=lambda s: (s.top, s.left))
        for bi, bs in enumerate(bubble_shapes):
            if bi < len(quotes):
                _set_text_preserve_format(bs.text_frame, _truncate(quotes[bi], 90))
            else:
                _set_text_preserve_format(bs.text_frame, "")

        return slide

    else:
        slide = _clone_slide(prs, T_CLOSER_LOOK_3)
        shapes = _find_text_shapes(slide)

        # Use unified dimensions (same 4 for all segments, different answers)
        unified_cl3 = segment.get("_unified_cl3", {})
        card_data = _get_closer_look_3_cards(segment, unified_cl3)

        # Build card text: "67% prefer yoga and hiking..." format
        # Max ~90 chars per card to prevent overflow in narrow card columns
        cards_text = []
        for cd in card_data:
            desc = cd['description'][:75]
            cards_text.append(f"{cd['stat_pct']}% {desc}")

        for s in shapes:
            text = s.text_frame.text.strip()
            if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
                _set_text_preserve_format(s.text_frame, _truncate(title, 55))

        # Find the 4 bottom text boxes (top > 4500000, sorted by left)
        bottom_shapes = sorted(
            [s for s in shapes if s.top > 4500000 and s.text_frame.text.strip() and "base" not in s.text_frame.text.lower()],
            key=lambda s: s.left,
        )
        for idx, s in enumerate(bottom_shapes[:4]):
            if idx < len(cards_text) and cards_text[idx]:
                detail = _truncate(cards_text[idx], 150)
                _set_bold_pct_text(s.text_frame, detail)
            else:
                _set_text_preserve_format(s.text_frame, "")

        # Replace 4 card images with topic-relevant images from LLM prompts
        _replace_closer_look_3_images_v2(slide, card_data, name)

        return slide


def _replace_closer_look_3_images(slide, cards: list[str], segment_name: str):
    """Legacy — kept for compatibility. Use _replace_closer_look_3_images_v2 instead."""
    pass


def _replace_closer_look_3_images_v2(slide, card_data: list[dict], segment_name: str):
    """Replace the 4 card images on Closer Look 3 using LLM-provided image prompts.

    card_data: list of dicts with 'image_prompt' key from _llm_generate_closer_look_3_cards.
    """
    import io as _io
    import re as _re
    import tempfile
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000
    pic_shapes = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            area_pct = (shape.width * shape.height) / SLIDE_AREA
            if 0.03 < area_pct < 0.25:
                pic_shapes.append(shape)

    if not pic_shapes:
        return

    pic_shapes.sort(key=lambda s: s.left)

    from pipeline.image_gen import generate_image
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    for idx, pic in enumerate(pic_shapes[:4]):
        if idx >= len(card_data):
            continue
        cd = card_data[idx]
        prompt = cd.get("image_prompt", "Professional lifestyle photography, editorial quality")

        safe_name = _re.sub(r'[^a-z0-9]', '_', segment_name.lower())[:15]
        out_path = Path(tempfile.mkstemp(suffix=f"_cl3_{safe_name}_{idx}.png")[1])

        result = generate_image(prompt, output_path=out_path, backend="gpt-image",
                                size="1024x1024", quality="standard")
        if result and result.exists():
            try:
                img_part, rId = slide.part.get_or_add_image_part(
                    _io.BytesIO(result.read_bytes()))
                for blip in pic._element.iter(f"{ns_a}blip"):
                    blip.set(f"{ns_r}embed", rId)
                print(f"[closer_look_3] Replaced card {idx+1} image for {segment_name}: {cd.get('dimension', '')}")
            except Exception as e:
                print(f"[closer_look_3] Image replacement failed: {e}")
            finally:
                try:
                    out_path.unlink(missing_ok=True)
                except Exception:
                    pass
        else:
            print(f"[closer_look_3] Image generation failed for card {idx+1}")


def _build_segment_behavioral_summary(prs, segment):
    """Build a behavioral summary slide for a segment.

    Uses T_CONTENT template. Shows 3-4 mini_table categories as structured
    text with bold headers and data items. Uses websearch for real brand
    product images instead of AI-generated.
    """
    from lxml import etree
    from pptx.oxml.ns import qn
    import copy

    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)
    name = segment.get("name", "SEGMENT")
    mini = _normalize_mini_tables(segment)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame,
            _truncate(f"{name.upper()} – BEHAVIORAL SUMMARY", 55))

    # Build structured summary with bold headers using proper PPTX formatting
    mini_items = list(mini.items()) if mini else []
    seen_items = set()
    sections = []  # list of (header, items_text)

    for label, rows in mini_items[:4]:
        if not isinstance(rows, list) or not rows:
            continue
        header = label.upper()[:35]
        items = []
        for r in rows[:3]:
            if isinstance(r, dict) and r.get("item"):
                item_text = str(r["item"])
                if item_text.lower() in seen_items:
                    continue
                seen_items.add(item_text.lower())
                pct = r.get("pct", "")
                line = f"{item_text} ({pct}%)" if pct else str(item_text)
                items.append(line[:50])
        if items:
            sections.append((header, " · ".join(items)))

    # Add what_premium_means
    wpm = segment.get("what_premium_means", "")
    if wpm:
        sections.append(("WHAT PREMIUM MEANS", wpm[:100]))

    # Render into text frame with bold headers using XML
    if len(shapes) >= 2:
        tf = shapes[1].text_frame
        paras = list(tf.paragraphs)
        template_p = paras[0] if paras else None
        template_rPr = None
        template_pPr = None
        font_sz = "1600"  # 16pt for clear readability

        if template_p and template_p.runs:
            rpr = template_p.runs[0]._r.find(qn("a:rPr"))
            if rpr is not None:
                template_rPr = rpr
                sz = rpr.get("sz")
                if sz and int(sz) >= 1600:
                    font_sz = sz  # only inherit if >= 16pt
            template_pPr = template_p._p.find(qn("a:pPr"))

        txBody = tf._txBody
        for p in list(txBody.findall(qn("a:p"))):
            txBody.remove(p)

        for si, (header, body) in enumerate(sections):
            p_elem = etree.SubElement(txBody, qn("a:p"))

            if template_pPr is not None:
                new_pPr = copy.deepcopy(template_pPr)
                p_elem.insert(0, new_pPr)
                # Generous spacing between sections for breathing room
                if si > 0:
                    spcBef = new_pPr.find(qn("a:spcBef"))
                    if spcBef is None:
                        spcBef = etree.SubElement(new_pPr, qn("a:spcBef"))
                    else:
                        for child in list(spcBef):
                            spcBef.remove(child)
                    spc_pts = etree.SubElement(spcBef, qn("a:spcPts"))
                    spc_pts.set("val", "1400")  # ~14pt gap for clear separation

            # Bold header run
            h_r = etree.SubElement(p_elem, qn("a:r"))
            h_rPr = copy.deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(h_r, qn("a:rPr"))
            h_r.insert(0, h_rPr)
            h_rPr.set("sz", font_sz)
            h_rPr.set("b", "1")
            h_t = etree.SubElement(h_r, qn("a:t"))
            h_t.text = header + ": "

            # Regular body run
            b_r = etree.SubElement(p_elem, qn("a:r"))
            b_rPr = copy.deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(b_r, qn("a:rPr"))
            b_r.insert(0, b_rPr)
            b_rPr.set("sz", font_sz)
            b_rPr.set("b", "0")
            b_t = etree.SubElement(b_r, qn("a:t"))
            b_t.text = body

    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, "")

    # Replace template image with real brand product photo via websearch
    brand_name = segment.get("_brand_name", "")
    img_shapes = [s for s in slide.shapes if s.shape_type == 13]
    if img_shapes and brand_name:
        import io as _io
        import tempfile
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

        tmp_img = Path(tempfile.mkstemp(suffix="_behav_ws.png")[1])
        result = _websearch_brand_image(brand_name, "product collection", tmp_img)
        if result and result.exists():
            try:
                img_part, rId = slide.part.get_or_add_image_part(
                    _io.BytesIO(result.read_bytes()))
                for blip in img_shapes[0]._element.iter(f"{ns_a}blip"):
                    blip.set(f"{ns_r}embed", rId)
                print(f"[behavioral_summary] Replaced with brand image for {name}")
            except Exception as e:
                print(f"[behavioral_summary] Brand image replacement failed: {e}")
            finally:
                try:
                    tmp_img.unlink(missing_ok=True)
                except Exception:
                    pass
        else:
            # Fallback: generate product image
            from pipeline.image_gen import generate_image
            prompt = (
                f"Professional product photography of {brand_name} products, "
                f"clean modern setting, warm lighting, editorial quality, no text overlays."
            )
            result = generate_image(prompt, output_path=tmp_img, backend="gpt-image",
                                   size="1024x1024", quality="standard")
            if result and result.exists():
                try:
                    img_part, rId = slide.part.get_or_add_image_part(
                        _io.BytesIO(result.read_bytes()))
                    for blip in img_shapes[0]._element.iter(f"{ns_a}blip"):
                        blip.set(f"{ns_r}embed", rId)
                    print(f"[behavioral_summary] AI image for {name}")
                except Exception as e:
                    print(f"[behavioral_summary] Image replacement failed: {e}")
                finally:
                    try:
                        tmp_img.unlink(missing_ok=True)
                    except Exception:
                        pass
    elif img_shapes:
        # No brand name — remove template image
        for s in img_shapes:
            s._element.getparent().remove(s._element)

    return slide


def _build_segment_social_media(prs, segment):
    """Build a social media & lifestyle signals slide for a segment.

    Sample PPTs show 4 lifestyle cards with social media platform usage,
    music preferences, car brand affinities, and lifestyle values.
    Uses T_CLOSER_LOOK_3 (slide 56) template with 4 image slots + 4 text boxes.
    """
    slide = _clone_slide(prs, T_CLOSER_LOOK_3)
    shapes = _find_text_shapes(slide)
    name = segment.get("name", "SEGMENT")

    # Update title
    for s in shapes:
        text = s.text_frame.text.strip()
        if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
            _set_text_preserve_format(s.text_frame,
                _truncate(f"{name.upper()} – SOCIAL MEDIA & LIFESTYLE", 55))

    # Build exactly 4 lifestyle signal cards from segment data
    social = segment.get("social_media", [])
    music = segment.get("music_preferences", "")
    car_brand = segment.get("car_brand_affinities", "")
    lifestyle_vals = segment.get("lifestyle_values", [])
    channels = segment.get("channels", [])
    touchpoints = segment.get("media_touchpoints", [])

    # Each card: (label, icon_filename)
    cards = []
    if social:
        social_str = ", ".join(str(s) for s in social[:3]) if isinstance(social, list) else str(social)
        cards.append((f"Top platforms: {social_str}", "phone_icon.png"))
    if music:
        cards.append((f"Music: {music}", "mic_icon.png"))
    if car_brand:
        cards.append((f"Car style: {car_brand}", "compass_icon.png"))
    if lifestyle_vals:
        vals = ", ".join(str(v) for v in lifestyle_vals[:3]) if isinstance(lifestyle_vals, list) else str(lifestyle_vals)
        cards.append((f"Values: {vals}", "heart_icon.png"))

    # Fallback cards from lifestyle_signals
    if len(cards) < 4:
        lifestyle = segment.get("lifestyle_signals", [])
        icon_fallbacks = ["star_icon.png", "globe_icon.png", "target_icon.png", "sparkles_icon.png"]
        for ls in lifestyle:
            if len(cards) >= 4:
                break
            detail = ls.get("detail", str(ls)) if isinstance(ls, dict) else str(ls)
            if detail and not any(detail in c[0] for c in cards):
                cards.append((detail, icon_fallbacks[len(cards) % len(icon_fallbacks)]))

    # Further fallback from channels/touchpoints
    if len(cards) < 4 and channels:
        cards.append((f"Channels: {', '.join(str(c) for c in channels[:3])}", "shopping_cart_icon.png"))
    if len(cards) < 4 and touchpoints:
        cards.append((f"Media: {', '.join(str(t) for t in touchpoints[:3])}", "video_icon.png"))

    while len(cards) < 4:
        cards.append(("", ""))

    # Find bottom text boxes — include EMPTY ones too (the bug was filtering them out)
    bottom_shapes = sorted(
        [s for s in shapes if s.top > 4500000
         and "base" not in s.text_frame.text.lower()
         and ("CLOSER LOOK" not in s.text_frame.text.upper())],
        key=lambda s: s.left,
    )
    for idx, s in enumerate(bottom_shapes[:4]):
        if idx < len(cards) and cards[idx][0]:
            _set_text_preserve_format(s.text_frame, _truncate(cards[idx][0], 150))
        else:
            _set_text_preserve_format(s.text_frame, "")

    # Replace images with matching icons
    img_shapes = [s for s in slide.shapes if s.shape_type == 13]
    for idx, img_s in enumerate(sorted(img_shapes, key=lambda s: s.left)[:4]):
        icon_name = cards[idx][1] if idx < len(cards) else ""
        if icon_name:
            icon_path = ASSETS_DIR / icon_name
            if icon_path.exists():
                try:
                    import io
                    img_part, rId = slide.part.get_or_add_image_part(io.BytesIO(icon_path.read_bytes()))
                    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                    for blip in img_s._element.iter(f"{ns_a}blip"):
                        blip.set(f"{ns_r}embed", rId)
                except Exception:
                    pass

    return slide


def _derive_segment_demographics(segment: dict) -> dict:
    """Derive structured demographic breakdowns from segment data.

    Returns dict with: generation, gender, marital, income, ethnicity —
    each as {categories: [...], values: [...]}.
    Derives from segment demographics strings + adjusts for age/income profile.
    """
    import re as _re
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    age_skew = demo.get("age_skew", "")
    gender_split = demo.get("gender_split", "")
    income_desc = demo.get("income", "")

    result = {}

    # --- Generation breakdown ---
    gen_cats = ["Gen Z (18-28)", "Millennials (29-44)", "Gen X (45-60)", "Boomers (61+)"]
    gen_vals = [10, 35, 35, 20]  # defaults
    if age_skew:
        # Split by comma to parse each fragment independently
        gz = ml = gx = 0
        for part in age_skew.split(","):
            part = part.strip()
            pct_m = _re.search(r'(\d+)%', part)
            if not pct_m:
                continue
            pct = int(pct_m.group(1))
            if _re.search(r'Gen Z|18–|14–|younger', part, _re.IGNORECASE):
                gz = pct
            elif _re.search(r'Millennial|Mill', part, _re.IGNORECASE):
                ml = pct
            elif _re.search(r'Gen X|Boomer', part, _re.IGNORECASE):
                gx = pct
        bm = max(0, 100 - gz - ml - gx) if (gz + ml + gx) > 0 else 0
        if gz + ml + gx > 0:
            gen_vals = [gz, ml, gx, bm]
    result["generation"] = {"categories": gen_cats, "values": gen_vals}

    # --- Gender ---
    male_pct, female_pct = 35, 65
    if gender_split:
        fm = _re.search(r'(\d+)%\s*female', gender_split.lower())
        mm = _re.search(r'(\d+)%\s*male', gender_split.lower())
        if fm:
            female_pct = int(fm.group(1))
        if mm:
            male_pct = int(mm.group(1))
    result["gender"] = {"male_pct": male_pct, "female_pct": female_pct}

    # --- Marital status (derived from age profile) ---
    # Younger segments → more single; older → more married
    if gen_vals[0] > 40:  # Gen Z dominant
        married, single, divorced = 20, 68, 12
    elif gen_vals[2] > 30:  # Gen X dominant
        married, single, divorced = 62, 22, 16
    elif gen_vals[1] > 40:  # Millennial dominant
        married, single, divorced = 45, 42, 13
    else:
        married, single, divorced = 48, 38, 14
    result["marital"] = {"married": married, "single": single, "divorced": divorced}

    # --- Income (derived from income description) ---
    inc_cats = ["Low income", "Low middle", "Upper middle", "High income"]
    if "150k" in income_desc.lower() or "high" in income_desc.lower():
        inc_vals = [3, 12, 35, 50]
    elif "100k" in income_desc.lower() or "upper" in income_desc.lower():
        inc_vals = [5, 18, 42, 35]
    elif "55k" in income_desc.lower() or "middle" in income_desc.lower():
        inc_vals = [8, 28, 38, 26]
    elif "25k" in income_desc.lower() or "low" in income_desc.lower():
        inc_vals = [22, 35, 28, 15]
    else:
        inc_vals = [10, 25, 35, 30]
    result["income"] = {"categories": inc_cats, "values": inc_vals}

    # --- Ethnicity (use US census-adjusted for segment) ---
    eth_cats = ["White/Caucasian", "Black/African American", "Hispanic/Latino",
                "Asian", "Other"]
    eth_vals = [52, 18, 17, 8, 5]  # US census baseline
    result["ethnicity"] = {"categories": eth_cats, "values": eth_vals}

    return result


def _build_segment_profile(prs, segment):
    """Clone the respondent profile slide for a segment.

    Template slide 52 layout (matching CozyFit reference Image #85):
      - Title: "[SEGMENT] – RESPONDENT PROFILE"
      - Generation: horizontal bar chart (4 bars)
      - Gender: group with male/female icons + percentages
      - Marital Status: 3 icons (married/single/widowed) + percentages
      - Race/Ethnicity: pie chart
      - Household Income: 4 money icons + percentages
      - Footer: "Segment size: X% of audience"
    """
    import re as _re
    from pptx.chart.data import CategoryChartData

    slide = _clone_slide_with_charts(prs, T_SEGMENT_PROFILE)
    name = segment.get("name", "SEGMENT")
    size_pct = segment.get("size_pct", "")

    # Derive structured demographics for this segment
    seg_demo = _derive_segment_demographics(segment)
    gen_data = seg_demo["generation"]
    gender = seg_demo["gender"]
    marital = seg_demo["marital"]
    income = seg_demo["income"]
    eth = seg_demo["ethnicity"]

    shapes = _find_text_shapes(slide)

    # --- Update title and footer ---
    for s in shapes:
        text = s.text_frame.text.strip()
        if "RESPONDENT PROFILE" in text.upper() or "ENDURANCE" in text.upper():
            _set_text_preserve_format(s.text_frame, f"{name.upper()} – RESPONDENT PROFILE")
        elif "base" in text.lower() and "n" in text.lower():
            _set_text_preserve_format(s.text_frame,
                f"Segment size: {size_pct}% of audience" if size_pct else "")

    # --- Update Marital Status percentages ---
    marital_pcts = [
        (1.5, 2.5, str(marital["married"]) + "%"),
        (2.5, 3.8, str(marital["single"]) + "%"),
        (3.8, 5.0, str(marital["divorced"]) + "%"),
    ]
    for s in shapes:
        text = s.text_frame.text.strip()
        if _re.match(r'^\d{1,3}%$', text):
            y_in = s.top / 914400
            for y_min, y_max, new_val in marital_pcts:
                if y_min < y_in < y_max:
                    _set_text_preserve_format(s.text_frame, new_val)
                    break

    # --- Update Income percentages ---
    inc_pct_shapes = []
    for s in shapes:
        text = s.text_frame.text.strip()
        y_in = s.top / 914400
        if _re.match(r'^\d{1,3}%$', text) and y_in > 5.5:
            inc_pct_shapes.append(s)
    inc_pct_shapes.sort(key=lambda s: s.left)
    for i, s in enumerate(inc_pct_shapes[:4]):
        if i < len(income["values"]):
            _set_text_preserve_format(s.text_frame, f"{income['values'][i]}%")

    # --- Update Gender percentages in group shape ---
    for s in slide.shapes:
        if s.shape_type == 6:  # GROUP shape
            try:
                for gs in s.shapes:
                    if gs.has_text_frame:
                        t = gs.text_frame.text.strip()
                        if _re.match(r'^\d{1,3}%$', t):
                            if gs.left < s.width // 2:
                                _set_text_preserve_format(gs.text_frame,
                                    f"{gender['male_pct']}%")
                            else:
                                _set_text_preserve_format(gs.text_frame,
                                    f"{gender['female_pct']}%")
            except Exception:
                pass

    # --- Update native Generation bar chart (BAR_CLUSTERED) ---
    # Template has native charts preserved by _clone_slide_with_charts
    chart_shapes = [s for s in slide.shapes if hasattr(s, 'has_chart') and s.has_chart]
    for cs in chart_shapes:
        try:
            ct = cs.chart.chart_type
            if ct == 57:  # BAR_CLUSTERED — generation breakdown
                cd = CategoryChartData()
                cd.categories = gen_data["categories"]
                cd.add_series("Generation", [float(v) / 100 for v in gen_data["values"]])
                cs.chart.replace_data(cd)
                # Fix number format: show as "29%" not "0.29"
                _set_chart_pct_format(cs.chart)
                print(f"[segment_profile] Native bar chart updated for {name}")
            elif ct == 5:  # PIE — ethnicity breakdown
                cd = CategoryChartData()
                cd.categories = eth["categories"]
                cd.add_series("Ethnicity", [float(v) / 100 for v in eth["values"]])
                cs.chart.replace_data(cd)
                # Fix number format + set orange color scheme
                _set_chart_pct_format(cs.chart)
                _set_pie_orange_colors(cs.chart)
                print(f"[segment_profile] Native pie chart updated for {name}")
        except Exception as e:
            print(f"[segment_profile] Chart update failed: {e}")

    return slide


_unified_pain_cache: dict | None = None


def _collect_unified_pain_items(segments: list, category: str = "consumer products") -> dict:
    """Generate unified pain point items with per-segment percentages via LLM.

    CozyFit pattern: ALL segments share the SAME 8 pain point items,
    each segment has different percentages, sorted by their own ranking.

    Returns dict: {segment_name: [(item_text, pct_int), ...]} with 8 items per segment.
    Also includes "__items__" key with the ordered item list.
    """
    global _unified_pain_cache
    if _unified_pain_cache is not None:
        return _unified_pain_cache

    import json as _json

    # Gather context from all segments
    seg_context = []
    for seg in segments:
        name = seg.get("name", "?")
        mini = _normalize_mini_tables(seg)
        tagline = seg.get("tagline", "")
        items_summary = []
        for label, rows in mini.items():
            if isinstance(rows, list):
                for r in rows:
                    if isinstance(r, dict) and r.get("item"):
                        items_summary.append(f"{r['item']} ({r.get('pct', '?')}%)")
        seg_context.append(
            f"- {name} ({seg.get('size_pct', '?')}%): {tagline}\n"
            f"  Data: {', '.join(items_summary[:6])}"
        )

    system = (
        "You are a quantitative consumer researcher designing survey analysis for a brand discovery project. "
        f"Generate realistic survey pain point data specific to {category}. "
        "Pain items must be concrete product frustrations — the kind real consumers write in open-ended survey fields. "
        "Never use generic items like 'quality issues' or 'poor customer service'. "
        "Output valid JSON only, no markdown fences."
    )

    seg_names = [s.get("name", "?") for s in segments]
    user = f"""Generate unified pain point survey results for these {len(segments)} consumer segments in the {category.upper()} category.

SEGMENTS:
{chr(10).join(seg_context)}

TASK: Create exactly 8 pain point items that ALL segments were asked about (like a real survey).
Each pain item must be a specific, tangible frustration about {category} — things like sizing inconsistency, durability loss after washing, limited color options, poor lid seal, etc.
For each segment, assign a realistic percentage (10-55%) based on their persona and priorities.
Different segments should rank items differently based on their lifestyle and values.

Make pain items specific to {category} — a reader should immediately know what product category this survey covers.

Output JSON format:
{{
  "items": ["pain item 1", "pain item 2", ...],
  "segments": {{
    "{seg_names[0]}": [44, 36, 29, 27, 27, 27, 13, 10],
    "{seg_names[1]}": [41, 25, 25, 25, 22, 20, 20, 10],
    ...
  }}
}}

Each segment's array has 8 numbers (percentages) in the SAME ORDER as the items array.
Make percentages realistic and varied — not all similar. Higher % = bigger pain for that segment."""

    result = _llm_generate_text(system, user, max_tokens=600, temperature=0.6)
    if result:
        try:
            text = _clean_llm_json(result)
            data = _json.loads(text)
            items = data.get("items", [])
            seg_data = data.get("segments", {})

            if len(items) >= 6 and seg_data:
                cache = {"__items__": items}
                for seg in segments:
                    sname = seg.get("name", "?")
                    pcts = seg_data.get(sname, [])
                    if not pcts:
                        # Try fuzzy match
                        for k, v in seg_data.items():
                            if k.lower() in sname.lower() or sname.lower() in k.lower():
                                pcts = v
                                break
                    rows = []
                    for i, item in enumerate(items):
                        pct = int(pcts[i]) if i < len(pcts) else 15
                        rows.append((item, pct))
                    cache[sname] = rows
                _unified_pain_cache = cache
                print(f"[challenges] LLM generated {len(items)} unified pain items for {len(seg_data)} segments")
                return cache
        except Exception as e:
            print(f"[challenges] LLM pain items parse failed: {e}")

    # Fallback: use existing mini_table items
    print("[challenges] Using fallback pain items from mini_tables")
    from collections import Counter
    item_counter = Counter()
    for seg in segments:
        mini = _normalize_mini_tables(seg)
        for label, rows in mini.items():
            if not isinstance(rows, list):
                continue
            for item in rows:
                if isinstance(item, dict) and item.get("item"):
                    item_counter[str(item["item"])] += 1
    top_items = [item for item, _ in item_counter.most_common(8)]
    cache = {"__items__": top_items}
    for seg in segments:
        sname = seg.get("name", "?")
        mini = _normalize_mini_tables(seg)
        # Build lookup of this segment's percentages
        seg_pcts = {}
        for _label, rows in mini.items():
            if isinstance(rows, list):
                for r in rows:
                    if isinstance(r, dict) and r.get("item"):
                        seg_pcts[str(r["item"]).lower()] = int(r.get("pct", 0))
        rows = [(item, seg_pcts.get(item.lower(), 0)) for item in top_items]
        cache[sname] = rows
    _unified_pain_cache = cache
    return cache


def _build_segment_challenges(prs, segment, unified_pain_items=None):
    """Clone the challenges & pain points slide for a segment.

    Template idx 55 has two tables:
      - 16×1 verbatim quote table (right side, 12pt)
      - 8×2 pain-point + percentage table (left side, 11pt, green gradient fill)

    CozyFit reference slide 55: left panel shows top pain points with green
    heat-map fills, right panel shows customer-voice verbatim quotes.
    """
    slide = _clone_slide(prs, T_CHALLENGES)
    name = segment.get("name", "SEGMENT")
    mini = _normalize_mini_tables(segment)

    shapes = _find_text_shapes(slide)
    for s in shapes:
        text = s.text_frame.text.strip()
        if "CHALLENGES" in text.upper() or "ENDURANCE" in text.upper():
            _set_text_preserve_format(s.text_frame, f"{name.upper()} – KEY NEEDS & CHALLENGES")
        elif text == "Pain Points":
            _set_text_preserve_format(s.text_frame, "Top Needs")

    # --- LEFT TABLE: unified pain point items, same across ALL segments ---
    # unified_pain_items is a dict from _collect_unified_pain_items():
    #   {segment_name: [(item, pct), ...], "__items__": [...]}
    # CozyFit pattern: same items, different %, sorted by each segment's ranking.
    if isinstance(unified_pain_items, dict) and name in unified_pain_items:
        need_rows_raw = unified_pain_items[name]  # [(item, pct), ...]
    else:
        # Fallback: gather from this segment's mini_tables
        need_rows_raw = []
        seen_items = set()
        for label, rows in mini.items():
            if not isinstance(rows, list):
                continue
            for item in rows:
                if isinstance(item, dict) and item.get("item"):
                    item_text = str(item["item"])
                    if item_text.lower() in seen_items:
                        continue
                    seen_items.add(item_text.lower())
                    pct = int(item["pct"]) if item.get("pct") else 0
                    need_rows_raw.append((item_text[:50], pct))

    # Sort by percentage descending (CozyFit reference: highest at top)
    need_rows_raw_sorted = sorted(need_rows_raw, key=lambda x: x[1], reverse=True)
    need_rows = [(item, f"{pct}%" if pct else "") for item, pct in need_rows_raw_sorted[:9]]

    # --- RIGHT TABLE: LLM-generated verbatim quotes matching segment personality ---
    quote_lines = _llm_generate_challenge_quotes(segment, need_rows_raw_sorted[:9])

    # Fill tables — preserve template font formatting (12pt verbatim, 11pt needs)
    for sh in slide.shapes:
        if not (hasattr(sh, "has_table") and sh.has_table):
            continue
        t = sh.table
        cols = len(t.columns)
        rows_count = len(t.rows)
        if cols == 1 and rows_count > 4:
            # Right side: verbatim quotes table (16x1, 12pt)
            # Consistent question header across ALL segments
            _cat = segment.get("_category", "consumer products")
            header = f"What frustrates you most about buying or using {_cat} today?"
            # Same question for ALL segments — only the answers differ
            for ri in range(rows_count):
                cell = t.cell(ri, 0)
                new_text = ""
                if ri == 0:
                    new_text = header
                elif ri - 1 < len(quote_lines):
                    q = quote_lines[ri - 1]
                    # Truncate long quotes to prevent page overflow
                    if len(q) > 85:
                        # Trim inside the quotes: "text..."
                        inner = q.strip('"').strip()
                        inner = inner[:78] + '…'
                        q = f'"{inner}"'
                    new_text = q
                # Preserve font: set text via first run if available
                _set_cell_text_preserve_font(cell, new_text)
        elif cols == 2:
            # Left side: pain points with percentages (8x2, 11pt, green fill)
            for ri in range(rows_count):
                if ri < len(need_rows):
                    _set_cell_text_preserve_font(t.cell(ri, 0), need_rows[ri][0])
                    _set_cell_text_preserve_font(t.cell(ri, 1), need_rows[ri][1])
                else:
                    _set_cell_text_preserve_font(t.cell(ri, 0), "")
                    _set_cell_text_preserve_font(t.cell(ri, 1), "")

    return slide


def _set_cell_text_preserve_font(cell, new_text):
    """Set table cell text while preserving template font size, bold, and color."""
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        # Preserve formatting from first run
        run = tf.paragraphs[0].runs[0]
        run.text = new_text
        # Clear extra runs
        for extra_run in tf.paragraphs[0].runs[1:]:
            extra_run.text = ""
        # Clear extra paragraphs
        for p in list(tf.paragraphs)[1:]:
            p._p.getparent().remove(p._p)
    else:
        cell.text = new_text


def _build_brand_metrics_def(prs):
    """Clone brand metrics definitions boilerplate slide."""
    slide = _clone_slide(prs, T_CHART_TABLE)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "BRAND METRICS DEFINITIONS")

    definitions = (
        "Aided Awareness: % of consumers who recognize your brand when prompted.\n"
        "Purchase: % of consumers who have bought your brand's products within a specific timeframe.\n"
        "Satisfaction: % of purchasers who report being satisfied with the product.\n"
        "Recommendation: % of purchasers likely to recommend the brand to others."
    )

    _remove_chart_shapes(slide, clean_region=True)
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Emu(419100), Emu(1400000), Emu(11353800), Emu(4000000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in definitions.split("\n"):
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x29, 0x25, 0x24)

    return slide


def _build_segmentation_intro(prs):
    """Clone segmentation benefits boilerplate slide."""
    slide = _clone_slide(prs, T_SEGMENTATION_INTRO)
    return slide


def _build_focusing_segments(prs, segments, project_id: int = 0):
    """Clone 'FOCUSING ON THE MOST DOMINANT SEGMENTS…' slide (slide 50).

    Shows segments with percentages and taglines, highlighting the dominant ones.
    Template has: title + 5 percentage auto-shapes + 5 name text boxes + 5 taglines
    + 5 persona circle images + gray overlay on the last (smallest) segment.
    Replaces template images with the same AI-generated personas from overview slide.
    """
    slide = _clone_slide(prs, T_FOCUSING_SEGMENTS)

    # Categorize shapes by position
    pct_shapes = []
    name_shapes = []
    tag_shapes = []

    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        text = s.text_frame.text.strip()
        if "FOCUSING" in text.upper() or "DOMINANT" in text.upper():
            continue  # title — keep as-is
        elif text.endswith("%") and len(text) <= 4:
            pct_shapes.append(s)
        elif s.top > 4600000:
            tag_shapes.append(s)
        elif s.top > 4000000 and len(text) < 30:
            name_shapes.append(s)

    # Sort by left position
    pct_shapes.sort(key=lambda s: s.left)
    name_shapes.sort(key=lambda s: s.left)
    tag_shapes.sort(key=lambda s: s.left)

    for i, seg in enumerate(segments[:5]):
        if i < len(pct_shapes):
            _set_text_preserve_format(pct_shapes[i].text_frame, f"{seg.get('size_pct', '?')}%")
        if i < len(name_shapes):
            _set_text_preserve_format(name_shapes[i].text_frame, seg.get("name", f"Segment {i+1}"))
        if i < len(tag_shapes):
            _set_text_preserve_format(tag_shapes[i].text_frame, _truncate(seg.get("tagline", ""), 70))

    # Clear excess shapes
    for shapes_list in [pct_shapes, name_shapes, tag_shapes]:
        for i in range(len(segments), len(shapes_list)):
            _set_text_preserve_format(shapes_list[i].text_frame, "")

    # Shift taglines down to prevent overlap with bold segment names
    for ts in tag_shapes:
        ts.top = max(ts.top, 4850000)

    # Replace persona images (same as overview slide — reuse cached images)
    img_shapes = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left
    )
    for i, img_s in enumerate(img_shapes[:5]):
        if i >= len(segments):
            break
        seg = segments[i]
        name = seg.get("name", "")
        safe_name = "".join(c if c.isalnum() or c == "_" else "_" for c in name.lower())
        persona_path = OUTPUT_DIR / f"project_{project_id}" / "images" / f"persona_{safe_name}.png"
        if persona_path.exists():
            try:
                import io
                img_part, rId = slide.part.get_or_add_image_part(
                    io.BytesIO(persona_path.read_bytes()))
                ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                for blip in img_s._element.iter(f"{ns_a}blip"):
                    blip.set(f"{ns_r}embed", rId)
            except Exception:
                pass

    return slide


def _build_why_not_segments(prs, deprioritized, brand_name, category: str = "consumer products"):
    """Build 'WHY NOT PRIORITIZE OTHER SEGMENTS (FOR NOW)' slide.

    Uses T_CONTENT template with structured deprioritization rationale.
    """
    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "WHY NOT PRIORITIZE OTHER SEGMENTS (FOR NOW)")

    bullets = []
    for dep in deprioritized[:3]:
        name = dep.get("name", "Segment")
        reason = dep.get("reason", "Not the right fit for now")
        bullets.append(f"{name} ({dep.get('size_pct', '?')}%): {_truncate(reason, 120)}")

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, bullets or ["All segments show potential"])
    if len(shapes) >= 3:
        closing = f"Building long-term authority requires focus. {brand_name} must anchor the brand before expanding."
        _set_text_preserve_format(shapes[2].text_frame, _truncate(closing, 85))

    # Replace template image with persona-style image (matching target section slides)
    import tempfile as _tmpf
    from pipeline.image_gen import generate_image as _gen_img
    _tmp = Path(_tmpf.mkstemp(suffix="_whynot.png")[1])
    _result = _gen_img(
        f"Professional lifestyle photograph of a diverse group of people in a casual "
        f"urban setting, each with a different style {category}. Selective focus on "
        f"one person contemplating their choice. Warm editorial photography, natural lighting. "
        f"No text, no logos, no watermarks. Landscape orientation 16:9.",
        output_path=_tmp, backend="gpt-image", size="1536x1024", quality="high"
    )
    if _result and _result.exists():
        _replace_slide_image(slide, _result)
        print(f"[target] AI-generated persona image for WHY NOT slide")
    else:
        _replace_slide_image_websearch(slide, brand_name, f"{category} lifestyle photography")

    return slide


def _build_competitive_fares(prs, fares_data, brand_name, category: str = "consumer products"):
    """Build 'HOW [BRAND] FARES AGAINST THE COMPETITION' slide.

    Shows competitive positioning: what each brand wins on, the compromise forced,
    and the strategic question.
    """
    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame,
            _truncate(f"HOW {brand_name.upper()} FARES AGAINST THE COMPETITION", 55))

    brand_strengths = fares_data.get("brand_strengths", "")
    compromise = fares_data.get("category_compromise", "")
    opportunity = fares_data.get("strategic_opportunity", "")

    # Ensure string values (DB may store as lists)
    if isinstance(brand_strengths, list):
        brand_strengths = "; ".join(str(s) for s in brand_strengths)
    if isinstance(compromise, list):
        compromise = "; ".join(str(s) for s in compromise)
    if isinstance(opportunity, list):
        opportunity = "; ".join(str(s) for s in opportunity)

    bullets = []
    if brand_strengths:
        bullets.append(_truncate(str(brand_strengths), 85))
    if compromise:
        bullets.append(_truncate(str(compromise), 85))
    if opportunity:
        bullets.append(_truncate(str(opportunity), 85))

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, bullets or ["Competitive analysis in progress"])
    if len(shapes) >= 3:
        question = fares_data.get("strategic_question",
            f"What would it look like to build a brand that doesn't force that compromise?")
        _set_text_preserve_format(shapes[2].text_frame, _truncate(question, 85))

    # Replace template image with persona-style competitive image (matching target section)
    import tempfile as _tmpf2
    from pipeline.image_gen import generate_image as _gen_img2
    _tmp2 = Path(_tmpf2.mkstemp(suffix="_fares.png")[1])
    _result2 = _gen_img2(
        f"Professional lifestyle photograph of a young person {category} shopping, "
        f"comparing different options on a store shelf or browsing online. "
        f"Thoughtful expression, modern retail or home environment. "
        f"Warm editorial photography, natural lighting. "
        f"No text, no logos, no watermarks. Landscape orientation 16:9.",
        output_path=_tmp2, backend="gpt-image", size="1536x1024", quality="high"
    )
    if _result2 and _result2.exists():
        _replace_slide_image(slide, _result2)
        print(f"[target] AI-generated persona image for FARES slide")
    else:
        _replace_slide_image_websearch(slide, brand_name, f"{category} lifestyle photography")

    return slide


def _llm_generate_target_section(segments: list, brand_name: str, category: str = "") -> dict:
    """Use LLM to generate the entire target recommendation section.

    Selects the primary target segment and generates all supporting content:
    - Comparison table metrics
    - Primary target rationale bullets
    - Why this segment reasoning
    - Enables / does not decide
    - Consumer summary
    - Final summary columns
    """
    # Build segment context
    seg_summaries = []
    for s in segments[:5]:
        demo = s.get("demographics", {})
        if not isinstance(demo, dict):
            demo = {}
        mini = _normalize_mini_tables(s)
        top_items = []
        for label, rows in mini.items():
            if isinstance(rows, list):
                for r in rows[:2]:
                    if isinstance(r, dict) and r.get("item") and r.get("pct"):
                        top_items.append(f"{r['item']} ({r['pct']}%)")
        seg_summaries.append(
            f"SEGMENT: {s.get('name')} ({s.get('size_pct', '?')}%)\n"
            f"  Tagline: {s.get('tagline', '')}\n"
            f"  Demographics: age={demo.get('age_skew','?')}, gender={demo.get('gender_split','?')}, income={demo.get('income','?')}\n"
            f"  What premium means: {s.get('what_premium_means','')[:100]}\n"
            f"  Key data: {', '.join(top_items[:5])}"
        )

    system = (
        "You are a senior brand strategist at a top consultancy presenting your target audience recommendation "
        f"to {brand_name}'s leadership team. Your language is authoritative and strategic — you state conclusions, "
        "not suggestions. Never use hedge words like 'likely', 'potentially', 'could', 'may', 'seems', or 'appears to'. "
        "Every statement is backed by research. Frame the primary target as a strategic decision that unlocks "
        "positioning, not just the biggest segment. Output valid JSON only, no markdown."
    )

    user = f"""Analyze these consumer segments for {brand_name} ({category or 'consumer products'}) and select the PRIMARY TARGET.

{chr(10).join(seg_summaries)}

Generate a complete target recommendation section. Output JSON:
{{
  "primary_segment": "segment name",
  "selecting_title": "SELECTING {brand_name.upper()}'S TARGET AUDIENCE – [PRIMARY SEGMENT NAME]",
  "comparison_rows": [
    {{"metric": "Primary Use Case", "values": ["value for seg1", "value for seg2", ...]}},
    {{"metric": "Age Profile", "values": [...]}},
    {{"metric": "Buying Motivation", "values": [...]}},
    {{"metric": "% willing to pay premium", "values": ["62%", "48%", ...]}},
    {{"metric": "Average annual spend", "values": ["$180", "$95", ...]}},
    {{"metric": "Brand loyalty indicator", "values": [...]}}
  ],
  "primary_title": "PRIMARY TARGET: [SEGMENT NAME IN CAPS]",
  "rationale_bullets": [
    "Description of who they are and their core behavior",
    "View [product] as essential [function], not [alternative framing]",
    "Prioritize [key needs] through [context]",
    "Represent [X]% of the market and the highest [metric] segment"
  ],
  "target_insight": "This segment [key strategic importance statement].",
  "why_title": "WHY [SEGMENT] IS THE RIGHT FOCUS",
  "why_bullets": [
    "Defines quality: [reasoning about market trust]",
    "Highest value: [spending power and willingness]",
    "Strong product fit: [alignment with brand strengths]",
    "Natural channel fit: [distribution alignment]"
  ],
  "why_insight": "For this segment, 'premium' means [specific definition] — not [alternative].",
  "enables": [
    "A clear decision filter for [specific area]",
    "A credible path to [growth strategy]",
    "Natural spillover to adjacent segments who value [shared attributes]"
  ],
  "does_not_decide": [
    "Final brand positioning or tone",
    "The future role of [sub-brand] versus a new brand",
    "Pricing architecture or promotional strategy"
  ],
  "enables_insight": "Research shows this segment values [key insight] — creating a clear path to [outcome].",
  "consumer_summary": "[Segment name] [key behaviors and market influence — one powerful sentence].",
  "summary_capabilities": "[Brand] is [current state], now facing the need to [strategic challenge].",
  "summary_competition": "The [category] market is [state], with leading brands succeeding by [strategy].",
  "summary_consumer": "[Target segment] [value proposition — why they matter most].",
  "summary_closing": "Building on these insights, we will define a clear and differentiated brand position for {brand_name} — one that resonates with its most demanding customers and scales credibly across the broader market."
}}

RULES:
- Choose the segment with the strongest combination of: market size, spending power, brand alignment, and strategic value
- All content must be specific to {brand_name} and {category or 'this product category'} — no generic text
- Keep bullets concise (under 90 chars each)
- comparison_rows values must follow the SAME order as segments listed above
- NEVER use hedge words: 'likely', 'potentially', 'could be', 'may', 'seems to', 'appears to', 'tends to'
- State everything as an authoritative conclusion from completed research
- rationale_bullets should read like a strategist presenting to a CEO, not an analyst hedging
- consumer_summary must be a single powerful sentence that makes the strategic case undeniable
- JSON ONLY, no markdown"""

    result = _llm_generate_text(system, user, max_tokens=2000, temperature=0.6)
    if result:
        import json
        try:
            clean = _clean_llm_json(result)
            data = json.loads(clean)
            if isinstance(data, dict) and data.get("primary_segment"):
                print(f"[target] LLM selected primary target: {data['primary_segment']}")
                return data
        except (json.JSONDecodeError, ValueError) as e:
            print(f"[target] JSON parse error: {e}")

    # Fallback: pick largest segment
    print("[target] Using fallback target selection")
    primary = max(segments[:5], key=lambda s: float(s.get("size_pct", 0)))
    pname = primary.get("name", "Primary Segment")
    return {
        "primary_segment": pname,
        "selecting_title": f"SELECTING {brand_name.upper()}'S TARGET AUDIENCE – {pname.upper()}",
        "primary_title": f"PRIMARY TARGET: {pname.upper()}",
        "rationale_bullets": [
            f"Largest segment at {primary.get('size_pct','')}% of the market",
            f"Strong alignment with {brand_name}'s core product strengths",
            f"Highest engagement with the product category",
            f"Clear path to sustainable brand growth",
        ],
        "target_insight": f"This segment defines the standard for the category.",
        "why_title": f"WHY {pname.upper()} IS THE RIGHT FOCUS",
        "why_bullets": [
            f"Defines quality: Their standards set expectations across the market",
            f"Highest value: Spend more and are willing to pay for proven performance",
            f"Strong product fit: Needs align with {brand_name}'s strengths",
            f"Natural channel fit: Already engage through {brand_name}'s key channels",
        ],
        "why_insight": f"For this segment, premium means proof of performance — not prestige.",
        "enables": [
            "A clear decision filter for product quality standards",
            f"A credible path to brand elevation for {brand_name}",
            "Natural spillover to adjacent segments",
        ],
        "does_not_decide": [
            "Final brand positioning or tone",
            "Pricing architecture or promotional strategy",
            "Channel expansion timeline",
        ],
        "enables_insight": "Research shows this segment values proven performance — creating a clear path to credibility.",
        "consumer_summary": f"{pname} consumers define what quality means in this category — making them the most valuable and influential segment.",
        "summary_capabilities": f"{brand_name} has strong product fundamentals and growing market presence.",
        "summary_competition": f"The {category or 'product'} market rewards brands with clear positioning.",
        "summary_consumer": f"{pname} consumers are the most valuable segment in the market.",
        "summary_closing": f"Building on these insights, we will define a clear brand position for {brand_name}.",
        "comparison_rows": [],
    }


def _build_selecting_target(prs, segments: list, target_data: dict, brand_name: str):
    """Build SELECTING TARGET AUDIENCE comparison table slide (CozyFit slide 75).

    Clones template slide 75 with 5-col x 8-row table. First col = metrics,
    cols 2-5 = top 4 segments. Primary segment column highlighted with red border.
    """
    slide = _clone_slide(prs, T_SELECTING_TARGET)
    shapes = _find_text_shapes(slide)

    primary_name = target_data.get("primary_segment", "")
    title = target_data.get("selecting_title",
        f"SELECTING {brand_name.upper()}'S TARGET AUDIENCE – {primary_name.upper()}")

    # Set title
    for s in shapes:
        txt = s.text_frame.text.strip()
        if "SELECTING" in txt.upper() or "TARGET" in txt.upper():
            _set_text_preserve_format(s.text_frame, _truncate(title, 70))
            break

    # Fill comparison table
    for sh in slide.shapes:
        if not (hasattr(sh, "has_table") and sh.has_table):
            continue
        t = sh.table
        cols = len(t.columns)
        rows_count = len(t.rows)

        # Row 0: segment names (cols 1-4)
        top_segments = sorted(segments[:5], key=lambda s: float(s.get("size_pct", 0)), reverse=True)[:4]
        for ci in range(1, min(cols, len(top_segments) + 1)):
            seg = top_segments[ci - 1]
            seg_name = seg.get("name", f"Segment {ci}")
            _set_cell_text_preserve_font(t.cell(0, ci), seg_name)

        # Fill metrics from LLM-generated comparison_rows
        comp_rows = target_data.get("comparison_rows", [])
        # Build fallback from segment data if LLM didn't provide
        if not comp_rows:
            comp_rows = _build_comparison_rows_fallback(top_segments, brand_name)

        for ri in range(1, rows_count):
            if ri - 1 < len(comp_rows):
                cr = comp_rows[ri - 1]
                _set_cell_text_preserve_font(t.cell(ri, 0), str(cr.get("metric", ""))[:50])
                values = cr.get("values", [])
                for ci in range(1, min(cols, len(values) + 1)):
                    _set_cell_text_preserve_font(t.cell(ri, ci), str(values[ci - 1])[:40])
            else:
                for ci in range(cols):
                    _set_cell_text_preserve_font(t.cell(ri, ci), "")

    print(f"[target] Built selecting target table: {primary_name}")
    return slide


def _build_comparison_rows_fallback(segments: list, brand_name: str) -> list[dict]:
    """Build comparison table rows from segment data when LLM doesn't provide them."""
    rows = []
    # Row 1: Age profile
    row_vals = []
    for s in segments[:4]:
        demo = s.get("demographics", {})
        if isinstance(demo, dict):
            row_vals.append(str(demo.get("age_skew", "N/A"))[:35])
        else:
            row_vals.append("N/A")
    rows.append({"metric": "Age Profile", "values": row_vals})

    # Row 2: Gender split
    row_vals = []
    for s in segments[:4]:
        demo = s.get("demographics", {})
        if isinstance(demo, dict):
            row_vals.append(str(demo.get("gender_split", "N/A"))[:35])
        else:
            row_vals.append("N/A")
    rows.append({"metric": "Gender Split", "values": row_vals})

    # Row 3: Income
    row_vals = []
    for s in segments[:4]:
        demo = s.get("demographics", {})
        if isinstance(demo, dict):
            row_vals.append(str(demo.get("income", "N/A"))[:35])
        else:
            row_vals.append("N/A")
    rows.append({"metric": "Household Income", "values": row_vals})

    # Row 4: Top purchase driver
    row_vals = []
    for s in segments[:4]:
        mini = _normalize_mini_tables(s)
        found = False
        for label, mrows in mini.items():
            if isinstance(mrows, list) and mrows and any(w in label.lower() for w in ("driver", "purchase", "important")):
                top = mrows[0]
                if isinstance(top, dict):
                    row_vals.append(f"{top.get('item','')} ({top.get('pct','')}%)"[:35])
                    found = True
                    break
        if not found:
            row_vals.append("N/A")
    rows.append({"metric": "Top Purchase Driver", "values": row_vals})

    # Row 5: What premium means
    row_vals = [s.get("what_premium_means", "N/A")[:35] for s in segments[:4]]
    rows.append({"metric": "What Premium Means", "values": row_vals})

    # Row 6: Segment size
    row_vals = [f"{s.get('size_pct', '?')}%" for s in segments[:4]]
    rows.append({"metric": "Market Share", "values": row_vals})

    # Row 7: Channels
    row_vals = []
    for s in segments[:4]:
        ch = s.get("channels", [])
        row_vals.append(", ".join(str(c) for c in ch[:2])[:35] if ch else "N/A")
    rows.append({"metric": "Primary Channels", "values": row_vals})

    return rows


def _extract_persona_descriptors(demographics: dict) -> tuple[str, str]:
    """Extract gender_desc and age_desc from demographics for image prompts.

    Returns (gender_desc, age_desc) like ("young woman", "in their early 20s").
    """
    import re as _re
    gender = str(demographics.get("gender_split", "mixed")).lower()
    age = str(demographics.get("age_skew", "25-35")).lower()

    # Extract dominant gender by finding percentages
    gender_desc = "person"
    fem_match = _re.search(r'female[^0-9]*(\d+)', gender) or _re.search(r'(\d+)[^0-9]*female', gender)
    male_match = _re.search(r'(?<!fe)male[^0-9]*(\d+)', gender) or _re.search(r'(\d+)[^0-9]*(?<!fe)male', gender)
    fem_pct = int(fem_match.group(1)) if fem_match else 0
    male_pct = int(male_match.group(1)) if male_match else 0
    if fem_pct >= 60:
        gender_desc = "young woman"
    elif male_pct >= 60:
        gender_desc = "young man"

    # Extract age descriptor
    if "gen z" in age or "18-24" in age or "18-25" in age:
        age_desc = "in their early 20s"
    elif "millennial" in age or "25-34" in age or "25-35" in age:
        age_desc = "in their early 30s"
    elif "gen x" in age or "35-44" in age or "35-50" in age:
        age_desc = "in their early 40s"
    elif "boomer" in age or "50+" in age or "55+" in age:
        age_desc = "in their 50s"
    else:
        age_desc = "young adult"

    return gender_desc, age_desc


def _generate_persona_image_for_target(segment: dict, brand_name: str, category: str = "consumer products", project_id: int = 0) -> "Path | None":
    """Generate a high-quality persona image for the target segment.

    Uses websearch + AI generation + LLM as judge to find the best image
    showing people that match the target segment's demographic profile.
    """
    import tempfile
    from pipeline.image_gen import generate_image

    name = segment.get("name", "Segment")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}
    tagline = segment.get("tagline", "")

    gender_desc, age_desc = _extract_persona_descriptors(demo)

    # Generate 2 candidates via AI
    candidates = []
    prompt = (
        f"Professional lifestyle portrait of a {gender_desc} {age_desc}, "
        f"using a {category}, active and healthy lifestyle, "
        f"warm natural lighting, editorial photography quality, "
        f"authentic and relatable, upper body shot, genuine smile. "
        f"Clean modern background, no text overlays."
    )

    for i in range(2):
        tmp_path = Path(tempfile.mkstemp(suffix=f"_persona_{i}.png")[1])
        result = generate_image(prompt, output_path=tmp_path, backend="gpt-image",
                               size="1024x1024", quality="high")
        if result and result.exists():
            candidates.append(result)

    if not candidates:
        return None

    # LLM as judge: pick the best image
    if len(candidates) >= 2:
        import asyncio
        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                # Already in async context — just pick first
                best = candidates[0]
            else:
                best_list = loop.run_until_complete(_llm_judge_images(
                    candidates,
                    f"Select the best portrait photo for a '{name}' consumer segment ({tagline}). "
                    f"The person should look {age_desc}, {gender_desc}, authentic and relatable. "
                    f"Pick the most professional and appealing image. Reply with just the image number.",
                    pick_count=1
                ))
                best = best_list[0] if best_list else candidates[0]
        except Exception:
            best = candidates[0]
    else:
        best = candidates[0]

    # Clean up non-selected candidates
    for c in candidates:
        if c != best:
            try:
                c.unlink(missing_ok=True)
            except Exception:
                pass

    print(f"[target] Generated persona image for {name}")
    return best


def _generate_target_section_image(segment: dict, brand_name: str, category: str,
                                    style: str = "lifestyle") -> "Path | None":
    """Generate an AI image for target section slides with varied styles.

    Styles:
    - "lifestyle": Person using product in daily life (for Primary Target)
    - "group": Multiple people or community context (for Why Target)
    - "product_hero": Clean product-focused shot, no people (for WHY NOT / FARES)
    - "scene": Person in a lifestyle setting with product (for Consumer Summary)
    """
    import tempfile
    from pipeline.image_gen import generate_image

    name = segment.get("name", "Segment")
    demo = segment.get("demographics", {})
    if not isinstance(demo, dict):
        demo = {}

    gender_desc, age_desc = _extract_persona_descriptors(demo)

    prompts = {
        "lifestyle": (
            f"Professional lifestyle photograph of a {gender_desc} {age_desc} "
            f"using a {category} during their daily routine. "
            f"Active, practical, on-the-go setting (office, commute, gym). "
            f"Warm natural lighting, editorial quality, genuine expression. "
            f"No text, no logos, no watermarks. Landscape orientation 16:9."
        ),
        "group": (
            f"Professional lifestyle photograph showing diverse young professionals "
            f"{age_desc} in a casual social setting, each with their own {category}. "
            f"Friends at a park, cafe, or outdoor activity. Warm, authentic feel. "
            f"Editorial quality, natural lighting. "
            f"No text, no logos, no watermarks. Landscape orientation 16:9."
        ),
        "product_hero": (
            f"Clean product photography of multiple colorful {category} "
            f"arranged aesthetically on a clean surface. Premium editorial style, "
            f"soft studio lighting, minimalist modern background. "
            f"No text, no logos, no people. Landscape orientation 16:9."
        ),
        "scene": (
            f"Professional lifestyle photograph of a {gender_desc} {age_desc} "
            f"in an outdoor setting with a {category}. "
            f"Hiking trail, beach, or park. Active lifestyle, vibrant colors. "
            f"Wide shot showing environment. Editorial photography quality. "
            f"No text, no logos, no watermarks. Landscape orientation 16:9."
        ),
    }

    prompt = prompts.get(style, prompts["lifestyle"])

    tmp_path = Path(tempfile.mkstemp(suffix=f"_target_{style}.png")[1])
    result = generate_image(prompt, output_path=tmp_path, backend="gpt-image",
                           size="1536x1024", quality="high")
    if result and result.exists():
        print(f"[target] AI-generated {style} image for {name}")
        return result
    return None


_brand_image_urls_cache: dict[str, list[str]] = {}


def _websearch_brand_image(brand_name: str, query_suffix: str, output_path) -> "Path | None":
    """Search web for real brand product image by scraping brand's official site.

    Strategy: Use Anthropic web_search to find brand's official site,
    then scrape product page HTML for high-quality CDN image URLs.
    """
    import httpx
    import re as _re

    cache_key = brand_name.lower()

    # Check cache first
    if cache_key not in _brand_image_urls_cache:
        try:
            from anthropic import Anthropic
            from config import ANTHROPIC_API_KEY
            client = Anthropic(api_key=ANTHROPIC_API_KEY)

            response = client.messages.create(
                model=MODEL_SONNET,
                max_tokens=512,
                messages=[{
                    "role": "user",
                    "content": f"Find the official {brand_name} website product page or collection page URL. Return ONLY the URL, nothing else."
                }],
                tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 2}],
            )

            page_url = None
            for block in response.content:
                if hasattr(block, "text") and block.text:
                    urls = _re.findall(r'https?://[^\s"\'<>]+', block.text.strip())
                    if urls:
                        page_url = urls[0]
                        break

            if page_url:
                r = httpx.get(page_url, timeout=15, follow_redirects=True,
                             headers={"User-Agent": "Mozilla/5.0"})
                if r.status_code == 200:
                    # Extract CDN image URLs from HTML
                    img_urls = _re.findall(
                        r'(https?://(?:cdn\.shopify\.com|images\.unsplash\.com|m\.media-amazon\.com|i\.imgur\.com)[^\s"\'<>]+\.(?:jpg|jpeg|png|webp))',
                        r.text
                    )
                    if not img_urls:
                        # Broader fallback: any https image URL
                        img_urls = _re.findall(
                            r'(https?://[^\s"\'<>]+\.(?:jpg|jpeg|png))',
                            r.text
                        )
                    # Dedupe, filter out icons/logos/nav/banners/accessories
                    unique = list(dict.fromkeys(img_urls))
                    _reject = {
                        # UI / navigation / chrome elements
                        "logo", "icon", "nav", "favicon", "badge", "seal",
                        "popup", "modal", "mobile-nav", "hamburger", "arrow",
                        "sprite", "checkout", "cart", "payment",
                        "shop-now", "buy-now", "add-to-cart",
                        # Thumbnail / tiny image markers
                        "thumbnail", "thumb",
                    }
                    filtered = [u for u in unique
                               if not any(rj in u.lower() for rj in _reject)
                               and len(u) > 50
                               # Prefer larger images (skip thumbnails)
                               and not any(sz in u for sz in ["_50x", "_100x", "_150x", "_200x", "w=50", "w=100"])]
                    _brand_image_urls_cache[cache_key] = filtered
                    print(f"[websearch] Cached {len(filtered)} product images for {brand_name}")
        except Exception as e:
            print(f"[websearch] Brand image search failed: {e}")

        if cache_key not in _brand_image_urls_cache:
            _brand_image_urls_cache[cache_key] = []

    # Pick image from cache, filtered by LLM vision judge
    urls = _brand_image_urls_cache.get(cache_key, [])
    if not urls:
        return None

    import random
    # Shuffle to get variety across slides
    random.shuffle(urls)

    # Download candidates, apply quality filters
    candidates = []
    for url in urls[:12]:
        try:
            r = httpx.get(url, timeout=15, follow_redirects=True,
                         headers={"User-Agent": "Mozilla/5.0"})
            if r.status_code != 200 or len(r.content) < 15000:
                continue
            try:
                from PIL import Image as _PILImg
                import io as _qio
                _qimg = _PILImg.open(_qio.BytesIO(r.content))
                w, h = _qimg.size
                _qimg.close()
                if w < 300 or h < 300:
                    continue
                ratio = w / h
                if ratio > 2.5 or ratio < 0.35:
                    continue
            except Exception:
                pass
            candidates.append((url, r.content))
            if len(candidates) >= 6:
                break
        except Exception:
            continue

    if not candidates:
        return None

    # LLM Vision judge: pick the best image for a professional PPT slide
    best_content = _llm_judge_brand_images(brand_name, query_suffix, candidates)
    if best_content:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(best_content)
        print(f"[websearch] LLM-selected brand image for {brand_name}")
        return output_path

    # Fallback: use first candidate if LLM judge fails
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_bytes(candidates[0][1])
    print(f"[websearch] Fallback brand image: {candidates[0][0][:80]}...")
    return output_path


def _llm_judge_brand_images(brand_name: str, context: str, candidates: list) -> bytes | None:
    """Use Claude Vision to select the best brand image for a PPT slide.

    Rejects: single product cutouts on white/plain backgrounds, accessories,
    non-flagship products, low-quality images.
    Prefers: hero product shots, lifestyle imagery, brand campaigns, collections.
    """
    try:
        from anthropic import Anthropic
        from config import ANTHROPIC_API_KEY
        client = Anthropic(api_key=ANTHROPIC_API_KEY)

        # Prepare image content blocks
        content = []
        for i, (url, img_bytes) in enumerate(candidates):
            try:
                from PIL import Image as _PILImg
                import io as _bio
                img = _PILImg.open(_bio.BytesIO(img_bytes))
                img.thumbnail((400, 400))
                buf = _bio.BytesIO()
                img.convert("RGB").save(buf, format="JPEG", quality=70)
                import base64
                b64 = base64.standard_b64encode(buf.getvalue()).decode()
                content.append({"type": "text", "text": f"Image {i+1}:"})
                content.append({
                    "type": "image",
                    "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
                })
            except Exception:
                continue

        if len(content) < 4:  # Need at least 2 images
            return None

        content.append({"type": "text", "text": f"""You are selecting a brand image for "{brand_name}" (context: {context}) to use in a professional brand strategy presentation slide. The image must look impressive and authoritative on stage.

Pick the SINGLE BEST image. Reply with just the image number (e.g. "3").

IDEAL — images that "镇场子" (command the room):
- Brand campaign / lifestyle photography showing the product in an aspirational real-world setting
- Hero shot of the brand's FLAGSHIP product(s) — the product the brand is best known for
- Product collection / lineup showing variety (multiple colors, sizes, styles)
- Editorial or environmental product photography with rich backgrounds and lighting

REJECT — never pick these:
- E-commerce cutout: single product floating on plain white/gray background with no context
- Accessories, replacement parts, or peripheral items — NOT the brand's core product
- Products clearly from a different category than what {brand_name} is primarily known for
- Tiny, blurry, heavily cropped, or watermarked images
- UI elements, icons, banners, or website chrome

If ALL images are bad, reply "none"."""})

        response = client.messages.create(
            model=MODEL_SONNET,
            max_tokens=32,
            messages=[{"role": "user", "content": content}],
        )

        answer = response.content[0].text.strip().lower()
        print(f"[llm_judge] Brand image selection for {brand_name}: {answer}")

        if "none" in answer:
            return None

        # Extract number
        import re as _re_judge
        nums = _re_judge.findall(r'\d+', answer)
        if nums:
            idx = int(nums[0]) - 1
            if 0 <= idx < len(candidates):
                return candidates[idx][1]

    except Exception as e:
        print(f"[llm_judge] Vision judge failed: {e}")

    return None


def _replace_slide_image_websearch(slide, brand_name: str, query_suffix: str = "products"):
    """Replace the first large image on slide with a websearch brand image."""
    img_shapes = [s for s in slide.shapes if s.shape_type == 13]
    if not img_shapes or not brand_name:
        return

    import io as _io
    import tempfile
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    tmp_img = Path(tempfile.mkstemp(suffix="_ws_brand.png")[1])
    result = _websearch_brand_image(brand_name, query_suffix, tmp_img)
    if result and result.exists():
        try:
            img_part, rId = slide.part.get_or_add_image_part(
                _io.BytesIO(result.read_bytes()))
            # Replace largest image shape
            largest = max(img_shapes, key=lambda s: s.width * s.height)
            for blip in largest._element.iter(f"{ns_a}blip"):
                blip.set(f"{ns_r}embed", rId)
            print(f"[target] Websearch image: {brand_name} {query_suffix}")
        except Exception as e:
            print(f"[target] Websearch image failed: {e}")
        finally:
            try:
                tmp_img.unlink(missing_ok=True)
            except Exception:
                pass
    else:
        # Fallback to AI-generated
        from pipeline.image_gen import generate_image
        prompt = f"Professional product photography of {brand_name}, clean modern setting, editorial quality"
        result = generate_image(prompt, output_path=tmp_img, backend="gpt-image",
                               size="1024x1024", quality="standard")
        if result and result.exists():
            try:
                img_part, rId = slide.part.get_or_add_image_part(
                    _io.BytesIO(result.read_bytes()))
                largest = max(img_shapes, key=lambda s: s.width * s.height)
                for blip in largest._element.iter(f"{ns_a}blip"):
                    blip.set(f"{ns_r}embed", rId)
            except Exception:
                pass
            finally:
                try:
                    tmp_img.unlink(missing_ok=True)
                except Exception:
                    pass


def _build_target_recommendation(prs, target):
    """Clone PRIMARY TARGET slide (slide 76 pattern).

    Shape 0: Title — "PRIMARY TARGET: [SEGMENT NAME]"
    Shape 1: Rationale bullets (4 bullets with bold labels)
    Shape 2: Image (right half) — websearch brand image
    Shape 3: Insight text (bottom)
    """
    slide = _clone_slide(prs, T_TARGET_RECOMMENDATION)
    shapes = _find_text_shapes(slide)

    # Support both old keys and LLM-generated keys
    title = target.get("primary_title") or target.get("title", "PRIMARY TARGET")

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 55))
    if len(shapes) >= 2:
        bullets = target.get("rationale_bullets", [])
        bullets = [_truncate(b, 90) for b in bullets[:4]]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        insight = target.get("target_insight") or target.get("insight", "")
        _set_text_preserve_format(shapes[2].text_frame, _truncate(insight, 100))

    return slide


def _build_why_target(prs, target):
    """Clone WHY [SEGMENT] slide (slide 77 pattern).

    Shape 0: Title
    Shape 1: Rationale bullets (left)
    Shape 2: Image (right)
    Shape 3: Insight (bottom)
    """
    slide = _clone_slide(prs, T_WHY_TARGET)
    shapes = _find_text_shapes(slide)

    segment_name = target.get("primary_segment", "THIS SEGMENT")
    title = target.get("why_title") or f"WHY {segment_name.upper()} IS THE RIGHT FOCUS"

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 55))
    if len(shapes) >= 2:
        bullets = target.get("why_bullets") or target.get("rationale_bullets", [])
        bullets = [_truncate(b, 85) for b in bullets[:4]]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        insight = target.get("why_insight") or target.get("insight", "")
        _set_text_preserve_format(shapes[2].text_frame, _truncate(insight, 85))

    return slide


def _build_enables_slide(prs, target):
    """Clone ENABLES slide (slide 78 pattern).

    Shape 0: Title
    Shape 1: "What This Does Not Decide Yet" (right column)
    Shape 2: "What Targeting [X] Unlocks" (left column)
    Shape 3: Closing insight (bottom)
    """
    slide = _clone_slide(prs, T_ENABLES)
    shapes = _find_text_shapes(slide)

    segment_name = target.get("primary_segment", "this segment")
    enables = target.get("enables", [])
    does_not = target.get("does_not_decide", [])

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "WHAT THIS CHOICE ENABLES (AND DOES NOT)")

    # Left column = enables, right column = does not decide
    enables_text = f"What Targeting {segment_name} Unlocks\n" + "\n".join(
        _truncate(e, 85) for e in enables[:3]
    )
    does_not_text = "What This Does Not Decide Yet\n" + "\n".join(
        _truncate(d, 85) for d in does_not[:3]
    )

    if len(shapes) >= 3:
        # shapes sorted by top,left — left column first
        _set_text_preserve_format(shapes[2].text_frame, enables_text)
        _set_text_preserve_format(shapes[1].text_frame, does_not_text)
    if len(shapes) >= 4:
        insight = target.get("enables_insight") or target.get("insight", "")
        _set_text_preserve_format(shapes[3].text_frame, _truncate(insight, 100))

    return slide


def _build_consumer_summary(prs, summary_text):
    """Clone consumer summary slide (slide 79 — half-text, half-image)."""
    slide = _clone_slide(prs, T_CONSUMER_SUMMARY)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "CONSUMER SUMMARY")
    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, _truncate(summary_text, 500))

    return slide


def _build_final_summary(prs, summary_data):
    """Clone three-column summary slide (slide 80 pattern).

    Shape 0: Title — "SUMMARY & NEXT STEPS"
    Shape 1-3: Column headers (Consumer, Capabilities, Competition)
    Shape 4-6: Column text paragraphs
    Shape 7: Closing insight (bottom)
    """
    slide = _clone_slide(prs, T_FINAL_SUMMARY)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "SUMMARY & NEXT STEPS")

    cap_text = summary_data.get("capabilities_column") or summary_data.get("summary_capabilities", "")
    comp_text = summary_data.get("competition_column") or summary_data.get("summary_competition", "")
    cons_text = summary_data.get("consumer_column") or summary_data.get("summary_consumer", "")
    closing = summary_data.get("closing_insight") or summary_data.get("summary_closing", "")

    # Find column header shapes (short text) and body shapes (long text)
    headers = []
    bodies = []
    for s in shapes[1:]:
        text = s.text_frame.text.strip()
        if len(text) < 20:
            headers.append(s)
        elif len(text) > 20:
            bodies.append(s)

    # Set headers
    header_labels = ["Capabilities", "Competition", "Consumer"]
    for i, label in enumerate(header_labels):
        if i < len(headers):
            _set_text_preserve_format(headers[i].text_frame, label)

    # Set body paragraphs (keep short to avoid overlapping closing insight)
    column_texts = [cap_text, comp_text, cons_text]
    for i, txt in enumerate(column_texts):
        if i < len(bodies):
            _set_text_preserve_format(bodies[i].text_frame, _truncate(txt, 220))

    # Closing insight (last shape with substantial width)
    closing_shapes = [s for s in shapes if s.width > 7000000 and s.top > 4000000]
    if closing_shapes:
        _set_text_preserve_format(closing_shapes[0].text_frame, _truncate(closing, 120))

    return slide


# ── Chart Slide Builders (Questionnaire Section) ──────────────

def _remove_chart_shapes(slide, clean_region=True):
    """Remove chart shapes and optionally all chart-region elements.

    After cloning, chart objects decompose into extra shapes (axis titles,
    connectors, sub-labels). With clean_region=True, we also remove these
    orphaned elements, keeping only the slide title (top < 600000) and
    base/sample text (top > 6000000).
    """
    for shape in list(slide.shapes):
        if shape.shape_type == 3:  # CHART
            shape._element.getparent().remove(shape._element)
        elif shape.shape_type == 9:  # LINE / CONNECTOR
            shape._element.getparent().remove(shape._element)

    if not clean_region:
        return

    for shape in list(slide.shapes):
        if shape.top < 600000 or shape.top > 6000000:
            continue
        shape._element.getparent().remove(shape._element)


def _insert_chart_image(slide, chart_path: Path, left=None, top=None, width=None, height=None):
    """Insert a chart PNG image onto a slide at specified position.

    If position not given, uses default chart area (centered, below title).
    """
    if not chart_path or not chart_path.exists():
        return
    if left is None:
        left = Emu(348906)
    if top is None:
        top = Emu(1879166)
    if width is None:
        width = Emu(11843094)
    if height is None:
        height = Emu(4114800)
    slide.shapes.add_picture(str(chart_path), left, top, width, height)


def _insert_asset_image(slide, asset_name: str, left, top, width, height):
    """Insert an asset image (gender_icon, etc.) at a specific position."""
    asset_path = ASSETS_DIR / asset_name
    if not asset_path.exists():
        return
    slide.shapes.add_picture(str(asset_path), left, top, width, height)


def _build_chart_divider(prs, template_idx, title_override=None):
    """Clone a section divider slide (Demographics, Shopping, Brand Eval)."""
    slide = _clone_slide(prs, template_idx)
    if title_override:
        shapes = _find_text_shapes(slide)
        if shapes:
            _set_text_preserve_format(shapes[0].text_frame, title_override)
    return slide


def _build_chart_slide(prs, chart_data: dict, chart_path: Path, template_idx=None):
    """Build a chart slide by cloning a template and replacing charts with rendered images.

    Args:
        chart_data: Chart metadata from analyzer (title, subtitle, chart_type, etc.)
        chart_path: Path to the rendered chart PNG
        template_idx: Which template slide to clone (auto-selected by chart_type if None)
    """
    chart_type = chart_data.get("chart_type", chart_data.get("type", "hbar"))

    if template_idx is None:
        template_idx = {
            "dual": T_CHART_DUAL,
            "donut": T_CHART_DUAL,
            "pie": T_CHART_SINGLE_HBAR,
            "hbar": T_CHART_SINGLE_HBAR,
            "vbar": T_CHART_SINGLE_VBAR,
            "stacked": T_CHART_STACKED,
            "funnel": T_CHART_SINGLE_HBAR,
            "grouped_bar": T_CHART_SINGLE_HBAR,
            "wordcloud": T_CHART_SINGLE_HBAR,
            "matrix": T_CHART_TABLE,
            "table": T_CHART_TABLE,
        }.get(chart_type, T_CHART_SINGLE_HBAR)

    slide = _clone_slide(prs, template_idx)

    _remove_chart_shapes(slide)

    # Update title (kept by _remove_chart_shapes — top < 600000)
    title = chart_data.get("title", "")
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip()
        if text and text.isupper() and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 60))
            break

    # Add subtitle/question as a new text box (old ones were removed with chart region)
    subtitle = chart_data.get("subtitle", "") or chart_data.get("question", "")
    if subtitle:
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(
            Emu(419100), Emu(1124334), Emu(11353800), Emu(500000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _truncate(subtitle, 120)
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x29, 0x25, 0x24)
        run.font.italic = True

    # Insert rendered chart image below subtitle — size by chart type
    if chart_type == "dual":
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    elif chart_type in ("donut", "pie"):
        # Full-width layout (donut left + bars right) to fill the slide
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    elif chart_type == "vbar":
        # Vertical bar: slightly narrower to avoid edge clipping
        _insert_chart_image(slide, chart_path,
                            left=Emu(600000), top=Emu(1750000),
                            width=Emu(10992000), height=Emu(4300000))
    elif chart_type == "matrix":
        # Matrix/table: full width, slightly taller for readability
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1650000),
                            width=Emu(11353800), height=Emu(4600000))
    elif chart_type == "wordcloud":
        # Word cloud: centered, large square area
        chart_w = Emu(8229600)   # 9.0 inches
        chart_h = Emu(4389120)   # 4.8 inches
        chart_left = Emu((12192000 - 8229600) // 2)
        _insert_chart_image(slide, chart_path,
                            left=chart_left, top=Emu(1750000),
                            width=chart_w, height=chart_h)
    elif chart_type == "grouped_bar":
        # Grouped bar: full width with small margins
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    else:
        # Default hbar: full width
        _insert_chart_image(slide, chart_path)

    return slide


def _build_respondent_profile(prs, generation_data: dict = None, ethnicity_data: dict = None,
                              gender_data: dict = None, sample_size: int = 200,
                              generation_chart: Path = None, ethnicity_chart: Path = None):
    """Build respondent profile slide matching CozyFit slide 28.

    Uses native PowerPoint CHART objects via _clone_slide_with_charts when
    structured data is provided. Falls back to PNG image insertion when only
    chart image paths are given.

    Args:
        generation_data: {"categories": [...], "values": [...]} (values as 0-1 decimals)
        ethnicity_data: {"categories": [...], "values": [...]} (values as 0-1 decimals)
        gender_data: {"male_pct": 30, "female_pct": 70}
        sample_size: Survey sample size
        generation_chart: Fallback PNG path for generation chart
        ethnicity_chart: Fallback PNG path for ethnicity chart
    """
    from pptx.chart.data import CategoryChartData

    use_native = generation_data or ethnicity_data

    if use_native:
        slide = _clone_slide_with_charts(prs, 28)
        # Find chart shapes — template has COLUMN_CLUSTERED (ethnicity) and BAR_CLUSTERED (generation)
        chart_shapes = [s for s in slide.shapes if hasattr(s, "has_chart") and s.has_chart]
        for cs in chart_shapes:
            ct = cs.chart.chart_type
            if ct == 57 and generation_data:  # BAR_CLUSTERED = generation hbar
                cd = CategoryChartData()
                cats = generation_data.get("categories", [])
                vals = generation_data.get("values", [])
                # Convert integer percentages to decimals if needed
                vals = [v / 100.0 if v > 1 else v for v in vals]
                cd.categories = cats
                cd.add_series("Series 1", tuple(vals))
                cs.chart.replace_data(cd)
            elif ct == 51 and ethnicity_data:  # COLUMN_CLUSTERED = ethnicity vbar
                cd = CategoryChartData()
                cats = ethnicity_data.get("categories", [])
                vals = ethnicity_data.get("values", [])
                vals = [v / 100.0 if v > 1 else v for v in vals]
                cd.categories = cats
                cd.add_series("Series 1", tuple(vals))
                cs.chart.replace_data(cd)
    else:
        slide = _clone_slide(prs, 28)
        for shape in list(slide.shapes):
            if shape.shape_type == 3:
                shape._element.getparent().remove(shape._element)
        if generation_chart:
            _insert_chart_image(slide, generation_chart,
                                left=Emu(3364800), top=Emu(1573910),
                                width=Emu(8229600), height=Emu(2103120))
        if ethnicity_chart:
            _insert_chart_image(slide, ethnicity_chart,
                                left=Emu(3364800), top=Emu(4232530),
                                width=Emu(8229600), height=Emu(2103120))

    # Update gender percentages inside the GROUP shape
    male_pct = gender_data.get("male_pct", 30) if gender_data else 30
    female_pct = gender_data.get("female_pct", 70) if gender_data else 70
    groups = [s for s in slide.shapes if s.shape_type == 6]
    for grp in groups:
        pct_children = []
        for child in grp.shapes:
            if child.has_text_frame:
                text = child.text_frame.text.strip()
                if "%" in text and len(text) <= 5:
                    pct_children.append(child)
        pct_children.sort(key=lambda c: c.left)
        if len(pct_children) >= 2:
            _set_text_preserve_format(pct_children[0].text_frame, f"{male_pct}%")
            _set_text_preserve_format(pct_children[1].text_frame, f"{female_pct}%")

    # Update title
    for s in _find_text_shapes(slide):
        if "RESPONDENT" in s.text_frame.text.strip().upper():
            _set_text_preserve_format(s.text_frame, "RESPONDENT PROFILE")
            break

    # Update sample size
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().lower()
        if "total sample" in text or "base n" in text:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            break

    return slide


def _build_respondent_profile_continued(prs, income_data: dict = None, sample_size: int = 200,
                                        marital_data: dict = None, income_chart: Path = None):
    """Build respondent profile continued slide matching CozyFit slide 29.

    Args:
        income_data: {"categories": [...], "values": [...]} (values as 0-1 decimals)
        sample_size: Survey sample size
        marital_data: {"married_pct": 52, "single_pct": 33, "divorced_pct": 15}
        income_chart: Fallback PNG path for income chart
    """
    from pptx.chart.data import CategoryChartData

    use_native = income_data is not None

    if use_native:
        slide = _clone_slide_with_charts(prs, 29)
        # Slide 29 has one COLUMN_CLUSTERED chart (income)
        for s in slide.shapes:
            if hasattr(s, "has_chart") and s.has_chart:
                cd = CategoryChartData()
                cats = income_data.get("categories", [])
                vals = income_data.get("values", [])
                vals = [v / 100.0 if v > 1 else v for v in vals]
                cd.categories = cats
                cd.add_series("Series 1", tuple(vals))
                s.chart.replace_data(cd)
                break
    else:
        slide = _clone_slide(prs, 29)
        for shape in list(slide.shapes):
            if shape.shape_type == 3:
                shape._element.getparent().remove(shape._element)
        if income_chart:
            _insert_chart_image(slide, income_chart,
                                left=Emu(4214124), top=Emu(1813653),
                                width=Emu(7506100), height=Emu(3133704))

    # Update marital status data in the 3 GROUP shapes
    if marital_data is None:
        marital_data = {"married_pct": 52, "single_pct": 33, "divorced_pct": 15}
    marital_items = [
        (marital_data.get("married_pct", 52), "Married or domestic partnership"),
        (marital_data.get("single_pct", 33), "Single, never married"),
        (marital_data.get("divorced_pct", 15), "Widowed, divorced or separated"),
    ]
    groups = sorted(
        [s for s in slide.shapes if s.shape_type == 6],
        key=lambda s: s.top,
    )
    for gi, grp in enumerate(groups[:3]):
        if gi >= len(marital_items):
            break
        pct, desc = marital_items[gi]
        for child in grp.shapes:
            if child.has_text_frame:
                text = child.text_frame.text.strip()
                if "%" in text:
                    _set_text_preserve_format(child.text_frame, f"{pct}%")
                elif len(text) > 3:
                    _set_text_preserve_format(child.text_frame, desc)

    # Update title
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().upper()
        if "RESPONDENT" in text:
            _set_text_preserve_format(s.text_frame, "RESPONDENT PROFILE (CONTINUED)")
            break

    # Update "Household Income" label
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip()
        if "Household" in text or "Income" in text:
            _set_text_preserve_format(s.text_frame, "Household Income")
            break

    # Update "Marital Status" label
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip()
        if "Marital" in text or "Status" in text:
            _set_text_preserve_format(s.text_frame, "Marital Status")
            break

    # Update sample size
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().lower()
        if "total sample" in text or "base n" in text:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            break

    return slide


def _build_occupation_slide(prs, occupation_data: dict = None, sample_size: int = 200):
    """Build occupation slide matching CozyFit slide 30.

    Template slide 30 has one COLUMN_CLUSTERED chart with occupation categories.

    Args:
        occupation_data: {"categories": [...], "values": [...]}
        sample_size: Survey sample size
    """
    from pptx.chart.data import CategoryChartData

    if not occupation_data:
        return None

    slide = _clone_slide_with_charts(prs, 30)
    for s in slide.shapes:
        if hasattr(s, "has_chart") and s.has_chart:
            cd = CategoryChartData()
            cats = occupation_data.get("categories", [])
            vals = occupation_data.get("values", [])
            vals = [v / 100.0 if v > 1 else v for v in vals]
            cd.categories = cats
            cd.add_series("Series 1", tuple(vals))
            s.chart.replace_data(cd)
            break

    # Update title
    title_text = occupation_data.get("title", "OCCUPATION AND WORK DETAILS")
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().upper()
        if "OCCUPATION" in text or "WORK" in text:
            _set_text_preserve_format(s.text_frame, title_text.upper())
            break

    # Update sample size
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().lower()
        if "total sample" in text or "base n" in text:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            break

    return slide


def _build_occupation_detail_slide(prs, detail_charts: list = None, sample_size: int = 200):
    """Build occupation detail slide matching CozyFit slide 31.

    Template slide 31 has: DOUGHNUT + 2x COLUMN_CLUSTERED charts.

    Args:
        detail_charts: list of up to 3 dicts, each {"categories": [...], "values": [...], "title": "..."}
                       First is mapped to DOUGHNUT, remaining to COLUMN_CLUSTERED.
        sample_size: Survey sample size
    """
    from pptx.chart.data import CategoryChartData

    if not detail_charts or len(detail_charts) < 1:
        return None

    slide = _clone_slide_with_charts(prs, 31)
    # Template chart order: DOUGHNUT (-4120), COLUMN_CLUSTERED (51), COLUMN_CLUSTERED (51)
    chart_shapes = [s for s in slide.shapes if hasattr(s, "has_chart") and s.has_chart]

    for i, cs in enumerate(chart_shapes):
        if i >= len(detail_charts):
            break
        data = detail_charts[i]
        cd = CategoryChartData()
        cats = data.get("categories", [])
        vals = data.get("values", [])
        vals = [v / 100.0 if v > 1 else v for v in vals]
        cd.categories = cats
        cd.add_series("Series 1", tuple(vals))
        cs.chart.replace_data(cd)

    # Update title
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().upper()
        if "OCCUPATION" in text:
            title = detail_charts[0].get("slide_title", "OCCUPATION AND WORK DETAILS (CONTINUED)")
            _set_text_preserve_format(s.text_frame, title.upper())
            break

    # Update sample size
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().lower()
        if "total sample" in text or "base n" in text:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            break

    return slide


def _build_social_media_slide(prs, social_media_data: dict = None, sample_size: int = 200,
                              chart_path: Path = None):
    """Build social media platform usage slide matching CozyFit slide 32.

    Args:
        social_media_data: {"categories": [...], "values": [...]} (values as 0-1 decimals)
        sample_size: Survey sample size
        chart_path: Fallback PNG path for social media chart
    """
    from pptx.chart.data import CategoryChartData

    use_native = social_media_data is not None

    if use_native:
        slide = _clone_slide_with_charts(prs, 32)
        # Slide 32 has one BAR_CLUSTERED chart (social media hbar)
        for s in slide.shapes:
            if hasattr(s, "has_chart") and s.has_chart:
                cd = CategoryChartData()
                cats = social_media_data.get("categories", [])
                vals = social_media_data.get("values", [])
                vals = [v / 100.0 if v > 1 else v for v in vals]
                cd.categories = cats
                cd.add_series("Series 1", tuple(vals))
                s.chart.replace_data(cd)
                break
    else:
        slide = _clone_slide(prs, 32)
        for shape in list(slide.shapes):
            if shape.shape_type == 3:
                shape._element.getparent().remove(shape._element)
        if chart_path:
            _insert_chart_image(slide, chart_path,
                                left=Emu(4719665), top=Emu(1600200),
                                width=Emu(6928881), height=Emu(4754880))

    # Update title
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().upper()
        if "SOCIAL MEDIA" in text:
            _set_text_preserve_format(s.text_frame, "SOCIAL MEDIA PLATFORM USAGE")
            break

    # Update sample size
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip().lower()
        if "total sample" in text or "base n" in text:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            break

    return slide


# ── Shopping / Survey Native Chart Builders ──────────────────


def _build_native_single_hbar(prs, title: str, subtitle: str, data: dict,
                               sample_size: int = 200, template_idx: int = 34):
    """Build a single full-width hbar slide using native PowerPoint chart.

    Clones template slide 34 (or 38) and replaces the BAR_CLUSTERED chart.

    Args:
        title: ALL-CAPS slide title
        subtitle: Survey question text (displayed under title)
        data: {"categories": [...], "values": [...]}
        sample_size: Survey sample size
        template_idx: Template slide index (34 or 38)
    """
    from pptx.chart.data import CategoryChartData

    slide = _clone_slide_with_charts(prs, template_idx)

    # Replace chart data
    for s in slide.shapes:
        if hasattr(s, "has_chart") and s.has_chart:
            cd = CategoryChartData()
            cats = data.get("categories", [])
            vals = data.get("values", [])
            vals = [v / 100.0 if v > 1 else v for v in vals]
            cd.categories = cats
            cd.add_series("Series 1", tuple(vals))
            s.chart.replace_data(cd)
            for plot in s.chart.plots:
                plot.data_labels.number_format = '0%'
                plot.data_labels.number_format_is_linked = False
                for series in plot.series:
                    series.data_labels.number_format = '0%'
                    series.data_labels.number_format_is_linked = False
            break

    # Update text shapes — replace known patterns, clear leftovers
    text_shapes = _find_text_shapes(slide)
    title_set = False
    subtitle_set = False
    for s in text_shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        text_upper = text.upper()
        if not title_set and text_upper == text and len(text) > 5 and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            title_set = True
        elif "total sample" in text.lower() or "base n" in text.lower():
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
        elif not subtitle_set and len(text) > 20 and text != title:
            _set_text_preserve_format(s.text_frame, _truncate(subtitle, 120))
            subtitle_set = True
        elif len(text) > 3 and text != title:
            # Clear any other template leftover text
            _set_text_preserve_format(s.text_frame, "")

    return slide


def _build_appendix_section(prs, charts: list, sample_size: int = 200,
                            brand_name: str = "", category: str = "",
                            segments: list = None):
    """Build Research Appendix: transition slide + segmented grouped bar charts.

    Each chart shows a survey question broken down by consumer segment,
    with each segment as a separate colored series.

    Args:
        charts: Full chart list from consumer analysis
        sample_size: Survey sample size (n=)
        brand_name: Brand name for context
        category: Product category for context
        segments: List of segment dicts with name, size_pct, demographics, etc.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.dml.color import RGBColor
    import json as _json

    if not segments:
        segments = []

    # 1. Transition slide
    slide = _clone_slide(prs, T_APPENDIX_TRANSITION)
    shapes = _find_text_shapes(slide)
    if shapes:
        _set_text_preserve_format(shapes[0].text_frame, "Research Appendix")
    print("[appendix] Research Appendix transition slide")

    # 2. Select charts with usable data
    _SKIP_TYPES = {"matrix", "grouped_bar", "stacked", "table", "wordcloud"}
    appendix_charts = []
    for c in charts:
        ctype = c.get("chart_type", c.get("type", ""))
        if ctype in _SKIP_TYPES:
            continue
        cats = c.get("categories", [])
        vals = c.get("values", [])
        data = c.get("data", [])
        if not cats and isinstance(data, list) and data:
            cats = [d.get("label", "") for d in data if isinstance(d, dict)]
            vals = [d.get("value", 0) for d in data if isinstance(d, dict)]
        if cats and vals and len(cats) >= 2:
            if not any(c2 for c2 in cats) or not any(v != 0 for v in vals):
                continue
            appendix_charts.append({
                "title": c.get("title", "Survey Question"),
                "categories": cats[:8],
                "values": vals[:8],
                "chart_type": ctype,
                "section": c.get("section", ""),
            })
        if len(appendix_charts) >= 12:
            break

    if not appendix_charts:
        print("[appendix] No suitable chart data for appendix")
        return 0

    # 3. Generate question text + per-segment data via LLM
    seg_names = [s.get("name", f"Segment {i+1}") for i, s in enumerate(segments[:5])]
    seg_sizes = [s.get("size_pct", 20) for s in segments[:5]]
    seg_n = [max(1, round(sample_size * pct / 100)) for pct in seg_sizes]

    appendix_data = _generate_appendix_segmented_data(
        appendix_charts, seg_names, seg_sizes, brand_name, category, sample_size
    )

    # 4. Segment colors matching CozyFit reference (slides 83-91)
    _SEG_COLORS = [
        RGBColor(0xE8, 0x6C, 0x00),  # dark orange (segment 1)
        RGBColor(0xF0, 0x9E, 0x4A),  # medium orange (segment 2)
        RGBColor(0xE8, 0xC8, 0x7A),  # tan/gold (segment 3)
        RGBColor(0xA0, 0xA8, 0xB8),  # gray-blue (segment 4)
        RGBColor(0xC8, 0xCC, 0xD4),  # light gray (segment 5)
    ]

    # 5. Build column-based bar chart slides (one column per segment)
    #    Layout matches CozyFit reference slides 83-91:
    #    - White background (cloned from T_APPENDIX_WHITE = slide 83)
    #    - Title text at top (from template placeholder)
    #    - "Column n / Column %" label top-left of chart area
    #    - Segment name + n as column headers
    #    - Response option labels on the left
    #    - Horizontal bars in each segment column, each with % label
    #    - "base n = X" at bottom

    # Layout constants (EMU)
    _LABEL_LEFT = 200000           # left edge of response labels
    _LABEL_WIDTH = 2100000         # width for response label column
    _CHART_LEFT = 2400000          # left edge of chart columns
    _CHART_RIGHT = 11800000        # right edge of chart columns
    _HEADER_TOP = 1450000          # top of column headers
    _HEADER_HEIGHT = 450000        # height of header row
    _DATA_TOP = 1950000            # top of data rows
    _DATA_BOTTOM = 6050000         # bottom of data rows
    _BAR_HEIGHT_RATIO = 0.55       # bar thickness relative to row height
    _COL_GAP = 50000               # gap between columns
    _GRAY = RGBColor(0x66, 0x66, 0x66)
    _DARK = RGBColor(0x33, 0x33, 0x33)
    _LINE_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

    n_segs = min(len(seg_names), 5)
    _col_total = _CHART_RIGHT - _CHART_LEFT
    _col_width = (_col_total - _COL_GAP * (n_segs - 1)) // n_segs

    for i, ac in enumerate(appendix_charts):
        seg_data = appendix_data.get(ac["title"])
        if not seg_data:
            continue

        question = seg_data.get("question", ac["title"])
        response_options = seg_data.get("responses", ac["categories"][:6])
        seg_values = seg_data.get("segments", {})

        # Clone white appendix template
        slide = _clone_slide(prs, T_APPENDIX_WHITE)

        # Update existing shapes
        for sh in slide.shapes:
            if sh.shape_type == 14:  # Title placeholder
                _set_text_preserve_format(sh.text_frame, question)
            elif sh.shape_type == 17 and sh.has_text_frame:
                if "base" in sh.text_frame.text.lower():
                    _set_text_preserve_format(sh.text_frame, f"base n = {sample_size}")
            elif sh.shape_type == 13:  # Picture — remove
                sh._element.getparent().remove(sh._element)

        n_rows = len(response_options)
        row_height = (_DATA_BOTTOM - _DATA_TOP) // n_rows
        bar_h = int(row_height * _BAR_HEIGHT_RATIO)
        bar_y_offset = (row_height - bar_h) // 2

        # "Column n / Column %" header label
        _add_appendix_text(slide, _LABEL_LEFT, _HEADER_TOP,
                           _LABEL_WIDTH, _HEADER_HEIGHT,
                           "Column n\nColumn %", Pt(8), _GRAY,
                           bold=False, alignment="center")

        # Column headers: segment name + n
        for si in range(n_segs):
            col_left = _CHART_LEFT + si * (_col_width + _COL_GAP)
            _add_appendix_text(slide, col_left, _HEADER_TOP,
                               _col_width, _HEADER_HEIGHT,
                               f"{seg_names[si]}\n{seg_n[si]}",
                               Pt(9), _DARK, bold=False, alignment="center")

        # Data rows
        for ri, resp in enumerate(response_options):
            row_top = _DATA_TOP + ri * row_height

            # Row separator line (thin gray)
            if ri > 0:
                line = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    Emu(_LABEL_LEFT), Emu(row_top),
                    Emu(_CHART_RIGHT - _LABEL_LEFT), Emu(12700),  # 1pt line
                )
                line.fill.solid()
                line.fill.fore_color.rgb = _LINE_GRAY
                line.line.fill.background()

            # Response option label (left column)
            _add_appendix_text(slide, _LABEL_LEFT, row_top,
                               _LABEL_WIDTH, row_height,
                               resp, Pt(8), _DARK,
                               bold=False, alignment="right", v_center=True)

            # Bars for each segment
            for si in range(n_segs):
                col_left = _CHART_LEFT + si * (_col_width + _COL_GAP)
                vals_for_seg = seg_values.get(seg_names[si], [])
                while len(vals_for_seg) < n_rows:
                    vals_for_seg.append(0.1)
                val = vals_for_seg[ri]
                if val > 1:
                    val = val / 100.0  # normalize

                # Bar rectangle
                bar_w = max(int(_col_width * 0.85 * val), Emu(25400))  # min 2pt
                bar_top = row_top + bar_y_offset
                bar = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    Emu(col_left), Emu(bar_top),
                    Emu(bar_w), Emu(bar_h),
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = _SEG_COLORS[si % len(_SEG_COLORS)]
                bar.line.fill.background()

                # Percentage label at end of bar
                pct_text = f"{int(round(val * 100))}%"
                pct_left = col_left + bar_w + 25000
                _add_appendix_text(slide, pct_left, bar_top,
                                   400000, bar_h,
                                   pct_text, Pt(8), _DARK,
                                   bold=False, alignment="left", v_center=True)

        print(f"[appendix] Segmented chart {i+1}/{len(appendix_charts)}: {ac['title'][:50]}")

    print(f"[appendix] Built {len(appendix_charts)} segmented appendix chart slides")
    return len(appendix_charts)  # number of chart slides (excludes transition)


def _add_appendix_text(slide, left, top, width, height, text, font_size,
                        color, bold=False, alignment="left", v_center=False):
    """Add a styled text box to an appendix slide."""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if v_center:
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        try:
            tf.auto_size = None
            txBox.text_frame._txBody.attrib[
                '{http://schemas.openxmlformats.org/drawingml/2006/main}anchor'
            ] = 'ctr' if v_center else 't'
        except Exception:
            pass

    p = tf.paragraphs[0]
    p.text = text
    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER
    elif alignment == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


def _generate_appendix_segmented_data(charts: list, seg_names: list,
                                       seg_sizes: list, brand_name: str,
                                       category: str, sample_size: int) -> dict:
    """Generate per-segment survey response data for appendix charts via LLM.

    Returns dict mapping chart title → {question, responses, segments: {name: [vals]}}.
    The data must be coherent with each segment's persona and support the narrative.
    """
    import json as _json

    chart_info = []
    for c in charts:
        chart_info.append({
            "title": c["title"],
            "response_options": c["categories"][:6],
            "overall_values": [round(v, 2) if v <= 1 else round(v / 100, 2) for v in c["values"][:6]],
        })

    seg_desc = ", ".join(f"{n} ({p}%)" for n, p in zip(seg_names, seg_sizes))

    system = (
        f"You are a quantitative research analyst generating per-segment survey data for a {category} brand study. "
        "Each segment must show distinct behavioral patterns that align with their persona — "
        "a value-conscious parent responds differently from a trend-driven Gen Z collector. "
        "Data must be internally consistent and support the brand strategy narrative. "
        "Output valid JSON only."
    )
    user_msg = (
        f"Brand: {brand_name}, Category: {category}.\n"
        f"Total sample: n={sample_size}\n"
        f"Segments: {seg_desc}\n\n"
        f"For each chart below, generate:\n"
        f"1. A natural survey question (2nd person, conversational)\n"
        f"2. Response option labels (use the provided ones, shorten if >40 chars)\n"
        f"3. Per-segment percentage values (as decimals 0.0-1.0) that:\n"
        f"   - Are realistic and internally consistent\n"
        f"   - Reflect each segment's persona and behavior patterns\n"
        f"   - Within each segment, values for a given question should sum to roughly 1.0 if the question is single-select\n"
        f"   - Hover near the overall values but with meaningful segment-level variation\n"
        f"   - Support the brand strategy narrative\n\n"
        f"Charts:\n{_json.dumps(chart_info, indent=2)}\n\n"
        f"Output format — JSON object where keys are chart titles:\n"
        f"{{\n"
        f'  "CHART TITLE": {{\n'
        f'    "question": "Survey question text?",\n'
        f'    "responses": ["Option A", "Option B", ...],\n'
        f'    "segments": {{\n'
        f'      "{seg_names[0]}": [0.27, 0.33, ...],\n'
        f'      "{seg_names[1]}": [0.35, 0.25, ...],\n'
        f"      ...\n"
        f"    }}\n"
        f"  }},\n"
        f"  ...\n"
        f"}}\n"
    )

    try:
        import anthropic
        from config import ANTHROPIC_API_KEY
        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        response = client.messages.create(
            model=MODEL_SONNET,
            max_tokens=8000,
            system=system,
            messages=[{"role": "user", "content": user_msg}],
        )
        result = response.content[0].text
        parsed = _json.loads(_clean_llm_json(result))
        print(f"[appendix] LLM generated segmented data for {len(parsed)} charts")
        return parsed
    except Exception as e:
        print(f"[appendix] Segmented data generation failed: {e}")
        # Fallback: distribute overall values with small random variation per segment
        import random
        fallback = {}
        for c in charts:
            segs = {}
            for sn in seg_names:
                segs[sn] = [
                    max(0.01, min(0.99, v + random.uniform(-0.08, 0.08)))
                    for v in (v2 / 100.0 if v2 > 1 else v2 for v2 in c["values"][:6])
                ]
            fallback[c["title"]] = {
                "question": c["title"],
                "responses": c["categories"][:6],
                "segments": segs,
            }
        return fallback


def _build_native_dual_donut_hbar(prs, title: str, subtitle: str,
                                   left_data: dict, right_data: dict,
                                   left_question: str = "", right_question: str = "",
                                   sample_size: int = 200, template_idx: int = 35):
    """Build a dual-chart slide: donut (left) + hbar (right).

    Clones template slide 35 (or 37) which has DOUGHNUT + BAR_CLUSTERED.

    Args:
        title: ALL-CAPS slide title
        subtitle: Time-period note (e.g., "All data reflect...")
        left_data: {"categories": [...], "values": [...]} for donut
        right_data: {"categories": [...], "values": [...]} for hbar
        left_question: Question text for left chart
        right_question: Question text for right chart
        sample_size: Survey sample size
        template_idx: Template slide index (35 or 37)
    """
    from pptx.chart.data import CategoryChartData

    slide = _clone_slide_with_charts(prs, template_idx)

    # Replace chart data — find donut and hbar charts
    for s in slide.shapes:
        if not (hasattr(s, "has_chart") and s.has_chart):
            continue
        ct = s.chart.chart_type
        if ct == -4120 and left_data:  # DOUGHNUT
            cd = CategoryChartData()
            cats = left_data.get("categories", [])
            vals = left_data.get("values", [])
            vals = [v / 100.0 if v > 1 else v for v in vals]
            cd.categories = cats
            cd.add_series("Column1", tuple(vals))
            s.chart.replace_data(cd)
            # Force percentage format on data labels
            for plot in s.chart.plots:
                plot.data_labels.number_format = '0%'
                plot.data_labels.number_format_is_linked = False
                for series in plot.series:
                    series.data_labels.number_format = '0%'
                    series.data_labels.number_format_is_linked = False
        elif ct == 57 and right_data:  # BAR_CLUSTERED (hbar)
            cd = CategoryChartData()
            cats = right_data.get("categories", [])
            vals = right_data.get("values", [])
            vals = [v / 100.0 if v > 1 else v for v in vals]
            cd.categories = cats
            cd.add_series("Series 1", tuple(vals))
            s.chart.replace_data(cd)
            for plot in s.chart.plots:
                plot.data_labels.number_format = '0%'
                plot.data_labels.number_format_is_linked = False
                for series in plot.series:
                    series.data_labels.number_format = '0%'
                    series.data_labels.number_format_is_linked = False

    # Update text shapes — replace known patterns, clear everything else
    text_shapes = _find_text_shapes(slide)
    title_set = False
    subtitle_set = False
    sample_set = False
    question_shapes = []
    extra_shapes = []  # template leftovers to clear
    for s in text_shapes:
        text = s.text_frame.text.strip()
        text_lower = text.lower()
        if not text:
            continue
        if ("total sample" in text_lower or "base n" in text_lower) and not sample_set:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            sample_set = True
        elif not title_set and text.upper() == text and len(text) > 5 and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            title_set = True
        elif not subtitle_set and "all data" in text_lower:
            _set_text_preserve_format(s.text_frame, _truncate(subtitle, 100))
            subtitle_set = True
        elif len(text) > 20:
            # Likely a question label
            question_shapes.append(s)
        else:
            # Short text like "$297", "Average Median Spend", "Legend" → clear
            extra_shapes.append(s)

    # Set left and right question labels (sorted by horizontal position)
    question_shapes.sort(key=lambda s: s.left)
    if left_question and len(question_shapes) >= 1:
        _set_text_preserve_format(question_shapes[0].text_frame, _truncate(left_question, 90))
    if right_question and len(question_shapes) >= 2:
        _set_text_preserve_format(question_shapes[1].text_frame, _truncate(right_question, 90))
    # Clear any additional question shapes beyond the 2 we need
    for qs in question_shapes[2:]:
        _set_text_preserve_format(qs.text_frame, "")
    # Clear template leftover text (dollar amounts, labels, etc.)
    for es in extra_shapes:
        _set_text_preserve_format(es.text_frame, "")

    return slide


def _build_native_dual_hbar(prs, title: str, subtitle: str,
                             left_data: dict, right_data: dict,
                             left_question: str = "", right_question: str = "",
                             sample_size: int = 200):
    """Build a dual-chart slide: hbar (left) + hbar (right).

    Clones template slide 36 which has 2x BAR_CLUSTERED.

    Args:
        title: ALL-CAPS slide title
        subtitle: Time-period note
        left_data: {"categories": [...], "values": [...]} for left hbar
        right_data: {"categories": [...], "values": [...]} for right hbar
        left_question: Question text for left chart
        right_question: Question text for right chart
        sample_size: Survey sample size
    """
    from pptx.chart.data import CategoryChartData

    slide = _clone_slide_with_charts(prs, 36)

    # Replace both hbar charts — sort by left position to determine left vs right
    chart_shapes = sorted(
        [s for s in slide.shapes if hasattr(s, "has_chart") and s.has_chart],
        key=lambda s: s.left
    )
    for i, cs in enumerate(chart_shapes):
        chart_data = left_data if i == 0 else right_data
        if chart_data:
            cd = CategoryChartData()
            cats = chart_data.get("categories", [])
            vals = chart_data.get("values", [])
            vals = [v / 100.0 if v > 1 else v for v in vals]
            cd.categories = cats
            cd.add_series("Series 1", tuple(vals))
            cs.chart.replace_data(cd)
            for plot in cs.chart.plots:
                plot.data_labels.number_format = '0%'
                plot.data_labels.number_format_is_linked = False
                for series in plot.series:
                    series.data_labels.number_format = '0%'
                    series.data_labels.number_format_is_linked = False

    # Update text shapes — replace known patterns, clear leftovers
    text_shapes = _find_text_shapes(slide)
    title_set = False
    subtitle_set = False
    sample_set = False
    question_shapes = []
    extra_shapes = []
    for s in text_shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        text_lower = text.lower()
        if ("total sample" in text_lower or "base n" in text_lower) and not sample_set:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            sample_set = True
        elif not title_set and text.upper() == text and len(text) > 5 and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            title_set = True
        elif not subtitle_set and "all data" in text_lower:
            _set_text_preserve_format(s.text_frame, _truncate(subtitle, 100))
            subtitle_set = True
        elif len(text) > 20:
            question_shapes.append(s)
        else:
            extra_shapes.append(s)

    question_shapes.sort(key=lambda s: s.left)
    if left_question and len(question_shapes) >= 1:
        _set_text_preserve_format(question_shapes[0].text_frame, _truncate(left_question, 90))
    if right_question and len(question_shapes) >= 2:
        _set_text_preserve_format(question_shapes[1].text_frame, _truncate(right_question, 90))
    for qs in question_shapes[2:]:
        _set_text_preserve_format(qs.text_frame, "")
    for es in extra_shapes:
        _set_text_preserve_format(es.text_frame, "")

    return slide


def _build_native_stacked_bar(prs, title: str, brands: list,
                               awareness: list, purchase: list,
                               sample_size: int = 200):
    """Build a stacked bar slide: brand awareness + purchase conversion.

    Clones template slide 42 (BAR_STACKED with 3 series).
    Series 0 = "Aware only" (orange), Series 1 = "Purchased previously" (gray),
    Series 2 = Total (awareness + purchase, shown as end label).

    Args:
        title: ALL-CAPS slide title
        brands: List of brand names (categories)
        awareness: List of awareness-only percentages (integers)
        purchase: List of purchase percentages (integers)
        sample_size: Survey sample size
    """
    from pptx.chart.data import CategoryChartData

    slide = _clone_slide_with_charts(prs, 42)

    for s in slide.shapes:
        if not (hasattr(s, "has_chart") and s.has_chart):
            continue
        cd = CategoryChartData()
        cd.categories = brands
        # Convert to decimals
        aware_dec = [a / 100.0 if a > 1 else a for a in awareness]
        purch_dec = [p / 100.0 if p > 1 else p for p in purchase]
        total_dec = [a + p for a, p in zip(aware_dec, purch_dec)]
        cd.add_series("Aware only", tuple(aware_dec))
        cd.add_series("Purchased previously", tuple(purch_dec))
        cd.add_series("Total", tuple(total_dec))
        s.chart.replace_data(cd)
        for plot in s.chart.plots:
            plot.data_labels.number_format = '0%'
            plot.data_labels.number_format_is_linked = False
            for series in plot.series:
                series.data_labels.number_format = '0%'
                series.data_labels.number_format_is_linked = False
        break

    # Update text shapes
    text_shapes = _find_text_shapes(slide)
    title_set = False
    sample_set = False
    q1_set = False
    q2_set = False
    for s in text_shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        tl = text.lower()
        if ("total sample" in tl or "base n" in tl) and not sample_set:
            _set_text_preserve_format(s.text_frame, f"Total sample; base n = {sample_size}")
            sample_set = True
        elif not title_set and text.upper() == text and len(text) > 5 and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            title_set = True
        elif ("which of the following" in tl or "heard of" in tl) and not q1_set:
            _set_text_preserve_format(
                s.text_frame,
                "Which of the following brands have you heard of? (Select all that apply)",
            )
            q1_set = True
        elif ("purchased" in tl or "among the brands" in tl) and not q2_set:
            _set_text_preserve_format(
                s.text_frame,
                "Among the brands you indicated you've heard of, which ones have you ever purchased from? (Select all that apply)",
            )
            q2_set = True

    return slide


def _build_native_brand_matrix(prs, title: str, question: str,
                                brands: list, attributes: list,
                                data: list, sample_size: int = 200,
                                row_n: list = None):
    """Build a brand association heat-map table slide.

    Clones template slide 44 (TABLE 8×14). Fills with brands × attributes
    percentages and applies green gradient heat map coloring.
    Matches CozyFit reference: col 0 = "% Row Responses", cols 1..N = brands,
    col N+1 = "None", col N+2 = "Row n" (sample counts).

    Args:
        title: ALL-CAPS slide title
        question: Survey question text
        brands: List of brand names (column headers)
        attributes: List of attribute names (row headers)
        data: 2D list [attribute_idx][brand_idx] of integer percentages
        sample_size: Survey sample size
        row_n: Optional list of per-row sample counts (len = len(attributes))
    """
    from pptx.dml.color import RGBColor

    slide = _clone_slide(prs, 44)

    # Update title and question text
    text_shapes = _find_text_shapes(slide)
    title_set = False
    question_set = False
    for s in text_shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        if not title_set and text.upper() == text and len(text) > 5 and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            title_set = True
        elif not question_set and len(text) > 20:
            _set_text_preserve_format(s.text_frame, _truncate(question, 120))
            question_set = True

    # Compute "None" percentages (100 - sum of brand %s) and row_n defaults
    none_vals = []
    for ri, row_data in enumerate(data):
        row_sum = sum(row_data)
        none_pct = max(0, 100 - row_sum)
        none_vals.append(none_pct)
    if not row_n:
        # Estimate row_n from sample_size with slight variation
        import random
        base = int(sample_size * 0.48)  # ~half respond to perception Qs
        row_n = [base - i * 2 for i in range(len(attributes))]

    # Find and fill the table
    for s in slide.shapes:
        if not s.has_table:
            continue
        table = s.table
        n_rows = len(table.rows)
        n_cols = len(table.columns)

        # Layout: col 0 = "% Row Responses", cols 1..n_brands = brand names,
        #          col n_brands+1 = "None", col n_brands+2 = "Row n"
        n_brands = len(brands)
        none_col = n_brands + 1  # column index for "None"
        rown_col = n_brands + 2  # column index for "Row n"

        # Header row (row 0)
        # Col 0: keep "% Row Responses" from template or set it
        _replace_cell_text(table.cell(0, 0), "% Row Responses")

        # Brand columns
        for ci, brand in enumerate(brands):
            if ci + 1 < n_cols:
                _replace_cell_text(table.cell(0, ci + 1), brand[:15])

        # "None" column header
        if none_col < n_cols:
            _replace_cell_text(table.cell(0, none_col), "None")

        # "Row n" column header
        if rown_col < n_cols:
            _replace_cell_text(table.cell(0, rown_col), "Row n")

        # Clear any remaining header columns beyond Row n
        for ci in range(rown_col + 1, n_cols):
            _hdr_cell = table.cell(0, ci)
            _replace_cell_text(_hdr_cell, "")
            try:
                _hdr_cell.fill.background()
            except Exception:
                pass

        # Fill data rows
        for ri, attr in enumerate(attributes):
            if ri + 1 >= n_rows:
                break
            row_idx = ri + 1
            # Row label (col 0)
            _replace_cell_text(table.cell(row_idx, 0), attr)

            row_data = data[ri] if ri < len(data) else []

            # Brand data cells (cols 1..n_brands)
            for ci in range(n_brands):
                if ci + 1 >= n_cols:
                    break
                cell = table.cell(row_idx, ci + 1)
                val = row_data[ci] if ci < len(row_data) else 0
                _replace_cell_text(cell, f"{val}%")
                # Apply heat map color: 0% = lightest, 30%+ = darkest green
                intensity = min(val / 30.0, 1.0)
                r = int(255 - intensity * 120)  # 255→135
                g = int(255 - intensity * 50)   # 255→205
                b = int(255 - intensity * 120)  # 255→135
                try:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = RGBColor(r, g, b)
                except Exception:
                    pass

            # "None" cell
            if none_col < n_cols:
                cell = table.cell(row_idx, none_col)
                nv = none_vals[ri] if ri < len(none_vals) else 0
                _replace_cell_text(cell, f"{nv}%")
                intensity = min(nv / 30.0, 1.0)
                r = int(255 - intensity * 120)
                g = int(255 - intensity * 50)
                b = int(255 - intensity * 120)
                try:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = RGBColor(r, g, b)
                except Exception:
                    pass

            # "Row n" cell (no heat map — plain number)
            if rown_col < n_cols:
                cell = table.cell(row_idx, rown_col)
                rn = row_n[ri] if ri < len(row_n) else sample_size
                _replace_cell_text(cell, str(rn))
                # Clear any fill on Row n cell
                try:
                    cell.fill.background()
                except Exception:
                    pass

            # Clear any columns beyond Row n
            for ci in range(rown_col + 1, n_cols):
                _extra_cell = table.cell(row_idx, ci)
                _replace_cell_text(_extra_cell, "")
                try:
                    _extra_cell.fill.background()
                except Exception:
                    pass

        # Clear unused rows entirely
        for ri in range(len(attributes) + 1, n_rows):
            for ci in range(n_cols):
                _unused_cell = table.cell(ri, ci)
                _replace_cell_text(_unused_cell, "")
                try:
                    _unused_cell.fill.background()
                except Exception:
                    pass

        # Fit table within slide: dynamically size columns to avoid overflow
        from pptx.util import Emu
        used_cols = rown_col + 1  # cols 0 through rown_col are used

        # Collapse unused columns first
        for ci in range(used_cols, n_cols):
            table.columns[ci].width = Emu(0)

        # Calculate available width for used columns
        shape_left = s.left
        right_margin = Emu(200000)  # ~0.22in right margin
        available_width = prs.slide_width - shape_left - right_margin

        # Col 0 (row labels) gets ~14% of space, rest split evenly
        label_width = int(available_width * 0.14)
        data_col_width = (available_width - label_width) // max(used_cols - 1, 1)

        table.columns[0].width = label_width
        for ci in range(1, used_cols):
            table.columns[ci].width = data_col_width

        # Also constrain the shape width to match
        s.width = available_width
        break

    return slide


def _build_native_verbatim_table(prs, title: str, subtitle: str, quotes: list,
                                  sample_size: int = 200):
    """Build a verbatim quotes table slide (CozyFit slide 39 pattern).

    Args:
        title: ALL-CAPS slide title
        subtitle: Column header text
        quotes: List of verbatim quote strings
    """
    slide = _clone_slide(prs, 39)

    text_shapes = _find_text_shapes(slide)
    if text_shapes:
        _set_text_preserve_format(text_shapes[0].text_frame, _truncate(title, 55))

    # Find table shape
    from pptx.table import Table as _Table
    for s in slide.shapes:
        if s.has_table:
            table = s.table
            # First row is header
            if table.rows and len(table.rows) > 0:
                header_cell = table.cell(0, 0)
                if header_cell.text_frame.paragraphs:
                    _replace_para_text(header_cell.text_frame.paragraphs[0], subtitle)

            # Fill remaining rows with quotes
            for ri in range(1, len(table.rows)):
                if ri - 1 < len(quotes):
                    cell = table.cell(ri, 0)
                    if cell.text_frame.paragraphs:
                        _replace_para_text(cell.text_frame.paragraphs[0], f'"{quotes[ri-1]}"')
                else:
                    # Clear unused rows
                    cell = table.cell(ri, 0)
                    if cell.text_frame.paragraphs:
                        _replace_para_text(cell.text_frame.paragraphs[0], "")
            break

    return slide


# ── Main Generator ───────────────────────────────────────────

async def generate_pptx(
    project_id: int,
    analysis: dict,
    brand_name: str,
    phase: str = "full",
    collected_images: dict = None,
    brand_url: str = "",
    competitor_names: list[str] = None,
) -> tuple[Path, list[dict]]:
    """Generate a Brand Discovery PPTX from analysis data.

    Clones slides from the CozyFit reference template and replaces
    text content with analysis results. If collected_images is provided,
    replaces template images with brand-specific images.

    Args:
        collected_images: Output from image_collector.collect_images()
            {"brand": [Path], "product": [Path], "lifestyle": [Path], "all": [Path]}

    Returns:
        (pptx_path, slide_previews)
    """
    # Reset caches from any prior run in the same process
    _reset_caches()

    # Auto-discover images from disk if none passed
    if not collected_images:
        img_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
        if img_dir.exists():
            all_imgs = sorted(img_dir.glob("*.png")) + sorted(img_dir.glob("*.jpg")) + sorted(img_dir.glob("*.webp"))
            # Categorize by filename prefix
            brand_imgs = [p for p in all_imgs if any(p.name.startswith(pfx) for pfx in ("httpx_", "search_", "brand_"))]
            product_imgs = [p for p in all_imgs if any(p.name.startswith(pfx) for pfx in ("ecom_", "amazon_", "product_"))]
            lifestyle_imgs = [p for p in all_imgs if p.name.startswith("lifestyle_")]
            # Filter out cropped/segment/persona/hero/topic from all
            usable = [p for p in all_imgs if not any(
                p.name.startswith(pfx) for pfx in ("_cropped", "_fitted", "segment_bg", "persona_", "hero_", "topic_", "comp_summary")
            )]
            collected_images = {
                "brand": brand_imgs or usable,
                "product": product_imgs or usable,
                "lifestyle": lifestyle_imgs or usable,
                "all": usable,
            }
            print(f"[ppt_gen] Auto-discovered {len(usable)} images from {img_dir}")

            # Vision-filter to remove irrelevant images
            try:
                from pipeline.image_collector import _filter_relevant_images
                from PIL import Image as _PILImage
                relevant = await _filter_relevant_images(usable, brand_name)
                if relevant is not None:
                    relevant_set = set(relevant)
                    # Always preserve wide-aspect-ratio images (website screenshots)
                    # — they may serve as evidence on "brand confusion" slides
                    for p in usable:
                        if p not in relevant_set:
                            try:
                                w, h = _PILImage.open(p).size
                                if w / h > 2.5:
                                    relevant.append(p)
                                    relevant_set.add(p)
                                    print(f"[ppt_gen] Preserved screenshot: {p.name[:50]}")
                            except Exception:
                                pass
                    collected_images = {
                        "brand": [p for p in collected_images["brand"] if p in relevant_set],
                        "product": [p for p in collected_images["product"] if p in relevant_set],
                        "lifestyle": [p for p in collected_images["lifestyle"] if p in relevant_set],
                        "all": relevant,
                    }
                    print(f"[ppt_gen] After Vision filter: {len(relevant)} images kept")
            except Exception as e:
                print(f"[ppt_gen] Vision filter skipped: {e}")

    img_pool = _ImagePool(collected_images)

    # Start from the reference PPTX to get its theme, layouts, fonts
    prs = Presentation(str(TEMPLATE_PATH))

    # Remove ALL existing slides — we'll clone fresh ones
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    slide_meta = []
    date_str = analysis.get("date", "2026")

    # ── 1. Cover ──────────────────────────────────────────────
    _build_cover(prs, brand_name, date_str)
    slide_meta.append({"type": "cover", "content": {"brand_name": brand_name}})

    # ── 2. Agenda ─────────────────────────────────────────────
    slide = _build_agenda(prs)
    # Use dedicated hero product image (web search + AI generated, portrait orientation)
    hero_img = await _get_hero_product_image(brand_name, project_id, brand_url=brand_url)
    if hero_img:
        _replace_slide_image(slide, hero_img)
    elif img_pool.has_images():
        _replace_slide_image(slide, img_pool.next_brand())
    slide_meta.append({"type": "agenda", "content": {}})

    # ── 3. Approach ───────────────────────────────────────────
    _build_approach(prs)
    slide_meta.append({"type": "approach", "content": {}})

    # ── 4. Step 1 – Discovery ─────────────────────────────────
    _build_step_divider(prs)
    slide_meta.append({"type": "step", "content": {"step": 1}})

    # ── Capabilities ──────────────────────────────────────────

    _build_section_header(prs, "capabilities")
    slide_meta.append({"type": "section", "content": {"section": "capabilities"}})

    cap = analysis.get("capabilities", {})

    # Content slides for each capability dimension
    content_keys = [
        "execution_summary", "product_offer", "product_fundamentals",
        "pricing_position", "channel_analysis",
    ]
    # Alternate between template slide 6 and 7 for visual variety
    template_pool = [T_CONTENT, T_CONTENT_ALT]

    # Collect all slide topics for content-aware image assignment
    all_cap_topics = []
    all_cap_sections = []
    for key in content_keys:
        section = cap.get(key)
        if section:
            all_cap_topics.append({
                "title": section.get("title", key.replace("_", " ").upper()),
                "key": key,
                "bullets": section.get("bullets", []),
            })
            all_cap_sections.append(("content", section, key))
    for ci, challenge in enumerate(cap.get("brand_challenges", [])):
        all_cap_topics.append({
            "title": challenge.get("title", "BRAND CHALLENGE"),
            "key": f"challenge_{ci}",
            "bullets": challenge.get("bullets", []),
        })
        all_cap_sections.append(("challenge", challenge, f"challenge_{ci}"))

    # Assign images to slides using Vision or keyword matching
    image_assignments = {}
    if img_pool.has_images():
        image_assignments = await _assign_images_to_slides(
            img_pool, all_cap_topics, brand_name, project_id=project_id
        )

    # Build the slides with assigned images
    for slide_i, (section_type, section, key) in enumerate(all_cap_sections):
        tmpl = template_pool[slide_i % len(template_pool)]
        slide = _build_content_slide(
            prs,
            title=section.get("title", key.replace("_", " ").upper()),
            bullets=section.get("bullets", []),
            insight_text=section.get("insight", ""),
            template_idx=tmpl,
        )
        # Use content-aware assigned image, fallback to pool (with dedup)
        if slide_i in image_assignments:
            _replace_slide_image(slide, image_assignments[slide_i])
        elif img_pool.has_images():
            # Skip images already used by content-aware assignments
            used_in_assignments = set(image_assignments.values())
            fallback_img = img_pool.next_brand()
            attempts = 0
            while fallback_img and fallback_img in used_in_assignments and attempts < 20:
                fallback_img = img_pool.next_brand()
                attempts += 1
            if fallback_img:
                _replace_slide_image(slide, fallback_img)
        slide_meta.append({"type": "insight", "content": section})

    # Clarity Scoring slide — after content slides, before summary
    clarity_data = analysis.get("clarity_scoring")
    if not clarity_data:
        clarity_data = cap.get("clarity_scoring")
    if clarity_data and isinstance(clarity_data, dict) and clarity_data.get("dimensions"):
        _build_clarity_scoring_slide(prs, clarity_data)
        slide_meta.append({"type": "clarity_scoring", "content": clarity_data})

    # Capabilities summary — pick best unused product image, avoid website screenshots
    cap_summary = cap.get("capabilities_summary", "")
    if cap_summary:
        slide = _build_summary_slide(prs, "CAPABILITIES SUMMARY", cap_summary)
        if img_pool.has_images():
            used_in_caps = set(image_assignments.values())
            best_summary_img = _pick_clean_product_image(img_pool, used_in_caps)
            if best_summary_img:
                _replace_slide_image(slide, best_summary_img)
            else:
                _replace_slide_image(slide, img_pool.next_brand())
        slide_meta.append({"type": "summary", "content": {"text": cap_summary}})

    # ── Competition (Phase 2+) ────────────────────────────────

    if phase in ("market_structure", "full") and analysis.get("competition"):
        _build_section_header(prs, "competition")
        slide_meta.append({"type": "section", "content": {"section": "competition"}})

        comp = analysis.get("competition", {})
        competitor_list = comp.get("competitor_analyses", [])
        category = analysis.get("capabilities", {}).get("category", "")
        # Infer category from market overview title if not explicitly set
        if not category:
            overview = comp.get("market_overview", {})
            overview_title = overview.get("title", "")
            # Try title first — e.g., "KEY PLAYERS SHAPING THE INSULATED BOTTLE CATEGORY"
            title_lower = overview_title.lower()
            # Extract category dynamically from overview title
            # Pattern: "KEY PLAYERS SHAPING THE {CATEGORY} MARKET/CATEGORY/INDUSTRY"
            import re as _re
            _cat_match = _re.search(
                r'(?:shaping|defining|in|of)\s+(?:the\s+)?(.+?)\s+(?:market|category|industry|space|landscape|sector)',
                title_lower
            )
            if _cat_match:
                category = _cat_match.group(1).strip().title()
            # Fallback: try to extract from "A WELL-DEFINED {CATEGORY} MARKET" pattern
            if not category:
                _cat_match2 = _re.search(r'well.defined\s+(.+?)\s+market', title_lower)
                if _cat_match2:
                    category = _cat_match2.group(1).strip().title()
            # Fallback: scan bullets for "X market" or "X category" pattern
            if not category:
                for bullet in overview.get("bullets", []):
                    bullet_lower = str(bullet).lower()
                    _bm = _re.search(r'the\s+(.{3,30}?)\s+(?:market|category|industry|space)', bullet_lower)
                    if _bm:
                        category = _bm.group(1).strip().title()
                        break
        if not category:
            category = "category"

        # Fetch real competitor logos and per-competitor images
        all_comp_names = [c.get("name", "Competitor") for c in competitor_list[:12]]
        # Use full competitor list (from DB column) if available
        full_comp_names = competitor_names or all_comp_names
        # Select the most strategically relevant competitors via websearch + LLM
        focused_competitors = []  # Always re-select via LLM for best results
        if len(all_comp_names) > 6:
            try:
                from config import ANTHROPIC_API_KEY
                if ANTHROPIC_API_KEY:
                    import anthropic
                    _client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

                    # Step 1: Web search to understand brand positioning and competitive landscape
                    web_context = ""
                    try:
                        _ws_resp = _client.messages.create(
                            model=MODEL_OPUS,
                            max_tokens=1200,
                            tools=[{"type": "web_search_20250305", "name": "web_search", "max_uses": 5}],
                            messages=[{"role": "user", "content": (
                                f"You are a senior brand strategist. Research {brand_name} in the {category} market.\n\n"
                                f"Find out:\n"
                                f"1. {brand_name}'s price positioning (value/mid/premium/luxury) and price range\n"
                                f"2. Target consumer segments (demographics, psychographics)\n"
                                f"3. Core brand proposition and differentiation strategy\n"
                                f"4. Distribution channels (DTC, Amazon, wholesale, retail partners)\n"
                                f"5. Who industry analysts and consumers consider their TOP direct competitors\n"
                                f"6. Recent competitive dynamics (market share shifts, new entrants, M&A)\n\n"
                                f"Be specific with data points. Provide a concise summary (under 400 words)."
                            )}],
                        )
                        for block in _ws_resp.content:
                            if hasattr(block, "text"):
                                web_context = block.text
                                break
                        print(f"[competitor] Web research complete ({len(web_context)} chars)")
                    except Exception as e:
                        print(f"[competitor] Web search failed ({e}), using analysis data only")

                    # Step 2: Build context from each competitor's analysis data
                    comp_summaries = []
                    for c in competitor_list:
                        name = c.get("name", "?")
                        pos = c.get("positioning", [])
                        if isinstance(pos, list) and pos:
                            parts = []
                            for p in pos[:2]:
                                if isinstance(p, dict):
                                    parts.append(f"{p.get('label','')}: {str(p.get('detail',''))[:80]}")
                                else:
                                    parts.append(str(p)[:100])
                            pos_text = "; ".join(parts)
                        else:
                            pos_text = "no positioning data"
                        comp_summaries.append(f"- {name}: {pos_text}")
                    comp_context = "\n".join(comp_summaries)

                    # Step 3: LLM selects 6 most relevant using web research + analysis data
                    llm_prompt = (
                        f"You are a senior brand strategist building a competitive landscape presentation for {brand_name} "
                        f"in the {category} category.\n\n"
                    )
                    if web_context:
                        llm_prompt += f"BRAND INTELLIGENCE (from web research):\n{web_context}\n\n"
                    llm_prompt += (
                        f"COMPETITOR PROFILES:\n{comp_context}\n\n"
                        f"Select exactly 6 competitors that would give {brand_name}'s leadership team the most actionable strategic insight. "
                        f"Your selection should cover:\n"
                        f"- 2-3 direct head-to-head competitors (same price tier, same consumer, fighting for the same purchase occasion)\n"
                        f"- 1-2 category leaders/incumbents whose scale and brand equity set the benchmark\n"
                        f"- 1 aspirational or disruptive brand whose strategy reveals a white-space opportunity or emerging threat\n\n"
                        f"Exclude brands that operate in adjacent categories unless they directly compete for {brand_name}'s core consumer.\n\n"
                        f"Reply with ONLY the 6 brand names, comma-separated."
                    )
                    _resp = _client.messages.create(
                        model=MODEL_OPUS,
                        max_tokens=150,
                        messages=[{"role": "user", "content": llm_prompt}],
                    )
                    answer = _resp.content[0].text.strip()
                    print(f"[competitor] LLM focused selection: {answer}")
                    # Parse names — match against actual competitor names
                    all_lower_map = {n.lower(): n for n in all_comp_names}
                    for part in answer.split(","):
                        part = part.strip().strip('"').strip("'")
                        matched = all_lower_map.get(part.lower())
                        if matched and matched not in focused_competitors:
                            focused_competitors.append(matched)
                    focused_competitors = focused_competitors[:6]
            except Exception as e:
                print(f"[competitor] LLM focused selection failed: {e}")

        if not focused_competitors:
            focused_competitors = all_comp_names[:6]

        comp_img_dir = OUTPUT_DIR / f"project_{project_id}" / "images"
        comp_img_dir.mkdir(parents=True, exist_ok=True)

        # Fetch logos for all competitors (async — we're inside async generate_pptx)
        logos = await _fetch_competitor_logos(full_comp_names, comp_img_dir)
        print(f"[competitor] Logos ready: {len(logos)}/{len(full_comp_names)}")

        # Collect per-competitor images only for focused competitors (deep-dive slides)
        competitor_images = {}  # {name: [Path, ...]}
        focused_for_imgs = focused_competitors or all_comp_names[:6]
        for competitor in competitor_list:
            cname = competitor.get("name", "Competitor")
            if cname.lower() not in {n.lower() for n in focused_for_imgs}:
                continue
            imgs = await _collect_competitor_images(cname, comp_img_dir, num_images=4, category=category)
            competitor_images[cname] = imgs
        print(f"[competitor] Images collected for {len(competitor_images)} competitors")

        # Build full competitor dicts for grid (include those not in analyses)
        grid_competitors = list(competitor_list)  # start with analyzed ones
        analyzed_names_lower = {c.get("name", "").lower() for c in competitor_list}
        for name in full_comp_names:
            if name.lower() not in analyzed_names_lower:
                grid_competitors.append({"name": name})

        # Reorder grid: focused competitors first for consistent positioning
        # (slide 16 template has a red box around the first 6 positions)
        if focused_competitors:
            focused_lower = {n.lower() for n in focused_competitors}
            focused_comps = [c for c in grid_competitors if c.get("name", "").lower() in focused_lower]
            other_comps = [c for c in grid_competitors if c.get("name", "").lower() not in focused_lower]
            grid_competitors = focused_comps + other_comps

        # 1. All-competitor overview grid (CozyFit slide 16 pattern)
        if len(grid_competitors) >= 2:
            slide = _build_competitor_overview(prs, grid_competitors, category, logos=logos, focused_names=focused_competitors)
            slide_meta.append({"type": "competitor_grid", "content": {"title": "market overview", "count": len(grid_competitors)}})

        # 2. Focused review grid — highlighted brands with red boxes
        if len(grid_competitors) >= 2:
            _build_competitor_focused(
                prs, grid_competitors, category,
                logos=logos, focused_names=focused_competitors,
            )
            slide_meta.append({"type": "competitor_focused", "content": {"count": len(grid_competitors)}})

        # 3. Competitor deep dives — only for the 6 focused competitors
        #    T_COMPETITOR (17): 1 wide banner image
        #    T_COMPETITOR_2IMG (18): 2 side-by-side images
        #    T_COMPETITOR_4IMG (19): 4 images in a row
        focused_lower_set = {n.lower() for n in focused_competitors} if focused_competitors else set()
        deep_dive_list = [
            c for c in competitor_list
            if c.get("name", "").lower() in focused_lower_set
        ] if focused_lower_set else competitor_list[:6]

        for ci, competitor in enumerate(deep_dive_list):
            comp_name = competitor.get("name", "Competitor")
            comp_imgs_available = len(competitor_images.get(comp_name, []))
            # Default to 2-image layout (1 lifestyle/ad + 1 product) for best visual impact
            # Only fall back to 1-image if we truly have just one
            if comp_imgs_available >= 2:
                tmpl = T_COMPETITOR_2IMG
            else:
                tmpl = T_COMPETITOR

            pos_bullets = [
                f"{p['label']}: {p['detail']}" if isinstance(p, dict) else str(p)
                for p in competitor.get("positioning", [])
            ]
            learn_bullets = [
                f"{k['label']}: {k['detail']}" if isinstance(k, dict) else str(k)
                for k in competitor.get("key_learnings", [])
            ]
            slide = _build_competitor_slide(
                prs,
                name=comp_name,
                positioning_bullets=pos_bullets,
                learnings_bullets=learn_bullets,
                template_idx=tmpl,
            )
            # Replace picture shapes with per-competitor images (cover mode for photos)
            comp_imgs = competitor_images.get(comp_name, [])
            if comp_imgs:
                from pptx.shapes.picture import Picture
                SLIDE_AREA = 12192000 * 6858000
                pics = [s for s in slide.shapes if isinstance(s, Picture)
                        and (s.width * s.height) / SLIDE_AREA < 0.9]
                for pi in range(min(len(pics), len(comp_imgs))):
                    if comp_imgs[pi].exists():
                        _replace_nth_picture(slide, comp_imgs[pi], pi, mode="cover")
            elif img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_product())
            slide_meta.append({"type": "competitor", "content": competitor})

        # Landscape summary
        landscape = comp.get("landscape_summary", {})
        if landscape:
            roles = landscape.get("market_roles", [])
            role_bullets = [
                f"{r['role']}: {', '.join(r.get('brands', []))} — {r.get('description', '')}"
                for r in roles[:4]
            ]
            white_space = landscape.get("white_space", "")
            sidebar = f"White Space Opportunity:\n{white_space}" if white_space else ""
            slide = _build_landscape_slide(
                prs,
                title="COMPETITIVE LANDSCAPE ROLES",
                bullets=role_bullets or ["No market roles identified"],
                sidebar_text=sidebar,
            )
            slide_meta.append({"type": "landscape", "content": landscape})

        # Competition summary — prefer clean product image, fall back to brand search
        comp_summary = comp.get("competition_summary", "")
        if comp_summary:
            slide = _build_summary_slide(prs, "COMPETITION SUMMARY", comp_summary, T_COMP_SUMMARY)
            # First try: clean product image from pool (no text overlays)
            comp_summary_img = _pick_clean_product_image(img_pool)
            if not comp_summary_img:
                # Second try: web search for brand advertising image
                try:
                    comp_summary_img = await _collect_brand_summary_image(
                        brand_name, comp_img_dir, category=category,
                    )
                except Exception as e:
                    print(f"[comp_summary] Brand image search failed: {e}")
            if comp_summary_img and Path(comp_summary_img).exists():
                _replace_slide_image(slide, comp_summary_img)
            else:
                best = img_pool.best_brand()
                if best:
                    _replace_slide_image(slide, best)
                else:
                    _replace_slide_image(slide, img_pool.next_product())
            slide_meta.append({"type": "summary", "content": {"text": comp_summary}})

    # ── Consumer (Full only) ──────────────────────────────────

    target = {}  # Will be populated in consumer section for use by final summary

    if phase == "full" and analysis.get("consumer"):
        _build_section_header(prs, "consumer")
        slide_meta.append({"type": "section", "content": {"section": "consumer"}})

        consumer = analysis.get("consumer", {})

        # Research approach — always include (use default if analysis doesn't provide one)
        research = consumer.get("research_approach", [])
        if not research:
            # Generate adaptive methodology from survey designer
            research = _generate_research_approach(brand_name, analysis, date_str)
        _build_research_approach(prs, research)
        slide_meta.append({"type": "research", "content": {"items": research}})

        # Key consumer insights as content slides
        for insight in consumer.get("key_insights", []):
            slide = _build_content_slide(
                prs,
                title=insight.get("title", "CONSUMER INSIGHT"),
                bullets=insight.get("bullets", []),
                insight_text=insight.get("insight", ""),
            )
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_lifestyle())
            slide_meta.append({"type": "insight", "content": insight})

        # ── Evidence Plan slide ───────────────────────────────
        evidence_plan = analysis.get("evidence_plan")
        if evidence_plan and isinstance(evidence_plan, dict):
            _build_evidence_plan_slide(prs, evidence_plan)
            slide_meta.append({"type": "evidence_plan", "content": evidence_plan})

        # ── Hypothesis Validation slide ──────────────────────
        hyp_val = analysis.get("hypothesis_validation")
        if hyp_val and isinstance(hyp_val, list) and len(hyp_val) > 0:
            _build_hypothesis_validation_slide(prs, hyp_val)
            slide_meta.append({"type": "hypothesis_validation", "content": hyp_val})

        # ── Questionnaire / Chart Slides ─────────────────────
        charts = consumer.get("charts", [])
        if charts:
            from pipeline.chart_renderer import render_chart
            chart_output_dir = OUTPUT_DIR / f"project_{project_id}" / "charts"

            # Demographics divider
            _build_chart_divider(prs, T_CHART_DIVIDER_DEMO, "Demographics &\nBackground")
            slide_meta.append({"type": "divider", "content": {"title": "Demographics & Background"}})

            # ── Demographics slides (CozyFit slides 28-32 pattern) ──
            # Detect generation, ethnicity, income, and social media charts
            # to render them using specialized slide templates.
            demographics_indices = set()
            gen_chart_idx = None
            eth_chart_idx = None
            combined_gender_eth_idx = None  # LLM may combine gender+ethnicity
            income_chart_idx = None
            social_chart_idx = None
            occupation_chart_indices = []  # category-specific: occupation/role charts
            for ci, cd in enumerate(charts):
                tl = cd.get("title", "").lower()
                ct = cd.get("chart_type", cd.get("type", "")).lower()
                if gen_chart_idx is None and any(kw in tl for kw in ("generation", "age group", "age distribution")):
                    gen_chart_idx = ci
                elif eth_chart_idx is None and any(kw in tl for kw in ("ethnicit", "race")):
                    # Check if this is a combined gender+ethnicity chart
                    if "gender" in tl or ct in ("dual_donut", "dual"):
                        combined_gender_eth_idx = ci
                    else:
                        eth_chart_idx = ci
                elif income_chart_idx is None and any(kw in tl for kw in ("income", "household income")):
                    income_chart_idx = ci
                elif social_chart_idx is None and any(kw in tl for kw in ("social media", "platform usage")):
                    social_chart_idx = ci
                elif any(kw in tl for kw in ("occupation", "work detail", "work environment",
                         "shift", "work setting", "activity level", "role")):
                    occupation_chart_indices.append(ci)

            # Extract sample size from research approach
            sample_n = 200
            for ra in consumer.get("research_approach", []):
                if isinstance(ra, dict):
                    detail = str(ra.get("detail", ""))
                    import re as _re
                    _n_match = _re.search(r'n\s*[=:]\s*(\d+)', detail)
                    if _n_match:
                        sample_n = int(_n_match.group(1))
                        break

            # Handle combined gender+ethnicity chart: split into gender_data + standalone ethnicity vbar
            gender_data = consumer.get("gender_data")
            if combined_gender_eth_idx is not None:
                demographics_indices.add(combined_gender_eth_idx)
                combo = charts[combined_gender_eth_idx]
                sub_charts = combo.get("charts", [])
                # Extract gender data from sub-chart
                if not gender_data:
                    for sc in sub_charts:
                        if sc.get("label", "").lower() == "gender":
                            gd = sc.get("data", [])
                            male = next((d["value"] for d in gd if "male" in d.get("label", "").lower() and "female" not in d.get("label", "").lower()), 30)
                            female = next((d["value"] for d in gd if "female" in d.get("label", "").lower()), 70)
                            gender_data = {"male_pct": male, "female_pct": female}
                            break
                # Extract ethnicity data → create standalone vbar chart for rendering
                for sc in sub_charts:
                    if sc.get("label", "").lower() in ("ethnicity", "race"):
                        eth_data = sc.get("data", [])
                        eth_standalone = {
                            "chart_type": "vbar",
                            "title": "RACE / ETHNICITY",
                            "categories": [d["label"] for d in eth_data],
                            "values": [d["value"] for d in eth_data],
                        }
                        eth_chart_idx = f"_synth_eth"  # synthetic index
                        break

            if not gender_data:
                gender_data = {"male_pct": 30, "female_pct": 70}

            # ── Helper: extract categories/values from LLM chart data ──
            def _extract_chart_data(cd):
                """Extract {categories, values} from various LLM chart formats."""
                cats = cd.get("categories", [])
                vals = cd.get("values", [])
                if not cats and "data" in cd:
                    data = cd["data"]
                    if isinstance(data, list) and data and isinstance(data[0], dict):
                        cats = [d.get("label", "") for d in data]
                        vals = [d.get("value", 0) for d in data]
                return {"categories": cats, "values": vals}

            # Slide 28: RESPONDENT PROFILE (generation hbar + ethnicity vbar + gender icon)
            gen_native = None
            eth_native = None
            if gen_chart_idx is not None:
                demographics_indices.add(gen_chart_idx)
                gen_native = _extract_chart_data(charts[gen_chart_idx])
            if eth_chart_idx == "_synth_eth":
                eth_native = {
                    "categories": eth_standalone.get("categories", []),
                    "values": eth_standalone.get("values", []),
                }
            elif eth_chart_idx is not None:
                demographics_indices.add(eth_chart_idx)
                eth_native = _extract_chart_data(charts[eth_chart_idx])
            if gen_native or eth_native:
                _build_respondent_profile(prs, generation_data=gen_native,
                                          ethnicity_data=eth_native,
                                          gender_data=gender_data, sample_size=sample_n)
                slide_meta.append({"type": "respondent_profile", "content": {
                    "gender_data": gender_data, "sample_size": sample_n}})

            # Slide 29: RESPONDENT PROFILE (CONTINUED) — marital status + household income
            if income_chart_idx is not None:
                demographics_indices.add(income_chart_idx)
                income_native = _extract_chart_data(charts[income_chart_idx])
                marital_data = consumer.get("marital_data", {"married_pct": 52, "single_pct": 33, "divorced_pct": 15})
                _build_respondent_profile_continued(prs, income_data=income_native,
                                                    sample_size=sample_n,
                                                    marital_data=marital_data)
                slide_meta.append({"type": "respondent_profile_continued", "content": {
                    "marital_data": marital_data, "sample_size": sample_n}})

            # Slides 30-31: OCCUPATION AND WORK DETAILS (category-specific)
            if occupation_chart_indices:
                for oi in occupation_chart_indices:
                    demographics_indices.add(oi)
                # First occupation chart → slide 30 (single vbar)
                first_occ = _extract_chart_data(charts[occupation_chart_indices[0]])
                first_occ["title"] = charts[occupation_chart_indices[0]].get("title", "OCCUPATION AND WORK DETAILS")
                _build_occupation_slide(prs, occupation_data=first_occ, sample_size=sample_n)
                slide_meta.append({"type": "occupation", "content": first_occ})
                # Remaining occupation charts → slide 31 (doughnut + 2 vbar)
                if len(occupation_chart_indices) > 1:
                    detail_list = []
                    for oi in occupation_chart_indices[1:4]:  # max 3 charts on slide 31
                        d = _extract_chart_data(charts[oi])
                        d["title"] = charts[oi].get("title", "")
                        detail_list.append(d)
                    _build_occupation_detail_slide(prs, detail_charts=detail_list, sample_size=sample_n)
                    slide_meta.append({"type": "occupation_detail", "content": {"charts": len(detail_list)}})

            # Slide 32: SOCIAL MEDIA PLATFORM USAGE (icons + hbar)
            if social_chart_idx is not None:
                demographics_indices.add(social_chart_idx)
                social_native = _extract_chart_data(charts[social_chart_idx])
                _build_social_media_slide(prs, social_media_data=social_native,
                                          sample_size=sample_n)
                slide_meta.append({"type": "social_media", "content": {
                    "sample_size": sample_n}})

            # ── Dynamic demographics charts from survey question_data ──
            # Render any Section 1 (demographics) questions that aren't covered
            # by the native chart builders above (generation, ethnicity, income,
            # social_media, gender, marital, occupation).
            _native_demo_keywords = {"generation", "age", "ethnic", "race", "income",
                                     "social media", "platform", "gender", "marital",
                                     "occupation", "work detail", "work environment",
                                     "shift", "work setting", "role"}
            survey_qdata = analysis.get("survey_question_data", {})
            from pipeline.chart_renderer import render_chart as _render_chart_demo
            _demo_chart_dir = OUTPUT_DIR / f"project_{project_id}" / "charts"
            if survey_qdata:
                _demo_extra_idx = 0
                for qid, qd in survey_qdata.items():
                    if qd.get("section", "").lower() != "demographics":
                        continue
                    # Skip if already handled natively
                    qt_lower = qd.get("chart_title", qd.get("question_text", "")).lower()
                    if any(kw in qt_lower for kw in _native_demo_keywords):
                        continue
                    cats = qd.get("categories", [])
                    vals = qd.get("values", [])
                    if not cats or not vals:
                        continue
                    chart_title = qd.get("chart_title", qd.get("question_text", qid)).upper()
                    # Use native hbar/vbar chart instead of matplotlib
                    t_idx = 34 if _demo_extra_idx % 2 == 0 else 38
                    _build_native_single_hbar(
                        prs, title=_truncate(chart_title, 55),
                        subtitle=qd.get("question_text", qd.get("chart_title", "")),
                        data={"categories": cats, "values": vals},
                        sample_size=sample_n, template_idx=t_idx,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": chart_title}})
                    print(f"[demographics] Native chart: {chart_title} (template {t_idx})")
                    _demo_extra_idx += 1

            # ── Collect survey questions by section for native chart rendering ──
            _shopping_qs = []  # (qid, qd) tuples for shopping section
            _drivers_qs = []
            _brand_qs = []
            _other_section_qs = []  # lifestyle etc
            if survey_qdata:
                for qid, qd in survey_qdata.items():
                    sec = qd.get("section", "").lower()
                    cats = qd.get("categories", [])
                    vals = qd.get("values", [])
                    if sec == "demographics" or not sec or (not cats and not vals):
                        continue
                    # Skip screener questions (Yes/No with extreme distribution)
                    if len(cats) <= 2 and vals and max(vals) >= 95:
                        print(f"[survey] Skip screener: {qd.get('chart_title', qid)} ({vals})")
                        continue
                    if sec == "shopping":
                        _shopping_qs.append((qid, qd))
                    elif sec == "drivers":
                        _drivers_qs.append((qid, qd))
                    elif sec == "brand":
                        _brand_qs.append((qid, qd))
                    else:
                        _other_section_qs.append((qid, qd))

            # ── Build native Shopping slides from survey data ──
            shopping_inserted = False
            native_shopping_built = False
            if _shopping_qs:
                _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Shopping Habits, Usage,\nAttitude and Image")
                slide_meta.append({"type": "divider", "content": {"title": "Shopping Habits"}})
                shopping_inserted = True

                # Classify questions: ≤5 categories → donut candidate, >5 → hbar
                donut_candidates = []
                hbar_candidates = []
                for qid, qd in _shopping_qs:
                    n_cats = len(qd.get("categories", []))
                    ct = qd.get("chart_type", "hbar").lower()
                    if ct == "donut" or (ct in ("vbar", "pie") and n_cats <= 6):
                        donut_candidates.append((qid, qd))
                    else:
                        hbar_candidates.append((qid, qd))

                # Strategy: always pair into dual slides for visual richness.
                # Prefer donut+hbar pairs; fallback to hbar+hbar; single for odd remainder.
                _slide_count = 0

                # Phase 1: Pair donut + hbar → dual donut+hbar slides
                while donut_candidates and hbar_candidates:
                    dq_id, dq = donut_candidates.pop(0)
                    hq_id, hq = hbar_candidates.pop(0)
                    t_idx = 35 if _slide_count % 2 == 0 else 37
                    slide_title = dq.get("chart_title", dq.get("question_text", "")).upper()
                    _build_native_dual_donut_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=f"All data reflect past 12 months  |  n = {sample_n}",
                        left_data={"categories": dq["categories"], "values": dq["values"]},
                        right_data={"categories": hq["categories"], "values": hq["values"]},
                        left_question=dq.get("question_text", dq.get("chart_title", "")),
                        right_question=hq.get("question_text", hq.get("chart_title", "")),
                        sample_size=sample_n, template_idx=t_idx,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[shopping] Native dual donut+hbar: {slide_title} (template {t_idx})")
                    _slide_count += 1

                # Phase 2: Pair remaining → dual hbar slides
                remaining = donut_candidates + hbar_candidates
                while len(remaining) >= 2:
                    lq_id, lq = remaining.pop(0)
                    rq_id, rq = remaining.pop(0)
                    slide_title = lq.get("chart_title", lq.get("question_text", "")).upper()
                    _build_native_dual_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=f"All data reflect past 12 months  |  n = {sample_n}",
                        left_data={"categories": lq["categories"], "values": lq["values"]},
                        right_data={"categories": rq["categories"], "values": rq["values"]},
                        left_question=lq.get("question_text", lq.get("chart_title", "")),
                        right_question=rq.get("question_text", rq.get("chart_title", "")),
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[shopping] Native dual hbar: {slide_title}")
                    _slide_count += 1

                # Phase 3: Odd remainder → single hbar slide
                for sq_id, sq in remaining:
                    t_idx = 34 if _slide_count % 2 == 0 else 38
                    slide_title = sq.get("chart_title", sq.get("question_text", "")).upper()
                    _build_native_single_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=sq.get("question_text", sq.get("chart_title", "")),
                        data={"categories": sq["categories"], "values": sq["values"]},
                        sample_size=sample_n, template_idx=t_idx,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[shopping] Native single hbar: {slide_title} (template {t_idx})")
                    _slide_count += 1

                # Phase 4: Verbatim table — from survey_verbatim or survey_qdata
                _verbatim_quotes = []
                _verbatim_title = "CHALLENGES WHEN PURCHASING"
                _verbatim_subtitle = "What frustrates you most about buying in this category?"
                # Check survey_verbatim first
                survey_verbatim = analysis.get("survey_verbatim", {})
                for vkey, vdata in survey_verbatim.items():
                    if vdata.get("section", "").lower() == "shopping" or "challenge" in vkey.lower() or "frustrat" in vkey.lower():
                        _verbatim_quotes = vdata.get("quotes", [])
                        _verbatim_title = vdata.get("chart_title", _verbatim_title).upper()
                        _verbatim_subtitle = vdata.get("question_text", _verbatim_subtitle)
                        break
                if not _verbatim_quotes:
                    # Fallback: check question_data for open_ended/verbatim type
                    for qid, qd in survey_qdata.items():
                        if qd.get("section", "").lower() != "shopping":
                            continue
                        if qd.get("chart_type", "").lower() in ("verbatim", "open_ended", "table"):
                            _verbatim_quotes = qd.get("quotes", qd.get("values", []))
                            _verbatim_title = qd.get("chart_title", _verbatim_title).upper()
                            _verbatim_subtitle = qd.get("question_text", _verbatim_subtitle)
                            break
                if _verbatim_quotes and isinstance(_verbatim_quotes[0], str):
                    _build_native_verbatim_table(
                        prs,
                        title=_truncate(_verbatim_title, 55),
                        subtitle=_verbatim_subtitle,
                        quotes=_verbatim_quotes[:26],
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "verbatim", "content": {"title": _verbatim_title}})
                    print(f"[shopping] Native verbatim table: {_verbatim_title}")
                    _slide_count += 1

                native_shopping_built = True
                print(f"[shopping] Built {_slide_count} native chart slides from survey data")

            # ── Build native Purchase Drivers slides from survey data ──
            native_drivers_built = False
            if _drivers_qs:
                _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Purchase Drivers\n& Barriers")
                slide_meta.append({"type": "divider", "content": {"title": "Purchase Drivers"}})

                # Phase 1: Pair drivers questions into dual hbar slides, singles for remainder
                _drv_slide_count = 0
                drv_list = list(_drivers_qs)
                while len(drv_list) >= 2:
                    lq_id, lq = drv_list.pop(0)
                    rq_id, rq = drv_list.pop(0)
                    slide_title = lq.get("chart_title", lq.get("question_text", "")).upper()
                    _build_native_dual_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=f"All data reflect past 12 months  |  n = {sample_n}",
                        left_data={"categories": lq["categories"], "values": lq["values"]},
                        right_data={"categories": rq["categories"], "values": rq["values"]},
                        left_question=lq.get("question_text", lq.get("chart_title", "")),
                        right_question=rq.get("question_text", rq.get("chart_title", "")),
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[drivers] Native dual hbar: {slide_title}")
                    _drv_slide_count += 1

                for sq_id, sq in drv_list:
                    t_idx = 34 if _drv_slide_count % 2 == 0 else 38
                    slide_title = sq.get("chart_title", sq.get("question_text", "")).upper()
                    _build_native_single_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=sq.get("question_text", sq.get("chart_title", "")),
                        data={"categories": sq["categories"], "values": sq["values"]},
                        sample_size=sample_n, template_idx=t_idx,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[drivers] Native single hbar: {slide_title}")
                    _drv_slide_count += 1

                # Phase 2: Drivers verbatim (wishlist feedback)
                survey_verbatim = analysis.get("survey_verbatim", {})
                for vkey, vdata in survey_verbatim.items():
                    if vdata.get("section", "").lower() in ("drivers", "brand") and "wishlist" in vkey.lower():
                        vquotes = vdata.get("quotes", [])
                        if vquotes and isinstance(vquotes[0], str):
                            _build_native_verbatim_table(
                                prs,
                                title=_truncate(vdata.get("chart_title", "BRAND WISHLIST FEEDBACK").upper(), 55),
                                subtitle=vdata.get("question_text", "What would you change?"),
                                quotes=vquotes[:26],
                                sample_size=sample_n,
                            )
                            slide_meta.append({"type": "verbatim", "content": {"title": vdata.get("chart_title", "")}})
                            print(f"[drivers] Native verbatim: {vdata.get('chart_title', 'wishlist')}")
                            _drv_slide_count += 1
                            break

                native_drivers_built = True
                print(f"[drivers] Built {_drv_slide_count} native chart slides from survey data")

            # ── Build native Brand Evaluation slides from survey data ──
            native_brand_built = False
            brand_awareness_data = analysis.get("survey_brand_awareness", {})
            brand_matrix_data = analysis.get("survey_brand_matrix", {})

            if _brand_qs or brand_awareness_data or brand_matrix_data:
                _build_chart_divider(prs, T_CHART_DIVIDER_BRAND, "Brand Evaluation &\nCompetitor Analysis")
                slide_meta.append({"type": "divider", "content": {"title": "Brand Evaluation"}})
                _brand_metrics_def_done = False
                _brd_slide_count = 0

                # Phase 1: Brand Awareness stacked bar
                if brand_awareness_data:
                    _build_native_stacked_bar(
                        prs,
                        title="BRAND AWARENESS AND PURCHASE CONVERSION",
                        brands=brand_awareness_data["brands"],
                        awareness=brand_awareness_data["awareness"],
                        purchase=brand_awareness_data["purchase"],
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": "BRAND AWARENESS AND PURCHASE CONVERSION"}})
                    print(f"[brand] Native stacked bar: BRAND AWARENESS ({len(brand_awareness_data['brands'])} brands)")
                    _brd_slide_count += 1

                # Phase 2: Brand association matrix heat map
                if brand_matrix_data:
                    # Derive category-specific title (e.g., "WATER BOTTLE BRAND ASSOCIATION")
                    _matrix_title = brand_matrix_data.get("title", "").upper()
                    if not _matrix_title:
                        # Try to extract category from question or brand context
                        _q_text = brand_matrix_data.get("question", "")
                        _matrix_title = f"{brand_name.upper()} CATEGORY BRAND ASSOCIATION"
                    _build_native_brand_matrix(
                        prs,
                        title=_truncate(_matrix_title, 55),
                        question=brand_matrix_data.get("question",
                            "Based on your perception, which brand best fits each description?"),
                        brands=brand_matrix_data["brands"],
                        attributes=brand_matrix_data["attributes"],
                        data=brand_matrix_data["data"],
                        sample_size=sample_n,
                        row_n=brand_matrix_data.get("row_n"),
                    )
                    slide_meta.append({"type": "brand_matrix", "content": {"title": _matrix_title}})
                    print(f"[brand] Native matrix: {len(brand_matrix_data['attributes'])}×{len(brand_matrix_data['brands'])}")
                    _brd_slide_count += 1

                # Phase 3: Brand survey questions
                # Strategy: pair first two into a donut+hbar combo (like CozyFit
                # slide 43 employer required/recommended), rest as single hbar
                brd_list = list(_brand_qs)
                if len(brd_list) >= 2:
                    lq_id, lq = brd_list.pop(0)
                    rq_id, rq = brd_list.pop(0)
                    slide_title = lq.get("chart_title", lq.get("question_text", "")).upper()
                    _build_native_dual_donut_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=f"All data reflect past 12 months  |  n = {sample_n}",
                        left_data={"categories": lq["categories"], "values": lq["values"]},
                        right_data={"categories": rq["categories"], "values": rq["values"]},
                        left_question=lq.get("question_text", lq.get("chart_title", "")),
                        right_question=rq.get("question_text", rq.get("chart_title", "")),
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[brand] Native donut+hbar: {slide_title}")
                    _brd_slide_count += 1

                # Remaining brand questions as single hbar
                for sq_id, sq in brd_list:
                    slide_title = sq.get("chart_title", sq.get("question_text", "")).upper()
                    _build_native_single_hbar(
                        prs, title=_truncate(slide_title, 55),
                        subtitle=sq.get("question_text", sq.get("chart_title", "")),
                        data={"categories": sq["categories"], "values": sq["values"]},
                        sample_size=sample_n, template_idx=34,
                    )
                    slide_meta.append({"type": "native_chart", "content": {"title": slide_title}})
                    print(f"[brand] Native single hbar: {slide_title}")
                    _brd_slide_count += 1

                # Phase 4: Brand wishlist verbatim table (skip if already built in drivers Phase 2)
                _brand_verbatim = analysis.get("survey_verbatim", {}).get("brand_wishlist", {})
                _wishlist_already_built = any(
                    m.get("type") == "verbatim" and "wishlist" in m.get("content", {}).get("title", "").lower()
                    for m in slide_meta
                )
                if _brand_verbatim and _brand_verbatim.get("quotes") and not _wishlist_already_built:
                    _bv_title = _brand_verbatim.get("chart_title", "BRAND WISHLIST FEEDBACK").upper()
                    _bv_question = _brand_verbatim.get("question_text",
                        f"What do you wish {brand_name} or other brands in this category did better or differently?")
                    _build_native_verbatim_table(
                        prs,
                        title=_bv_title,
                        subtitle=_bv_question,
                        quotes=_brand_verbatim["quotes"],
                        sample_size=sample_n,
                    )
                    slide_meta.append({"type": "verbatim", "content": {"title": _bv_title}})
                    print(f"[brand] Native verbatim: {_bv_title}")
                    _brd_slide_count += 1

                native_brand_built = True
                print(f"[brand] Built {_brd_slide_count} native chart slides from survey data")

            # ── Fallback: main chart loop for any sections not handled natively ──
            if not native_drivers_built or not native_brand_built:
                drivers_inserted = native_drivers_built
                brand_inserted = native_brand_built
                for ci, chart_data in enumerate(charts):
                    if ci in demographics_indices:
                        continue

                    chart_type = chart_data.get("chart_type", chart_data.get("type", "hbar"))
                    title_lower = chart_data.get("title", "").lower()
                    section_lower = chart_data.get("section", "").lower()

                    is_shopping = section_lower in ("shopping habits", "shopping") or any(kw in title_lower for kw in ("shopping", "habit", "frequency", "spend", "channel", "purchase channel", "occasion", "pre-purchase", "usage"))
                    is_driver = section_lower in ("purchase drivers", "drivers") or any(kw in title_lower for kw in ("driver", "matters most", "premium", "willingness", "pay", "wordcloud", "say about", "pain point"))
                    is_brand = section_lower in ("brand evaluation", "brand") or any(kw in title_lower for kw in ("brand", "awareness", "metric", "association", "matrix", "perception", "favorite", "likelihood", "switching"))

                    if native_shopping_built and is_shopping:
                        continue
                    if native_drivers_built and is_driver:
                        continue
                    if native_brand_built and is_brand:
                        continue

                    chart_path = render_chart(chart_data, chart_output_dir, ci)

                    if not shopping_inserted and ci >= 1 and (is_shopping or ci == 4):
                        _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Shopping Habits, Usage,\nAttitude and Image")
                        slide_meta.append({"type": "divider", "content": {"title": "Shopping Habits"}})
                        shopping_inserted = True
                    elif not drivers_inserted and (shopping_inserted or native_shopping_built) and (is_driver or ci >= 8):
                        _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Purchase Drivers\n& Barriers")
                        slide_meta.append({"type": "divider", "content": {"title": "Purchase Drivers"}})
                        drivers_inserted = True
                    elif not brand_inserted and (shopping_inserted or drivers_inserted or native_shopping_built) and (is_brand or ci >= len(charts) - 3):
                        _build_chart_divider(prs, T_CHART_DIVIDER_BRAND, "Brand Evaluation &\nCompetitor Analysis")
                        slide_meta.append({"type": "divider", "content": {"title": "Brand Evaluation"}})
                        brand_inserted = True

                    if chart_path is None:
                        continue

                    _build_chart_slide(prs, chart_data, chart_path)
                    slide_meta.append({"type": "chart", "content": chart_data})

        # Segmentation divider + intro boilerplate
        segments = consumer.get("segments", [])
        if segments:
            # Sort segments by size descending (largest segment first, like CozyFit reference)
            segments = sorted(segments, key=lambda s: s.get("size_pct", 0), reverse=True)

            _clone_slide(prs, T_SEGMENT_DIVIDER)
            slide_meta.append({"type": "divider", "content": {"title": "Market Segmentation"}})

            _build_segmentation_intro(prs)
            slide_meta.append({"type": "boilerplate", "content": {"title": "Benefits of Segmentation"}})

            # Segment overview (all segments at a glance)
            _build_segment_overview(prs, segments, project_id=project_id)
            slide_meta.append({"type": "segment_overview", "content": {"segments": [s.get("name") for s in segments]}})

            # "FOCUSING ON THE MOST DOMINANT SEGMENTS…"
            _build_focusing_segments(prs, segments, project_id=project_id)
            slide_meta.append({"type": "focusing", "content": {"segments": [s.get("name") for s in segments]}})

            # Individual segment slides: 8-slide pattern per segment
            # (matching sample PPT structure discovered in analysis)
            # 1. Meet Segment (full-bleed intro with persona quote)
            # 2. Respondent Profile (demographics)
            # 3. Closer Look 1 (premium/driver data with group shapes)
            # 4. Closer Look 2 (brand awareness + verbatim quotes)
            # 5. Behavioral Summary (4-column: social/drivers/pain/pre-purchase)
            # 6. Challenges Table (pain points + unmet needs)
            # 7. Social Media & Lifestyle (platform icons + lifestyle signals)
            # 8. Closer Look 3 (4 lifestyle signal cards)
            # Pre-compute unified pain point list across ALL segments
            # CozyFit pattern: same items for every segment, different percentages
            _cat = _extract_category(analysis, brand_name)
            unified_pain_items = _collect_unified_pain_items(segments[:5], category=_cat)

            # Inject brand context into segments for LLM card generation
            # Pre-compute unified Closer Look 3 dimensions (same 4 for all segments)
            _unified_cl3 = _collect_unified_closer_look_3(segments[:5], brand_name, _cat)
            for seg in segments[:5]:
                seg["_brand_name"] = brand_name
                seg["_category"] = _cat
                seg["_unified_cl3"] = _unified_cl3

            for seg in segments[:5]:
                slide = _build_meet_segment(prs, seg, project_id=project_id)
                slide_meta.append({"type": "meet_segment", "content": seg})

                _build_segment_profile(prs, seg)
                slide_meta.append({"type": "segment_profile", "content": {"segment": seg.get("name")}})

                _build_segment_closer_look(prs, seg, slide_num=1)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 1}})

                _build_segment_closer_look(prs, seg, slide_num=2)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 2}})

                _build_segment_behavioral_summary(prs, seg)
                slide_meta.append({"type": "behavioral_summary", "content": {"segment": seg.get("name")}})

                _build_segment_challenges(prs, seg, unified_pain_items=unified_pain_items)
                slide_meta.append({"type": "challenges", "content": {"segment": seg.get("name")}})

                slide = _build_segment_closer_look(prs, seg, slide_num=3)
                # Card images now generated inside _build_segment_closer_look (topic-relevant)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 3}})

        # ── Conflict Matrix slide (after segments, before target) ──
        conflict_data = analysis.get("conflict_matrix")
        if conflict_data and isinstance(conflict_data, dict) and conflict_data.get("conflicts"):
            _build_conflict_matrix_slide(prs, conflict_data)
            slide_meta.append({"type": "conflict_matrix", "content": conflict_data})

        # Target recommendation
        target = consumer.get("target_recommendation", {})

        # If no target data from DB, generate via LLM
        if not target and segments:
            _cat = _extract_category(analysis, brand_name)
            target = _llm_generate_target_section(segments[:5], brand_name, _cat)
            print(f"[target] LLM-generated target section for {brand_name}")

        if target:
            # Find the primary segment object for persona image generation
            primary_name = target.get("primary_segment", "")
            primary_seg = next(
                (s for s in segments if s.get("name", "").lower() == primary_name.lower()),
                segments[0] if segments else {}
            )

            # Generate varied AI images for target section slides
            _cat = _extract_category(analysis, brand_name)
            persona_img = _generate_persona_image_for_target(primary_seg, brand_name, category=_cat)
            group_img = _generate_target_section_image(primary_seg, brand_name, _cat, "group")
            scene_img = _generate_target_section_image(primary_seg, brand_name, _cat, "scene")

            # Slide 75: Selecting Target Audience comparison table
            _build_selecting_target(prs, segments, target, brand_name)
            slide_meta.append({"type": "selecting_target", "content": target})

            # Slide 76: Primary Target — persona with product
            slide = _build_target_recommendation(prs, target)
            if persona_img:
                _replace_slide_image(slide, persona_img)
            else:
                _replace_slide_image_websearch(slide, brand_name, f"{_cat} lifestyle photography")
            slide_meta.append({"type": "target", "content": target})

            # Slide 77: Why Target — group/community image
            slide = _build_why_target(prs, target)
            if group_img:
                _replace_slide_image(slide, group_img)
            else:
                _replace_slide_image_websearch(slide, brand_name, f"{_cat} collection lifestyle")
            slide_meta.append({"type": "why_target", "content": target})

            # Slide 78: Enables
            _build_enables_slide(prs, target)
            slide_meta.append({"type": "enables", "content": target})

            # "WHY NOT PRIORITIZE OTHER SEGMENTS"
            deprioritized = consumer.get("deprioritized_segments", [])
            if not deprioritized:
                deprioritized = target.get("deprioritized_segments", [])
            if not deprioritized and len(segments) > 1:
                # Generate segment-specific reasons based on their characteristics
                deprioritized = []
                for s in segments:
                    if s.get("name") == primary_name:
                        continue
                    sz = s.get("size_pct", "?")
                    tagline = s.get("tagline", "")
                    prem = s.get("what_premium_means", "")[:80]
                    reason = (
                        f"While {tagline.lower() if tagline else 'relevant'}, "
                        f"this {sz}% segment has different purchase drivers "
                        f"that don't align with {brand_name}'s core positioning"
                    )
                    deprioritized.append({"name": s.get("name", ""), "size_pct": sz, "reason": reason})
                deprioritized = deprioritized[:3]
            if deprioritized:
                _build_why_not_segments(prs, deprioritized, brand_name, category=_cat)
                slide_meta.append({"type": "why_not", "content": {"segments": deprioritized}})

            # "HOW [BRAND] FARES AGAINST THE COMPETITION"
            fares = consumer.get("competitive_fares", {})
            if fares:
                _build_competitive_fares(prs, fares, brand_name, category=_cat)
                slide_meta.append({"type": "fares", "content": fares})

            # Slide 79: Consumer Summary — lifestyle scene image
            cons_summary = target.get("consumer_summary") or consumer.get("consumer_summary", "")
            if cons_summary:
                slide = _build_consumer_summary(prs, cons_summary)
                if scene_img:
                    _replace_slide_image(slide, scene_img)
                else:
                    # Fallback: generate product hero shot (no text/banners)
                    _product_img = _generate_target_section_image(
                        primary_seg, brand_name, _cat, "product_hero")
                    if _product_img:
                        _replace_slide_image(slide, _product_img)
                    else:
                        _replace_slide_image_websearch(slide, brand_name, f"{_cat} product")
                slide_meta.append({"type": "consumer_summary", "content": {"text": cons_summary}})

    # ── Final Summary & Next Steps ───────────────────────────

    # Use LLM-generated summary from target section, or fall back to DB data
    summary_data = analysis.get("summary_and_next_steps", {})
    if not summary_data and target:
        summary_data = {
            "capabilities_column": target.get("summary_capabilities", ""),
            "competition_column": target.get("summary_competition", ""),
            "consumer_column": target.get("summary_consumer", ""),
            "closing_insight": target.get("summary_closing", ""),
        }
    if summary_data and phase == "full":
        _build_final_summary(prs, summary_data)
        slide_meta.append({"type": "final_summary", "content": summary_data})

    # ── Research Appendix ─────────────────────────────────────

    if phase == "full" and analysis.get("consumer", {}).get("charts"):
        _appendix_charts = analysis["consumer"]["charts"]
        # Extract sample size from research approach (same logic as consumer section)
        _appendix_n = 200
        for _ra in analysis.get("consumer", {}).get("research_approach", []):
            if isinstance(_ra, dict):
                import re as _re2
                _nm = _re2.search(r'n\s*[=:]\s*(\d+)', str(_ra.get("detail", "")))
                if _nm:
                    _appendix_n = int(_nm.group(1))
                    break
        _appendix_cat = _extract_category(analysis, brand_name)
        _n_appendix = _build_appendix_section(
            prs, _appendix_charts,
            sample_size=_appendix_n,
            brand_name=brand_name,
            category=_appendix_cat,
            segments=segments[:5] if segments else [],
        ) or 0
        # Transition slide
        slide_meta.append({"type": "appendix", "content": {"charts": _n_appendix}})
        # Individual chart slides
        for _ai in range(_n_appendix):
            slide_meta.append({"type": "appendix_chart", "content": {"index": _ai + 1}})

    # ── Thank You ─────────────────────────────────────────────

    _build_thank_you(prs)
    slide_meta.append({"type": "thank_you", "content": {}})

    # ── Save ──────────────────────────────────────────────────

    # Fix CJK fonts across all slides before saving
    _fix_cjk_fonts(prs)

    output_path = OUTPUT_DIR / f"project_{project_id}" / f"{brand_name}_Brand_Discovery.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Re-save to fix XML structure issues from slide cloning (LibreOffice compat)
    prs2 = Presentation(str(output_path))
    prs2.save(str(output_path))

    # Post-generation QA
    qa_result = _qa_check(output_path, phase=phase)
    slide_meta.append({"type": "_qa", "content": qa_result})

    # Generate previews
    preview_paths = _generate_previews(output_path, project_id)
    for i, pp in enumerate(preview_paths):
        if i < len(slide_meta):
            slide_meta[i]["preview_path"] = str(pp)

    return output_path, slide_meta


# ── Post-Generation QA ────────────────────────────────────────

def _qa_check(pptx_path: Path, phase: str = "full") -> dict:
    """Run automated quality checks on the generated PPTX.

    Checks:
    1. Text overflow — text that likely won't fit its shape
    2. Empty content shapes — shapes that should have text but don't
    3. Slide count sanity — too few or too many slides (phase-aware)
    4. Missing images — image placeholders with no image fill

    Returns:
        {"passed": bool, "score": int (0-100), "issues": [{"slide": n, "severity": str, "detail": str}]}
    """
    prs = Presentation(str(pptx_path))
    issues = []
    total_slides = len(prs.slides)

    # Phase-aware slide count thresholds
    _min_slides = {"brand_reality": 8, "market_structure": 15, "full": 25}
    _max_slides = {"brand_reality": 25, "market_structure": 50, "full": 120}
    min_expected = _min_slides.get(phase, 15)
    max_expected = _max_slides.get(phase, 120)

    if total_slides < min_expected:
        issues.append({"slide": 0, "severity": "high", "detail": f"Only {total_slides} slides for {phase} — likely missing sections (expected ≥{min_expected})"})
    elif total_slides > max_expected:
        issues.append({"slide": 0, "severity": "medium", "detail": f"{total_slides} slides — unusually long for {phase} (expected ≤{max_expected})"})

    empty_content_count = 0
    overflow_count = 0
    missing_image_count = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Check text overflow
            if shape.has_text_frame:
                tf = shape.text_frame
                full_text = tf.text.strip()

                if not full_text:
                    # Skip shapes that are likely decorative (very small or at edges)
                    if shape.width and shape.height:
                        w_in = shape.width / 914400  # EMU to inches
                        h_in = shape.height / 914400
                        if w_in > 2 and h_in > 0.5:
                            empty_content_count += 1
                            if empty_content_count <= 5:  # Only report first few
                                issues.append({
                                    "slide": slide_idx,
                                    "severity": "low",
                                    "detail": f"Empty text shape ({w_in:.1f}\" x {h_in:.1f}\")"
                                })
                    continue

                # Estimate text overflow: compare text length vs shape area
                if shape.width and shape.height:
                    w_in = shape.width / 914400
                    h_in = shape.height / 914400
                    area_sq_in = w_in * h_in

                    # Estimate average font size from first paragraph
                    avg_font_pt = 12
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                avg_font_pt = run.font.size.pt
                                break
                        if avg_font_pt != 12:
                            break

                    # Characters that fit: roughly (width / char_width) * (height / line_height)
                    char_width_in = avg_font_pt / 72 * 0.6  # approx char width
                    line_height_in = avg_font_pt / 72 * 1.4  # approx line spacing
                    if char_width_in > 0 and line_height_in > 0:
                        chars_per_line = max(1, int(w_in / char_width_in))
                        max_lines = max(1, int(h_in / line_height_in))
                        approx_capacity = chars_per_line * max_lines

                        # Large fonts (>24pt) often have auto-size/shrink — use looser threshold
                        overflow_ratio = 2.0 if avg_font_pt >= 24 else 1.5
                        if len(full_text) > approx_capacity * overflow_ratio:
                            overflow_count += 1
                            if overflow_count <= 8:
                                issues.append({
                                    "slide": slide_idx,
                                    "severity": "medium" if avg_font_pt < 24 else "low",
                                    "detail": (f"Potential text overflow: {len(full_text)} chars in "
                                               f"{w_in:.1f}\" x {h_in:.1f}\" at {avg_font_pt}pt "
                                               f"(~{approx_capacity} capacity). Text: \"{full_text[:60]}...\"")
                                })

            # Check for picture placeholders with no fill
            if shape.shape_type is not None:
                try:
                    shape_type_val = int(shape.shape_type)
                    if shape_type_val in (13, 14):  # Picture, Placeholder
                        # Check if image fill exists
                        sp = shape._element
                        blip_fills = sp.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
                        if not blip_fills:
                            missing_image_count += 1
                except (ValueError, TypeError):
                    pass

    if missing_image_count > 3:
        issues.append({
            "slide": 0,
            "severity": "medium",
            "detail": f"{missing_image_count} image placeholders appear empty"
        })

    if empty_content_count > 10:
        issues.append({
            "slide": 0,
            "severity": "medium",
            "detail": f"{empty_content_count} large text shapes are empty — possible unpopulated sections"
        })

    # Score: start at 100, deduct per issue
    score = 100
    for issue in issues:
        if issue["severity"] == "high":
            score -= 15
        elif issue["severity"] == "medium":
            score -= 5
        elif issue["severity"] == "low":
            score -= 2
    score = max(0, score)

    passed = score >= 60 and not any(i["severity"] == "high" for i in issues)

    # Log summary
    print(f"[qa] PPTX QA: {score}/100, {len(issues)} issues, {'PASS' if passed else 'WARN'}")
    for issue in issues[:10]:
        print(f"[qa]   [{issue['severity']}] Slide {issue['slide']}: {issue['detail'][:120]}")
    if len(issues) > 10:
        print(f"[qa]   ... and {len(issues) - 10} more issues")

    return {"passed": passed, "score": score, "issues": issues, "slide_count": total_slides}


# ── Preview Generation ───────────────────────────────────────

def _generate_previews(pptx_path: Path, project_id: int) -> list[Path]:
    """Convert PPTX slides to PNG previews via LibreOffice + PyMuPDF."""
    import subprocess
    import tempfile

    preview_dir = PREVIEW_DIR / f"project_{project_id}"
    preview_dir.mkdir(parents=True, exist_ok=True)

    # Clean old previews
    for old in preview_dir.glob("*.png"):
        old.unlink()

    # Step 1: PPTX -> PDF via LibreOffice
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run([
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", tmpdir, str(pptx_path)
            ], capture_output=True, timeout=120)

            pdf_files = list(Path(tmpdir).glob("*.pdf"))
            if not pdf_files:
                return _generate_placeholder_previews(pptx_path, preview_dir)

            # Step 2: PDF -> per-page PNG via PyMuPDF
            import fitz
            doc = fitz.open(str(pdf_files[0]))
            paths = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                png_path = preview_dir / f"slide_{i:03d}.png"
                pix.save(str(png_path))
                paths.append(png_path)
            doc.close()
            return paths

    except (FileNotFoundError, subprocess.TimeoutExpired, ImportError):
        pass

    return _generate_placeholder_previews(pptx_path, preview_dir)


def _generate_placeholder_previews(pptx_path: Path, preview_dir: Path) -> list[Path]:
    """Simple placeholder previews when LibreOffice unavailable."""
    from PIL import Image, ImageDraw

    prs = Presentation(str(pptx_path))
    paths = []

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (1280, 720), "#FAFAF9")
        draw = ImageDraw.Draw(img)
        draw.text((20, 20), f"Slide {i + 1}", fill="#E8652D")

        y = 80
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text[:100]
                if text.strip():
                    draw.text((40, y), text, fill="#292524")
                    y += 30
                    if y > 650:
                        break

        path = preview_dir / f"slide_{i:03d}.png"
        img.save(str(path))
        paths.append(path)

    return paths
