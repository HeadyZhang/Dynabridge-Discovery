"""File content extractor for DynaBridge case files.

Extracts structured text, tables, and metadata from
PPTX, PDF, DOCX, XLSX, and image files.
"""
from pathlib import Path


def extract_file(file_path: str) -> dict:
    """Extract structured content from a file.

    Args:
        file_path: Absolute or relative path to the file.

    Returns:
        {
            "source_file": str,
            "file_type": str,
            "content": {
                "slides": [...] | "paragraphs": [...] | "sheets": [...],
                "tables": [...],
                "raw_text": str,
            },
            "metadata": {
                "page_count": int,
                "word_count": int,
                "language_hint": str,  # "en", "zh", "en+zh"
                "quality": str,  # "high", "medium", "low"
            },
        }
    """
    path = Path(file_path)
    if not path.exists():
        return _error_result(file_path, f"File not found: {file_path}")

    ext = path.suffix.lower()

    extractors = {
        ".pptx": _extract_pptx,
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".doc": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xlsx,
    }

    image_exts = {".jpg", ".jpeg", ".png", ".svg", ".ai", ".psd", ".tif", ".tiff"}

    if ext in extractors:
        try:
            return extractors[ext](path)
        except Exception as e:
            return _error_result(file_path, str(e))
    elif ext in image_exts:
        return _extract_image(path)
    else:
        return _error_result(file_path, f"Unsupported file type: {ext}")


def _extract_pptx(path: Path) -> dict:
    """Extract text and tables from PowerPoint files."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    all_text_parts = []
    tables = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    tables.append({
                        "slide": i + 1,
                        "headers": rows[0] if rows else [],
                        "rows": rows[1:] if len(rows) > 1 else [],
                    })

        slide_text = "\n".join(slide_texts)
        slides.append({
            "slide_number": i + 1,
            "text": slide_text,
        })
        all_text_parts.append(slide_text)

    raw_text = "\n\n".join(all_text_parts)

    return {
        "source_file": str(path),
        "file_type": "pptx",
        "content": {
            "slides": slides,
            "tables": tables,
            "raw_text": raw_text[:50000],
        },
        "metadata": _build_metadata(raw_text, len(slides)),
    }


def _extract_pdf(path: Path) -> dict:
    """Extract text and tables from PDF files."""
    try:
        import pdfplumber
        return _extract_pdf_plumber(path)
    except ImportError:
        pass

    try:
        import fitz  # PyMuPDF
        return _extract_pdf_pymupdf(path)
    except ImportError:
        return _error_result(str(path), "No PDF library available")


def _extract_pdf_plumber(path: Path) -> dict:
    import pdfplumber

    pages = []
    tables = []
    all_text_parts = []

    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages[:100]):
            text = page.extract_text() or ""
            pages.append({
                "page_number": i + 1,
                "text": text,
            })
            all_text_parts.append(text)

            page_tables = page.extract_tables() or []
            for t in page_tables:
                if t and len(t) > 1:
                    tables.append({
                        "page": i + 1,
                        "headers": [str(c) for c in t[0]] if t[0] else [],
                        "rows": [[str(c) for c in row] for row in t[1:]],
                    })

    raw_text = "\n\n".join(all_text_parts)

    return {
        "source_file": str(path),
        "file_type": "pdf",
        "content": {
            "pages": pages,
            "tables": tables,
            "raw_text": raw_text[:50000],
        },
        "metadata": _build_metadata(raw_text, len(pages)),
    }


def _extract_pdf_pymupdf(path: Path) -> dict:
    import fitz

    doc = fitz.open(str(path))
    pages = []
    all_text_parts = []

    for i, page in enumerate(doc[:100]):
        text = page.get_text()
        pages.append({
            "page_number": i + 1,
            "text": text,
        })
        all_text_parts.append(text)

    doc.close()
    raw_text = "\n\n".join(all_text_parts)

    return {
        "source_file": str(path),
        "file_type": "pdf",
        "content": {
            "pages": pages,
            "tables": [],
            "raw_text": raw_text[:50000],
        },
        "metadata": _build_metadata(raw_text, len(pages)),
    }


def _extract_docx(path: Path) -> dict:
    """Extract text and tables from Word documents."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = []
    tables = []
    all_text_parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append({
                "text": text,
                "style": para.style.name if para.style else "",
            })
            all_text_parts.append(text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            tables.append({
                "headers": rows[0],
                "rows": rows[1:] if len(rows) > 1 else [],
            })

    raw_text = "\n\n".join(all_text_parts)

    return {
        "source_file": str(path),
        "file_type": "docx",
        "content": {
            "paragraphs": paragraphs,
            "tables": tables,
            "raw_text": raw_text[:50000],
        },
        "metadata": _build_metadata(raw_text, len(paragraphs)),
    }


def _extract_xlsx(path: Path) -> dict:
    """Extract sheet structure and sample data from Excel files."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    all_text_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 50), values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(cells)

        headers = rows[0] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else []

        sheet_text = " ".join(
            " ".join(row) for row in rows
        )
        all_text_parts.append(sheet_text)

        sheets.append({
            "name": sheet_name,
            "headers": headers,
            "row_count": ws.max_row or 0,
            "col_count": ws.max_column or 0,
            "sample_rows": data_rows[:10],
        })

    wb.close()
    raw_text = "\n\n".join(all_text_parts)

    return {
        "source_file": str(path),
        "file_type": "xlsx",
        "content": {
            "sheets": sheets,
            "tables": [],
            "raw_text": raw_text[:50000],
        },
        "metadata": _build_metadata(raw_text, len(sheets)),
    }


def _extract_image(path: Path) -> dict:
    """Record image metadata without extracting content."""
    size = path.stat().st_size
    return {
        "source_file": str(path),
        "file_type": "image",
        "content": {
            "raw_text": "",
        },
        "metadata": {
            "page_count": 1,
            "word_count": 0,
            "language_hint": "unknown",
            "quality": "low",
            "file_size_bytes": size,
        },
    }


def _error_result(file_path: str, error: str) -> dict:
    return {
        "source_file": file_path,
        "file_type": "error",
        "content": {"raw_text": "", "error": error},
        "metadata": {
            "page_count": 0,
            "word_count": 0,
            "language_hint": "unknown",
            "quality": "low",
        },
    }


def _build_metadata(raw_text: str, page_count: int) -> dict:
    """Build metadata from extracted text."""
    word_count = len(raw_text.split())
    lang = _detect_language(raw_text)

    if word_count > 500 and page_count > 5:
        quality = "high"
    elif word_count > 100:
        quality = "medium"
    else:
        quality = "low"

    return {
        "page_count": page_count,
        "word_count": word_count,
        "language_hint": lang,
        "quality": quality,
    }


def _detect_language(text: str) -> str:
    """Simple language detection based on character ranges."""
    if not text:
        return "unknown"

    sample = text[:5000]
    chinese_chars = sum(1 for c in sample if '\u4e00' <= c <= '\u9fff')
    ascii_chars = sum(1 for c in sample if c.isascii() and c.isalpha())
    total = chinese_chars + ascii_chars

    if total == 0:
        return "unknown"

    cn_ratio = chinese_chars / total

    if cn_ratio > 0.5:
        return "zh"
    elif cn_ratio > 0.1:
        return "en+zh"
    else:
        return "en"
