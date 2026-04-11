"""
Build the Brand Discovery master PPTX template — matching CozyFit reference.

Fonts: Montserrat (Light / Regular / SemiBold / Bold / Italic / SemiBoldItalic)
       Microsoft YaHei for Chinese text

Run standalone:  python build_template.py
Produces:        ../templates/brand_discovery_master.pptx
                 ../output/sample/CozyFit_Brand_Discovery_SAMPLE.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Brand Colors (exact from CozyFit PDF) ─────────────────────
ORANGE       = RGBColor(0xE8, 0x65, 0x2D)
ORANGE_LIGHT = RGBColor(0xF0, 0x9E, 0x72)  # lighter orange for chart series
BLUE         = RGBColor(0x00, 0xA8, 0xB5)   # teal/blue accent
DARK         = RGBColor(0x29, 0x25, 0x24)    # near-black text
GRAY_TEXT    = RGBColor(0x80, 0x80, 0x7A)
GRAY_BG      = RGBColor(0x58, 0x57, 0x52)   # section divider bg
GRAY_LIGHT   = RGBColor(0xD6, 0xD3, 0xD1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL     = RGBColor(0x33, 0x33, 0x30)   # thank-you bg

# Donut chart palette (CozyFit uses orange → brown tones)
DONUT_COLORS = [
    ORANGE,
    RGBColor(0xC5, 0x8B, 0x5F),   # brown/tan
    RGBColor(0x8B, 0x6D, 0x52),   # dark brown
    RGBColor(0xD4, 0xB8, 0x9C),   # light tan
    RGBColor(0xF5, 0xD5, 0xC3),   # very light peach
    GRAY_LIGHT,
    RGBColor(0xA0, 0x8B, 0x75),
    RGBColor(0xE8, 0xCF, 0xB8),
]

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Fonts — matching CozyFit PDF exactly
FONT         = "Montserrat"
FONT_LIGHT   = "Montserrat Light"
FONT_CN      = "Microsoft YaHei"

ASSETS_DIR = Path(__file__).parent.parent / "templates" / "assets"


# ── Low-level helpers ─────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_image(slide, img_path, left, top, width=None, height=None):
    """Add image. If width/height omitted, uses natural size."""
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(img_path), left, top, **kwargs)


def add_text(slide, left, top, width, height, text, size=Pt(16),
             color=DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             font_name=None, line_spacing=None, space_before=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name or FONT
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if space_before:
        p.space_before = Pt(space_before)
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=Pt(14),
                  color=DARK, bold=False, font_name=None, align=PP_ALIGN.LEFT,
                  line_spacing=None, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {text}" if bullet else text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name or FONT
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        p.space_after = Pt(4)
    return txBox


def add_rich_text(slide, left, top, width, height, runs):
    """runs: list of (text, size, color, bold, italic, font_name) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for i, item in enumerate(runs):
        text, size, color, bold, italic = item[:5]
        fname = item[5] if len(item) > 5 else FONT
        run = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
        if i == 0 and not p.runs:
            run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = fname
    return txBox


def add_circle_outline(slide, left, top, size, line_color, line_width=Pt(1.5)):
    """Outline-only circle (no fill) — matches CozyFit Venn style."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.background()  # no fill
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_dotted_line_v(slide, x, y_start, y_end, color=WHITE):
    """Vertical dotted line using small circles."""
    dot_spacing = Inches(0.12)
    dot_size = Inches(0.04)
    y = y_start
    while y < y_end:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dot_size, dot_size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        y += dot_spacing


def add_dotted_line_h(slide, x_start, x_end, y, color=WHITE):
    """Horizontal dotted line using small circles."""
    dot_spacing = Inches(0.12)
    dot_size = Inches(0.04)
    x = x_start
    while x < x_end:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dot_size, dot_size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        x += dot_spacing


def add_sample_footer(slide, text="Total sample; base n = 201"):
    """Italic footnote bottom-left — as in CozyFit data slides."""
    add_text(slide, Inches(0.5), Inches(6.9), Inches(4), Inches(0.4),
             text, Pt(10), DARK, italic=True, font_name=FONT)


# ── Venn Diagram (outline style, matching CozyFit p.3) ────────

def add_venn(slide, left, top, size, highlight=None):
    """Three overlapping outline circles. highlight: 'capabilities'/'competition'/'consumer'/None."""
    r = int(size * 0.38)
    cx = left + size // 2
    cy = top + size // 2

    positions = {
        'competition':  (cx, cy - int(r * 0.55)),
        'capabilities': (cx - int(r * 0.55), cy + int(r * 0.35)),
        'consumer':     (cx + int(r * 0.55), cy + int(r * 0.35)),
    }
    labels_map = {
        'competition':  ("Competition", BLUE),
        'capabilities': ("Capabilities", BLUE),
        'consumer':     ("Consumer", BLUE),
    }

    for key, (ox, oy) in positions.items():
        is_hl = (key == highlight)
        line_c = ORANGE if is_hl else RGBColor(0x5A, 0x7A, 0x8A)
        lw = Pt(2.5) if is_hl else Pt(1.5)
        add_circle_outline(slide, ox - r, oy - r, r * 2, line_c, lw)

        label, default_color = labels_map[key]
        text_color = ORANGE if is_hl else default_color
        add_text(slide, ox - int(r * 0.8), oy - Pt(7),
                 int(r * 1.6), int(r * 0.5),
                 label, Pt(12), text_color, bold=is_hl,
                 align=PP_ALIGN.CENTER, font_name=FONT)


# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_cover(prs, brand_name="[Brand Name]", subtitle="Brand Discovery",
                date="[DATE]"):
    """Cover: extracted CozyFit page as full-bleed bg, brand name + date overlaid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full-bleed background from extracted PDF page
    bg_path = ASSETS_DIR / "fixed_cover.png"
    if bg_path.exists():
        add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)

    # Brand name — overlaid on the orange right side
    add_text(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(1.5),
             brand_name, Pt(60), WHITE, font_name=FONT_LIGHT)

    # "Brand Discovery"
    add_text(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(1),
             subtitle, Pt(38), WHITE, font_name=FONT_LIGHT)

    # Date
    add_text(slide, Inches(7), Inches(5), Inches(5), Inches(0.5),
             date, Pt(15), WHITE, bold=True, font_name=FONT)

    return slide


def build_agenda(prs, items=None):
    """Agenda: plain text list left, image placeholder right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if items is None:
        items = [
            "Overview",
            "Our Approach",
            "Capabilities Review",
            "Competition Analysis",
            "Consumer Insights",
            "Next Steps",
            "Research Appendix",
        ]

    # Title
    add_text(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.5),
             "TODAY'S AGENDA", Pt(20), ORANGE, bold=True)

    # Items — Montserrat Light, large, stacked
    y = Inches(1.8)
    for item in items:
        add_text(slide, Inches(0.6), y, Inches(5.5), Inches(0.55),
                 item, Pt(26), DARK, font_name=FONT_LIGHT)
        y += Inches(0.65)

    # Right half: light gray bg + image placeholder
    add_rect(slide, Inches(6.666), 0, Inches(6.667), SLIDE_H,
             RGBColor(0xE8, 0xE5, 0xE0))
    add_text(slide, Inches(8), Inches(3.2), Inches(3.5), Inches(0.5),
             "[Brand / Product Image]", Pt(14), GRAY_TEXT,
             align=PP_ALIGN.CENTER)

    return slide


def build_approach(prs):
    """Our Brand Building Process — exact copy from CozyFit page 3."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg_path = ASSETS_DIR / "fixed_approach.png"
    if bg_path.exists():
        add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)

    return slide


def build_step_divider(prs, step_num=1, step_label="DISCOVERY", highlight=None):
    """Step marker — exact copy from CozyFit page 4."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg_path = ASSETS_DIR / "fixed_step1.png"
    if bg_path.exists():
        add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)

    return slide


def build_section_header(prs, text="A closer look at the\nbrand capabilities",
                         highlight='capabilities'):
    """Chapter divider — exact copy from CozyFit section pages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Pick the matching extracted page
    asset_map = {
        'capabilities': 'fixed_capabilities.png',
        'competition':  'fixed_competition.png',
        'consumer':     'fixed_consumer.png',
    }
    bg_file = asset_map.get(highlight, 'fixed_capabilities.png')
    bg_path = ASSETS_DIR / bg_file
    if bg_path.exists():
        add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)

    return slide


def build_insight_slide(prs, title="INSIGHT TITLE", subtitle_text=None,
                        bullets=None, insight_text=None, has_image=False):
    """Standard content slide: orange title, bullets, optional blue insight."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.55),
             title, Pt(20), ORANGE, bold=True)

    # Optional subtitle (italic, like data context)
    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.35),
                 subtitle_text, Pt(11), DARK, italic=True)

    if bullets is None:
        bullets = ["Key finding here", "Supporting evidence"]

    content_w = Inches(5.8) if has_image else Inches(11.5)

    # Bullets
    y = Inches(1.6) if subtitle_text else Inches(1.3)
    for b in bullets:
        add_text(slide, Inches(0.8), y, content_w, Inches(0.7),
                 f"•  {b}", Pt(14), DARK, line_spacing=20)
        y += Inches(0.75)

    # Image placeholder (right side)
    if has_image:
        add_rect(slide, Inches(7.2), Inches(1), Inches(5.5), Inches(4.5),
                 RGBColor(0xE8, 0xE5, 0xE0))
        add_text(slide, Inches(8.5), Inches(3), Inches(2.5), Inches(0.4),
                 "[Image]", Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

    # Blue insight box
    if insight_text:
        add_text(slide, Inches(0.8), Inches(6), Inches(11.5), Inches(0.8),
                 insight_text, Pt(14), BLUE, italic=True)

    return slide


def build_competitor_slide(prs, name="COMPETITOR", positioning=None,
                           key_learnings=None):
    """Competitor: image banner + dual column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if positioning is None:
        positioning = [("Target", "Detail"), ("Price Point", "$XX–$XX")]
    if key_learnings is None:
        key_learnings = [("Strength", "Detail"), ("Opportunity", "Detail")]

    # Title
    add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.55),
             f"{name} — POSITIONING & KEY LEARNINGS", Pt(18), ORANGE, bold=True)

    # Image banner placeholder
    add_rect(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(2.2),
             RGBColor(0xE8, 0xE5, 0xE0))
    add_text(slide, Inches(4.5), Inches(2), Inches(4), Inches(0.4),
             f"[{name} Brand Images]", Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

    # POSITIONING column
    col_l = Inches(0.6)
    add_text(slide, col_l, Inches(3.7), Inches(5.5), Inches(0.35),
             "POSITIONING", Pt(12), DARK, bold=True)

    y = Inches(4.2)
    for label, detail in positioning:
        add_rich_text(slide, col_l, y, Inches(5.5), Inches(0.55), [
            (f"{label}:  ", Pt(12), ORANGE, True, False),
            (detail, Pt(12), DARK, False, False),
        ])
        y += Inches(0.5)

    # KEY LEARNINGS column
    col_r = Inches(7)
    add_text(slide, col_r, Inches(3.7), Inches(5.5), Inches(0.35),
             "KEY LEARNINGS", Pt(12), DARK, bold=True)

    y = Inches(4.2)
    for label, detail in key_learnings:
        add_rich_text(slide, col_r, y, Inches(5.5), Inches(0.55), [
            (f"{label}:  ", Pt(12), BLUE, True, False),
            (detail, Pt(12), DARK, False, False),
        ])
        y += Inches(0.5)

    return slide


def build_summary_slide(prs, title="SUMMARY", text="", has_image=True):
    """Summary: paragraph left + optional image right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.55),
             title, Pt(20), ORANGE, bold=True)

    text_w = Inches(5.5) if has_image else Inches(11.5)
    add_text(slide, Inches(0.6), Inches(1.5), text_w, Inches(5),
             text, Pt(15), DARK, line_spacing=24, font_name=FONT)

    if has_image:
        add_rect(slide, Inches(7.2), Inches(1), Inches(5.5), Inches(5.2),
                 RGBColor(0xE8, 0xE5, 0xE0))
        add_text(slide, Inches(8.5), Inches(3.3), Inches(2.5), Inches(0.4),
                 "[Summary Visual]", Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

    return slide


def build_research_approach(prs, rows=None):
    """Research approach: dark bg + labeled rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.55),
             "RESEARCH APPROACH", Pt(20), WHITE, bold=True)

    if rows is None:
        rows = [
            ("Format", "Online survey"),
            ("Sample", "N=500"),
            ("Participants", "Target demographic"),
            ("Analysis", "Quantitative with cross-tabulation"),
            ("Timing", "Q1 2026"),
        ]

    y = Inches(1.8)
    for label, value in rows:
        # Orange label box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.6), y, Inches(2.2), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ORANGE
        shape.line.fill.background()
        shape.adjustments[0] = 0.15

        add_text(slide, Inches(0.75), y + Pt(2), Inches(2), Inches(0.4),
                 label, Pt(13), WHITE, bold=True)
        add_text(slide, Inches(3.1), y + Pt(2), Inches(9), Inches(0.4),
                 value, Pt(13), WHITE)
        y += Inches(0.75)

    return slide


def build_subsection_divider(prs, text="Demographics &\nBackground"):
    """Sub-section divider: gray bg + white text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRAY_BG)
    add_text(slide, Inches(0.8), Inches(2.5), Inches(8), Inches(2),
             text, Pt(36), WHITE, bold=True, font_name=FONT)
    return slide


def build_dual_chart_slide(prs, title="TITLE", subtitle_text=None,
                           left_title="", left_categories=None, left_values=None,
                           right_title="", right_categories=None, right_values=None,
                           left_type="donut", right_type="bar"):
    """Two charts side-by-side (CozyFit data page pattern)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.55),
             title, Pt(20), ORANGE, bold=True)

    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.3),
                 subtitle_text, Pt(11), DARK, italic=True)

    # Left sub-question
    if left_title:
        add_text(slide, Inches(0.3), Inches(1.5), Inches(6), Inches(0.8),
                 left_title, Pt(16), DARK, align=PP_ALIGN.CENTER, font_name=FONT)

    # Right sub-question
    if right_title:
        add_text(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.8),
                 right_title, Pt(16), DARK, align=PP_ALIGN.CENTER, font_name=FONT)

    # Left chart
    if left_categories and left_values:
        cd = CategoryChartData()
        cd.categories = left_categories
        cd.add_series("", left_values)
        ct = XL_CHART_TYPE.DOUGHNUT if left_type == "donut" else XL_CHART_TYPE.COLUMN_CLUSTERED
        cf = slide.shapes.add_chart(ct, Inches(0.3), Inches(2.5), Inches(5.8), Inches(4), cd)
        _style_chart(cf.chart, left_type, left_categories)

    # Right chart
    if right_categories and right_values:
        cd = CategoryChartData()
        cd.categories = right_categories
        cd.add_series("", right_values)
        ct = XL_CHART_TYPE.BAR_CLUSTERED if right_type == "hbar" else XL_CHART_TYPE.COLUMN_CLUSTERED
        cf = slide.shapes.add_chart(ct, Inches(6.8), Inches(2.5), Inches(6), Inches(4), cd)
        _style_chart(cf.chart, right_type, right_categories)

    add_sample_footer(slide)
    return slide


def build_bar_chart_slide(prs, title="TITLE", subtitle_text=None,
                          question="", categories=None, values=None,
                          is_horizontal=True):
    """Single chart data slide — horizontal or vertical bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.55),
             title, Pt(20), ORANGE, bold=True)

    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.3),
                 subtitle_text, Pt(11), DARK, italic=True)

    if question:
        add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.7),
                 question, Pt(18), DARK, align=PP_ALIGN.CENTER, font_name=FONT)

    if categories is None:
        categories = ["A", "B", "C"]
    if values is None:
        values = [50, 30, 20]

    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("", values)

    ct = XL_CHART_TYPE.BAR_CLUSTERED if is_horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_top = Inches(2.4) if question else Inches(1.5)
    cf = slide.shapes.add_chart(ct, Inches(0.8), chart_top,
                                 Inches(11.5), Inches(4.5), cd)
    _style_chart(cf.chart, "hbar" if is_horizontal else "bar", categories)

    add_sample_footer(slide)
    return slide


def build_donut_chart_slide(prs, title="TITLE", subtitle_text=None,
                            left_title="", left_categories=None, left_values=None,
                            right_title="", right_categories=None, right_values=None):
    """Donut left + horizontal bar right (CozyFit shopping habits pattern)."""
    return build_dual_chart_slide(
        prs, title, subtitle_text,
        left_title, left_categories, left_values,
        right_title, right_categories, right_values,
        left_type="donut", right_type="hbar"
    )


def build_next_steps(prs, steps=None):
    """Recommended next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.55),
             "RECOMMENDED NEXT STEPS", Pt(20), ORANGE, bold=True)

    if steps is None:
        steps = ["Step 1", "Step 2", "Step 3"]

    y = Inches(1.8)
    for i, step in enumerate(steps, 1):
        # Number
        add_text(slide, Inches(0.6), y, Inches(0.5), Inches(0.45),
                 str(i), Pt(18), ORANGE, bold=True)
        # Text
        add_text(slide, Inches(1.2), y, Inches(11), Inches(0.45),
                 step, Pt(15), DARK)
        y += Inches(0.7)

    return slide


def build_thank_you(prs, phone="", email="contact@dynabridge.com",
                    website="https://www.dynabridge.cn/"):
    """Thank You: exact CozyFit page 92 as background, contact info overlaid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full-bleed background from extracted PDF page
    bg_path = ASSETS_DIR / "fixed_thankyou.png"
    if bg_path.exists():
        add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)

    # Overlay contact info (positioned to match original layout)
    contact_lines = []
    if phone:
        contact_lines.append(f"电话/微信: {phone}")
    contact_lines.append(f"邮箱: {email}")
    contact_lines.append(f"网站: {website}")

    add_multiline(slide, Inches(0.8), Inches(5.2), Inches(6), Inches(1.5),
                  contact_lines, Pt(13), WHITE, font_name=FONT,
                  line_spacing=22)

    return slide


# ── Chart styling helper ──────────────────────────────────────

def _style_chart(chart, chart_type, categories):
    """Apply CozyFit styling to a chart."""
    plot = chart.plots[0]
    series = plot.series[0]

    if chart_type == "donut":
        # Color each slice
        for i in range(len(categories)):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = DONUT_COLORS[i % len(DONUT_COLORS)]
        series.has_data_labels = True
        series.data_labels.font.size = Pt(12)
        series.data_labels.font.color.rgb = DARK
        series.data_labels.font.bold = True
        series.data_labels.font.name = FONT
        series.data_labels.number_format = '0"%"'
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = FONT
        chart.legend.include_in_layout = False

    elif chart_type in ("bar", "hbar"):
        plot.gap_width = 60
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ORANGE
        series.has_data_labels = True
        dl = series.data_labels
        dl.font.size = Pt(11)
        dl.font.color.rgb = DARK
        dl.font.bold = True
        dl.font.name = FONT
        dl.number_format = '0"%"'
        dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

        chart.has_legend = False
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(11)
        cat_axis.tick_labels.font.color.rgb = DARK
        cat_axis.tick_labels.font.name = FONT
        cat_axis.major_tick_mark = 2  # none
        cat_axis.has_major_gridlines = False

        val_axis = chart.value_axis
        val_axis.visible = False
        val_axis.has_major_gridlines = False


# ═══════════════════════════════════════════════════════════════
#  MASTER TEMPLATE + SAMPLE GENERATOR
# ═══════════════════════════════════════════════════════════════

def build_master_template(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    build_agenda(prs)
    build_approach(prs)
    build_step_divider(prs, 1, "DISCOVERY")
    build_section_header(prs, "A closer look at the\nbrand capabilities", "capabilities")
    build_insight_slide(prs, "INSIGHT PLACEHOLDER", bullets=["Bullet 1", "Bullet 2"])
    build_competitor_slide(prs)
    build_summary_slide(prs)
    build_section_header(prs, "A closer look at the\nconsumer", "consumer")
    build_research_approach(prs)
    build_subsection_divider(prs)
    build_bar_chart_slide(prs, "BAR CHART")
    build_dual_chart_slide(prs, "DUAL CHART")
    build_next_steps(prs)
    build_thank_you(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✓ Master template: {output_path}")


def generate_sample():
    """Generate sample deck with CozyFit-style mock data."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    build_cover(prs, "CozyFit", "Brand Discovery", "FEBRUARY 2026")

    # 2. Agenda
    build_agenda(prs, [
        "Overview", "Our Approach", "Capabilities Review",
        "Competition Analysis", "Consumer Insights",
        "Next Steps", "Research Appendix",
    ])

    # 3. Approach
    build_approach(prs)

    # 4. Step 1
    build_step_divider(prs, 1, "DISCOVERY")

    # 5. Capabilities
    build_section_header(prs, "A closer look at the\nbrand capabilities", "capabilities")

    build_insight_slide(prs, "EXECUTION SUMMARY", bullets=[
        "CozyFit has built a strong foundation in medical scrubs with growing DTC presence",
        "Product line focuses on comfort and functionality with antimicrobial fabrics",
        "Pricing positioned in the mid-premium range ($45–85 per piece)",
        "Primary channel is e-commerce with emerging wholesale partnerships",
    ], insight_text="CozyFit has solid product-market fit but needs to differentiate more clearly from established players.",
       has_image=True)

    build_insight_slide(prs, "PRODUCT OFFERING", bullets=[
        "Core line: 12 SKUs across scrub tops, pants, and jackets",
        "Materials: proprietary CoolBreeze™ antimicrobial fabric blend",
        "Size range: XXS to 4XL with inclusive fit options",
        "Colors: 8 core colors + seasonal limited editions",
    ], insight_text="Strong technical fabric story but color range lags behind competitors like FIGS (20+ colors).")

    build_insight_slide(prs, "BRAND CHALLENGES", bullets=[
        "Limited brand awareness outside core nursing demographic",
        "No clear brand story or emotional positioning vs. competitors",
        "Heavy reliance on functional benefits without lifestyle appeal",
        "Social media presence lacks consistency and engagement strategy",
    ], insight_text="The brand needs to evolve from a product company to a lifestyle brand to command premium pricing.")

    build_summary_slide(prs, "CAPABILITIES SUMMARY",
        "CozyFit has established a solid product foundation with technically "
        "superior fabrics and inclusive sizing. However, the brand lacks a clear "
        "emotional positioning and storytelling framework that would differentiate "
        "it in an increasingly crowded market.\n\n"
        "Key areas for development include brand narrative, visual identity "
        "consistency, and lifestyle marketing approach.")

    # Competition
    build_section_header(prs, "A closer look at the\ncompetition", "competition")

    for name, pos, learn in [
        ("CARHARTT",
         [("Target", "Blue-collar workers expanding to lifestyle"),
          ("Price Point", "$25–$120"),
          ("Key Differentiator", "Heritage brand with 130+ year legacy"),
          ("Channel", "Wholesale dominant + growing DTC")],
         [("Strength", "Authentic workwear heritage creates trust"),
          ("Opportunity", "Slow to adapt to medical/healthcare segment"),
          ("Threat", "Brand extension power could enter medical easily")]),
        ("FIGS",
         [("Target", "Young healthcare professionals"),
          ("Price Point", "$38–$90"),
          ("Key Differentiator", "Fashion-forward medical apparel"),
          ("Channel", "100% DTC e-commerce")],
         [("Strength", "Created lifestyle category in medical wear"),
          ("Opportunity", "Premium pricing faces pushback in downturn"),
          ("Insight", "Ambassador program drives 40% of new customers")]),
        ("MED COUTURE",
         [("Target", "Fashion-conscious nurses"),
          ("Price Point", "$22–$55"),
          ("Key Differentiator", "Trend-driven designs at accessible prices"),
          ("Channel", "Wholesale + retail partnerships")],
         [("Strength", "Fast fashion approach keeps designs fresh"),
          ("Opportunity", "Lower quality perception vs. premium brands"),
          ("Insight", "Strong in print/pattern — an underserved niche")]),
    ]:
        build_competitor_slide(prs, name, pos, learn)

    build_summary_slide(prs, "COMPETITION SUMMARY",
        "The medical scrubs market is dominated by FIGS in premium DTC, "
        "with traditional players holding wholesale share. Carhartt represents "
        "a lateral threat with brand extension potential.\n\n"
        "The white space for CozyFit lies in combining FIGS-level brand appeal "
        "with superior technical fabrics and more inclusive positioning.")

    # Consumer
    build_section_header(prs, "A closer look at the\nconsumer", "consumer")

    build_research_approach(prs, [
        ("Format", "Online survey via Qualtrics"),
        ("Sample", "N=201 healthcare professionals"),
        ("Participants", "Nurses, doctors, technicians across US hospital systems"),
        ("Analysis", "Quantitative analysis with cross-tabulation by role/age/region"),
        ("Timing", "January 2026"),
    ])

    build_subsection_divider(prs, "Demographics &\nBackground")

    # Shopping habits (donut + hbar)
    build_donut_chart_slide(prs, "SCRUBS SHOPPING HABITS",
        subtitle_text="All data reflect a period of the last 12 months prior to January 2026",
        left_title="In the past 12 months, how often have you\npurchased scrubs?",
        left_categories=["Monthly or more", "Every 2–3 months", "2–3 times per year",
                          "Once per year", "Only when items wear out"],
        left_values=[18, 42, 27, 6, 6],
        right_title="Where have you purchased scrubs in the\npast 12 months?",
        right_categories=["Amazon", "Specialty uniform stores", "Walmart",
                           "Brand websites (DTC)", "Target",
                           "Employer-provided", "Secondhand / resale", "Other"],
        right_values=[59, 51, 51, 41, 26, 25, 12, 2],
    )

    # Feature importance
    build_bar_chart_slide(prs, "SCRUBS SHOPPING HABITS (CONTINUED)",
        subtitle_text="All data reflect a period of the last 12 months prior to January 2026",
        question="What matter most to you in scrubs?",
        categories=[
            "All-day comfort", "Stretch and flexibility",
            "Durability after repeated washing", "Breathability / moisture-wicking",
            "Easy care (wrinkle-resistant, quick-dry)",
            "Smart storage solutions (pockets)", "Fluid resistance",
            "Soft hand feel", "Odor resistance / no chemical smell",
        ],
        values=[61, 42, 40, 28, 27, 23, 18, 17, 13],
        is_horizontal=True,
    )

    # Consumer insights
    build_insight_slide(prs, "KEY CONSUMER INSIGHTS", bullets=[
        "Comfort is the #1 purchase driver, far ahead of brand loyalty (61% vs 13%)",
        "Amazon dominates as purchase channel (59%), but DTC has strong showing (41%)",
        "73% are willing to pay premium ($50+) for proven comfort and durability",
        "Social media (Instagram, TikTok) influences 56% of purchase decisions",
    ], insight_text="Clear opportunity to win on comfort + style without competing on brand heritage.")

    # Next steps
    build_next_steps(prs, [
        "Develop brand positioning centered on 'comfort innovation' narrative",
        "Conduct focus groups with Millennial nurses to validate brand concepts",
        "Design visual identity system balancing clinical trust with lifestyle appeal",
        "Create social media strategy targeting healthcare influencer partnerships",
        "Prototype seasonal color program to compete with FIGS' color range",
    ])

    # Thank you
    build_thank_you(prs,
                    phone="13736758116",
                    email="contact@dynabridge.com",
                    website="https://www.dynabridge.cn/")

    output_dir = Path(__file__).parent.parent / "output" / "sample"
    output_dir.mkdir(parents=True, exist_ok=True)
    out = output_dir / "CozyFit_Brand_Discovery_SAMPLE.pptx"
    prs.save(str(out))
    print(f"✓ Sample deck: {out}")
    return out


if __name__ == "__main__":
    base = Path(__file__).parent.parent
    build_master_template(base / "templates" / "brand_discovery_master.pptx")
    sample = generate_sample()
    print(f"\n  open '{sample}'")
